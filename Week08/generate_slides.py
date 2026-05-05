#!/usr/bin/env python3
"""Generate Week 8 Dashboard Tab slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ──────────────────────────────────────────────────────────
GREEN      = RGBColor(0x1B, 0xB8, 0x89)   # teal-green (P&L positive)
RED        = RGBColor(0xE5, 0x47, 0x47)   # soft red (P&L negative)
BLUE       = RGBColor(0x28, 0x7D, 0xFA)   # calendar blue
PURPLE     = RGBColor(0x8E, 0x44, 0xAD)   # journal purple
ORANGE     = RGBColor(0xF3, 0x96, 0x20)   # quick-actions orange
DARK       = RGBColor(0x1A, 0x1A, 0x2E)   # slide background
DARK2      = RGBColor(0x16, 0x21, 0x3E)   # card background
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT_BLUE = RGBColor(0xA8, 0xD8, 0xFF)

W = 13.33   # slide width  (inches, 16:9)
H = 7.5     # slide height (inches, 16:9)


# ── Helper: add a slide with solid dark background ──────────────────────────
def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide  = prs.slides.add_slide(layout)
    # Remove all placeholder shapes
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)
    # Solid dark background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


# ── Helper: add a text box ──────────────────────────────────────────────────
def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


# ── Helper: add a filled rounded rect ───────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ── Helper: bullet list ──────────────────────────────────────────────────────
def add_bullets(slide, items, x, y, w, h, size=16, color=WHITE,
                bullet_color=GREEN, indent=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("    • " if indent else "• ") + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color


# ── Helper: code block ───────────────────────────────────────────────────────
def add_code(slide, code_lines, x, y, w, h, size=11):
    bg = add_rect(slide, x, y, w, h, RGBColor(0x0D, 0x1B, 0x2A))
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.1),
        Inches(w - 0.3),  Inches(h - 0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = LIGHT_BLUE
        run.font.name = "Courier New"


# ════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)

# ── Slide 1: Title ──────────────────────────────────────────────────────────
slide = add_slide(prs)

# Accent bar
add_rect(slide, 0, 0, 0.08, H, GREEN)

# Course label
add_text(slide, "SmartFarmer Assistant · iOS SwiftUI Course",
         0.3, 0.3, 8, 0.4, size=13, color=GREY)

# Title
add_text(slide, "Week 8", 0.3, 0.9, 12, 1.0,
         size=52, bold=True, color=GREEN)
add_text(slide, "Dashboard Tab",
         0.3, 1.75, 12, 1.0, size=40, bold=True, color=WHITE)

# Subtitle
add_text(slide,
         "Aggregating data across Finance, Calendar & Journal\n"
         "into a single overview screen",
         0.3, 2.85, 10, 1.2, size=20, color=GREY)

# Module pills
pill_data = [
    (0.3,  4.3, GREEN,  "💰  Finance"),
    (2.6,  4.3, BLUE,   "📅  Calendar"),
    (4.9,  4.3, PURPLE, "📖  Journal"),
    (7.2,  4.3, ORANGE, "⚡  Quick Actions"),
]
for px, py, pc, pt in pill_data:
    add_rect(slide, px, py, 2.1, 0.5, pc)
    add_text(slide, pt, px + 0.1, py + 0.05, 2.0, 0.4, size=14, bold=True)

# ── Slide 2: Learning Objectives ────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.1, DARK2)
add_text(slide, "🎯  Learning Objectives", 0.4, 0.2, 12, 0.7,
         size=28, bold=True, color=GREEN)

objectives = [
    "Fetch data from multiple Core Data entities in one view",
    "Compute monthly profit/loss using pure ViewModel functions",
    "Build a gradient summary card with conditional colour",
    "Create a reusable DashboardSection<Content: View> component",
    "Wire Quick Action buttons to open Add sheets from the Dashboard",
    "Add a new tab to MainTabView without breaking existing deep links",
]
add_bullets(slide, objectives, 0.5, 1.3, 12.5, 5.5, size=19)

# ── Slide 3: Dashboard Overview ─────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "📊  Dashboard — What It Shows",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

sections = [
    (GREEN,  "Monthly P&L Card",
             "Net profit/loss · Income · Expense · Gradient colour (green = profit)"),
    (GREEN,  "Recent Transactions",
             "Last 3 transactions — type icon, note/category, date, amount"),
    (BLUE,   "Upcoming Activities",
             "Next 3 pending tasks — type icon, title, date, reminder bell"),
    (PURPLE, "Latest Journal Entry",
             "Most recent entry — weather icon, title, snippet, photo count"),
    (ORANGE, "Quick Actions",
             "3 buttons: Add Transaction · Add Activity · Add Journal Entry"),
]
for i, (color, title, desc) in enumerate(sections):
    y = 1.15 + i * 1.1
    add_rect(slide, 0.3, y, 0.06, 0.75, color)
    add_text(slide, title, 0.55, y,      12.5, 0.4, size=17, bold=True, color=color)
    add_text(slide, desc,  0.55, y + 0.38, 12.5, 0.4, size=14, color=GREY)

# ── Slide 4: Architecture ────────────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🏛  Architecture — Cross-Module Data Flow",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Left column: data sources
add_text(slide, "Core Data Stores", 0.4, 1.1, 4, 0.5,
         size=15, bold=True, color=GREY, align=PP_ALIGN.CENTER)
stores = [
    (GREEN,  "Transaction"),
    (BLUE,   "FarmActivity"),
    (PURPLE, "JournalEntry"),
]
for i, (c, name) in enumerate(stores):
    y = 1.7 + i * 1.25
    add_rect(slide, 0.3, y, 3.6, 0.7, c)
    add_text(slide, name, 0.35, y + 0.15, 3.5, 0.4,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow column
for i in range(3):
    y = 1.7 + i * 1.25 + 0.2
    add_text(slide, "→", 4.1, y, 1.0, 0.4, size=22, bold=True, color=GREY)

# Middle column: FetchRequest
add_text(slide, "@FetchRequest", 5.1, 1.1, 3.2, 0.5,
         size=15, bold=True, color=GREY, align=PP_ALIGN.CENTER)
for i, (c, req) in enumerate([
    (GREEN,  "@FetchRequest<Transaction>"),
    (BLUE,   "@FetchRequest<FarmActivity>"),
    (PURPLE, "@FetchRequest<JournalEntry>"),
]):
    y = 1.7 + i * 1.25
    add_rect(slide, 5.0, y, 3.3, 0.7, RGBColor(0x16, 0x21, 0x3E))
    add_text(slide, req, 5.05, y + 0.15, 3.2, 0.4, size=12, color=c)

# Arrow column 2
for i in range(3):
    y = 1.7 + i * 1.25 + 0.2
    add_text(slide, "→", 8.5, y, 0.8, 0.4, size=22, bold=True, color=GREY)

# Right column: ViewModel
add_text(slide, "DashboardViewModel", 9.3, 1.1, 3.5, 0.5,
         size=15, bold=True, color=GREY, align=PP_ALIGN.CENTER)
vm_methods = [
    "monthlyIncome/Expense()",
    "upcomingActivities()",
    "latestEntry()",
]
for i, (c, m) in enumerate(zip([GREEN, BLUE, PURPLE], vm_methods)):
    y = 1.7 + i * 1.25
    add_rect(slide, 9.2, y, 3.8, 0.7, RGBColor(0x0D, 0x1B, 0x2A))
    add_text(slide, m, 9.25, y + 0.15, 3.7, 0.4, size=13, color=c,
             italic=True)

add_text(slide,
         "Dashboard is read-only — it never writes to Core Data directly.",
         0.4, 6.6, 12.5, 0.5, size=14, color=ORANGE, italic=True)

# ── Slide 5: DashboardViewModel ─────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "⚙️  DashboardViewModel — Pure Functions",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

code = [
    "class DashboardViewModel: ObservableObject {",
    "",
    "    func monthlyIncome(_ transactions: [Transaction]) -> Double {",
    "        currentMonthTransactions(transactions)",
    '            .filter { $0.type == "income" }',
    "            .reduce(0) { $0 + $1.amount }",
    "    }",
    "",
    "    func upcomingActivities(_ activities: [FarmActivity], limit: Int = 3)",
    "        -> [FarmActivity] {",
    "        let startOfToday = Calendar.current.startOfDay(for: Date())",
    "        return activities",
    "            .filter { !$0.isCompleted && ($0.date ?? .distantPast) >= startOfToday }",
    "            .sorted { ($0.date ?? .distantPast) < ($1.date ?? .distantPast) }",
    "            .prefix(limit).map { $0 }",
    "    }",
    "}",
]
add_code(slide, code, 0.3, 1.1, 7.8, 5.9)

# Right panel: key points
add_text(slide, "Key Design Rules", 8.4, 1.1, 4.6, 0.5,
         size=17, bold=True, color=GREEN)
points = [
    "Receives plain [Entity] arrays\n→ no Core Data context in VM",
    "@FetchRequest belongs in the\nView, not the ViewModel",
    "Pure functions = easy to\nunit test without Core Data",
    "currentMonthTransactions() is\nprivate — only VM uses it",
    "limit: Int = 3 default\n→ callers can override",
]
for i, pt in enumerate(points):
    y = 1.75 + i * 1.1
    add_rect(slide, 8.3, y, 4.7, 0.85, DARK2)
    add_text(slide, pt, 8.45, y + 0.08, 4.5, 0.75, size=13, color=WHITE)

# ── Slide 6: Monthly P&L Card ────────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "💳  Monthly Profit / Loss Card",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Mock card — profit state
card_bg = slide.shapes.add_shape(1, Inches(0.3), Inches(1.15), Inches(5.6), Inches(2.8))
card_bg.fill.solid()
card_bg.fill.fore_color.rgb = RGBColor(0x12, 0x99, 0x77)
card_bg.line.fill.background()

add_text(slide, "ចំណេញ / ខាតប្រចាំខែ", 0.5, 1.25, 4, 0.4, size=13, color=WHITE)
add_text(slide, "ខែឧសភា ២០២៦",        0.5, 1.55, 4, 0.4, size=11, color=GREY)
add_text(slide, "$1,240.00",            0.5, 1.9,  4, 0.6, size=32, bold=True, color=WHITE)
add_text(slide, "ចំណូល  $2,100",       0.5, 2.6,  2.5, 0.4, size=14, color=WHITE)
add_text(slide, "ចំណាយ  $860",         3.0, 2.6,  2.5, 0.4, size=14, color=WHITE)

# Mock card — loss state
card_r = slide.shapes.add_shape(1, Inches(6.2), Inches(1.15), Inches(5.6), Inches(2.8))
card_r.fill.solid()
card_r.fill.fore_color.rgb = RGBColor(0xC0, 0x3B, 0x3B)
card_r.line.fill.background()

add_text(slide, "ចំណេញ / ខាតប្រចាំខែ", 6.4, 1.25, 4, 0.4, size=13, color=WHITE)
add_text(slide, "ខែឧសភា ២០២៦",        6.4, 1.55, 4, 0.4, size=11, color=GREY)
add_text(slide, "-$320.00",             6.4, 1.9,  4, 0.6, size=32, bold=True, color=WHITE)
add_text(slide, "ចំណូល  $540",         6.4, 2.6,  2.5, 0.4, size=14, color=WHITE)
add_text(slide, "ចំណាយ  $860",         9.0, 2.6,  2.5, 0.4, size=14, color=WHITE)

# Labels
add_text(slide, "✅ Profit", 1.5, 4.1, 3, 0.5, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(slide, "❌ Loss",   7.7, 4.1, 3, 0.5, size=16, bold=True, color=RED,   align=PP_ALIGN.CENTER)

# Key point
add_text(slide,
         "isProfit flag drives both the gradient colour AND the arrow SF Symbol — computed once, used twice.",
         0.3, 4.75, 12.7, 0.6, size=15, color=GREY, italic=True)

# Code snippet
code2 = [
    "let isProfit = pl >= 0",
    "",
    "LinearGradient(",
    "    colors: isProfit",
    "        ? [Color(red:0.12,green:0.7,blue:0.55), ...]",
    "        : [Color(red:0.9,green:0.28,blue:0.28), ...],",
    "    startPoint: .topLeading, endPoint: .bottomTrailing",
    ")",
]
add_code(slide, code2, 0.3, 5.4, 12.7, 1.9)

# ── Slide 7: DashboardSection Component ─────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🧩  Reusable DashboardSection<Content: View>",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

code3 = [
    "struct DashboardSection<Content: View>: View {",
    "    let title: String",
    "    let icon: String",
    "    let color: Color",
    "    @ViewBuilder let content: () -> Content  // ← any SwiftUI content",
    "",
    "    var body: some View {",
    "        VStack(alignment: .leading, spacing: 0) {",
    "            HStack(spacing: 8) {",
    "                Image(systemName: icon).foregroundColor(color)",
    "                Text(title).font(.headline)",
    "            }",
    "            .padding(.horizontal, 16).padding(.vertical, 12)",
    "            Divider()",
    "            content()",
    "        }",
    "        .background(Color(.systemBackground))",
    "        .cornerRadius(12)",
    "        .shadow(color: .black.opacity(0.06), radius: 8, x: 0, y: 2)",
    "    }",
    "}",
]
add_code(slide, code3, 0.3, 1.1, 7.8, 5.8)

add_text(slide, "Why @ViewBuilder?", 8.4, 1.1, 4.5, 0.4,
         size=16, bold=True, color=ORANGE)
add_text(slide,
         "Without @ViewBuilder, the closure\n"
         "can only have ONE expression.\n\n"
         "With @ViewBuilder, callers can write\n"
         "if, ForEach, multiple rows —\n"
         "full SwiftUI DSL syntax.",
         8.4, 1.65, 4.5, 2.2, size=14, color=WHITE)

add_text(slide, "Usage:", 8.4, 4.0, 4.5, 0.4, size=15, bold=True, color=GREEN)
code_usage = [
    'DashboardSection(',
    '    title: "ប្រតិបត្តិការ",',
    '    icon: "dollarsign.circle.fill",',
    '    color: .green',
    ') {',
    '    ForEach(recent) { tx in',
    '        TransactionDashboardRow(tx)',
    '    }',
    '}',
]
add_code(slide, code_usage, 8.3, 4.45, 4.7, 2.5, size=12)

# ── Slide 8: Quick Actions ───────────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "⚡  Quick Actions",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Mock quick-action buttons
button_data = [
    (0.4,  GREEN,  "plus.circle.fill", "ចំណូល/ចំណាយ"),
    (4.6,  BLUE,   "calendar.badge.plus", "សកម្មភាព"),
    (8.8,  PURPLE, "square.and.pencil", "កំណត់ហេតុ"),
]
for bx, bc, _, bt in button_data:
    add_rect(slide, bx, 1.3, 3.8, 2.0, DARK2)
    icon_bg = slide.shapes.add_shape(
        1, Inches(bx + 1.2), Inches(1.5), Inches(1.4), Inches(1.0)
    )
    icon_bg.fill.solid()
    icon_bg.fill.fore_color.rgb = bc
    icon_bg.line.fill.background()
    add_text(slide, bt, bx + 0.1, 2.65, 3.6, 0.5, size=14,
             bold=True, color=bc, align=PP_ALIGN.CENTER)

add_text(slide,
         "Each button sets a @State Bool flag → .sheet(isPresented:) opens the correct Add view.",
         0.4, 3.5, 12.5, 0.5, size=15, color=GREY, italic=True)

code4 = [
    "@State private var showAddTransaction = false",
    "@State private var showAddActivity    = false",
    "@State private var showAddJournal     = false",
    "",
    "QuickActionButton(title: \"ចំណូល/ចំណាយ\", icon: \"plus.circle.fill\", color: .green) {",
    "    showAddTransaction = true",
    "}",
    "",
    ".sheet(isPresented: $showAddTransaction) {",
    "    AddTransactionView().environment(\\.managedObjectContext, viewContext)",
    "}",
]
add_code(slide, code4, 0.3, 4.1, 12.7, 3.0)

# ── Slide 9: Updating MainTabView ────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🔧  Updating MainTabView — Adding Tab 0",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Before / After comparison
add_text(slide, "Before (Week 7)",  0.3, 1.1, 5.8, 0.4, size=15, bold=True, color=RED)
add_text(slide, "After (Week 8)",   7.1, 1.1, 5.8, 0.4, size=15, bold=True, color=GREEN)

before = [
    "FinanceTabView()    .tag(0)",
    "CalendarTabView()   .tag(1)",
    "PestGuideTabView()  .tag(2)",
    "JournalTabView()    .tag(3)",
    "",
    "// notification deep link:",
    "selectedTab = 1   // Calendar",
]
after = [
    "DashboardTabView()  .tag(0)  ← NEW",
    "FinanceTabView()    .tag(1)",
    "CalendarTabView()   .tag(2)",
    "PestGuideTabView()  .tag(3)",
    "JournalTabView()    .tag(4)",
    "",
    "// notification deep link:",
    "selectedTab = 2   // Calendar moved",
]
add_code(slide, before, 0.3, 1.6, 6.4, 3.2)
add_code(slide, after,  6.8, 1.6, 6.2, 3.2)

add_rect(slide, 0.3, 5.05, 12.7, 0.05, GREEN)

add_text(slide, "⚠️  Rule: always use explicit .tag() — implicit indices break silently when tabs are reordered.",
         0.4, 5.2, 12.5, 0.5, size=14, color=ORANGE)
add_text(slide, "⚠️  Update ALL hardcoded tab indices (here: the notification deep link).",
         0.4, 5.75, 12.5, 0.5, size=14, color=ORANGE)

# ── Slide 10: Common Mistakes ────────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🐛  Common Mistakes",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

mistakes = [
    ("Forget .tag() after reordering",
     "Notification jumps to wrong tab silently",
     "Always use explicit integer tags on every tab"),
    ("P&L logic in body computed property",
     "Business logic buried in the view",
     "Move to DashboardViewModel as a pure function"),
    ("Hold FetchedResults<T> in ViewModel",
     "@FetchRequest only works inside a View",
     "Pass Array(results) to the ViewModel method"),
    ("Miss @ViewBuilder on DashboardSection",
     "Compiler error: cannot use result builder",
     "Add @ViewBuilder to the closure parameter"),
    ("Omit .buttonStyle(PlainButtonStyle())",
     "Entire button row turns blue on press",
     "Add .buttonStyle(PlainButtonStyle()) to QuickActionButton"),
]
for i, (mistake, problem, fix) in enumerate(mistakes):
    y = 1.15 + i * 1.15
    add_rect(slide, 0.3, y, 0.06, 0.95, RED)
    add_text(slide, f"❌ {mistake}", 0.55, y,       8.5, 0.42, size=14, bold=True, color=RED)
    add_text(slide, f"→ {problem}", 0.55, y + 0.4,  8.5, 0.38, size=12, color=GREY)
    add_rect(slide, 9.0, y, 0.06, 0.95, GREEN)
    add_text(slide, f"✅ {fix}",    9.15, y + 0.2,  3.8, 0.6,  size=12, color=GREEN)

# ── Slide 11: UI / UX Summary ────────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🎨  UI / UX Summary",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

ux_items = [
    (GREEN,  "square.grid.2x2.fill", "Tab icon", "ផ្ទាំងគ្រប់គ្រង"),
    (GREEN,  "arrow.up.right.circle.fill", "P&L positive", "Green gradient + up arrow"),
    (RED,    "arrow.down.right.circle.fill", "P&L negative", "Red gradient + down arrow"),
    (GREEN,  "dollarsign.circle.fill", "Finance section", "Green accent"),
    (BLUE,   "calendar.badge.clock", "Calendar section", "Blue accent"),
    (PURPLE, "book.fill", "Journal section", "Purple accent"),
    (ORANGE, "bolt.fill", "Quick Actions header", "Orange accent"),
]

col_w = 5.8
for i, (c, sym, label, desc) in enumerate(ux_items):
    col = i % 2
    row = i // 2
    x = 0.3 + col * col_w + 0.1
    y = 1.1 + row * 1.35
    add_rect(slide, x, y, col_w - 0.1, 1.1, DARK2)
    add_rect(slide, x, y, 0.06, 1.1, c)
    add_text(slide, label, x + 0.2, y + 0.08, col_w - 0.4, 0.4,
             size=15, bold=True, color=c)
    add_text(slide, desc,  x + 0.2, y + 0.55, col_w - 0.4, 0.45,
             size=13, color=GREY)

# Last row odd item centred
if len(ux_items) % 2 == 1:
    c, sym, label, desc = ux_items[-1]
    add_rect(slide, 0.3, 1.1 + (len(ux_items)//2) * 1.35,
             col_w - 0.1, 1.1, DARK2)

# ── Slide 12: Summary & Next Steps ───────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "✅  Week 8 Summary",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

achieved = [
    "Built DashboardViewModel with pure functions for monthly P&L, recent items & upcoming tasks",
    "Created a conditional-colour LinearGradient hero card",
    "Factored DashboardSection<Content: View> to eliminate repeated card frames",
    "Wrote three row views — one per data source — keeping detail level appropriate for a dashboard",
    "Wired Quick Action buttons to AddTransactionView, AddActivityView, AddJournalEntryView",
    "Updated MainTabView: inserted Dashboard as tab 0, updated notification deep link to tab 2",
]
add_text(slide, "What we built:", 0.4, 1.1, 12.5, 0.45, size=18, bold=True, color=GREEN)
add_bullets(slide, achieved, 0.5, 1.65, 12.3, 3.5, size=16)

add_rect(slide, 0.3, 5.3, 12.7, 0.05, GREY)

add_text(slide, "Extension ideas:", 0.4, 5.5, 12.5, 0.45, size=18, bold=True, color=ORANGE)
extensions = [
    "Tap a Recent Transaction row → navigate directly to TransactionDetailView",
    "Add a simple bar chart for last 6 months P&L using SwiftUI Path/Shape",
    "Show a 'streak' counter for consecutive days with journal entries",
]
add_bullets(slide, extensions, 0.5, 6.05, 12.3, 1.3, size=15, color=GREY)

# ── Save ──────────────────────────────────────────────────────────────────────
output = os.path.join(os.path.dirname(__file__), "Week08_Dashboard_Tab.pptx")
prs.save(output)
print(f"Saved: {output}")
