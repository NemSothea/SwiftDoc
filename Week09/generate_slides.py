#!/usr/bin/env python3
"""Generate Week 9 Advanced UI & Animations slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ──────────────────────────────────────────────────────────
GREEN   = RGBColor(0x1B, 0xB8, 0x89)
BLUE    = RGBColor(0x28, 0x7D, 0xFA)
PURPLE  = RGBColor(0x8E, 0x44, 0xAD)
ORANGE  = RGBColor(0xF3, 0x96, 0x20)
PINK    = RGBColor(0xE9, 0x4F, 0x97)
TEAL    = RGBColor(0x00, 0xC9, 0xC8)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
DARK2   = RGBColor(0x16, 0x21, 0x3E)
CARD    = RGBColor(0x0F, 0x2A, 0x45)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xAA, 0xAA, 0xBB)
YELLOW  = RGBColor(0xFF, 0xD7, 0x00)

W = 13.33
H = 7.5


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None, radius=0):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, items, x, y, w, h, size=15, color=WHITE,
                bullet_color=GREEN, indent=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        dot = p.add_run()
        dot.text = "  • " if indent else "• "
        dot.font.size = Pt(size)
        dot.font.color.rgb = bullet_color
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_code(slide, code_lines, x, y, w, h, size=11):
    add_rect(slide, x, y, w, h, CARD)
    tb = slide.shapes.add_textbox(
        Inches(x + 0.18), Inches(y + 0.15),
        Inches(w - 0.36), Inches(h - 0.3)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = TEAL
        run.font.name = "Courier New"


def slide_number(slide, n, total=12):
    add_text(slide, f"{n} / {total}", W - 1.2, H - 0.42, 1.0, 0.35,
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
def slide_01(prs):
    s = add_slide(prs)

    # accent bar left
    add_rect(s, 0, 0, 0.12, H, GREEN)

    # badge
    add_rect(s, 1.0, 1.2, 2.0, 0.5, GREEN)
    add_text(s, "WEEK 09", 1.0, 1.2, 2.0, 0.5,
             size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    add_text(s, "UI កម្រិតខ្ពស់", 1.0, 1.9, 10.5, 1.0,
             size=46, bold=True, color=WHITE)
    add_text(s, "& Animations", 1.0, 2.75, 10.5, 1.0,
             size=46, bold=True, color=GREEN)

    add_text(s, "SmartFarmerAssistant — តុបតែង iOS App ឱ្យស្អាតជាក្រុមហ៊ុន",
             1.0, 3.8, 10.5, 0.5, size=16, italic=True, color=GREY)

    # icon row
    icons = [("🎨", "ViewModifiers"), ("🃏", "FarmCard"),
             ("✨", "Animations"), ("🔄", "Pull-to-Refresh"),
             ("🌙", "Dark Mode"), ("♿", "ភាពងាយស្រួលប្រើ")]
    for i, (icon, label) in enumerate(icons):
        bx = 1.0 + i * 2.05
        add_rect(s, bx, 4.7, 1.85, 1.35, DARK2)
        add_text(s, icon, bx, 4.75, 1.85, 0.6,
                 size=22, align=PP_ALIGN.CENTER)
        add_text(s, label, bx, 5.35, 1.85, 0.45,
                 size=10, color=GREY, align=PP_ALIGN.CENTER)

    slide_number(s, 1)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
def slide_02(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "📋  របៀបវារៈ — Week 9", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=GREEN)

    topics = [
        ("01", GREEN,  "Custom ViewModifiers", "រចនាប័ទ្មស្អាតឆ្លងកាត់គ្រប់ Views"),
        ("02", BLUE,   "Components ប្រើម្ដងហើយម្ដងទៀត",   "FarmCard · PrimaryButton · SectionHeader"),
        ("03", ORANGE, "Animations",            "Fade-in lists · Scale-press buttons · Slide-in"),
        ("04", PURPLE, "Pull-to-Refresh",        "Async refreshable + skeleton loading states"),
        ("05", TEAL,   "Dark Mode",              "Adaptive colors · systemBackground semantics"),
        ("06", PINK,   "ភាពងាយស្រួលប្រើ",          "VoiceOver labels · Dynamic Type · A11y hints"),
    ]

    for i, (num, color, title, sub) in enumerate(topics):
        row_y = 1.1 + i * 1.0
        add_rect(s, 0.4, row_y, 0.55, 0.62, color)
        add_text(s, num, 0.4, row_y, 0.55, 0.62,
                 size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.1, row_y, 5.5, 0.38,
                 size=16, bold=True, color=WHITE)
        add_text(s, sub, 1.1, row_y + 0.36, 5.5, 0.3,
                 size=12, color=GREY)

    slide_number(s, 2)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What are ViewModifiers?
# ════════════════════════════════════════════════════════════════════════════
def slide_03(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🎨  Custom ViewModifiers", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=GREEN)

    add_text(s, "ViewModifier គឺជាអ្វី?", 0.4, 1.05, 12.5, 0.45,
             size=19, bold=True, color=WHITE)
    add_bullets(s, [
        "ប្លុក transformations ប្រើម្ដងហើយម្ដងទៀត (background, shadow, corner radius…)",
        "ប្រើ .modifier(MyModifier()) ឬ extension .myModifier() ខ្លីជាង",
        "View code DRY — ផ្លាស់ប្ដូរម្ដង ធ្វើបច្ចុប្បន្នភាព cards ទាំងអស់",
        "អាចមាន @State, @Environment ហើយ respond ទៅ dark mode / accessibility",
    ], 0.4, 1.55, 8.0, 2.0, size=14, bullet_color=GREEN)

    add_code(s, [
        "struct FarmCardModifier: ViewModifier {",
        "    func body(content: Content) -> some View {",
        "        content",
        "            .background(Color(.systemBackground))",
        "            .cornerRadius(16)",
        "            .shadow(color: .black.opacity(0.08),",
        "                    radius: 10, x: 0, y: 4)",
        "    }",
        "}",
        "",
        "extension View {",
        '    func farmCard() -> some View { modifier(FarmCardModifier()) }',
        "}",
    ], 8.55, 0.95, 4.55, 3.5, size=10)

    add_text(s, "✅  ការប្រើប្រាស់", 0.4, 3.65, 8.0, 0.35,
             size=14, bold=True, color=TEAL)
    add_code(s, [
        "VStack { ... }",
        "    .farmCard()   // applies everywhere consistently",
    ], 0.4, 4.05, 8.0, 0.9, size=12)

    slide_number(s, 3)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FadeIn & ScalePress Modifiers
# ════════════════════════════════════════════════════════════════════════════
def slide_04(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "✨  FadeIn & ScalePress Modifiers", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=ORANGE)

    # FadeIn
    add_rect(s, 0.3, 1.0, 6.1, 2.85, CARD)
    add_text(s, "FadeInModifier — entrance list ជំហាន",
             0.5, 1.05, 5.8, 0.38, size=13, bold=True, color=ORANGE)
    add_code(s, [
        "struct FadeInModifier: ViewModifier {",
        "    let delay: Double",
        "    @State private var opacity: Double = 0",
        "",
        "    func body(content: Content) -> some View {",
        "        content",
        "            .opacity(opacity)",
        "            .onAppear {",
        "                withAnimation(.easeOut(duration: 0.5)",
        "                              .delay(delay)) {",
        "                    opacity = 1",
        "                }",
        "            }",
        "    }",
        "}",
    ], 0.3, 1.45, 6.1, 2.35, size=9.5)

    # ScalePress
    add_rect(s, 6.6, 1.0, 6.4, 2.85, CARD)
    add_text(s, "ScalePressModifier — feedback button ពេលចុច",
             6.8, 1.05, 6.0, 0.38, size=13, bold=True, color=BLUE)
    add_code(s, [
        "struct ScalePressModifier: ViewModifier {",
        "    @GestureState private var isPressed = false",
        "",
        "    func body(content: Content) -> some View {",
        "        content",
        "            .scaleEffect(isPressed ? 0.96 : 1.0)",
        "            .animation(.spring(response: 0.3,",
        "                       dampingFraction: 0.6),",
        "                       value: isPressed)",
        "            .simultaneousGesture(",
        "                DragGesture(minimumDistance: 0)",
        "                    .updating($isPressed) {",
        "                        _, state, _ in state = true",
        "                    }",
        "            )",
        "    }",
        "}",
    ], 6.6, 1.45, 6.4, 2.35, size=9.5)

    # Usage
    add_text(s, "ការប្រើប្រាស់ — stagger cards ក្នុង Dashboard:", 0.4, 3.95, 8.5, 0.35,
             size=13, bold=True, color=TEAL)
    add_code(s, [
        "monthlyProfitLossCard.fadeIn(delay: 0.1)",
        "recentTransactionsSection.fadeIn(delay: 0.2)",
        "QuickActionButton().scalePress()",
    ], 0.4, 4.35, 12.5, 0.95, size=12)

    slide_number(s, 4)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FarmCard Component
# ════════════════════════════════════════════════════════════════════════════
def slide_05(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🃏  FarmCard — Container ទូទៅ", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=BLUE)

    add_code(s, [
        "struct FarmCard<Content: View>: View {",
        "    let title: String",
        "    let icon: String",
        "    let iconColor: Color",
        "    @ViewBuilder let content: () -> Content",
        "",
        "    var body: some View {",
        "        VStack(alignment: .leading, spacing: 0) {",
        "            HStack(spacing: 10) {",
        "                Image(systemName: icon)",
        "                    .foregroundColor(iconColor)",
        "                Text(title).font(.headline)",
        "                Spacer()",
        "            }",
        "            .padding(.horizontal, 16).padding(.vertical, 14)",
        "            Divider().padding(.horizontal, 16)",
        "            content()",
        "                .padding(.horizontal, 16).padding(.vertical, 12)",
        "        }",
        "        .farmCard()",
        "    }",
        "}",
    ], 0.3, 0.98, 6.5, 5.8, size=10)

    # Visual card preview
    add_rect(s, 7.0, 1.1, 6.0, 2.2, RGBColor(0x1E, 0x1E, 0x2E))
    add_rect(s, 7.0, 1.1, 6.0, 0.55, RGBColor(0x28, 0x28, 0x38))
    add_text(s, "💰  ចំណូល / ចំណាយ", 7.15, 1.15, 5.7, 0.42,
             size=13, bold=True, color=WHITE)
    add_rect(s, 7.05, 1.63, 5.9, 0.02, GREY)
    add_text(s, "ចំណូល", 7.2, 1.72, 2.5, 0.3, size=12, color=GREY)
    add_text(s, "$1,200", 10.5, 1.72, 2.35, 0.3, size=12, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    add_text(s, "ចំណាយ", 7.2, 2.08, 2.5, 0.3, size=12, color=GREY)
    add_text(s, "$850", 10.5, 2.08, 2.35, 0.3, size=12, bold=True, color=RGBColor(0xE5, 0x47, 0x47), align=PP_ALIGN.RIGHT)
    add_text(s, "← FarmCard render content ណាមួយ", 7.2, 3.45, 5.6, 0.3, size=10, italic=True, color=GREY)

    add_text(s, "ការប្រើប្រាស់:", 7.0, 3.85, 6.0, 0.35, size=13, bold=True, color=TEAL)
    add_code(s, [
        'FarmCard(title: "ចំណូល / ចំណាយ",',
        '         icon: "dollarsign.circle.fill",',
        '         iconColor: .green) {',
        "    // SwiftUI view ណាមួយ",
        "}",
    ], 7.0, 4.25, 6.0, 1.9, size=11)

    slide_number(s, 5)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PrimaryButton & SectionHeader
# ════════════════════════════════════════════════════════════════════════════
def slide_06(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🔘  PrimaryButton & SectionHeader", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=PURPLE)

    # PrimaryButton
    add_text(s, "PrimaryButton", 0.4, 1.0, 6.5, 0.38, size=16, bold=True, color=PURPLE)
    add_code(s, [
        "struct PrimaryButton: View {",
        "    let title: String",
        "    var icon: String? = nil",
        "    var color: Color = .green",
        "    let action: () -> Void",
        "",
        "    var body: some View {",
        "        Button(action: action) {",
        "            HStack(spacing: 8) {",
        "                if let icon { Image(systemName: icon) }",
        "                Text(title).font(.system(size: 17, weight: .semibold))",
        "            }",
        "            .foregroundColor(.white)",
        "            .frame(maxWidth: .infinity).frame(height: 54)",
        "            .background(color).cornerRadius(14)",
        "        }",
        "        .scalePress()",
        "        .accessibilityLabel(title)",
        "    }",
        "}",
    ], 0.3, 1.4, 6.6, 5.6, size=9.5)

    # SectionHeader
    add_text(s, "SectionHeader", 7.1, 1.0, 6.0, 0.38, size=16, bold=True, color=TEAL)
    add_code(s, [
        "struct SectionHeader: View {",
        "    let title: String",
        "    let icon: String",
        "    let color: Color",
        "    var actionTitle: String? = nil",
        "    var action: (() -> Void)? = nil",
        "",
        "    var body: some View {",
        "        HStack(spacing: 8) {",
        "            Image(systemName: icon)",
        "                .foregroundColor(color)",
        "            Text(title).font(.headline).bold()",
        "            Spacer()",
        "            if let t = actionTitle, let a = action {",
        "                Button(action: a) {",
        "                    Text(t).font(.subheadline)",
        "                        .foregroundColor(color)",
        "                }",
        "            }",
        "        }",
        "    }",
        "}",
    ], 7.0, 1.4, 6.1, 5.6, size=9.5)

    slide_number(s, 6)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Slide Animations
# ════════════════════════════════════════════════════════════════════════════
def slide_07(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🎬  Slide & List Animations", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=ORANGE)

    add_text(s, "Pattern animation ០៣ ប្រើក្នុង Week 9:", 0.4, 1.0, 12.5, 0.4,
             size=15, color=GREY)

    patterns = [
        ("Fade-In (Staggered)", ORANGE,
         ["ប្រើលើ Dashboard section នីមួយៗ នៅ .onAppear",
          "delay: 0.1s ក្នុង section → បង្កើត waterfall effect",
          ".easeOut(duration: 0.5).delay(delay)"],
         [".fadeIn(delay: 0.1)  // P&L card",
          ".fadeIn(delay: 0.2)  // transactions",
          ".fadeIn(delay: 0.3)  // activities"]),
        ("Scale Press", BLUE,
         ["ប្រើ @GestureState + DragGesture(minimumDistance:0)",
          "scaleEffect 1.0 → 0.96 ពេលចុច ត្រឡប់វិញ",
          "ប្រើលើ QuickActionButton, PrimaryButton, toolbar"],
         [".scalePress()   // on any tappable view",
          "// spring(response:0.3, dampingFraction:0.6)"]),
        ("List Rows (ក្នុងrow)", GREEN,
         ["PestGuideTabView: enumerated ForEach",
          "Each row: .fadeIn(delay: Double(index) * 0.05)",
          "Max ~20 rows × 0.05 = 1.0 វ stagger សរុប"],
         ["ForEach(Array(pests.enumerated()), ...) {",
          "    PestRowView(pest: pest)",
          "        .fadeIn(delay: Double(index) * 0.05)"]),
    ]

    for i, (title, color, bullets, code) in enumerate(patterns):
        bx = 0.3 + i * 4.35
        add_rect(s, bx, 1.5, 4.15, 5.5, CARD)
        add_text(s, title, bx + 0.15, 1.58, 3.85, 0.38,
                 size=13, bold=True, color=color)
        add_rect(s, bx + 0.1, 1.94, 3.95, 0.02, color)
        for j, b in enumerate(bullets):
            add_text(s, "• " + b, bx + 0.15, 2.05 + j * 0.45, 3.85, 0.42,
                     size=10.5, color=WHITE)
        add_code(s, code, bx + 0.05, 3.6, 4.05, 1.35, size=9)

    slide_number(s, 7)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Pull-to-Refresh & Loading States
# ════════════════════════════════════════════════════════════════════════════
def slide_08(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🔄  Pull-to-Refresh & Loading States", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=TEAL)

    # Left: refreshable
    add_rect(s, 0.3, 1.0, 6.2, 3.2, CARD)
    add_text(s, ".refreshable — SwiftUI native (iOS 15+)",
             0.5, 1.08, 5.8, 0.38, size=13, bold=True, color=TEAL)
    add_code(s, [
        "ScrollView {",
        "    VStack { ... }",
        "}",
        ".refreshable {",
        "    isRefreshing = true",
        "    try? await Task.sleep(",
        "        nanoseconds: 1_500_000_000)",
        "    // re-fetch data here",
        "    isRefreshing = false",
        "}",
    ], 0.3, 1.5, 6.2, 2.65, size=11)

    # Right: LoadingRowView
    add_rect(s, 6.75, 1.0, 6.25, 3.2, CARD)
    add_text(s, "LoadingRowView — skeleton shimmer",
             6.95, 1.08, 5.85, 0.38, size=13, bold=True, color=ORANGE)
    add_code(s, [
        "struct LoadingRowView: View {",
        "    @State private var pulse = false",
        "",
        "    var body: some View {",
        "        VStack(alignment: .leading, spacing: 10) {",
        "            RoundedRectangle(cornerRadius: 6)",
        "                .fill(Color(.systemGray4))",
        "                .frame(width: 180, height: 14)",
        "            // ... 2 more skeleton lines",
        "        }",
        "        .opacity(pulse ? 1.0 : 0.4)",
        "        .onAppear {",
        "            withAnimation(.easeInOut(duration: 0.9)",
        "                          .repeatForever(autoreverses: true))",
        "            { pulse = true }",
        "        }",
        "    }",
        "}",
    ], 6.75, 1.5, 6.25, 2.65, size=9.5)

    # Trigger logic
    add_text(s, "បង្ហាញ skeleton ពេល refresh ក្នុង PestGuideTabView:", 0.4, 4.3, 12.5, 0.35,
             size=13, bold=True, color=WHITE)
    add_code(s, [
        "if isLoading {",
        "    List { ForEach(0..<5, id: \\.self) { _ in LoadingRowView() } }",
        "} else {",
        "    List { /* real pest rows */ }",
        "        .refreshable { isLoading = true; await Task.sleep(...); isLoading = false }",
        "}",
    ], 0.3, 4.72, 12.75, 1.95, size=11)

    slide_number(s, 8)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Dark Mode Support
# ════════════════════════════════════════════════════════════════════════════
def slide_09(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🌙  ការគាំទ្រ Dark Mode", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=PURPLE)

    add_text(s, "ប្រព័ន្ធ adaptive color SwiftUI — dark mode ស្វ័យប្រវត្តិ:", 0.4, 1.0, 12.5, 0.38,
             size=14, color=GREY)

    # Color table
    headers = ["ការប្រើ Color", "SwiftUI Idiom", "ស្ដាប់ Dynamic Mode"]
    col_x = [0.4, 5.2, 10.0]
    col_w = [4.6, 4.6, 3.0]

    add_rect(s, 0.35, 1.45, 12.65, 0.42, DARK2)
    for j, h in enumerate(headers):
        add_text(s, h, col_x[j], 1.48, col_w[j], 0.35,
                 size=12, bold=True, color=TEAL)

    rows = [
        ["Card background", "Color(.systemBackground)", "✅ Yes"],
        ["Grouped screen bg", "Color(.systemGroupedBackground)", "✅ Yes"],
        ["Skeleton shimmer", "Color(.systemGray4/5)", "✅ Yes"],
        ["Primary text", ".primary / .foregroundColor(.primary)", "✅ Yes"],
        ["Shadows", ".black.opacity(0.08)", "Subtle — stays subtle"],
        ["Gradient card", "Custom RGB — keep low saturation", "⚠️ Test manually"],
    ]

    for i, row in enumerate(rows):
        row_y = 1.92 + i * 0.52
        bg = CARD if i % 2 == 0 else DARK2
        add_rect(s, 0.35, row_y, 12.65, 0.5, bg)
        for j, cell in enumerate(row):
            cell_color = GREEN if cell.startswith("✅") else (ORANGE if cell.startswith("⚠️") else WHITE)
            add_text(s, cell, col_x[j] + 0.1, row_y + 0.08, col_w[j] - 0.2, 0.36,
                     size=11, color=cell_color)

    add_text(s, "ច្បាប់សំខាន់:", 0.4, 5.15, 2.0, 0.35, size=13, bold=True, color=YELLOW)
    add_text(s,
             "ប្រើឈ្មោះ semantic UIColor (systemBackground, label) ជាជំនួស hex hardcoded។ SwiftUI ផ្លាស់ប្ដូរដោយស្វ័យប្រវត្តិ",
             2.4, 5.15, 10.6, 0.55, size=12, color=WHITE)
    add_text(s, "Test dark mode: Simulator → Features → Toggle Appearance  (Cmd+Shift+A)",
             0.4, 5.8, 12.5, 0.38, size=11, italic=True, color=GREY)

    slide_number(s, 9)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Accessibility
# ════════════════════════════════════════════════════════════════════════════
def slide_10(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "♿  ភាពងាយស្រួលប្រើ", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=PINK)

    # Left column — modifiers
    add_text(s, "Accessibility Modifiers ដែលប្រើ", 0.4, 1.0, 6.0, 0.38,
             size=15, bold=True, color=WHITE)
    modifiers = [
        (".accessibilityLabel(\"...\")", "ជំនួស label លំនាំដើម សម្រាប់ VoiceOver"),
        (".accessibilityHint(\"...\")", "ពិពណ៌នា អ្វីដែលកើតឡើងពេល activate"),
        (".accessibilityElement(children: .combine)", "ដាក់ sub-elements ក្នុង unit ដែល focus បាន"),
        (".accessibilityLabel(pest.name ?? \"\")", "label ក្នុង row នីមួយៗ PestGuide list"),
        ("NavigationView label", "បន្ថែមទៅ nav កម្រិតខ្ពស់ ដើម្បី identity screen"),
    ]
    for i, (mod, desc) in enumerate(modifiers):
        row_y = 1.48 + i * 0.88
        add_rect(s, 0.35, row_y, 6.1, 0.82, CARD)
        add_text(s, mod, 0.5, row_y + 0.05, 5.8, 0.32,
                 size=11, bold=True, color=TEAL, wrap=False)
        add_text(s, desc, 0.5, row_y + 0.38, 5.8, 0.35,
                 size=11, color=GREY)

    # Right column — code sample
    add_text(s, "AccessibilityCardModifier (custom)", 6.7, 1.0, 6.2, 0.38,
             size=15, bold=True, color=WHITE)
    add_code(s, [
        "struct AccessibilityCardModifier: ViewModifier {",
        "    let label: String",
        "    let hint: String",
        "",
        "    func body(content: Content) -> some View {",
        "        content",
        "            .accessibilityElement(children: .combine)",
        "            .accessibilityLabel(label)",
        "            .accessibilityHint(hint)",
        "    }",
        "}",
        "",
        "// P&L card in Dashboard:",
        "monthlyProfitLossCard",
        "    .accessibilityElement(children: .combine)",
        '    .accessibilityLabel("ចំណេញ / ខាត \\(pl.formattedCurrency)")',
    ], 6.65, 1.45, 6.4, 4.0, size=10)

    add_text(s, "Dynamic Type:", 0.4, 6.0, 2.5, 0.32,
             size=12, bold=True, color=YELLOW)
    add_text(s,
             "ប្រើ .font(.headline), .font(.caption) — ជំនួស Pt fixed — ឱ្យ text scale ទៅតាម preferences",
             2.85, 6.0, 10.0, 0.38, size=12, color=WHITE)

    slide_number(s, 10)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Before / After visual comparison
# ════════════════════════════════════════════════════════════════════════════
def slide_11(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "📊  មុន vs ក្រោយ — ការតុបតែង Week 9", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=YELLOW)

    # Before column
    add_rect(s, 0.3, 1.0, 6.0, 6.1, CARD)
    add_rect(s, 0.3, 1.0, 6.0, 0.45, RGBColor(0xE5, 0x47, 0x47))
    add_text(s, "❌  មុន (Week 8)", 0.4, 1.04, 5.8, 0.38,
             size=14, bold=True, color=WHITE)
    before_items = [
        "Inline styling ដដែលៗរៀងរាល់ View",
        "Buttons ឆ្លើយតបភ្លាមៗ — គ្មាន tactile feedback",
        "Lists លោតចេញតែមួយដង — ខ្ទាស",
        "គ្មាន pull-to-refresh លើ screen ណាមួយ",
        "Empty states បង្ហាញ text ប្រផេះធម្មតា",
        "VoiceOver អានតែ elements ដើម",
        "White backgrounds hardcoded ខូចក្នុង dark mode",
        "គ្មាន loading indicator ពេលទាញទិន្នន័យ",
    ]
    add_bullets(s, before_items, 0.45, 1.55, 5.7, 5.3,
                size=12, color=GREY, bullet_color=RGBColor(0xE5, 0x47, 0x47))

    # After column
    add_rect(s, 6.8, 1.0, 6.2, 6.1, CARD)
    add_rect(s, 6.8, 1.0, 6.2, 0.45, GREEN)
    add_text(s, "✅  ក្រោយ (Week 9)", 6.9, 1.04, 5.9, 0.38,
             size=14, bold=True, color=DARK)
    after_items = [
        ".farmCard() modifier ប្រមូលរំដោះ card styling ទាំងអស់",
        ".scalePress() ឱ្យ animation ចុចដូចស្ព្រីង",
        ".fadeIn(delay:) ធ្វើ rows លេចចេញម្ដងៗ",
        ".refreshable {} លើ Dashboard, Finance, PestGuide",
        "EmptyStateView ជាមួយ icon + title + button",
        "accessibilityLabel លើ elements focus-able ទាំងអស់",
        "Color(.systemBackground) ស្ដាប់ light/dark ដោយស្វ័យ",
        "LoadingRowView skeleton ពេល refresh",
    ]
    add_bullets(s, after_items, 6.95, 1.55, 5.9, 5.3,
                size=12, color=WHITE, bullet_color=GREEN)

    slide_number(s, 11)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary & What's Next
# ════════════════════════════════════════════════════════════════════════════
def slide_12(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🎯  សង្ខេប & ជំហានបន្ទាប់", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=GREEN)

    add_text(s, "លទ្ធផល Week 9", 0.4, 1.0, 6.2, 0.38,
             size=16, bold=True, color=WHITE)
    deliverables = [
        "Components/ViewModifiers.swift — modifiers 5 + View extensions",
        "Components/FarmCard.swift — generic card container",
        "Components/PrimaryButton.swift — button variants ០៣",
        "Components/SectionHeader.swift — header + skeleton + empty state",
        "Dashboard: fade-in stagger, pull-to-refresh, accessibility labels",
        "PestGuide: fade ក្នុង row, skeleton loading, refreshable list",
        "Finance: summary cards staggered, pull-to-refresh",
    ]
    add_bullets(s, deliverables, 0.4, 1.45, 6.3, 4.2,
                size=12.5, bullet_color=GREEN)

    add_text(s, "ខាងមុខ — Week 10", 7.0, 1.0, 5.9, 0.38,
             size=16, bold=True, color=WHITE)
    next_items = [
        "CloudKit sync — ចែករំលែកទិន្នន័យចម្ការ",
        "WidgetKit — សង្ខេបចម្ការ Home Screen",
        "Checklist ដាក់ App Store",
        "TestFlight distribution ខាងក្នុង",
        "Performance profiling ជាមួយ Instruments",
    ]
    add_bullets(s, next_items, 7.0, 1.45, 5.9, 3.5,
                size=12.5, bullet_color=BLUE)

    # Key takeaway
    add_rect(s, 0.3, 5.85, 12.75, 1.25, GREEN)
    add_text(s, "💡  ចំណាំសំខាន់",
             0.55, 5.9, 12.3, 0.38, size=14, bold=True, color=DARK)
    add_text(s,
             "UX ស្អាតមិនមែនអំពីការបន្ថែម complexity — វាអំពីការប្រើ patterns ស៊ីសង្វាក់ "
             "(ViewModifiers, reusable components) ដែលធ្វើ interactions ទាំងអស់ ចេះចង់",
             0.55, 6.28, 12.3, 0.65, size=12.5, color=DARK)

    slide_number(s, 12)
    return s


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)

    out = "Week09_Advanced_UI_Animations.pptx"
    prs.save(out)
    print(f"✅  Saved → {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
