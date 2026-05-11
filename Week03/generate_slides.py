#!/usr/bin/env python3
"""Generate Week 03 — Navigation & Tab Coordination slide deck (Khmer)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

PURPLE = RGBColor(0x8E,0x44,0xAD)
GREEN  = RGBColor(0x1B,0xB8,0x89)
BLUE   = RGBColor(0x28,0x7D,0xFA)
ORANGE = RGBColor(0xF3,0x96,0x20)
TEAL   = RGBColor(0x00,0xC9,0xC8)
DARK   = RGBColor(0x1A,0x1A,0x2E)
DARK2  = RGBColor(0x16,0x21,0x3E)
CARD   = RGBColor(0x0F,0x2A,0x45)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GREY   = RGBColor(0xAA,0xAA,0xBB)
RED    = RGBColor(0xE5,0x47,0x47)
YELLOW = RGBColor(0xFF,0xD7,0x00)
W,H    = 13.33,7.5

def add_slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    for ph in s.placeholders: ph._element.getparent().remove(ph._element)
    f=s.background.fill; f.solid(); f.fore_color.rgb=DARK; return s

def add_text(s,text,x,y,w,h,size=16,bold=False,italic=False,
             color=WHITE,align=PP_ALIGN.LEFT,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; return tb

def add_rect(s,x,y,w,h,fc,lc=None):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if lc: sh.line.color.rgb=lc
    else: sh.line.fill.background(); return sh

def add_bullets(s,items,x,y,w,h,size=14,color=WHITE,bc=PURPLE):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r1=p.add_run(); r1.text="• "; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=bc
        r2=p.add_run(); r2.text=item; r2.font.size=Pt(size); r2.font.color.rgb=color

def add_code(s,lines,x,y,w,h,size=10):
    add_rect(s,x,y,w,h,CARD)
    tb=s.shapes.add_textbox(Inches(x+.18),Inches(y+.15),Inches(w-.36),Inches(h-.3))
    tf=tb.text_frame; tf.word_wrap=False
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.color.rgb=TEAL; r.font.name="Courier New"

def hdr(s,title,accent=BLUE):
    add_rect(s,0,0,W,.9,DARK2); add_rect(s,0,0,.06,.9,accent)
    add_text(s,title,.25,.1,12.5,.7,size=22,bold=True,color=accent)

def snum(s,n,total=12):
    add_text(s,f"{n}/{total}",12.3,7.1,.9,.3,size=10,color=GREY,align=PP_ALIGN.RIGHT)

# ── Slide 1 — Title ──────────────────────────────────────────────────────────
def s1(prs):
    s=add_slide(prs)
    add_rect(s,0,0,.12,H,BLUE)
    add_rect(s,.4,.4,2.2,.42,BLUE)
    add_text(s,"Week 03 · SmartFarmer",.4,.4,2.2,.42,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(s,"ការរុករក & ភ្ជាប់ Tab",.4,1.0,12.0,1.0,size=40,bold=True,color=BLUE)
    add_text(s,"Navigation & Tab Coordination",.4,1.95,12.0,.7,size=28,bold=True,color=WHITE)
    add_text(s,"iOS 13+  ·  NavigationView  ·  NavigationLink  ·  Coordinator Pattern",.4,2.85,12.0,.45,size=14,italic=True,color=GREY)
    add_rect(s,.4,3.45,10.5,.04,BLUE)
    icons=[("🗺️","NavigationView"),("🔗","NavigationLink"),("📋","List→Detail"),("🎯","Coordinator"),("🔀","Programmatic"),("🔔","Deep Link")]
    for i,(e,l) in enumerate(icons):
        cx=.4+i*2.1; add_rect(s,cx,3.65,1.9,1.0,CARD)
        add_text(s,e,cx,3.7,1.9,.45,size=20,align=PP_ALIGN.CENTER)
        add_text(s,l,cx,4.18,1.9,.4,size=10,color=GREY,align=PP_ALIGN.CENTER)
    add_text(s,"SmartFarmer Assistant · ភ្នំពេញ 2026",.4,6.9,12.0,.4,size=11,italic=True,color=GREY)
    snum(s,1)

# ── Slide 2 — Agenda ─────────────────────────────────────────────────────────
def s2(prs):
    s=add_slide(prs)
    hdr(s,"📋  មាតិកា Week 03 · Agenda")
    topics=[
        (BLUE,  "01","NavigationView Basics","container សម្រាប់ push navigation · iOS 13+"),
        (GREEN, "02","NavigationLink","List → Detail · tap-to-push pattern"),
        (PURPLE,"03","NavigationCoordinator","ObservableObject · centralize nav state"),
        (ORANGE,"04","tag/selection","programmatic navigation · @Binding"),
        (TEAL,  "05","Deep Link Simulation","jump to screen from outside the tab"),
        (BLUE,  "06","Mini-Project","Navigation ពេញ + deep link demo"),
    ]
    for i,(c,n,kh,en) in enumerate(topics):
        col=i%2; row=i//2; cx=.35+col*6.5; cy=1.1+row*2.0
        add_rect(s,cx,cy,6.1,1.75,CARD); add_rect(s,cx,cy,.08,1.75,c)
        add_rect(s,cx+.18,cy+.38,.6,.6,c)
        add_text(s,n,cx+.18,cy+.38,.6,.6,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        add_text(s,kh,cx+.95,cy+.25,4.9,.5,size=15,bold=True,color=WHITE)
        add_text(s,en,cx+.95,cy+.85,4.9,.45,size=11,color=GREY)
    snum(s,2)

# ── Slide 3 — NavigationView Basics ──────────────────────────────────────────
def s3(prs):
    s=add_slide(prs)
    hdr(s,"🗺️  NavigationView · Container សម្រាប់ Push Navigation")
    add_bullets(s,[
        "NavigationView = container ដែល wrap root view របស់ tab នីមួយៗ",
        "Tab នីមួយៗ ត្រូវការ NavigationView ផ្ទាល់ខ្លួន — មិនចែករំលែក",
        ".navigationTitle() — ដាក់ title ក្នុង navigation bar",
        "iOS 16+ ប្រើ NavigationStack — ប៉ុន្តែ iOS 13+ ប្រើ NavigationView ✅",
    ],.35,1.1,6.3,2.8,size=13,bc=BLUE)
    add_code(s,[
        "// ✅ iOS 13+ — NavigationView",
        "struct FinanceTabView: View {",
        "    var body: some View {",
        "        NavigationView {",
        "            List { ... }",
        '                .navigationTitle("ហិរញ្ញវត្ថុ")',
        "        }",
        "    }",
        "}",
        "",
        "// ❌ iOS 16+ only — NavigationStack",
        "struct FinanceTabView: View {",
        "    var body: some View {",
        "        NavigationStack { ... }",
        "    }",
        "}",
    ],.35,4.0,6.3,3.4,size=10.5)
    rows=[
        ("NavigationStack",      "❌ iOS 16+",  "✅ NavigationView"),
        ("NavigationPath",       "❌ iOS 16+",  "✅ @State var selectedID: UUID?"),
        (".navigationDestination","❌ iOS 16+", "✅ NavigationLink(destination:)"),
        ("@Environment(\\.dismiss)","❌ iOS 16+","✅ @Environment(\\.presentationMode)"),
    ]
    add_text(s,"⚠️  iOS 13+ Rules",7.0,1.05,6.1,.4,size=13,bold=True,color=ORANGE)
    hdrs=["Feature","❌ iOS 16+","✅ iOS 13+"]
    ws=[2.0,1.8,2.1]; xs=[7.05,9.15,11.05]
    for c,(h,cw,cx) in enumerate(zip(hdrs,ws,xs)):
        add_rect(s,cx,1.55,cw,.48,ORANGE if c==0 else RED if c==1 else GREEN)
        add_text(s,h,cx+.08,1.6,cw-.12,.38,size=11,bold=True,color=DARK if c>0 else WHITE)
    for r,(p,w,g) in enumerate(rows):
        cy=2.13+r*.8; bg=CARD if r%2==0 else DARK2
        for c,(cell,cw,cx) in enumerate(zip([p,w,g],ws,xs)):
            add_rect(s,cx,cy,cw,.72,bg)
            col=GREY if c==0 else (RED if c==1 else GREEN)
            add_text(s,cell,cx+.06,cy+.12,cw-.1,.48,size=9,color=col)
    snum(s,3)

# ── Slide 4 — NavigationLink ─────────────────────────────────────────────────
def s4(prs):
    s=add_slide(prs)
    hdr(s,"🔗  NavigationLink · Push a Detail View")
    add_bullets(s,[
        "NavigationLink = wrapper ជុំវិញ row ដើម្បី push detail view",
        "destination: — view ដែល push ទៅ",
        "label: (trailing closure) — អ្វីដែល user ឃើញ និង tap",
        "Back button ផ្ដល់ដោយ NavigationView ដោយស្វ័យប្រវត្តិ",
        "Style A (tap-only): NavigationLink(destination:)",
        "Style B (programmatic): NavigationLink(tag:selection:destination:)",
    ],.35,1.1,6.3,3.5,size=13,bc=BLUE)
    add_code(s,[
        "// Style A — tap only (simple)",
        "NavigationLink(",
        "    destination: TransactionDetailView(",
        "        transaction: transaction)",
        ") {",
        "    TransactionRowView(transaction: transaction)",
        "}",
        "",
        "// Style B — programmatic (deep link ready)",
        "NavigationLink(",
        "    destination: TransactionDetailView(",
        "        transaction: transaction),",
        "    tag: transaction.id ?? UUID(),",
        "    selection: $coordinator.selectedTransactionID",
        ") {",
        "    TransactionRowView(transaction: transaction)",
        "}",
    ],.35,4.75,6.3,2.65,size=10.5)
    notes=[
        (BLUE,  "Style A","Simple · tap-to-push\nប្រើ list ធម្មតា"),
        (ORANGE,"Style B","Programmatic · tag/selection\nប្រើ deep link / notification"),
        (GREEN, "destination:","View ដែល push\nទទួល data ពី parent"),
        (PURPLE,"label:","Row ដែល user tap\nrow view ឬ text"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,6.8,cy,6.25,1.38,CARD); add_rect(s,6.8,cy,.08,1.38,c)
        add_text(s,t,7.03,cy+.1,5.7,.4,size=12,bold=True,color=c)
        add_text(s,d,7.03,cy+.55,5.7,.72,size=11,color=WHITE)
    snum(s,4)

# ── Slide 5 — List → Detail Pattern ──────────────────────────────────────────
def s5(prs):
    s=add_slide(prs)
    hdr(s,"📋  List → Detail Pattern · រចនាបទ Navigation")
    add_bullets(s,[
        "List — បង្ហាញ records ទាំងអស់ · browse + delete",
        "Detail — បង្ហាញ record តែមួយ · មាន Edit button",
        "Edit sheet — sheet សម្រាប់ modify record",
        "TransactionDetailView ទទួល Transaction object ពី list",
        "Edit sheet បើក ពី detail view ដោយ .sheet(isPresented:)",
    ],.35,1.1,6.3,3.0,size=13,bc=BLUE)
    add_code(s,[
        "struct TransactionDetailView: View {",
        "    @EnvironmentObject private var viewModel: FarmViewModel",
        "    @Environment(\\.managedObjectContext) private var ctx",
        "    let transaction: Transaction",
        "    @State private var showingEdit = false",
        "",
        "    var body: some View {",
        "        Form {",
        '            Section(header: Text("ចំនួនទឹកប្រាក់")) {',
        "                Text(viewModel.formatCurrency(transaction.amount))",
        "                    .foregroundColor(",
        "                        transaction.isExpense ? .red : .green)",
        "            }",
        '            Section(header: Text("ប្រភេទ")) {',
        "                HStack {",
        '                    Text("ប្រភេទ").foregroundColor(.gray)',
        "                    Spacer()",
        "                    Text(transaction.isExpense ? \"ចំណាយ\" : \"ចំណូល\")",
        "                }",
        "            }",
        "        }",
        '        .navigationTitle("ព័ត៌មានប្រតិបត្តិការ")',
        '        .navigationBarItems(trailing: Button("កែប្រែ") {',
        "            showingEdit = true",
        "        })",
        "        .sheet(isPresented: $showingEdit) {",
        "            EditTransactionView(transaction: transaction)",
        "                .environment(\\.managedObjectContext, ctx)",
        "                .environmentObject(viewModel)",
        "        }",
        "    }",
        "}",
    ],.35,4.2,6.3,3.25,size=9)
    add_text(s,"Navigation Stack Diagram",7.0,1.05,6.1,.4,size=13,bold=True,color=BLUE)
    layers=[
        (BLUE,  "FinanceTabView (list)","root screen · ForEach transactions"),
        (GREEN, "  └── NavigationLink","tap row → push"),
        (ORANGE,"      └── TransactionDetailView","detail · Edit button"),
        (PURPLE,"            └── EditTransactionView","sheet · modify record"),
    ]
    for i,(c,title,desc) in enumerate(layers):
        cy=1.6+i*1.35
        add_rect(s,7.0,cy,6.1,1.2,CARD); add_rect(s,7.0,cy,.08,1.2,c)
        add_text(s,title,7.22,cy+.1,5.7,.45,size=12,bold=True,color=c)
        add_text(s,desc,7.22,cy+.62,5.7,.45,size=10,color=GREY)
    snum(s,5)

# ── Slide 6 — NavigationCoordinator ──────────────────────────────────────────
def s6(prs):
    s=add_slide(prs)
    hdr(s,"🎯  NavigationCoordinator · Single Source of Truth")
    add_code(s,[
        "// ❌ Without coordinator — nav state scattered",
        "struct FinanceTabView: View {",
        "    @State private var selectedTransaction: Transaction?",
        "    @State private var showingAdd = false",
        "    // Hard to trigger from outside",
        "}",
        "",
        "// ✅ With coordinator — centralized",
        "class FinanceCoordinator: ObservableObject {",
        "    // nil = list shown",
        "    // UUID = NavigationLink activates → detail pushed",
        "    @Published var selectedTransactionID: UUID? = nil",
        "",
        "    func navigate(to transaction: Transaction) {",
        "        selectedTransactionID = transaction.id",
        "    }",
        "",
        "    func reset() {",
        "        selectedTransactionID = nil",
        "    }",
        "}",
    ],.35,1.1,6.5,6.25,size=10)
    notes=[
        (BLUE,  "Problem","nav state ខ្ចាត់ខ្ចាយ\nhard to trigger from outside\n❌ notification/deep link"),
        (GREEN, "Solution","ObservableObject class\nselectedTransactionID: UUID?\n✅ one source of truth"),
        (ORANGE,"navigate(to:)","set selectedTransactionID\n→ NavigationLink activates\n→ detail pushes"),
        (PURPLE,"reset()","set selectedTransactionID = nil\n→ pop back to list\n→ any caller can use"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,7.0,cy,6.05,1.38,CARD); add_rect(s,7.0,cy,.08,1.38,c)
        add_text(s,t,7.22,cy+.1,5.6,.4,size=12,bold=True,color=c)
        add_text(s,d,7.22,cy+.55,5.6,.72,size=11,color=WHITE)
    snum(s,6)

# ── Slide 7 — Wiring Coordinator ─────────────────────────────────────────────
def s7(prs):
    s=add_slide(prs)
    hdr(s,"🔌  Wiring Coordinator · @StateObject vs @EnvironmentObject")
    add_code(s,[
        "// Step 1 — Create in MainTabView (root owner)",
        "struct MainTabView: View {",
        "    @StateObject private var viewModel: FarmViewModel",
        "    // ← @StateObject: create & own here",
        "    @StateObject private var financeCoordinator",
        "                         = FinanceCoordinator()",
        "    @State private var selectedTab = 0",
        "",
        "    var body: some View {",
        "        TabView(selection: $selectedTab) {",
        "            FinanceTabView()",
        '                .tabItem { Label("ហិរញ្ញវត្ថុ",',
        '                    systemImage: "dollarsign.circle") }',
        "                .tag(0)",
        "        }",
        "        .environmentObject(viewModel)",
        "        // ← inject coordinator to all children",
        "        .environmentObject(financeCoordinator)",
        "    }",
        "}",
        "",
        "// Step 2 — Read in child view",
        "struct FinanceTabView: View {",
        "    // ← @EnvironmentObject: read injected copy",
        "    @EnvironmentObject private var coordinator: FinanceCoordinator",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            FilteredTransactionList(",
        "                selectedTransactionID:",
        "                    $coordinator.selectedTransactionID",
        "            )",
        "        }",
        "    }",
        "}",
    ],.35,1.1,7.5,6.25,size=9)
    notes=[
        (BLUE,  "@StateObject","creates + owns the object\nOnly ONE place per coordinator\n→ MainTabView"),
        (GREEN, "@EnvironmentObject","reads injected copy\nAll child views\n→ FinanceTabView, etc."),
        (ORANGE,".environmentObject()","inject on TabView\nall children receive it\nmust be present or CRASH"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*2.1
        add_rect(s,8.0,cy,5.1,1.9,CARD); add_rect(s,8.0,cy,.08,1.9,c)
        add_text(s,t,8.22,cy+.12,4.6,.45,size=13,bold=True,color=c)
        add_text(s,d,8.22,cy+.65,4.6,.9,size=12,color=WHITE)
    snum(s,7)

# ── Slide 8 — tag/selection Programmatic Nav ─────────────────────────────────
def s8(prs):
    s=add_slide(prs)
    hdr(s,"🔀  tag/selection · Programmatic Navigation")
    add_code(s,[
        "// NavigationLink(tag:selection:) activates",
        "// when the @Binding matches the tag value",
        "ForEach(transactions, id: \\.self) { transaction in",
        "    NavigationLink(",
        "        destination: TransactionDetailView(",
        "            transaction: transaction),",
        "        tag: transaction.id ?? UUID(),",
        "        selection: $selectedTransactionID",
        "    ) {",
        "        TransactionRowView(transaction: transaction,",
        "                           viewModel: viewModel)",
        "    }",
        "}",
    ],.35,1.1,6.5,3.5,size=10.5)
    rows=[
        ("nil",         "list shown · nothing pushed"),
        ("matches tag", "NavigationLink activates → detail pushed"),
        ("set by code", "same effect → enables deep linking"),
    ]
    add_text(s,"selectedTransactionID value → Result",0.35,4.75,6.5,.4,size=12,bold=True,color=BLUE)
    hdrs=["selectedTransactionID","Result"]
    ws=[2.5,3.8]; xs=[.35,2.95]
    for c,(h,cw,cx) in enumerate(zip(hdrs,ws,xs)):
        add_rect(s,cx,5.2,cw,.48,BLUE)
        add_text(s,h,cx+.08,5.25,cw-.12,.38,size=11,bold=True,color=WHITE)
    for r,(v,res) in enumerate(rows):
        cy=5.78+r*.62; bg=CARD if r%2==0 else DARK2
        add_rect(s,.35,cy,2.5,.58,bg); add_rect(s,2.95,cy,3.8,.58,bg)
        add_text(s,v,.45,cy+.1,2.3,.38,size=10,color=TEAL)
        add_text(s,res,3.05,cy+.1,3.6,.38,size=10,color=GREY)
    notes=[
        (BLUE,  "tag:","unique id for each link\n= transaction.id\n→ UUID"),
        (GREEN, "selection:","@Binding<UUID?>\nfrom coordinator\n→ shared state"),
        (ORANGE,"activates when","selection == tag\nworks from code too\n→ deep link ✅"),
        (PURPLE,"nil = reset","set nil → pop back\ncoordinator.reset()\n→ back to list"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,7.0,cy,6.05,1.38,CARD); add_rect(s,7.0,cy,.08,1.38,c)
        add_text(s,t,7.22,cy+.1,5.6,.4,size=12,bold=True,color=c)
        add_text(s,d,7.22,cy+.55,5.6,.72,size=11,color=WHITE)
    snum(s,8)

# ── Slide 9 — Deep Linking ───────────────────────────────────────────────────
def s9(prs):
    s=add_slide(prs)
    hdr(s,"🔔  Deep Link Simulation · Jump from Outside")
    add_code(s,[
        "// In MainTabView — deep link helper",
        "func deepLink(to transactionID: UUID) {",
        "    // 1. Switch to Finance tab",
        "    selectedTab = 0",
        "    // 2. Activate the NavigationLink",
        "    financeCoordinator.selectedTransactionID = transactionID",
        "}",
        "",
        "// Call from anywhere:",
        "Button(\"Deep Link Demo\") {",
        "    if let first = transactions.first {",
        "        deepLink(to: first.id ?? UUID())",
        "    }",
        "}",
    ],.35,1.1,6.5,4.0,size=10.5)
    add_text(s,"Data Flow Diagram",0.35,5.2,6.5,.4,size=12,bold=True,color=BLUE)
    flow=[
        "MainTabView  (@StateObject financeCoordinator)",
        "  ↓  .environmentObject(financeCoordinator)",
        "FinanceTabView  (@EnvironmentObject coordinator)",
        "  ↓  $coordinator.selectedTransactionID",
        "NavigationLink(tag: uuid, selection: $selectedTransactionID)",
        "  ↓  activates when selectedTransactionID == uuid",
        "TransactionDetailView",
    ]
    add_code(s,flow,.35,5.65,6.5,1.8,size=9.5)
    notes=[
        (BLUE,  "Step 1","selectedTab = 0\n→ switch to Finance tab\n→ tab becomes visible"),
        (GREEN, "Step 2","coordinator.selectedTransactionID = id\n→ NavigationLink matches tag\n→ detail view pushes"),
        (ORANGE,"Why it works","coordinator = @StateObject\nchanges flow to @EnvironmentObject\n→ child updates automatically"),
        (PURPLE,"Real deep links","URL → AppDelegate/SceneDelegate\n→ call deepLink(to:)\n→ same mechanism"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,7.0,cy,6.05,1.38,CARD); add_rect(s,7.0,cy,.08,1.38,c)
        add_text(s,t,7.22,cy+.1,5.6,.4,size=12,bold=True,color=c)
        add_text(s,d,7.22,cy+.55,5.6,.72,size=11,color=WHITE)
    snum(s,9)

# ── Slide 10 — Common Mistakes ────────────────────────────────────────────────
def s10(prs):
    s=add_slide(prs)
    hdr(s,"⚠️  Common Mistakes · ❌ vs ✅")
    rows=[
        ("NavigationView","NavigationStack (iOS 16+)","NavigationView (iOS 13+)"),
        ("Multiple NavViews","One shared NavigationView for all tabs","Each tab wraps own NavigationView"),
        ("Sheet wrap","NavigationView { EditView() } in sheet","EditView() directly — no extra wrap"),
        ("@StateObject child","@StateObject coordinator in FinanceTabView","@EnvironmentObject coordinator"),
        ("Missing inject","Forgot .environmentObject(coordinator)","Add .environmentObject on TabView"),
        ("Style A always","NavigationLink(destination:) for deep link","NavigationLink(tag:selection:) instead"),
    ]
    hdrs=["Pattern","❌  Wrong","✅  Correct (iOS 13+)"]
    widths=[2.8,4.5,5.1]; starts=[.35,3.3,7.95]
    for c,(h,cw,cs) in enumerate(zip(hdrs,widths,starts)):
        add_rect(s,cs,1.1,cw,.55,BLUE if c==2 else RED if c==1 else DARK2)
        add_text(s,h,cs+.1,1.15,cw-.15,.45,size=12,bold=True,color=WHITE)
    for r,(p,w,g) in enumerate(rows):
        cy=1.75+r*.67; bg=CARD if r%2==0 else DARK2
        for c,(cell,cw,cs) in enumerate(zip([p,w,g],widths,starts)):
            add_rect(s,cs,cy,cw,.62,bg)
            col=GREY if c==0 else (RED if c==1 else GREEN)
            add_text(s,cell,cs+.1,cy+.1,cw-.15,.42,size=10,color=col)
    snum(s,10)

# ── Slide 11 — Mini-Project ───────────────────────────────────────────────────
def s11(prs):
    s=add_slide(prs)
    hdr(s,"🏠  Mini-Project · Navigation ពេញ + Deep Link Demo")
    cols=[
        ("📋  Requirements",BLUE,[
            "List → Detail navigation ✓",
            "TransactionDetailView (Form)",
            "Edit button → EditTransactionView sheet",
            "FinanceCoordinator class",
            "tag/selection NavigationLink",
            "Deep Link Demo button",
        ]),
        ("✅  Checklist",GREEN,[
            "NavigationView wraps FinanceTabView ✓",
            "NavigationLink(tag:selection:) ✓",
            "FinanceCoordinator @StateObject ✓",
            ".environmentObject() on TabView ✓",
            "Detail shows: amount/type/category ✓",
            "Edit + save updates detail view ✓",
        ]),
        ("🎯  Grading",ORANGE,[
            "List → Detail navigation: 25%",
            "Coordinator pattern: 25%",
            "tag/selection programmatic: 20%",
            "Detail view full fields: 15%",
            "Edit from detail works: 15%",
            "Bonus: Deep Link Demo button",
        ]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*4.35
        add_rect(s,cx,1.1,4.1,5.8,CARD); add_rect(s,cx,1.1,4.1,.55,c)
        add_text(s,t,cx+.12,1.15,3.85,.45,size=13,bold=True,color=WHITE)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.15,cy,3.7,.72,DARK2)
            add_text(s,item,cx+.28,cy+.1,3.4,.52,size=11,color=WHITE)
    add_rect(s,.35,7.0,12.6,.42,BLUE)
    add_text(s,"🗺️  Week 03 ចប់! — Navigation coordinator ready · Next: Finance Module Complete →",
             .5,7.02,12.3,.38,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    snum(s,11)

# ── Slide 12 — Cheat Sheet ───────────────────────────────────────────────────
def s12(prs):
    s=add_slide(prs)
    hdr(s,"📌  សង្ខេប Week 03 · Cheat Sheet",accent=BLUE)
    cols=[
        ("🗺️  Navigation",BLUE,["NavigationView { }","NavigationLink(dest:)","NavigationLink(tag:sel:)","navigationTitle()"]),
        ("🎯  Coordinator",GREEN,["class FC: ObservableObject","@Published selectedID: UUID?","func navigate(to:)","func reset()"]),
        ("🔌  Wiring",ORANGE,["@StateObject in MainTabView","@EnvironmentObject in child",".environmentObject(fc)","$coordinator.selectedID"]),
        ("🔔  Deep Link",PURPLE,["selectedTab = 0","coordinator.selectedID = id","tag matches → push","nil → pop to list"]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*3.25
        add_rect(s,cx,1.1,3.0,5.8,CARD); add_rect(s,cx,1.1,3.0,.55,c)
        add_text(s,t,cx+.1,1.15,2.8,.45,size=12,bold=True,color=WHITE)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.1,cy,2.8,.72,DARK2)
            add_text(s,item,cx+.18,cy+.1,2.6,.52,size=10,color=TEAL)
            s.shapes[-1].text_frame.paragraphs[0].runs[0].font.name="Courier New"
    add_rect(s,.35,7.0,12.6,.42,BLUE)
    add_text(s,"🗺️  Coordinator = single source of truth · tag/selection = deep link ready!",
             .5,7.02,12.3,.38,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    snum(s,12)

def build():
    prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H)
    s1(prs);s2(prs);s3(prs);s4(prs);s5(prs);s6(prs)
    s7(prs);s8(prs);s9(prs);s10(prs);s11(prs);s12(prs)
    out=os.path.join(os.path.dirname(__file__),"Week03_Navigation_KH.pptx")
    prs.save(out); print(f"✅  Saved → {out}")

if __name__=="__main__": build()
