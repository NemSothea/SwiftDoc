#!/usr/bin/env python3
"""Generate Week 12 Final Polish & TestFlight Distribution slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour palette ──────────────────────────────────────────────────────────
GREEN      = RGBColor(0x1B, 0xB8, 0x89)
RED        = RGBColor(0xE5, 0x47, 0x47)
BLUE       = RGBColor(0x28, 0x7D, 0xFA)
PURPLE     = RGBColor(0x8E, 0x44, 0xAD)
ORANGE     = RGBColor(0xF3, 0x96, 0x20)
TEAL       = RGBColor(0x00, 0xC9, 0xC8)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
DARK2      = RGBColor(0x16, 0x21, 0x3E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT_BLUE = RGBColor(0xA8, 0xD8, 0xFF)

W = 13.33
H = 7.5


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide  = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, items, x, y, w, h, size=16, color=WHITE,
                bullet_color=GREEN, indent=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("    • " if indent else "• ") + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color


def add_code(slide, code_lines, x, y, w, h, size=11):
    add_rect(slide, x, y, w, h, RGBColor(0x0D, 0x1B, 0x2A))
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.1),
        Inches(w - 0.3),  Inches(h - 0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = LIGHT_BLUE
        run.font.name = "Courier New"


# ════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)


# ── Slide 1: Title ──────────────────────────────────────────────────────────
slide = add_slide(prs)

add_rect(slide, 0, 0, 0.08, H, TEAL)

add_text(slide, "SmartFarmer Assistant · iOS SwiftUI Course",
         0.3, 0.3, 10, 0.4, size=13, color=GREY)

add_text(slide, "Week 12", 0.3, 0.9, 12, 1.0, size=52, bold=True, color=TEAL)
add_text(slide, "ការតុបតែងចុងក្រោយ & TestFlight",
         0.3, 1.75, 12, 1.0, size=32, bold=True, color=WHITE)

add_text(slide,
         "រៀបចំសម្រាប់ការសាកល្បងកសិករពិត —\n"
         "App icon, launch screen, TestFlight & ការណែនាំ",
         0.3, 2.85, 10, 1.2, size=20, color=GREY)

pill_data = [
    (0.3,  4.3, TEAL,   "🎨  ការតុបតែង App"),
    (2.8,  4.3, BLUE,   "🚀  TestFlight"),
    (5.3,  4.3, GREEN,  "📖  ការណែនាំ"),
    (7.8,  4.3, ORANGE, "📣  មតិ"),
    (10.3, 4.3, PURPLE, "🏆  Demo ចុងក្រោយ"),
]
for px, py, pc, pt in pill_data:
    add_rect(slide, px, py, 2.2, 0.5, pc)
    add_text(slide, pt, px + 0.1, py + 0.05, 2.1, 0.4, size=13, bold=True)

add_text(slide, "Week 12 ក្នុង 12", W - 2.5, H - 0.5, 2.3, 0.4,
         size=12, color=GREY, align=PP_ALIGN.RIGHT)


# ── Slide 2: Learning Objectives ────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.1, DARK2)
add_text(slide, "🎯  គោលបំណងរៀន", 0.4, 0.2, 12, 0.7,
         size=28, bold=True, color=TEAL)

objectives = [
    "Design និង export icon set app ប្រធានបទចម្ការ (ទំហំទាំងអស់) ដោយ SF Symbols / Figma",
    "Configure LaunchScreen.storyboard ជាមួយ splash screen ម៉ាក",
    "កំណត់ bundle ID, display name, version/build numbers ត្រឹមត្រូវ",
    "Archive ហើយ upload build ទៅ App Store Connect ដោយ Xcode Organizer",
    "បង្កើត TestFlight beta group ហើយអញ្ជើញកសិករពិតជា testers",
    "បង្កើតផ្ទាំងណែនាំ in-app ដែលពន្យល់ features ដល់ users ថ្មី",
    "Design ផែនការ feedback ដែលស្ថានភាពប្រមូលទិន្នន័យ ពីកសិករ testers",
    "សរសេរ README ចុងក្រោយ ឯកសារ architecture, features, ការ run project",
]
add_bullets(slide, objectives, 0.5, 1.3, 12.5, 5.8, size=17)


# ── Slide 3: App Icon & Launch Screen ───────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🎨  App Icon & Launch Screen",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Left: App Icon sizes table
add_text(slide, "ទំហំ Icon ដែលត្រូវការ", 0.3, 1.1, 5.8, 0.4,
         size=15, bold=True, color=TEAL)

icon_sizes = [
    ("iPhone Notification",  "20 pt  → 40×40 / 60×60 px"),
    ("iPhone Settings",      "29 pt  → 58×58 / 87×87 px"),
    ("iPhone Spotlight",     "40 pt  → 80×80 / 120×120 px"),
    ("iPhone App",           "60 pt  → 120×120 / 180×180 px"),
    ("App Store",            "1024 pt → 1024×1024 px"),
]
for i, (label, sizes) in enumerate(icon_sizes):
    y = 1.6 + i * 0.9
    add_rect(slide, 0.3, y, 6.0, 0.75, DARK2)
    add_text(slide, label, 0.5,  y + 0.1, 3.0, 0.35, size=13, bold=True, color=TEAL)
    add_text(slide, sizes, 0.5,  y + 0.4, 5.8, 0.3,  size=11, color=GREY)

# Right: Launch Screen
add_text(slide, "LaunchScreen.storyboard", 6.7, 1.1, 6.3, 0.4,
         size=15, bold=True, color=ORANGE)

# Mock phone frame
add_rect(slide, 7.5, 1.55, 3.5, 5.5, RGBColor(0x0D, 0x1B, 0x2A))
add_rect(slide, 7.55, 1.6, 3.4, 5.4, RGBColor(0x12, 0x5C, 0x3F))
add_text(slide, "🌾", 8.8, 2.8, 1.2, 1.2, size=36, align=PP_ALIGN.CENTER)
add_text(slide, "SmartFarmer", 7.65, 4.1, 3.2, 0.5,
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "ជំនួយស្មាតសម្រាប់\nកសិករ", 7.65, 4.65, 3.2, 0.7,
         size=11, color=GREY, align=PP_ALIGN.CENTER)

add_bullets(slide, [
    "ប្រើ UIImageView សម្រាប់ logo នៅចំកណ្ដាល",
    "កំណត់ background color ឱ្យស្របតាម theme",
    "បន្ថែម CFBundleDisplayName ក្នុង Info.plist ចំពោះឈ្មោះខ្មែរ",
    "គ្មាន animations — LaunchScreen ត្រូវ static",
], 6.7, 2.3, 6.2, 3.2, size=13, color=WHITE)


# ── Slide 4: App Name, Bundle ID & Version ───────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "⚙️  Identity App — Bundle ID, ឈ្មោះ & Version",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Config table
config_items = [
    ("Display Name",    "SmartFarmer Assistant",         "បង្ហាញពីក្រោម app icon ផ្ទាំងដើម"),
    ("Bundle ID",       "com.yourname.smartfarmer",      "តែមួយ — ត្រូវដូចគ្នា App Store Connect"),
    ("Version",         "1.0.0",                         "Version marketing បង្ហាញដល់ users"),
    ("Build",           "1",                             "Build number ខាងក្នុង — បង្កើនរៀងរាល់ upload"),
    ("Min Deployment",  "iOS 13.0",                      "គ្របដណ្ដប់ 99%+ iPhones ដំណើរការ"),
    ("Device Family",   "iPhone only",                   "កំណត់ iPad បើចង់ universal"),
    ("Orientations",    "Portrait only",                 "Lock ដើម្បីប្រើ — ការពារ rotation ចៃដន្យ"),
]
for i, (key, val, tip) in enumerate(config_items):
    y = 1.15 + i * 0.87
    add_rect(slide, 0.3, y, 2.8, 0.72, RGBColor(0x0D, 0x1B, 0x2A))
    add_rect(slide, 3.2, y, 4.0, 0.72, DARK2)
    add_rect(slide, 7.3, y, 5.7, 0.72, RGBColor(0x10, 0x1A, 0x30))
    add_text(slide, key, 0.45, y + 0.18, 2.6, 0.4, size=13, bold=True, color=TEAL)
    add_text(slide, val, 3.35, y + 0.18, 3.8, 0.4, size=13, color=WHITE)
    add_text(slide, tip, 7.45, y + 0.18, 5.5, 0.4, size=12, color=GREY)

# Column headers
add_text(slide, "ការកំណត់", 0.45, 0.88, 2.6, 0.3, size=12, bold=True, color=GREY)
add_text(slide, "តម្លៃ",    3.35, 0.88, 3.8, 0.3, size=12, bold=True, color=GREY)
add_text(slide, "ដំបូន្មាន", 7.45, 0.88, 5.5, 0.3, size=12, bold=True, color=GREY)


# ── Slide 5: TestFlight Setup ────────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🚀  ការរៀបចំ TestFlight — ជំហានម្ដងម្ដង",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

steps = [
    (TEAL,   "ជំហាន 1",
     "បង្កើត App Record ក្នុង App Store Connect",
     "appstoreconnect.apple.com → My Apps → + → New App\n"
     "បំពេញ: Platform (iOS), Name, Bundle ID, Primary Language, SKU"),
    (BLUE,   "ជំហាន 2",
     "Archive ក្នុង Xcode",
     "ជ្រើស 'Any iOS Device (arm64)'\n"
     "Product → Archive → Xcode Organizer បើកស្វ័យ"),
    (GREEN,  "ជំហាន 3",
     "Upload ទៅ App Store Connect",
     "ក្នុង Organizer: Distribute App → Upload\n"
     "បើក 'Upload your app's symbols'"),
    (ORANGE, "ជំហាន 4",
     "បន្ថែម Beta Testers",
     "App Store Connect → TestFlight → Group\n"
     "អញ្ជើញ email — កសិករទទួល link ដំឡើង"),
    (PURPLE, "ជំហាន 5",
     "ត្រួតពិនិត្យ & ប្រមូល Feedback",
     "TestFlight បង្ហាញ crash logs, screenshots feedback\n"
     "Testers long-press → Send Feedback ជាមួយ screenshot"),
]
for i, (color, label, title, desc) in enumerate(steps):
    y = 1.15 + i * 1.22
    add_rect(slide, 0.3, y, 1.0, 0.95, color)
    add_text(slide, label, 0.32, y + 0.27, 0.96, 0.4,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.4, y, 11.6, 0.95, DARK2)
    add_text(slide, title, 1.55, y + 0.05, 11.3, 0.38, size=14, bold=True, color=color)
    add_text(slide, desc,  1.55, y + 0.48, 11.3, 0.42, size=12, color=GREY)


# ── Slide 6: In-App User Guide ───────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "📖  ផ្ទាំងណែនាំ In-App",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Left: code
code = [
    "struct UserGuideView: View {",
    "    let sections: [GuideSection] = [",
    '        .init(icon: "dollarsign.circle.fill",',
    '              color: .green,',
    '              title: "ការគ្រប់គ្រងហិរញ្ញវត្ថុ",',
    '              tips: ["ចុច + ដើម្បីបន្ថែមប្រតិបត្តិការ",',
    '                     "ជ្រើសប្រភេទ: ចំណូល ឬ ចំណាយ",',
    '                     "ប្រើតម្រងដើម្បីមើលតាមប្រភេទ"]),',
    '        .init(icon: "calendar.badge.plus",',
    '              color: .blue,',
    '              title: "ប្រតិទិន & ការជូនដំណឹង",',
    '              tips: ["កំណត់ការរំលឹកចម្ការ",',
    '                     "ការជូនដំណឹងមកមុន ១ ថ្ងៃ"]),',
    "        // ... PestGuide, Journal sections",
    "    ]",
    "    var body: some View {",
    "        NavigationView {",
    "            List(sections) { section in",
    "                GuideSectionRow(section: section)",
    "            }",
    '            .navigationTitle("ការណែនាំប្រើប្រាស់")',
    "        }",
    "    }",
    "}",
]
add_code(slide, code, 0.3, 1.1, 7.8, 5.9, size=11)

# Right: design tips
add_text(slide, "ការអនុវត្ត UX ល្អ", 8.4, 1.1, 4.6, 0.4,
         size=16, bold=True, color=GREEN)

tips = [
    ("រក្សាភ្លឺ", "tips 3-4 per feature\n— កសិករ មិនអានអត្ថបទវែង"),
    ("ប្រើ SF Symbols", "ផ្ដល់ icons ដូចក្នុង app\nដើម្បីកសិករស្គាល់"),
    ("ខ្មែរជាមុន", "សរសេរ ការណែនាំ ជាខ្មែរ;\nបន្ថែម English ក្នុងវង់ក្រចក"),
    ("Screenshot cards", "បន្ថែម screenshot ពិត\nចំពោះ feature នីមួយៗ"),
    ("ចូលតាម ⚙️", "បន្ថែម ការណែនាំ ក្នុង Settings tab\nឱ្យ​ស្វែងរកបានជានិច្ច"),
]
for i, (title, desc) in enumerate(tips):
    y = 1.65 + i * 1.12
    add_rect(slide, 8.3, y, 4.7, 0.95, DARK2)
    add_text(slide, title, 8.45, y + 0.05, 4.5, 0.35, size=13, bold=True, color=GREEN)
    add_text(slide, desc,  8.45, y + 0.44, 4.5, 0.48, size=12, color=GREY)


# ── Slide 7: Feedback Collection Plan ───────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "📣  ផែនការប្រមូល Feedback កសិករ",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Left column: questions to ask
add_text(slide, "សំណួរសំខាន់ ដែលត្រូវសួរកសិករ", 0.3, 1.1, 6.2, 0.4,
         size=15, bold=True, color=ORANGE)
questions = [
    ("💰", "Finance",  "Did you record all income & expenses this week?\nWere any categories missing?"),
    ("📅", "Calendar", "Did you receive reminders before farm tasks?\nWere times correct?"),
    ("🌿", "Pest Guide","Did you find your crop's pest in the guide?\nWas the treatment clear enough?"),
    ("📓", "Journal",  "Did you add daily notes & photos?\nWas the weather picker useful?"),
    ("🏠", "Dashboard","Could you see today's summary at a glance?\nWas anything confusing?"),
]
for i, (icon, label, q) in enumerate(questions):
    y = 1.65 + i * 1.1
    add_rect(slide, 0.3, y, 6.0, 0.95, DARK2)
    add_rect(slide, 0.3, y, 0.06, 0.95, ORANGE)
    add_text(slide, f"{icon} {label}", 0.5, y + 0.05, 5.8, 0.35, size=13, bold=True, color=ORANGE)
    add_text(slide, q,                 0.5, y + 0.45, 5.8, 0.45, size=11, color=WHITE)

# Right column: methods & rating scale
add_text(slide, "វិធីប្រមូល", 6.7, 1.1, 6.2, 0.4,
         size=15, bold=True, color=TEAL)

methods = [
    ("TestFlight built-in",
     "កសិករ long-press → screenshot + comment\n→ ផ្ញើដោយផ្ទាល់ App Store Connect"),
    ("Form Feedback In-App",
     "Form សាមញ្ញ: rating 1-5 ⭐ ក្នុង module\n+ text field ប្រឹក្សា"),
    ("WhatsApp / Telegram Group",
     "បង្កើត group beta testers — កសិករ\nចែករំលែក voice messages (ងាយជាង type)"),
    ("ការសង្កេត",
     "មើល 3-5 កសិករប្រើ app 30 នាទី\nដោយមិនជួយ — កត់ចំណុចវង្វេង"),
]
for i, (method, desc) in enumerate(methods):
    y = 1.65 + i * 1.35
    add_rect(slide, 6.7, y, 6.2, 1.15, DARK2)
    add_rect(slide, 6.7, y, 0.06, 1.15, TEAL)
    add_text(slide, method, 6.9, y + 0.08, 5.9, 0.38, size=13, bold=True, color=TEAL)
    add_text(slide, desc,   6.9, y + 0.5,  5.9, 0.6,  size=12, color=WHITE)


# ── Slide 8: In-App Feedback Form (Code) ────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "⚙️  Form Feedback In-App — ឧទាហរណ៍",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

code2 = [
    "struct FeedbackView: View {",
    "    @State private var rating: Int = 0",
    '    @State private var comment: String = ""',
    "    @State private var submitted = false",
    "",
    "    var body: some View {",
    "        Form {",
    '            Section("វាយតម្លៃកម្មវិធី") {',
    "                HStack {",
    "                    ForEach(1...5, id: \\.self) { star in",
    "                        Image(systemName: star <= rating",
    '                            ? "star.fill" : "star")',
    "                            .foregroundColor(.yellow)",
    "                            .font(.title2)",
    "                            .onTapGesture { rating = star }",
    "                    }",
    "                }",
    "            }",
    '            Section("យោបល់") {',
    '                TextEditor(text: $comment)',
    "                    .frame(minHeight: 100)",
    "            }",
    '            Button("ផ្ញើ") { submitted = true }',
    "                .disabled(rating == 0)",
    "        }",
    '        .navigationTitle("ផ្ញើមតិ")',
    "    }",
    "}",
]
add_code(slide, code2, 0.3, 1.1, 7.8, 5.9, size=11)

# Right: mock UI
add_text(slide, "Preview", 8.4, 1.1, 4.5, 0.4, size=16, bold=True, color=GREY)  # keep as-is
add_rect(slide, 8.3, 1.6, 4.7, 5.4, RGBColor(0x0D, 0x1B, 0x2A))
add_text(slide, "វាយតម្លៃកម្មវិធី",   8.5, 1.75, 4.3, 0.4, size=13, bold=True, color=TEAL)
add_text(slide, "★ ★ ★ ★ ☆",         8.5, 2.25, 4.3, 0.5, size=22, color=RGBColor(0xFF, 0xD7, 0x00))
add_text(slide, "យោបល់",             8.5, 2.9,  4.3, 0.4, size=13, bold=True, color=TEAL)
add_rect(slide, 8.4, 3.35, 4.5, 1.2, RGBColor(0x16, 0x21, 0x3E))
add_text(slide, "ផ្នែកប្រតិទិនល្អណាស់!\nចង់បានប្រភេទពាណិជ្ជកម្ម\nបន្ថែម...",
         8.5, 3.4, 4.3, 1.1, size=11, color=GREY)
add_rect(slide, 8.8, 4.75, 3.1, 0.55, TEAL)
add_text(slide, "ផ្ញើ", 8.85, 4.83, 3.0, 0.4,
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "⭐ Ask farmers to rate EACH module\n    separately for granular feedback.",
         8.3, 5.5, 4.7, 0.8, size=12, color=ORANGE, italic=True)


# ── Slide 9: Archive & Upload Live Coding ───────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🔧  Live Coding — Archive & Upload to TestFlight",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Checklist
add_text(slide, "Pre-Archive Checklist", 0.3, 1.1, 6.2, 0.4,
         size=15, bold=True, color=TEAL)
checklist = [
    ("✅", GREEN,  "Scheme set to Release (not Debug)"),
    ("✅", GREEN,  "Version: 1.0.0 · Build: 1"),
    ("✅", GREEN,  "Signing: Automatic, your Team selected"),
    ("✅", GREEN,  "All DEBUG code removed (print statements)"),
    ("✅", GREEN,  "App icon set — no missing sizes"),
    ("✅", GREEN,  "LaunchScreen configured"),
    ("✅", GREEN,  "Destination: Any iOS Device (arm64)"),
]
for i, (icon, color, text) in enumerate(checklist):
    y = 1.65 + i * 0.77
    add_rect(slide, 0.3, y, 6.0, 0.62, DARK2)
    add_text(slide, icon, 0.45, y + 0.12, 0.5, 0.38, size=14, color=color)
    add_text(slide, text, 1.0,  y + 0.12, 5.2, 0.38, size=13, color=WHITE)

# Steps
add_text(slide, "Upload Steps", 6.7, 1.1, 6.3, 0.4,
         size=15, bold=True, color=ORANGE)
upload_steps = [
    "Product → Archive  (wait 2-5 min for build)",
    "Organizer window opens automatically",
    "Select archive → Distribute App",
    "Choose: App Store Connect → Next",
    "Upload (not Export) → keep all checkboxes on",
    "Wait for processing email (5-30 min)",
    "App Store Connect → TestFlight → Select build",
    "Add Internal Testers → Send Invitations",
]
for i, step in enumerate(upload_steps):
    y = 1.65 + i * 0.75
    add_rect(slide, 6.7, y, 0.5, 0.58, ORANGE)
    add_text(slide, str(i+1), 6.72, y + 0.1, 0.46, 0.38,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 7.3, y, 5.7, 0.58, DARK2)
    add_text(slide, step, 7.45, y + 0.12, 5.5, 0.38, size=12, color=WHITE)


# ── Slide 10: Common Mistakes ────────────────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🐛  Common Mistakes & How to Fix Them",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

mistakes = [
    ("Wrong destination selected",
     "Archive greys out or produces simulator build",
     "Always select 'Any iOS Device (arm64)' before archiving"),
    ("Bundle ID mismatch",
     "Upload fails: 'No app found with bundle ID'",
     "Bundle ID in Xcode must exactly match App Store Connect record"),
    ("Missing icon sizes",
     "Upload fails: 'Invalid App Store Icon'",
     "Use an icon generator — fill ALL required slots in Assets.xcassets"),
    ("Build number not incremented",
     "Upload fails: 'Build already exists'",
     "Increment CURRENT_PROJECT_VERSION before every archive"),
    ("LaunchScreen uses code/animations",
     "App rejected or slow launch",
     "LaunchScreen.storyboard must be static — no Swift code, no animation"),
    ("Testers can't find TestFlight invite",
     "Farmers say they received no email",
     "Check spam folder; use Internal Testing for first beta (no review wait)"),
]
for i, (mistake, problem, fix) in enumerate(mistakes):
    y = 1.15 + i * 1.06
    add_rect(slide, 0.3, y, 0.06, 0.88, RED)
    add_text(slide, f"❌ {mistake}", 0.5, y,       8.3, 0.38, size=13, bold=True, color=RED)
    add_text(slide, f"→ {problem}", 0.5, y + 0.36, 8.3, 0.35, size=11, color=GREY)
    add_rect(slide, 9.0, y, 0.06, 0.88, GREEN)
    add_text(slide, f"✅ {fix}",    9.15, y + 0.18, 3.9, 0.55, size=11, color=GREEN)


# ── Slide 11: Final Project Checklist ───────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.0, DARK2)
add_text(slide, "🏆  Final Project Checklist",
         0.4, 0.15, 12, 0.7, size=26, bold=True, color=WHITE)

# Left column
add_text(slide, "Polish & Submission", 0.3, 1.1, 6.0, 0.4,
         size=15, bold=True, color=TEAL)
polish = [
    ("App icon (all sizes)",           TEAL),
    ("LaunchScreen with farm theme",   TEAL),
    ("Bundle ID & version configured", TEAL),
    ("Release scheme build succeeds",  TEAL),
    ("App uploaded to TestFlight",     TEAL),
    ("3+ real farmers invited",        TEAL),
    ("In-app User Guide screen done",  TEAL),
    ("Feedback form implemented",      TEAL),
]
for i, (item, color) in enumerate(polish):
    y = 1.65 + i * 0.72
    add_rect(slide, 0.3, y, 6.0, 0.58, DARK2)
    add_text(slide, "☐ " + item, 0.5, y + 0.12, 5.8, 0.35, size=13, color=WHITE)

# Right column
add_text(slide, "Final Presentation & Docs", 6.7, 1.1, 6.3, 0.4,
         size=15, bold=True, color=ORANGE)
docs = [
    ("README.md with architecture overview", ORANGE),
    ("Feature list with screenshots",        ORANGE),
    ("How to run / build instructions",      ORANGE),
    ("CoreData schema diagram",              ORANGE),
    ("Known limitations section",            ORANGE),
    ("Live demo prepared (10 min)",          ORANGE),
    ("Farmer feedback summary written",      ORANGE),
    ("Next steps / roadmap documented",      ORANGE),
]
for i, (item, color) in enumerate(docs):
    y = 1.65 + i * 0.72
    add_rect(slide, 6.7, y, 6.3, 0.58, DARK2)
    add_text(slide, "☐ " + item, 6.9, y + 0.12, 6.1, 0.35, size=13, color=WHITE)

add_text(slide,
         "🎯 Goal: Farmers can install and use the app independently without any help from the developer.",
         0.3, 7.0, 12.7, 0.42, size=14, color=TEAL, italic=True)


# ── Slide 12: Course Completion & Summary ────────────────────────────────────
slide = add_slide(prs)
add_rect(slide, 0, 0, W, 1.1, DARK2)
add_text(slide, "🎓  Course Complete — What You Built",
         0.4, 0.15, 12, 0.8, size=28, bold=True, color=TEAL)

modules = [
    (GREEN,  "💰", "Finance Tracker",       "Income/expense, profit reports, categories"),
    (BLUE,   "📅", "Calendar & Reminders",  "Activity scheduling, local notifications"),
    (RED,    "🌿", "Pest & Disease Guide",  "Offline reference library with search"),
    (PURPLE, "📓", "Daily Journal",         "Notes with photos and weather"),
    (TEAL,   "🏠", "Dashboard",             "Unified view of all farm data"),
    (ORANGE, "📊", "Reports & Charts",      "Visual profit/loss analysis"),
    (GREEN,  "🚀", "TestFlight",            "Real-world testing with farmers"),
]
col_w = 5.9
for i, (color, icon, title, desc) in enumerate(modules):
    col = i % 2
    row = i // 2
    x = 0.3 + col * (col_w + 0.2)
    y = 1.25 + row * 1.1
    add_rect(slide, x, y, col_w, 0.92, DARK2)
    add_rect(slide, x, y, 0.06, 0.92, color)
    add_text(slide, f"{icon}  {title}", x + 0.18, y + 0.07, col_w - 0.3, 0.38,
             size=14, bold=True, color=color)
    add_text(slide, desc,               x + 0.18, y + 0.52, col_w - 0.3, 0.35,
             size=12, color=GREY)

# Odd last item (7th) centred bottom
c, icon, title, desc = modules[6]
x = 0.3 + 3 * 0 * 0  # centred
y = 1.25 + 3 * 1.1
add_rect(slide, 3.5, y, col_w, 0.92, DARK2)
add_rect(slide, 3.5, y, 0.06, 0.92, c)
add_text(slide, f"{icon}  {title}", 3.68, y + 0.07, col_w - 0.3, 0.38,
         size=14, bold=True, color=c)
add_text(slide, desc,               3.68, y + 0.52, col_w - 0.3, 0.35,
         size=12, color=GREY)

add_rect(slide, 0.3, 6.5, 12.7, 0.06, TEAL)
add_text(slide,
         "🌾  You have built a complete, production-ready agricultural app for Cambodian farmers. "
         "Ship it. Collect feedback. Iterate.",
         0.3, 6.65, 12.7, 0.7, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
output = os.path.join(os.path.dirname(__file__), "Week12_Final_Polish_TestFlight.pptx")
prs.save(output)
print(f"Saved: {output}")
