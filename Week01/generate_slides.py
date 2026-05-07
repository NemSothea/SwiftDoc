#!/usr/bin/env python3
"""Generate Week 01 — Project Setup & MVVM Architecture slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour palette ───────────────────────────────────────────────────────────
GREEN  = RGBColor(0x1B, 0xB8, 0x89)   # Week01 accent
BLUE   = RGBColor(0x28, 0x7D, 0xFA)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
ORANGE = RGBColor(0xF3, 0x96, 0x20)
TEAL   = RGBColor(0x00, 0xC9, 0xC8)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
DARK2  = RGBColor(0x16, 0x21, 0x3E)
CARD   = RGBColor(0x0F, 0x2A, 0x45)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xAA, 0xAA, 0xBB)
RED    = RGBColor(0xE5, 0x47, 0x47)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

W, H = 13.33, 7.5


# ── Helpers ──────────────────────────────────────────────────────────────────
def add_slide(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, items, x, y, w, h,
                size=15, color=WHITE, bullet_color=GREEN):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = "• "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_code(slide, code_lines, x, y, w, h, size=11):
    add_rect(slide, x, y, w, h, CARD)
    tb = slide.shapes.add_textbox(
        Inches(x + 0.18), Inches(y + 0.15),
        Inches(w - 0.36), Inches(h - 0.3))
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = TEAL
        run.font.name = "Courier New"
    return tb


def slide_number(slide, n, total=12):
    add_text(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.3,
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


def header_bar(slide, title, accent=GREEN):
    add_rect(slide, 0, 0, W, 0.9, DARK2)
    add_rect(slide, 0, 0, 0.06, 0.9, accent)
    add_text(slide, title, 0.25, 0.1, 12.5, 0.7,
             size=24, bold=True, color=accent)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = add_slide(prs)
    add_rect(slide, 0, 0, 0.12, H, GREEN)

    add_rect(slide, 0.4, 0.4, 2.2, 0.42, GREEN)
    add_text(slide, "Week 01 · SmartFarmer", 0.4, 0.4, 2.2, 0.42,
             size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    add_text(slide, "រៀបចំ Project", 0.4, 1.0, 12.0, 1.0,
             size=52, bold=True, color=GREEN)
    add_text(slide, "& MVVM Architecture", 0.4, 1.9, 12.0, 0.85,
             size=34, bold=True, color=WHITE)

    add_text(slide, "iOS 13+  ·  SwiftUI  ·  ObservableObject  ·  CoreData",
             0.4, 2.9, 12.0, 0.5, size=15, italic=True, color=GREY)

    add_rect(slide, 0.4, 3.5, 10.5, 0.04, GREEN)

    icons = [
        ("🏗️", "Architecture"), ("📁", "Folder Setup"),
        ("📦", "Data Models"),  ("🧠", "ViewModel"),
        ("📱", "TabView"),      ("🌾", "Khmer UI"),
    ]
    for i, (emoji, label) in enumerate(icons):
        cx = 0.4 + i * 2.1
        add_rect(slide, cx, 3.7, 1.9, 1.0, CARD)
        add_text(slide, emoji, cx, 3.75, 1.9, 0.45, size=22, align=PP_ALIGN.CENTER)
        add_text(slide, label, cx, 4.22, 1.9, 0.4,
                 size=11, color=GREY, align=PP_ALIGN.CENTER)

    add_text(slide, "SmartFarmer Assistant · ភ្នំពេញ 2026",
             0.4, 6.9, 12.0, 0.4, size=12, italic=True, color=GREY)
    slide_number(slide, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
def slide_agenda(prs):
    slide = add_slide(prs)
    header_bar(slide, "📋  មាតិកា Week 01 · Agenda")

    topics = [
        (GREEN,  "01", "MVVM Architecture",        "Model · View · ViewModel pattern"),
        (BLUE,   "02", "Xcode Project Setup",       "Create project, folder structure"),
        (PURPLE, "03", "Data Models",               "Transaction, FarmActivity, Pest, JournalEntry"),
        (ORANGE, "04", "FarmViewModel",             "ObservableObject + @Published properties"),
        (TEAL,   "05", "MainTabView & Environment", "4-tab app + environmentObject()"),
    ]

    for i, (color, num, kh, en) in enumerate(topics):
        col = i % 2
        row = i // 2
        cx = 0.35 + col * 6.5
        cy = 1.1 + row * 2.0
        if i == 4:
            cx = 0.35 + 0 * 6.5

        add_rect(slide, cx, cy, 6.1, 1.7, CARD)
        add_rect(slide, cx, cy, 0.08, 1.7, color)

        add_rect(slide, cx + 0.18, cy + 0.35, 0.6, 0.6, color)
        add_text(slide, num, cx + 0.18, cy + 0.35, 0.6, 0.6,
                 size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        add_text(slide, kh, cx + 0.95, cy + 0.22, 4.9, 0.5,
                 size=17, bold=True, color=WHITE)
        add_text(slide, en, cx + 0.95, cy + 0.8, 4.9, 0.45,
                 size=12, color=GREY)

    slide_number(slide, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is MVVM?
# ════════════════════════════════════════════════════════════════════════════
def slide_mvvm(prs):
    slide = add_slide(prs)
    header_bar(slide, "🏗️  MVVM Architecture គឺជាអ្វី?")

    # 3-layer diagram
    layers = [
        (PURPLE, "Model",     "Data structures\nTransaction, FarmActivity,\nPest, JournalEntry"),
        (GREEN,  "ViewModel", "Business Logic\nFarmViewModel\n@Published properties"),
        (BLUE,   "View",      "SwiftUI UI\nFinanceTabView,\nCalendarTabView..."),
    ]
    for i, (color, title, desc) in enumerate(layers):
        cx = 0.35 + i * 4.3
        add_rect(slide, cx, 1.1, 3.9, 2.8, CARD)
        add_rect(slide, cx, 1.1, 0.08, 2.8, color)
        add_text(slide, title, cx + 0.2, 1.22, 3.5, 0.55,
                 size=20, bold=True, color=color)
        add_text(slide, desc, cx + 0.2, 1.9, 3.5, 1.8,
                 size=13, color=WHITE)
        if i < 2:
            add_text(slide, "↔", cx + 3.75, 2.3, 0.7, 0.5,
                     size=22, bold=True, color=GREY, align=PP_ALIGN.CENTER)

    add_bullets(slide, [
        "Model — Pure data, no UI logic. Ex: @Model class Transaction { var amount: Double }",
        "ViewModel — Connects Model + View. Holds @Published state, business rules.",
        "View — SwiftUI struct. Reads ViewModel via @EnvironmentObject. Never mutates Model directly.",
        "ហេតុអ្វីប្រើ MVVM? — Testable, reusable, clean separation of concerns",
    ], 0.35, 4.15, 12.6, 2.9, size=13, bullet_color=GREEN)

    slide_number(slide, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Xcode Project Setup
# ════════════════════════════════════════════════════════════════════════════
def slide_xcode_setup(prs):
    slide = add_slide(prs)
    header_bar(slide, "⚙️  Xcode Project Setup · រៀបចំ Project")

    steps = [
        ("01", "File → New → Project",       "Choose iOS → App template"),
        ("02", "Configure Project",           "Product Name: SmartFarmerAssistant\nInterface: SwiftUI · Language: Swift"),
        ("03", "Deployment Target",           "Set Minimum Deployments → iOS 13.0"),
        ("04", "Uncheck Core Data checkbox",  "We create CoreDataManager manually (better control)"),
    ]
    for i, (num, title, detail) in enumerate(steps):
        cy = 1.1 + i * 1.55
        add_rect(slide, 0.35, cy, 6.2, 1.35, CARD)
        add_rect(slide, 0.35, cy, 0.08, 1.35, GREEN)
        add_rect(slide, 0.55, cy + 0.37, 0.55, 0.55, GREEN)
        add_text(slide, num, 0.55, cy + 0.37, 0.55, 0.55,
                 size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.25, cy + 0.1, 5.1, 0.45,
                 size=14, bold=True, color=WHITE)
        add_text(slide, detail, 1.25, cy + 0.58, 5.1, 0.65,
                 size=11, color=GREY)

    add_text(slide, "Project Folder Structure", 6.85, 1.1, 6.1, 0.4,
             size=13, bold=True, color=GREEN)
    add_code(slide, [
        "SmartFarmerAssistant/",
        "├── App/",
        "│   └── SmartFarmerAssistantApp.swift",
        "├── Models/",
        "│   ├── Transaction+CoreData.swift",
        "│   ├── FarmActivity+CoreData.swift",
        "│   ├── Pest+CoreData.swift",
        "│   └── JournalEntry+CoreData.swift",
        "├── ViewModels/",
        "│   └── FarmViewModel.swift",
        "├── Views/",
        "│   ├── Finance/",
        "│   ├── Calendar/",
        "│   ├── PestGuide/",
        "│   └── Journal/",
        "├── Utilities/",
        "│   ├── CoreDataManager.swift",
        "│   └── Constants.swift",
        "└── Resources/",
        "    └── Assets.xcassets",
    ], 6.85, 1.6, 6.1, 5.7, size=10)

    slide_number(slide, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Folder Structure Purpose
# ════════════════════════════════════════════════════════════════════════════
def slide_folder(prs):
    slide = add_slide(prs)
    header_bar(slide, "📁  Folder Structure · គោលបំណងនៃ Folder")

    folders = [
        (GREEN,  "App/",        "Entry point, @main struct, CoreData container setup"),
        (BLUE,   "Models/",     "NSManagedObject subclasses — pure data, no UI"),
        (PURPLE, "ViewModels/", "ObservableObject classes — business logic + CRUD"),
        (ORANGE, "Views/",      "SwiftUI structs — sub-folders per feature tab"),
        (TEAL,   "Utilities/",  "CoreDataManager, Constants, Extensions"),
        (GREEN,  "Resources/",  "Assets.xcassets, JSON files, .xcdatamodeld"),
    ]

    for i, (color, name, purpose) in enumerate(folders):
        col = i % 2
        row = i // 2
        cx = 0.35 + col * 6.5
        cy = 1.1 + row * 1.95
        add_rect(slide, cx, cy, 6.1, 1.7, CARD)
        add_rect(slide, cx, cy, 0.08, 1.7, color)
        add_code(slide, [name], cx + 0.18, cy + 0.1, 3.0, 0.55, size=14)
        add_text(slide, purpose, cx + 0.2, cy + 0.75, 5.7, 0.8,
                 size=12, color=GREY)

    slide_number(slide, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Data Models Overview
# ════════════════════════════════════════════════════════════════════════════
def slide_models_overview(prs):
    slide = add_slide(prs)
    header_bar(slide, "📦  Data Models Overview · គំរូទិន្នន័យ")

    models = [
        (GREEN,  "Transaction",   "💰", ["amount: Double", "date: Date", "note: String", "type: String", "category: String"]),
        (BLUE,   "FarmActivity",  "📅", ["title: String", "activityType: String", "date: Date", "isCompleted: Bool", "reminderEnabled: Bool"]),
        (PURPLE, "Pest",          "🐛", ["name: String", "pestType: String", "symptoms: String", "treatment: String", "isFavorite: Bool"]),
        (ORANGE, "JournalEntry",  "📖", ["date: Date", "content: String", "weather: String", "photoData: Data?", "location: String?"]),
    ]

    for i, (color, name, emoji, attrs) in enumerate(models):
        col = i % 2
        row = i // 2
        cx = 0.35 + col * 6.5
        cy = 1.1 + row * 3.1
        add_rect(slide, cx, cy, 6.1, 2.85, CARD)
        add_rect(slide, cx, cy, 0.08, 2.85, color)
        add_text(slide, f"{emoji}  {name}", cx + 0.2, cy + 0.1, 5.7, 0.5,
                 size=16, bold=True, color=color)
        add_text(slide, "NSManagedObject subclass", cx + 0.2, cy + 0.6, 5.7, 0.3,
                 size=10, italic=True, color=GREY)
        for j, attr in enumerate(attrs):
            add_code(slide, [attr], cx + 0.2, cy + 0.98 + j * 0.38, 5.7, 0.33, size=10)

    slide_number(slide, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Transaction Model Code
# ════════════════════════════════════════════════════════════════════════════
def slide_transaction_model(prs):
    slide = add_slide(prs)
    header_bar(slide, "💰  Transaction Model · CoreData @NSManaged")

    add_text(slide, "Transaction+CoreDataClass.swift", 0.35, 1.0, 6.1, 0.4,
             size=12, bold=True, color=GREEN)
    add_code(slide, [
        "// Models/Transaction+CoreDataClass.swift",
        "@objc(Transaction)",
        "public class Transaction: NSManagedObject { }",
        "",
        "// Models/Transaction+CoreDataProperties.swift",
        "extension Transaction {",
        "    @nonobjc public class func fetchRequest()",
        "        -> NSFetchRequest<Transaction> {",
        "        return NSFetchRequest<Transaction>(",
        '            entityName: "Transaction")',
        "    }",
        "    @NSManaged public var amount:   Double",
        "    @NSManaged public var date:     Date?",
        "    @NSManaged public var note:     String?",
        "    @NSManaged public var type:     String?",
        "    @NSManaged public var category: String?",
        "    @NSManaged public var id:       UUID?",
        "}",
        "extension Transaction: Identifiable {}",
    ], 0.35, 1.5, 6.1, 5.85, size=10)

    add_text(slide, "TransactionType Enums", 6.65, 1.0, 6.3, 0.4,
             size=12, bold=True, color=BLUE)
    add_code(slide, [
        "// Models/TransactionType.swift",
        "enum ExpenseCategory: String, CaseIterable {",
        '    case seeds      = "គ្រាប់ពូជ"',
        '    case fertilizer = "ជី"',
        '    case labor      = "កម្លាំងពលកម្ម"',
        '    case tools      = "ឧបករណ៍"',
        '    case other      = "ផ្សេងៗ"',
        "}",
        "",
        "enum IncomeCategory: String, CaseIterable {",
        '    case vegetable = "បន្លែ"',
        '    case fruit     = "ផ្លែឈើ"',
        '    case grain     = "ស្រូវ-ដំណាំ"',
        '    case livestock = "សត្វ"',
        '    case other     = "ផ្សេងៗ"',
        "}",
    ], 6.65, 1.5, 6.3, 3.8, size=10)

    add_rect(slide, 6.65, 5.55, 6.3, 1.35, CARD)
    add_rect(slide, 6.65, 5.55, 0.08, 1.35, YELLOW)
    add_text(slide, "💡  CoreData ប្រើ @NSManaged  (ជំនួស stored property)\n"
             "    Codegen → Manual/None (ដើម្បីចៀសវាង build error)\n"
             "    id: UUID → Identifiable conformance",
             6.88, 5.6, 5.9, 1.2, size=11, color=WHITE)

    slide_number(slide, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FarmActivity & Pest Models
# ════════════════════════════════════════════════════════════════════════════
def slide_other_models(prs):
    slide = add_slide(prs)
    header_bar(slide, "📅🐛  FarmActivity & Pest Models")

    add_text(slide, "FarmActivity+CoreDataProperties.swift", 0.35, 1.0, 6.1, 0.35,
             size=12, bold=True, color=BLUE)
    add_code(slide, [
        "extension FarmActivity {",
        "    @NSManaged public var id:              UUID?",
        "    @NSManaged public var title:           String?",
        "    @NSManaged public var activityType:    String?",
        "    @NSManaged public var date:            Date?",
        "    @NSManaged public var notes:           String?",
        "    @NSManaged public var isCompleted:     Bool",
        "    @NSManaged public var reminderEnabled: Bool",
        "}",
        "",
        "enum ActivityType: String, CaseIterable {",
        '    case planting     = "ដាំ"',
        '    case watering     = "ស្រោចទឹក"',
        '    case fertilizing  = "ដាក់ជី"',
        '    case harvesting   = "ប្រមូលផល"',
        '    case pesticide    = "បាញ់ថ្នាំ"',
        '    case other        = "ផ្សេងៗ"',
        "}",
    ], 0.35, 1.45, 6.1, 5.9, size=10)

    add_text(slide, "Pest+CoreDataProperties.swift", 6.65, 1.0, 6.3, 0.35,
             size=12, bold=True, color=PURPLE)
    add_code(slide, [
        "extension Pest {",
        "    @NSManaged public var id:         UUID?",
        "    @NSManaged public var name:       String?",
        "    @NSManaged public var pestType:   String?",
        "    @NSManaged public var symptoms:   String?",
        "    @NSManaged public var treatment:  String?",
        "    @NSManaged public var prevention: String?",
        "    @NSManaged public var imageName:  String?",
        "    @NSManaged public var isFavorite: Bool",
        "}",
        "",
        "// JournalEntry+CoreDataProperties.swift",
        "extension JournalEntry {",
        "    @NSManaged public var id:        UUID?",
        "    @NSManaged public var date:      Date?",
        "    @NSManaged public var content:   String?",
        "    @NSManaged public var weather:   String?",
        "    @NSManaged public var photoData: Data?",
        "    @NSManaged public var location:  String?",
        "}",
    ], 6.65, 1.45, 6.3, 5.9, size=10)

    slide_number(slide, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FarmViewModel
# ════════════════════════════════════════════════════════════════════════════
def slide_viewmodel(prs):
    slide = add_slide(prs)
    header_bar(slide, "🧠  FarmViewModel · ObservableObject (iOS 13+)")

    add_code(slide, [
        "// ViewModels/FarmViewModel.swift",
        "import Foundation",
        "import SwiftUI",
        "import CoreData",
        "",
        "// ✅ ObservableObject for iOS 13+",
        "// ❌ Do NOT use @Observable (requires iOS 17)",
        "class FarmViewModel: ObservableObject {",
        "",
        "    private var context: NSManagedObjectContext",
        "",
        "    init(context: NSManagedObjectContext",
        "         = CoreDataManager.shared.context) {",
        "        self.context = context",
        "    }",
        "",
        "    // MARK: - Helper Methods",
        "    func formatCurrency(_ amount: Double) -> String {",
        "        let f = NumberFormatter()",
        "        f.numberStyle = .currency",
        "        f.locale = Locale(identifier: \"en_US\")",
        "        return f.string(from: NSNumber(value: amount))",
        "               ?? \"$0.00\"",
        "    }",
        "",
        "    func formatDate(_ date: Date) -> String {",
        "        let f = DateFormatter()",
        "        f.locale = Locale(identifier: \"km-KH\")",
        "        f.dateStyle = .medium",
        "        return f.string(from: date)",
        "    }",
        "}",
    ], 0.35, 1.1, 6.8, 6.25, size=10)

    # explanation cards
    notes = [
        (GREEN,  "ObservableObject",  "Base protocol — SwiftUI\nwatches this object for\nchanges (iOS 13+)"),
        (BLUE,   "@Published",        "MarkUp ← property ⟹\nview updates automatically\nwhen value changes"),
        (PURPLE, "@StateObject",      "Use in root view to OWN\nthe ViewModel instance"),
        (ORANGE, "@EnvironmentObject", "Use in child views to\nREAD the shared ViewModel\npassed via environment"),
    ]
    for i, (color, title, desc) in enumerate(notes):
        cy = 1.1 + i * 1.57
        add_rect(slide, 7.35, cy, 5.6, 1.4, CARD)
        add_rect(slide, 7.35, cy, 0.08, 1.4, color)
        add_text(slide, title, 7.6, cy + 0.1, 5.1, 0.45,
                 size=13, bold=True, color=color)
        add_text(slide, desc, 7.6, cy + 0.58, 5.1, 0.75,
                 size=11, color=WHITE)

    slide_number(slide, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — MainTabView
# ════════════════════════════════════════════════════════════════════════════
def slide_maintabview(prs):
    slide = add_slide(prs)
    header_bar(slide, "📱  MainTabView · 4-Tab App + environmentObject")

    add_code(slide, [
        "// Views/MainTabView.swift",
        "struct MainTabView: View {",
        "",
        "    // ✅ @StateObject — owns the ViewModel",
        "    @StateObject private var viewModel: FarmViewModel",
        "    @State private var selectedTab = 0",
        "    @Environment(\\.managedObjectContext)",
        "        private var viewContext",
        "",
        "    init() {",
        "        _viewModel = StateObject(wrappedValue:",
        "            FarmViewModel(",
        "                context: CoreDataManager.shared.context))",
        "    }",
        "",
        "    var body: some View {",
        "        TabView(selection: $selectedTab) {",
        "            FinanceTabView()",
        '                .tabItem { Label("ហិរញ្ញវត្ថុ",',
        '                    systemImage: "dollarsign.circle") }',
        "                .tag(0)",
        "                .environment(\\.managedObjectContext,",
        "                             viewContext)",
        "            // ... Calendar, Pest, Journal tabs",
        "        }",
        "        // ✅ Pass ViewModel to ALL child views",
        "        .environmentObject(viewModel)",
        "    }",
        "}",
    ], 0.35, 1.1, 7.0, 6.25, size=10)

    # tab cards
    tabs = [
        (GREEN,  "ហិរញ្ញវត្ថុ",  "dollarsign.circle", "Finance Tracker"),
        (BLUE,   "ប្រតិទិន",    "calendar",          "Calendar & Reminders"),
        (PURPLE, "សត្វល្អិត",   "bug",               "Pest & Disease Guide"),
        (ORANGE, "កំណត់ហេតុ",   "book",              "Daily Journal"),
    ]
    for i, (color, title, icon, desc) in enumerate(tabs):
        cy = 1.1 + i * 1.57
        add_rect(slide, 7.55, cy, 5.4, 1.35, CARD)
        add_rect(slide, 7.55, cy, 0.08, 1.35, color)
        add_text(slide, f"SF: {icon}", 7.75, cy + 0.1, 5.0, 0.35,
                 size=10, italic=True, color=GREY)
        add_text(slide, title, 7.75, cy + 0.45, 2.8, 0.45,
                 size=15, bold=True, color=color)
        add_text(slide, desc, 7.75, cy + 0.9, 5.0, 0.35,
                 size=11, color=GREY)

    slide_number(slide, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Khmer Localization & Constants
# ════════════════════════════════════════════════════════════════════════════
def slide_khmer(prs):
    slide = add_slide(prs)
    header_bar(slide, "🌾  Khmer Localization & Constants · ភាសាខ្មែរ")

    add_code(slide, [
        "// Utilities/Constants.swift",
        "struct AppStrings {",
        "    // Tab titles",
        '    static let financeTab  = "ហិរញ្ញវត្ថុ"',
        '    static let calendarTab = "ប្រតិទិន"',
        '    static let pestTab     = "សត្វល្អិត"',
        '    static let journalTab  = "កំណត់ហេតុ"',
        "",
        "    // Common buttons",
        '    static let save   = "រក្សាទុក"',
        '    static let cancel = "បោះបង់"',
        '    static let delete = "លុប"',
        '    static let edit   = "កែប្រែ"',
        '    static let add    = "បន្ថែម"',
        "}",
        "",
        "// Utilities/Extensions/Date+Khmer.swift",
        "extension Date {",
        "    func khmerFormatted() -> String {",
        "        let f = DateFormatter()",
        '        f.locale = Locale(identifier: "km-KH")',
        "        f.dateStyle = .full",
        "        return f.string(from: self)",
        "    }",
        "}",
    ], 0.35, 1.1, 6.3, 6.25, size=10)

    # currency card
    add_text(slide, "formatCurrency · ប្រព័ន្ធលុយ", 6.85, 1.1, 6.1, 0.4,
             size=13, bold=True, color=GREEN)
    add_code(slide, [
        "// en_US  → $1,234.56",
        "// km_KH  → 1.234,56 ៛",
        "",
        'let f = NumberFormatter()',
        'f.numberStyle = .currency',
        'f.locale = Locale(identifier: "en_US")',
        'let result = f.string(from: NSNumber(value: 1234.5))',
        '// → "$1,234.50"',
    ], 6.85, 1.6, 6.1, 2.6, size=11)

    # weather/activity labels
    add_text(slide, "Khmer Labels ក្នុង Enums", 6.85, 4.4, 6.1, 0.4,
             size=13, bold=True, color=BLUE)
    add_code(slide, [
        'case planting  = "ដាំ"        // Planting',
        'case watering  = "ស្រោចទឹក"  // Watering',
        'case harvesting = "ប្រមូលផល" // Harvesting',
        'case sunny = "ថ្ងៃរះ"        // Sunny',
        'case rainy = "ភ្លៀង"         // Rainy',
    ], 6.85, 4.9, 6.1, 2.45, size=11)

    slide_number(slide, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary & Assignment
# ════════════════════════════════════════════════════════════════════════════
def slide_summary(prs):
    slide = add_slide(prs)
    header_bar(slide, "📌  សង្ខេប Week 01 & Mini-Project", accent=GREEN)

    # 3-col summary
    cols = [
        ("✅  Built", GREEN, [
            "Xcode project structure",
            "4 CoreData models",
            "FarmViewModel class",
            "MainTabView + 4 tabs",
            "Khmer constants",
        ]),
        ("📋  Checklist", BLUE, [
            "Project builds ✓",
            "4 model files exist ✓",
            "ObservableObject VM ✓",
            "4 tabs with Khmer titles ✓",
            "environmentObject() set ✓",
        ]),
        ("🎯  Next: Week 02", PURPLE, [
            "CoreData stack setup",
            "NSPersistentContainer",
            "@FetchRequest",
            "CRUD operations",
            "Finance Tab complete",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        cx = 0.35 + i * 4.35
        add_rect(slide, cx, 1.1, 4.1, 5.7, CARD)
        add_rect(slide, cx, 1.1, 4.1, 0.6, color)
        add_text(slide, title, cx + 0.15, 1.15, 3.8, 0.5,
                 size=14, bold=True, color=DARK)
        for j, item in enumerate(items):
            cy_item = 1.82 + j * 0.9
            add_rect(slide, cx + 0.15, cy_item, 3.7, 0.72, DARK2)
            add_text(slide, item, cx + 0.3, cy_item + 0.1, 3.4, 0.52,
                     size=12, color=WHITE)

    add_rect(slide, 0.35, 6.98, 12.6, 0.42, GREEN)
    add_text(slide, "🌾  Week 01 ចប់! — រៀបចំ foundation រួចរាល់ · Ready for CoreData in Week 02!",
             0.5, 7.01, 12.3, 0.38,
             size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    slide_number(slide, 12)


# ── BUILD ────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    slide_title(prs)
    slide_agenda(prs)
    slide_mvvm(prs)
    slide_xcode_setup(prs)
    slide_folder(prs)
    slide_models_overview(prs)
    slide_transaction_model(prs)
    slide_other_models(prs)
    slide_viewmodel(prs)
    slide_maintabview(prs)
    slide_khmer(prs)
    slide_summary(prs)

    out = os.path.join(os.path.dirname(__file__), "Week01_MVVM_Setup.pptx")
    prs.save(out)
    print(f"✅  Saved → {out}")


if __name__ == "__main__":
    build()
