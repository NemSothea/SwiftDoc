#!/usr/bin/env python3
"""Generate Week 02 — Core Data Persistence slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour palette ───────────────────────────────────────────────────────────
BLUE   = RGBColor(0x28, 0x7D, 0xFA)   # Week02 accent
GREEN  = RGBColor(0x1B, 0xB8, 0x89)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
ORANGE = RGBColor(0xF3, 0x96, 0x20)
TEAL   = RGBColor(0x00, 0xC9, 0xC8)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
DARK2  = RGBColor(0x16, 0x21, 0x3E)
CARD   = RGBColor(0x0F, 0x2A, 0x45)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xAA, 0xAA, 0xBB)
RED    = RGBColor(0xE5, 0x47, 0x47)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

W, H = 13.33, 7.5


# ── Helpers ──────────────────────────────────────────────────────────────────
def add_slide(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, items, x, y, w, h,
                size=15, color=WHITE, bullet_color=BLUE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_code(slide, code_lines, x, y, w, h, size=11):
    add_rect(slide, x, y, w, h, CARD)
    tb = slide.shapes.add_textbox(
        Inches(x + 0.18), Inches(y + 0.15),
        Inches(w - 0.36), Inches(h - 0.3))
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = TEAL
        run.font.name = "Courier New"
    return tb


def slide_number(slide, n, total=12):
    add_text(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.3,
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


def header_bar(slide, title, accent=BLUE):
    add_rect(slide, 0, 0, W, 0.9, DARK2)
    add_rect(slide, 0, 0, 0.06, 0.9, accent)
    add_text(slide, title, 0.25, 0.1, 12.5, 0.7,
             size=24, bold=True, color=accent)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = add_slide(prs)
    add_rect(slide, 0, 0, 0.12, H, BLUE)

    add_rect(slide, 0.4, 0.4, 2.2, 0.42, BLUE)
    add_text(slide, "Week 02 · SmartFarmer", 0.4, 0.4, 2.2, 0.42,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "CoreData", 0.4, 1.0, 12.0, 1.0,
             size=52, bold=True, color=BLUE)
    add_text(slide, "Persistence & CRUD", 0.4, 1.9, 12.0, 0.85,
             size=34, bold=True, color=WHITE)

    add_text(slide, "iOS 13+  ·  NSPersistentContainer  ·  @FetchRequest  ·  CRUD",
             0.4, 2.9, 12.0, 0.5, size=15, italic=True, color=GREY)

    add_rect(slide, 0.4, 3.5, 10.5, 0.04, BLUE)

    icons = [
        ("🗄️", "CoreData"),   ("📋", "@FetchRequest"),
        ("➕", "Create"),     ("✏️", "Update"),
        ("🗑️", "Delete"),     ("💾", "Persist"),
    ]
    for i, (emoji, label) in enumerate(icons):
        cx = 0.4 + i * 2.1
        add_rect(slide, cx, 3.7, 1.9, 1.0, CARD)
        add_text(slide, emoji, cx, 3.75, 1.9, 0.45, size=22, align=PP_ALIGN.CENTER)
        add_text(slide, label, cx, 4.22, 1.9, 0.4,
                 size=11, color=GREY, align=PP_ALIGN.CENTER)

    add_text(slide, "SmartFarmer Assistant · ភ្នំពេញ 2026",
             0.4, 6.9, 12.0, 0.4, size=12, italic=True, color=GREY)
    slide_number(slide, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
def slide_agenda(prs):
    slide = add_slide(prs)
    header_bar(slide, "📋  មាតិកា Week 02 · Agenda")

    topics = [
        (BLUE,   "01", "CoreData vs SwiftData",       "ស្គាល់ភាពខុសគ្នា — iOS 13+ vs iOS 17+"),
        (GREEN,  "02", "NSPersistentContainer Setup",  "CoreDataManager singleton pattern"),
        (PURPLE, "03", "Entity Design",                "4 entities in .xcdatamodeld editor"),
        (ORANGE, "04", "@FetchRequest",                "Auto UI refresh on data change"),
        (TEAL,   "05", "CRUD Operations",              "Create · Read · Update · Delete"),
        (BLUE,   "06", "Finance Tab Complete",         "FinanceTabView + SummaryCard + filter"),
    ]

    for i, (color, num, title, sub) in enumerate(topics):
        col = i % 2
        row = i // 2
        cx = 0.35 + col * 6.5
        cy = 1.1 + row * 2.0
        add_rect(slide, cx, cy, 6.1, 1.7, CARD)
        add_rect(slide, cx, cy, 0.08, 1.7, color)
        add_rect(slide, cx + 0.18, cy + 0.35, 0.6, 0.6, color)
        add_text(slide, num, cx + 0.18, cy + 0.35, 0.6, 0.6,
                 size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(slide, title, cx + 0.95, cy + 0.22, 4.9, 0.5,
                 size=16, bold=True, color=WHITE)
        add_text(slide, sub, cx + 0.95, cy + 0.82, 4.9, 0.45,
                 size=12, color=GREY)

    slide_number(slide, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CoreData vs SwiftData
# ════════════════════════════════════════════════════════════════════════════
def slide_comparison(prs):
    slide = add_slide(prs)
    header_bar(slide, "⚖️  CoreData vs SwiftData · ភាពខុសគ្នា")

    rows = [
        ("Feature",           "SwiftData (iOS 17+)", "CoreData (iOS 13+)  ✅"),
        ("Model declaration", "@Model class",         "NSManagedObject subclass"),
        ("Query in view",     "@Query",               "@FetchRequest"),
        ("Container setup",   ".modelContainer()",    "NSPersistentContainer"),
        ("ViewModel",         "@Observable class",    "class: ObservableObject"),
        ("Pass to views",     ".environment(vm)",     ".environmentObject(vm)"),
        ("Read in views",     "@Environment(VM.self)","@EnvironmentObject var vm"),
        ("Navigation",        "NavigationStack",      "NavigationView"),
        ("Min iOS target",    "iOS 17+",              "iOS 13+  (ប្រើក្នុង course នេះ)"),
    ]

    col_widths = [3.2, 4.4, 4.6]
    col_starts = [0.35, 3.7, 8.25]
    header_colors = [DARK2, RED, GREEN]
    text_colors   = [GREY,  WHITE, WHITE]

    for r, row in enumerate(rows):
        cy = 1.1 + r * 0.67
        for c, (cell, cw, cs, hc, tc) in enumerate(
                zip(row, col_widths, col_starts, header_colors, text_colors)):
            bg = hc if r == 0 else (CARD if r % 2 == 1 else DARK2)
            add_rect(slide, cs, cy, cw, 0.62, bg)
            fc = (WHITE if r == 0 else tc)
            bold = (r == 0)
            add_text(slide, cell, cs + 0.12, cy + 0.08, cw - 0.2, 0.46,
                     size=11, bold=bold, color=fc)

    # bottom note
    add_rect(slide, 0.35, 7.1, 12.6, 0.32, BLUE)
    add_text(slide, "✅  ក្នុង course នេះ — ប្រើ CoreData ព្រោះ target iOS 13+  "
             "| SwiftData ត្រូវការ iOS 17+",
             0.5, 7.12, 12.2, 0.28, size=11, bold=True, color=DARK)

    slide_number(slide, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — NSPersistentContainer Setup
# ════════════════════════════════════════════════════════════════════════════
def slide_coredata_setup(prs):
    slide = add_slide(prs)
    header_bar(slide, "🗄️  NSPersistentContainer & CoreDataManager")

    add_code(slide, [
        "// Utilities/CoreDataManager.swift",
        "import CoreData",
        "",
        "class CoreDataManager {",
        "    // ✅ Singleton — ONE container for entire app",
        "    static let shared = CoreDataManager()",
        "",
        "    lazy var persistentContainer: NSPersistentContainer = {",
        "        // ⚠️ Name must EXACTLY match .xcdatamodeld filename",
        '        let c = NSPersistentContainer(name:',
        '                    "SmartFarmerAssistantFinish")',
        "        c.loadPersistentStores { _, error in",
        "            if let error = error {",
        "                fatalError(\"CoreData failed: \\(error)\")",
        "            }",
        "        }",
        "        return c",
        "    }()",
        "",
        "    var context: NSManagedObjectContext {",
        "        persistentContainer.viewContext",
        "    }",
        "",
        "    func saveContext() {",
        "        if context.hasChanges {",
        "            try? context.save()",
        "        }",
        "    }",
        "}",
    ], 0.35, 1.1, 6.5, 6.25, size=10)

    add_text(slide, "App Entry Point", 7.05, 1.1, 6.0, 0.4,
             size=13, bold=True, color=BLUE)
    add_code(slide, [
        "// App/SmartFarmerAssistantFinishApp.swift",
        "@main",
        "struct SmartFarmerAssistantFinishApp: App {",
        "",
        "    // ✅ Reuse shared container",
        "    let context = CoreDataManager.shared.context",
        "",
        "    var body: some Scene {",
        "        WindowGroup {",
        "            MainTabView()",
        "                .environment(",
        "                    \\.managedObjectContext,",
        "                    context)",
        "        }",
        "    }",
        "}",
    ], 7.05, 1.6, 6.0, 4.0, size=10)

    add_rect(slide, 7.05, 5.85, 6.0, 1.5, CARD)
    add_rect(slide, 7.05, 5.85, 0.08, 1.5, RED)
    add_text(slide, "⚠️  Common Mistake",
             7.25, 5.9, 5.7, 0.4, size=13, bold=True, color=RED)
    add_text(slide, "ប្រើ NSPersistentContainer ២ ដង (ក្នុង App + CoreDataManager)\n"
             "→ Crash: \"Multiple NSEntityDescriptions claim the subclass\"\n"
             "Fix: ប្រើ CoreDataManager.shared.context ជានិច្ច",
             7.25, 6.35, 5.7, 0.9, size=11, color=WHITE)

    slide_number(slide, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Entity Design
# ════════════════════════════════════════════════════════════════════════════
def slide_entities(prs):
    slide = add_slide(prs)
    header_bar(slide, "🗂️  Entity Design · .xcdatamodeld Editor")

    entities = [
        (GREEN,  "Transaction",  [
            "amount: Double", "date: Date", "note: String",
            "type: String", "category: String", "id: UUID"]),
        (BLUE,   "FarmActivity", [
            "title: String", "activityType: String", "date: Date",
            "notes: String", "isCompleted: Boolean", "reminderEnabled: Boolean", "id: UUID"]),
        (PURPLE, "Pest",         [
            "name: String", "pestType: String", "symptoms: String",
            "treatment: String", "prevention: String", "imageName: String",
            "isFavorite: Boolean", "id: UUID"]),
        (ORANGE, "JournalEntry", [
            "date: Date", "content: String", "weather: String",
            "photoData: Binary Data", "location: String", "id: UUID"]),
    ]

    for i, (color, name, attrs) in enumerate(entities):
        col = i % 2
        row = i // 2
        cx = 0.35 + col * 6.5
        cy = 1.1 + row * 3.15
        add_rect(slide, cx, cy, 6.1, 2.9, CARD)
        add_rect(slide, cx, cy, 6.1, 0.5, color)
        add_text(slide, name, cx + 0.15, cy + 0.07, 5.8, 0.38,
                 size=15, bold=True, color=DARK)
        for j, attr in enumerate(attrs):
            add_code(slide, [attr], cx + 0.15, cy + 0.6 + j * 0.36, 5.7, 0.3, size=10)

    add_rect(slide, 0.35, 7.1, 12.6, 0.32, DARK2)
    add_text(slide, "⚙️  ក្នុង Xcode: ចុចលើ Entity → Inspector → Codegen → Manual/None   "
             "→  Editor → Create NSManagedObject Subclass",
             0.5, 7.12, 12.2, 0.28, size=10, color=GREY)

    slide_number(slide, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — NSManagedObject Subclass
# ════════════════════════════════════════════════════════════════════════════
def slide_managed_object(prs):
    slide = add_slide(prs)
    header_bar(slide, "🔧  NSManagedObject Subclass · Transaction Example")

    add_code(slide, [
        "// Models/Transaction+CoreDataClass.swift",
        "import CoreData",
        "",
        "@objc(Transaction)",
        "public class Transaction: NSManagedObject { }",
        "",
        "// ─────────────────────────────────────────────",
        "",
        "// Models/Transaction+CoreDataProperties.swift",
        "import CoreData",
        "",
        "extension Transaction {",
        "",
        "    @nonobjc public class func fetchRequest()",
        "        -> NSFetchRequest<Transaction> {",
        '        return NSFetchRequest<Transaction>(entityName: "Transaction")',
        "    }",
        "",
        "    @NSManaged public var amount:   Double",
        "    @NSManaged public var date:     Date?",
        "    @NSManaged public var note:     String?",
        "    @NSManaged public var type:     String?",
        "    @NSManaged public var category: String?",
        "    @NSManaged public var id:       UUID?",
        "}",
        "",
        "extension Transaction: Identifiable {}",
    ], 0.35, 1.1, 6.5, 6.25, size=10)

    # explanation
    notes = [
        (BLUE,   "@objc(Transaction)",    "ឈ្មោះ class ត្រូវដូចនឹង\nEntity name ក្នុង .xcdatamodeld"),
        (GREEN,  "@NSManaged",            "CoreData គ្រប់គ្រង\nstorage ដោយខ្លួនឯង\n(ជំនួស stored property)"),
        (PURPLE, "Identifiable",          "ធ្វើឲ្យ ForEach SwiftUI\nអាច loop transaction\nដោយ id"),
        (ORANGE, "Manual/None Codegen",   "ចៀសវាង build error:\n\"Multiple commands produce\"\nវាចាំបាច់ណាស់!"),
    ]
    for i, (color, title, desc) in enumerate(notes):
        cy = 1.1 + i * 1.57
        add_rect(slide, 7.05, cy, 5.9, 1.4, CARD)
        add_rect(slide, 7.05, cy, 0.08, 1.4, color)
        add_text(slide, title, 7.28, cy + 0.1, 5.4, 0.4,
                 size=13, bold=True, color=color)
        add_text(slide, desc, 7.28, cy + 0.55, 5.4, 0.78,
                 size=11, color=WHITE)

    slide_number(slide, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — @FetchRequest
# ════════════════════════════════════════════════════════════════════════════
def slide_fetchrequest(prs):
    slide = add_slide(prs)
    header_bar(slide, "🔄  @FetchRequest · Auto UI Refresh")

    add_code(slide, [
        "// ✅ Basic @FetchRequest",
        "struct TransactionListView: View {",
        "    @FetchRequest(",
        "        entity: Transaction.entity(),",
        "        sortDescriptors: [",
        "            NSSortDescriptor(",
        "                keyPath: \\Transaction.date,",
        "                ascending: false)",
        "        ]",
        "    ) var transactions: FetchedResults<Transaction>",
        "",
        "    var body: some View {",
        "        List {",
        "            ForEach(transactions, id: \\.self) { t in",
        "                Text(t.note ?? \"\")",
        "            }",
        "        }",
        "    }",
        "}",
        "",
        "// ✅ Filtered @FetchRequest (NSPredicate)",
        "@FetchRequest(",
        "    entity: Transaction.entity(),",
        "    sortDescriptors: [],",
        '    predicate: NSPredicate(format: "type == %@",',
        '                           "expense")',
        ") var expenses: FetchedResults<Transaction>",
    ], 0.35, 1.1, 6.5, 6.25, size=10)

    # flow diagram
    add_text(slide, "Auto-refresh Flow", 7.05, 1.1, 6.0, 0.4,
             size=13, bold=True, color=BLUE)

    flow = [
        (BLUE,   "User adds/edits/deletes"),
        (GREEN,  "viewContext.save()"),
        (PURPLE, "@FetchRequest detects change"),
        (ORANGE, "SwiftUI re-renders List"),
    ]
    for i, (color, label) in enumerate(flow):
        cy = 1.65 + i * 1.1
        add_rect(slide, 7.05, cy, 5.9, 0.75, CARD)
        add_rect(slide, 7.05, cy, 0.08, 0.75, color)
        add_text(slide, f"{'①②③④'[i]}  {label}", 7.28, cy + 0.15, 5.5, 0.45,
                 size=13, color=WHITE)
        if i < 3:
            add_text(slide, "↓", 9.75, cy + 0.73, 0.5, 0.4,
                     size=16, bold=True, color=GREY, align=PP_ALIGN.CENTER)

    add_rect(slide, 7.05, 6.2, 5.9, 1.15, CARD)
    add_rect(slide, 7.05, 6.2, 0.08, 1.15, YELLOW)
    add_text(slide, "💡  NSPredicate vs Swift .filter{}",
             7.25, 6.25, 5.5, 0.38, size=13, bold=True, color=YELLOW)
    add_text(slide, "NSPredicate → filters IN the SQLite database (fast)\n"
             ".filter{} → loads ALL rows first, filters in memory (slow)",
             7.25, 6.65, 5.5, 0.62, size=11, color=WHITE)

    slide_number(slide, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CRUD Create & Delete
# ════════════════════════════════════════════════════════════════════════════
def slide_crud_cd(prs):
    slide = add_slide(prs)
    header_bar(slide, "➕🗑️  CRUD · Create & Delete")

    add_text(slide, "Create — addTransaction()", 0.35, 1.0, 6.1, 0.4,
             size=13, bold=True, color=GREEN)
    add_code(slide, [
        "func addTransaction(amount: Double,",
        "                    note: String,",
        "                    type: String,",
        "                    category: String) {",
        "    let t = Transaction(context: context)",
        "    t.amount   = amount",
        "    t.date     = Date()",
        "    t.note     = note",
        "    t.type     = type",
        "    t.category = category",
        "    t.id       = UUID()",
        "    saveContext()  // ← persists to SQLite",
        "}",
        "",
        "// Call from View:",
        "viewModel.addTransaction(",
        "    amount: 50.0,",
        '    note: "ទិញគ្រាប់ពូជ",',
        '    type: "expense",',
        '    category: "គ្រាប់ពូជ"',
        ")",
    ], 0.35, 1.5, 6.1, 5.85, size=10)

    add_text(slide, "Delete — deleteTransaction()", 6.65, 1.0, 6.3, 0.4,
             size=13, bold=True, color=RED)
    add_code(slide, [
        "func deleteTransaction(_ t: Transaction) {",
        "    context.delete(t)",
        "    saveContext()",
        "}",
        "",
        "// Swipe-to-delete in List:",
        ".onDelete(perform: deleteTransactions)",
        "",
        "private func deleteTransactions(",
        "    offsets: IndexSet) {",
        "    for index in offsets {",
        "        viewContext.delete(",
        "            displayedTransactions[index])",
        "    }",
        "    try? viewContext.save()",
        "}",
        "",
        "// Batch delete ALL:",
        "func deleteAllTransactions() {",
        "    let req: NSFetchRequest<NSFetchRequestResult>",
        "            = Transaction.fetchRequest()",
        "    let batch = NSBatchDeleteRequest(fetchRequest: req)",
        "    try? context.execute(batch)",
        "    saveContext()",
        "}",
    ], 6.65, 1.5, 6.3, 5.85, size=10)

    slide_number(slide, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CRUD Update
# ════════════════════════════════════════════════════════════════════════════
def slide_crud_update(prs):
    slide = add_slide(prs)
    header_bar(slide, "✏️  CRUD · Update — Edit Flow & State(initialValue:)")

    add_text(slide, "updateTransaction() in ViewModel", 0.35, 1.0, 6.5, 0.4,
             size=13, bold=True, color=BLUE)
    add_code(slide, [
        "func updateTransaction(",
        "    _ t: Transaction,",
        "    amount: Double,",
        "    note: String,",
        "    type: String,",
        "    category: String) {",
        "    t.amount   = amount",
        "    t.note     = note",
        "    t.type     = type",
        "    t.category = category",
        "    saveContext()  // ← write to disk",
        "}",
        "",
        "// Pre-fill EditTransactionView:",
        "struct EditTransactionView: View {",
        "    let transaction: Transaction",
        "",
        "    // ✅ State(initialValue:) to pre-fill",
        "    @State private var amount: String",
        "",
        "    init(transaction: Transaction) {",
        "        self.transaction = transaction",
        "        _amount = State(initialValue:",
        "            String(transaction.amount))",
        "    }",
        "}",
    ], 0.35, 1.5, 6.5, 5.85, size=10)

    # edit flow
    add_text(slide, "Edit Flow Diagram", 7.05, 1.0, 6.0, 0.4,
             size=13, bold=True, color=ORANGE)

    flow = [
        (BLUE,   "User taps a transaction row"),
        (GREEN,  "selectedTransaction = transaction"),
        (PURPLE, ".sheet(item: $selectedTransaction)"),
        (ORANGE, "EditTransactionView opens (pre-filled)"),
        (TEAL,   "User edits → taps Save"),
        (BLUE,   "viewModel.updateTransaction(...) → disk"),
        (GREEN,  "@FetchRequest → UI refreshes ✅"),
    ]
    for i, (color, label) in enumerate(flow):
        cy = 1.5 + i * 0.82
        add_rect(slide, 7.05, cy, 5.9, 0.68, CARD)
        add_rect(slide, 7.05, cy, 0.08, 0.68, color)
        add_text(slide, f"{'①②③④⑤⑥⑦'[i]}  {label}", 7.28, cy + 0.12, 5.5, 0.44,
                 size=12, color=WHITE)

    add_rect(slide, 7.05, 7.2, 5.9, 0.22, DARK2)
    add_text(slide, "💡  .sheet(item:) ดีกว่า .sheet(isPresented:) — ส่ง item ตรงเข้า closure",
             7.15, 7.22, 5.7, 0.2, size=9, italic=True, color=GREY)

    slide_number(slide, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Finance Tab Complete
# ════════════════════════════════════════════════════════════════════════════
def slide_finance_tab(prs):
    slide = add_slide(prs)
    header_bar(slide, "💰  Finance Tab · Complete View")

    add_code(slide, [
        "struct FinanceTabView: View {",
        "    // ✅ iOS 13+ patterns",
        "    @EnvironmentObject private var viewModel: FarmViewModel",
        "    @Environment(\\.managedObjectContext)",
        "        private var viewContext",
        "    @State private var filterType = \"all\"",
        "    @State private var selectedTransaction:",
        "                           Transaction? = nil",
        "",
        "    @FetchRequest(entity: Transaction.entity(),",
        "        sortDescriptors: [NSSortDescriptor(",
        "            keyPath: \\Transaction.date,",
        "            ascending: false)]",
        "    ) var allTransactions: FetchedResults<Transaction>",
        "",
        "    var body: some View {",
        "        NavigationView {          // ✅ iOS 13+",
        "            VStack(spacing: 0) {",
        "                // Summary cards",
        "                HStack {",
        "                    SummaryCard(title: \"ចំណូល\",",
        "                        amount: totalIncome, color: .green)",
        "                    SummaryCard(title: \"ចំណាយ\",",
        "                        amount: totalExpense, color: .red)",
        "                    SummaryCard(title: \"សមតុល្យ\",",
        "                        amount: balance, color: .blue)",
        "                }",
        "                // Filter Picker + List",
        "                Picker(\"Filter\", selection: $filterType) {...}",
        "            }",
        "        }",
        "    }",
        "}",
    ], 0.35, 1.1, 7.0, 6.25, size=9.5)

    # summary card design
    add_text(slide, "SummaryCard UI", 7.55, 1.1, 5.5, 0.4,
             size=13, bold=True, color=GREEN)

    cards = [
        (GREEN, "ចំណូល", "$2,400"),
        (RED,   "ចំណាយ", "$1,100"),
        (BLUE,  "សមតុល្យ", "$1,300"),
    ]
    for i, (color, label, val) in enumerate(cards):
        cx = 7.55 + i * 1.85
        add_rect(slide, cx, 1.65, 1.7, 1.2, CARD)
        add_rect(slide, cx, 1.65, 1.7, 0.08, color)
        add_text(slide, label, cx + 0.1, 1.8, 1.5, 0.35,
                 size=11, color=GREY)
        add_text(slide, val, cx + 0.1, 2.18, 1.5, 0.5,
                 size=16, bold=True, color=color)

    # filter picker mockup
    add_text(slide, "Filter Picker", 7.55, 3.05, 5.5, 0.35,
             size=12, bold=True, color=BLUE)
    add_rect(slide, 7.55, 3.48, 5.5, 0.6, DARK2)
    for j, (lbl, fw) in enumerate([("ទាំងអស់", True), ("ចំណូល", False), ("ចំណាយ", False)]):
        bx = 7.65 + j * 1.8
        add_rect(slide, bx, 3.55, 1.6, 0.45,
                 GREEN if fw else DARK)
        add_text(slide, lbl, bx, 3.57, 1.6, 0.41,
                 size=11, bold=fw, color=DARK if fw else GREY,
                 align=PP_ALIGN.CENTER)

    add_bullets(slide, [
        "SummaryCard ប្រើ @EnvironmentObject ដើម្បី format currency",
        ".sheet(item: $selectedTransaction) — open edit form",
        "@FetchRequest ២ — total calculation + filtered list",
    ], 7.55, 4.3, 5.5, 2.8, size=12, bullet_color=BLUE)

    slide_number(slide, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Common Mistakes
# ════════════════════════════════════════════════════════════════════════════
def slide_mistakes(prs):
    slide = add_slide(prs)
    header_bar(slide, "⚠️  Common Mistakes · ❌ vs ✅")

    rows = [
        ("ViewModel pattern",   "@Observable class VM",          "class VM: ObservableObject"),
        ("State ViewModel",     "@State var vm: VM",              "@StateObject var vm: VM"),
        ("Read ViewModel",      "@Environment(VM.self)",          "@EnvironmentObject var vm"),
        ("Pass ViewModel",      ".environment(vm)",               ".environmentObject(vm)"),
        ("Navigation",          "NavigationStack {}",             "NavigationView {}"),
        ("Dismiss sheet",       "@Environment(\\.dismiss)",       "@Environment(\\.presentationMode)"),
        ("Pre-fill @State",     "@State var x = \"\"",            "@State var x: String  +  init"),
        ("2 containers",        "NSPersistentContainer() in App", "CoreDataManager.shared.context"),
    ]

    headers = ["Pattern", "❌  iOS 17+ Only", "✅  iOS 13+ Correct"]
    for c, (h, cw, cs) in enumerate(zip(headers, [2.8, 4.5, 5.1], [0.35, 3.3, 7.95])):
        add_rect(slide, cs, 1.1, cw, 0.55, BLUE if c == 2 else RED if c == 1 else DARK2)
        add_text(slide, h, cs + 0.1, 1.15, cw - 0.15, 0.45,
                 size=12, bold=True, color=WHITE)

    for r, (pattern, wrong, right) in enumerate(rows):
        cy = 1.75 + r * 0.67
        bg = CARD if r % 2 == 0 else DARK2
        for c, (cell, cw, cs) in enumerate(
                zip([pattern, wrong, right], [2.8, 4.5, 5.1], [0.35, 3.3, 7.95])):
            add_rect(slide, cs, cy, cw, 0.62, bg)
            col = GREY if c == 0 else (RED if c == 1 else GREEN)
            add_text(slide, cell, cs + 0.1, cy + 0.1, cw - 0.15, 0.42,
                     size=10, color=col)

    slide_number(slide, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary & What's Next
# ════════════════════════════════════════════════════════════════════════════
def slide_summary(prs):
    slide = add_slide(prs)
    header_bar(slide, "📌  សង្ខេប Week 02 & What's Next", accent=BLUE)

    cols = [
        ("✅  Week 02 Built", BLUE, [
            "CoreDataManager singleton",
            "4 Entity models defined",
            "@FetchRequest in views",
            "CRUD for Transaction",
            "Finance Tab complete",
        ]),
        ("📋  Key Rules", GREEN, [
            "ONE NSPersistentContainer only",
            "saveContext() after every write",
            "NSPredicate → db-level filter",
            "State(initialValue:) for edit forms",
            ".sheet(item:) for edit sheets",
        ]),
        ("🎯  Week 03 Preview", PURPLE, [
            "NavigationView & NavigationLink",
            "Detail screen (Transaction detail)",
            "Navigation state management",
            "Deep linking simulation",
            "NavigationCoordinator pattern",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        cx = 0.35 + i * 4.35
        add_rect(slide, cx, 1.1, 4.1, 5.7, CARD)
        add_rect(slide, cx, 1.1, 4.1, 0.6, color)
        add_text(slide, title, cx + 0.15, 1.15, 3.8, 0.5,
                 size=14, bold=True, color=DARK)
        for j, item in enumerate(items):
            cy_item = 1.82 + j * 0.9
            add_rect(slide, cx + 0.15, cy_item, 3.7, 0.72, DARK2)
            add_text(slide, item, cx + 0.3, cy_item + 0.1, 3.4, 0.52,
                     size=12, color=WHITE)

    add_rect(slide, 0.35, 6.98, 12.6, 0.42, BLUE)
    add_text(slide, "💾  Week 02 ចប់! — Data ត្រូវបានរក្សា & App ស្វែងហា! · Ready for Navigation in Week 03!",
             0.5, 7.01, 12.3, 0.38,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number(slide, 12)


# ── BUILD ────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    slide_title(prs)
    slide_agenda(prs)
    slide_comparison(prs)
    slide_coredata_setup(prs)
    slide_entities(prs)
    slide_managed_object(prs)
    slide_fetchrequest(prs)
    slide_crud_cd(prs)
    slide_crud_update(prs)
    slide_finance_tab(prs)
    slide_mistakes(prs)
    slide_summary(prs)

    out = os.path.join(os.path.dirname(__file__), "Week02_CoreData.pptx")
    prs.save(out)
    print(f"✅  Saved → {out}")


if __name__ == "__main__":
    build()
