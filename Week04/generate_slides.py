#!/usr/bin/env python3
"""Generate Week 04 — Finance Module Complete slide deck (Khmer)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

PURPLE = RGBColor(0x8E,0x44,0xAD)
GREEN  = RGBColor(0x1B,0xB8,0x89)
BLUE   = RGBColor(0x28,0x7D,0xFA)
ORANGE = RGBColor(0xF3,0x96,0x20)
TEAL   = RGBColor(0x00,0xC9,0xC8)
DARK   = RGBColor(0x1A,0x1A,0x2E)
DARK2  = RGBColor(0x16,0x21,0x3E)
CARD   = RGBColor(0x0F,0x2A,0x45)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GREY   = RGBColor(0xAA,0xAA,0xBB)
RED    = RGBColor(0xE5,0x47,0x47)
YELLOW = RGBColor(0xFF,0xD7,0x00)
W,H    = 13.33,7.5

def add_slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    for ph in s.placeholders: ph._element.getparent().remove(ph._element)
    f=s.background.fill; f.solid(); f.fore_color.rgb=DARK; return s

def add_text(s,text,x,y,w,h,size=16,bold=False,italic=False,
             color=WHITE,align=PP_ALIGN.LEFT,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; return tb

def add_rect(s,x,y,w,h,fc,lc=None):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if lc: sh.line.color.rgb=lc
    else: sh.line.fill.background(); return sh

def add_bullets(s,items,x,y,w,h,size=14,color=WHITE,bc=GREEN):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r1=p.add_run(); r1.text="• "; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=bc
        r2=p.add_run(); r2.text=item; r2.font.size=Pt(size); r2.font.color.rgb=color

def add_code(s,lines,x,y,w,h,size=10):
    add_rect(s,x,y,w,h,CARD)
    tb=s.shapes.add_textbox(Inches(x+.18),Inches(y+.15),Inches(w-.36),Inches(h-.3))
    tf=tb.text_frame; tf.word_wrap=False
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.color.rgb=TEAL; r.font.name="Courier New"

def hdr(s,title,accent=GREEN):
    add_rect(s,0,0,W,.9,DARK2); add_rect(s,0,0,.06,.9,accent)
    add_text(s,title,.25,.1,12.5,.7,size=22,bold=True,color=accent)

def snum(s,n,total=12):
    add_text(s,f"{n}/{total}",12.3,7.1,.9,.3,size=10,color=GREY,align=PP_ALIGN.RIGHT)

# ── Slide 1 — Title ──────────────────────────────────────────────────────────
def s1(prs):
    s=add_slide(prs)
    add_rect(s,0,0,.12,H,GREEN)
    add_rect(s,.4,.4,2.2,.42,GREEN)
    add_text(s,"Week 04 · SmartFarmer",.4,.4,2.2,.42,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(s,"ប្រព័ន្ធហិរញ្ញវត្ថុ ពេញលេញ",.4,1.0,12.0,1.0,size=38,bold=True,color=GREEN)
    add_text(s,"Finance Module Complete",.4,1.95,12.0,.7,size=28,bold=True,color=WHITE)
    add_text(s,"iOS 13+  ·  CoreData  ·  FarmViewModel  ·  CRUD  ·  Categories  ·  Real-time Totals",.4,2.85,12.0,.45,size=14,italic=True,color=GREY)
    add_rect(s,.4,3.45,10.5,.04,GREEN)
    icons=[("💰","Transactions"),("➕","Add Form"),("✏️","Edit Form"),("📋","Detail View"),("🗂️","Categories"),("📊","Live Totals")]
    for i,(e,l) in enumerate(icons):
        cx=.4+i*2.1; add_rect(s,cx,3.65,1.9,1.0,CARD)
        add_text(s,e,cx,3.7,1.9,.45,size=20,align=PP_ALIGN.CENTER)
        add_text(s,l,cx,4.18,1.9,.4,size=10,color=GREY,align=PP_ALIGN.CENTER)
    add_text(s,"SmartFarmer Assistant · ភ្នំពេញ 2026",.4,6.9,12.0,.4,size=11,italic=True,color=GREY)
    snum(s,1)

# ── Slide 2 — Agenda ─────────────────────────────────────────────────────────
def s2(prs):
    s=add_slide(prs)
    hdr(s,"📋  មាតិកា Week 04 · Agenda")
    topics=[
        (GREEN, "01","FarmViewModel","CRUD methods + formatting helpers"),
        (BLUE,  "02","AddTransactionView","Form: amount / type / category / note"),
        (ORANGE,"03","Category System","expense vs income categories — Picker adapts"),
        (PURPLE,"04","EditTransactionView","pre-filled Form · update existing record"),
        (TEAL,  "05","Real-time Totals","income / expense / balance → live @FetchRequest"),
        (GREEN, "06","Mini-Project","Finance module ពេញ · CRUD + filter + totals"),
    ]
    for i,(c,n,kh,en) in enumerate(topics):
        col=i%2; row=i//2; cx=.35+col*6.5; cy=1.1+row*2.0
        add_rect(s,cx,cy,6.1,1.75,CARD); add_rect(s,cx,cy,.08,1.75,c)
        add_rect(s,cx+.18,cy+.38,.6,.6,c)
        add_text(s,n,cx+.18,cy+.38,.6,.6,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        add_text(s,kh,cx+.95,cy+.25,4.9,.5,size=15,bold=True,color=WHITE)
        add_text(s,en,cx+.95,cy+.85,4.9,.45,size=11,color=GREY)
    snum(s,2)

# ── Slide 3 — FarmViewModel ───────────────────────────────────────────────────
def s3(prs):
    s=add_slide(prs)
    hdr(s,"🧠  FarmViewModel · CRUD + Formatting Helpers")
    add_code(s,[
        "class FarmViewModel: ObservableObject {",
        "    private let manager = CoreDataManager.shared",
        "",
        "    // CREATE",
        "    func addTransaction(amount: Double, note: String,",
        "                        type: String, category: String) {",
        "        manager.addTransaction(amount: amount, note: note,",
        "                               type: type, category: category)",
        "    }",
        "",
        "    // UPDATE",
        "    func updateTransaction(_ t: Transaction, amount: Double,",
        "                           note: String, type: String,",
        "                           category: String) {",
        "        t.amount = amount; t.note = note",
        "        t.type = type;     t.category = category",
        "        manager.save()",
        "    }",
        "",
        "    // DELETE",
        "    func deleteTransaction(_ t: Transaction) {",
        "        manager.context.delete(t); manager.save()",
        "    }",
        "",
        "    // FORMATTING",
        "    func formatCurrency(_ amount: Double) -> String {",
        '        return String(format: "%.2f ៛", amount)',
        "    }",
        "    func formatDate(_ date: Date?) -> String {",
        '        guard let d = date else { return "" }',
        "        let f = DateFormatter()",
        '        f.dateStyle = .medium; f.locale = Locale(identifier: "km_KH")',
        "        return f.string(from: d)",
        "    }",
        "}",
    ],.35,1.1,7.5,6.25,size=9)
    notes=[
        (GREEN, "Single ViewModel","one ObservableObject\nfor all Finance operations\n→ injected via .environmentObject"),
        (BLUE,  "CoreDataManager","shared singleton\n.save() persists changes\n@FetchRequest auto-refreshes"),
        (ORANGE,"formatCurrency","%.2f ៛ — Khmer Riel\nformat amount consistently\nacross all views"),
        (PURPLE,"formatDate","DateFormatter km_KH\nKhmer locale date\n→ farmer-friendly"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,8.0,cy,5.1,1.38,CARD); add_rect(s,8.0,cy,.08,1.38,c)
        add_text(s,t,8.22,cy+.1,4.6,.4,size=12,bold=True,color=c)
        add_text(s,d,8.22,cy+.55,4.6,.72,size=11,color=WHITE)
    snum(s,3)

# ── Slide 4 — AddTransactionView ──────────────────────────────────────────────
def s4(prs):
    s=add_slide(prs)
    hdr(s,"➕  AddTransactionView · Form បញ្ចូលប្រតិបត្តិការ")
    add_code(s,[
        "struct AddTransactionView: View {",
        "    @Environment(\\.presentationMode) var pm",
        "    @EnvironmentObject private var viewModel: FarmViewModel",
        "    @State private var amountText = \"\"",
        "    @State private var note      = \"\"",
        "    @State private var type      = \"expense\"",
        "    @State private var category  = \"Seeds\"",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            Form {",
        '                Section(header: Text("ចំនួនទឹកប្រាក់")) {',
        "                    TextField(\"0.00\", text: $amountText)",
        "                        .keyboardType(.decimalPad)",
        "                }",
        '                Section(header: Text("ប្រភេទ")) {',
        "                    Picker(\"ប្រភេទ\", selection: $type) {",
        "                        Text(\"ចំណាយ\").tag(\"expense\")",
        "                        Text(\"ចំណូល\").tag(\"income\")",
        "                    }.pickerStyle(SegmentedPickerStyle())",
        "                }",
        '                Section(header: Text("ប្រភេទរង")) {',
        "                    Picker(\"ប្រភេទរង\", selection: $category) {",
        "                        ForEach(categories, id: \\.self) {",
        "                            Text($0).tag($0)",
        "                        }",
        "                    }",
        "                }",
        '                Section(header: Text("កំណត់ចំណាំ")) {',
        "                    TextField(\"optional\", text: $note)",
        "                }",
        "            }",
        '            .navigationTitle("ប្រតិបត្តិការថ្មី")',
        "            .navigationBarItems(",
        "                leading: Button(\"បោះបង់\") {",
        "                    pm.wrappedValue.dismiss() },",
        "                trailing: Button(\"រក្សាទុក\") { save() }",
        "            )",
        "        }",
        "    }",
        "}",
    ],.35,1.1,7.5,6.25,size=8.8)
    notes=[
        (GREEN, "Form sections","amount / type / category / note\n→ clean grouped UI\n→ FormStyle automatic"),
        (BLUE,  "Segmented Picker","expense vs income\n→ SegmentedPickerStyle\n→ drives category list"),
        (ORANGE,".keyboardType","decimalPad for amount\n→ numeric-only keyboard\n→ user-friendly"),
        (PURPLE,"Save action","Double(amountText)\nviewModel.addTransaction()\npm.wrappedValue.dismiss()"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,8.0,cy,5.1,1.38,CARD); add_rect(s,8.0,cy,.08,1.38,c)
        add_text(s,t,8.22,cy+.1,4.6,.4,size=12,bold=True,color=c)
        add_text(s,d,8.22,cy+.55,4.6,.72,size=11,color=WHITE)
    snum(s,4)

# ── Slide 5 — Category System ─────────────────────────────────────────────────
def s5(prs):
    s=add_slide(prs)
    hdr(s,"🗂️  Category System · Picker Adapts to Type")
    add_code(s,[
        "// Categories adapt based on transaction type",
        "var categories: [String] {",
        "    if type == \"expense\" {",
        "        return [\"Seeds\", \"Fertiliser\",",
        "                \"Equipment\", \"Labour\", \"Other\"]",
        "    } else {",
        "        return [\"Crop Sales\", \"Livestock\",",
        "                \"Subsidy\", \"Other\"]",
        "    }",
        "}",
        "",
        "// Reset category when type changes",
        ".onChange(of: type) { _ in",
        "    category = categories.first ?? \"Other\"",
        "}",
        "",
        "// Transaction stored in CoreData",
        "// transaction.type     = \"expense\" | \"income\"",
        "// transaction.category = \"Seeds\" | \"Crop Sales\" | ...",
        "// transaction.isExpense: Bool {",
        "//     return type == \"expense\"",
        "// }",
    ],.35,1.1,6.5,5.5,size=10.5)
    add_text(s,"Expense Categories",0.35,6.7,3.1,.4,size=11,bold=True,color=RED)
    exp=["Seeds · គ្រាប់ពូជ","Fertiliser · ជី","Equipment · ឧបករណ៍","Labour · ពលករ","Other · ផ្សេងៗ"]
    for i,e in enumerate(exp):
        add_rect(s,.35,7.15+i*.0,3.1,.0,CARD)
    add_code(s,exp,.35,7.15,3.1,.0,size=10)

    add_text(s,"Income Categories",3.6,6.7,3.0,.4,size=11,bold=True,color=GREEN)
    inc=["Crop Sales · លក់ដំណាំ","Livestock · សត្វ","Subsidy · ជំនួយ","Other · ផ្សេងៗ"]

    exp_box=[
        ("Seeds",     "គ្រាប់ពូជ"),
        ("Fertiliser","ជី"),
        ("Equipment", "ឧបករណ៍"),
        ("Labour",    "ពលករ"),
        ("Other",     "ផ្សេងៗ"),
    ]
    inc_box=[
        ("Crop Sales","លក់ដំណាំ"),
        ("Livestock", "សត្វ"),
        ("Subsidy",   "ជំនួយ"),
        ("Other",     "ផ្សេងៗ"),
    ]
    add_text(s,"💸 Expense",7.0,1.1,2.8,.45,size=14,bold=True,color=RED)
    for i,(en,kh) in enumerate(exp_box):
        cy=1.65+i*.98
        add_rect(s,7.0,cy,5.8,.82,CARD); add_rect(s,7.0,cy,.08,.82,RED)
        add_text(s,en,7.18,cy+.08,3.0,.35,size=12,bold=True,color=WHITE)
        add_text(s,kh,7.18,cy+.46,3.0,.3,size=10,color=GREY)
        add_text(s,"expense",10.15,cy+.2,2.5,.42,size=10,color=RED,align=PP_ALIGN.RIGHT)

    add_text(s,"💰 Income",7.0,6.8,2.8,.45,size=14,bold=True,color=GREEN)
    # show income as horizontal row
    notes2=[
        (GREEN,"Crop Sales","លក់ដំណាំ"),
        (GREEN,"Livestock","សត្វ"),
        (GREEN,"Subsidy","ជំនួយ"),
        (GREEN,"Other","ផ្សេងៗ"),
    ]
    # just list them
    snum(s,5)

# ── Slide 5 override — redo cleanly ──────────────────────────────────────────
def s5_clean(prs):
    s=add_slide(prs)
    hdr(s,"🗂️  Category System · Picker Adapts to Type")
    add_code(s,[
        "var categories: [String] {",
        "    if type == \"expense\" {",
        "        return [\"Seeds\", \"Fertiliser\",",
        "                \"Equipment\", \"Labour\", \"Other\"]",
        "    } else {",
        "        return [\"Crop Sales\", \"Livestock\",",
        "                \"Subsidy\", \"Other\"]",
        "    }",
        "}",
        "",
        "// Reset when type changes",
        ".onChange(of: type) { _ in",
        "    category = categories.first ?? \"Other\"",
        "}",
    ],.35,1.1,6.3,3.8,size=10.5)
    exp_items=[("Seeds","គ្រាប់ពូជ"),("Fertiliser","ជី"),("Equipment","ឧបករណ៍"),("Labour","ពលករ"),("Other","ផ្សេងៗ")]
    add_text(s,"💸 ចំណាយ — Expense",.35,5.05,6.3,.45,size=13,bold=True,color=RED)
    for i,(en,kh) in enumerate(exp_items):
        col=i%3; row=i//3; cx=.35+col*2.1; cy=5.6+row*.75
        add_rect(s,cx,cy,1.95,.62,CARD); add_rect(s,cx,cy,.07,.62,RED)
        add_text(s,en,cx+.15,cy+.05,1.7,.28,size=10,bold=True,color=WHITE)
        add_text(s,kh,cx+.15,cy+.35,1.7,.22,size=9,color=GREY)
    inc_items=[("Crop Sales","លក់ដំណាំ"),("Livestock","សត្វ"),("Subsidy","ជំនួយ"),("Other","ផ្សេងៗ")]
    add_text(s,"💰 ចំណូល — Income",6.8,1.1,6.2,.45,size=13,bold=True,color=GREEN)
    for i,(en,kh) in enumerate(inc_items):
        cy=1.65+i*.98
        add_rect(s,6.8,cy,6.2,.82,CARD); add_rect(s,6.8,cy,.08,.82,GREEN)
        add_text(s,en,7.02,cy+.08,4.5,.35,size=13,bold=True,color=WHITE)
        add_text(s,kh,7.02,cy+.46,4.5,.3,size=11,color=GREY)
    notes=[
        (GREEN, "Computed var","categories changes\nwhen type changes\n→ dynamic Picker list"),
        (ORANGE,".onChange","reset category\nwhen switching type\n→ avoid invalid combo"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=5.6+i*1.0
        add_rect(s,6.8,cy,6.2,.88,CARD); add_rect(s,6.8,cy,.08,.88,c)
        add_text(s,t,7.02,cy+.08,5.8,.32,size=12,bold=True,color=c)
        add_text(s,d,7.02,cy+.45,5.8,.38,size=10,color=WHITE)
    snum(s,5)
    return s

# ── Slide 6 — EditTransactionView ─────────────────────────────────────────────
def s6(prs):
    s=add_slide(prs)
    hdr(s,"✏️  EditTransactionView · Pre-filled Form")
    add_code(s,[
        "struct EditTransactionView: View {",
        "    @Environment(\\.presentationMode) var pm",
        "    @EnvironmentObject private var viewModel: FarmViewModel",
        "    let transaction: Transaction",
        "",
        "    // Pre-fill @State from existing record",
        "    @State private var amountText: String",
        "    @State private var note: String",
        "    @State private var type: String",
        "    @State private var category: String",
        "",
        "    init(transaction: Transaction) {",
        "        self.transaction = transaction",
        "        _amountText = State(initialValue:",
        "            String(transaction.amount))",
        "        _note     = State(initialValue: transaction.note ?? \"\")",
        "        _type     = State(initialValue: transaction.type ?? \"expense\")",
        "        _category = State(initialValue:",
        "            transaction.category ?? \"Other\")",
        "    }",
        "",
        "    func save() {",
        "        guard let amount = Double(amountText) else { return }",
        "        viewModel.updateTransaction(transaction,",
        "            amount: amount, note: note,",
        "            type: type, category: category)",
        "        pm.wrappedValue.dismiss()",
        "    }",
        "}",
    ],.35,1.1,7.5,6.25,size=9)
    notes=[
        (GREEN, "Pre-fill with init","_amountText = State(initialValue:)\n→ existing value shown\n→ user edits from there"),
        (BLUE,  "Same Form layout","reuse same fields as Add\n→ consistent UX\n→ only init differs"),
        (ORANGE,"updateTransaction","viewModel updates CoreData\n→ @FetchRequest refreshes\n→ detail view updates"),
        (PURPLE,"dismiss()","pm.wrappedValue.dismiss()\n→ close edit sheet\n→ return to detail view"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,8.0,cy,5.1,1.38,CARD); add_rect(s,8.0,cy,.08,1.38,c)
        add_text(s,t,8.22,cy+.1,4.6,.4,size=12,bold=True,color=c)
        add_text(s,d,8.22,cy+.55,4.6,.72,size=11,color=WHITE)
    snum(s,6)

# ── Slide 7 — Real-time Totals ────────────────────────────────────────────────
def s7(prs):
    s=add_slide(prs)
    hdr(s,"📊  Real-time Totals · ចំណូល / ចំណាយ / នៅសល់")
    add_code(s,[
        "struct FinanceTabView: View {",
        "    @FetchRequest(",
        "        entity: Transaction.entity(),",
        "        sortDescriptors: [NSSortDescriptor(",
        "            keyPath: \\Transaction.date,",
        "            ascending: false)]",
        "    ) var transactions: FetchedResults<Transaction>",
        "",
        "    // Computed live from @FetchRequest",
        "    var totalIncome: Double {",
        "        transactions",
        "            .filter { $0.type == \"income\" }",
        "            .reduce(0) { $0 + $1.amount }",
        "    }",
        "    var totalExpense: Double {",
        "        transactions",
        "            .filter { $0.type == \"expense\" }",
        "            .reduce(0) { $0 + $1.amount }",
        "    }",
        "    var balance: Double { totalIncome - totalExpense }",
        "",
        "    var body: some View {",
        "        HStack {",
        "            SummaryCard(title: \"ចំណូល\",",
        "                        amount: totalIncome, color: .green)",
        "            SummaryCard(title: \"ចំណាយ\",",
        "                        amount: totalExpense, color: .red)",
        "            SummaryCard(title: \"នៅសល់\",",
        "                        amount: balance,",
        "                        color: balance >= 0 ? .green : .red)",
        "        }",
        "    }",
        "}",
    ],.35,1.1,7.5,6.25,size=9)
    notes=[
        (GREEN, "Live @FetchRequest","CoreData change → @FetchRequest\n→ computed vars recalculate\n→ cards update instantly"),
        (BLUE,  ".filter { }","filter by type string\n→ \"income\" vs \"expense\"\n→ no extra fetch needed"),
        (ORANGE,".reduce(0)","sum all amounts\n→ one-liner total\n→ always in sync"),
        (PURPLE,"balance color","green if >= 0\nred if negative\n→ visual alert for farmer"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,8.0,cy,5.1,1.38,CARD); add_rect(s,8.0,cy,.08,1.38,c)
        add_text(s,t,8.22,cy+.1,4.6,.4,size=12,bold=True,color=c)
        add_text(s,d,8.22,cy+.55,4.6,.72,size=11,color=WHITE)
    snum(s,7)

# ── Slide 8 — Delete + Filter ─────────────────────────────────────────────────
def s8(prs):
    s=add_slide(prs)
    hdr(s,"🗑️  Delete + Filter · Swipe to Delete + Segmented Filter")
    add_code(s,[
        "// Delete — swipe left on row",
        "List {",
        "    ForEach(filteredTransactions, id: \\.self) { t in",
        "        NavigationLink(destination: TransactionDetailView(",
        "            transaction: t)) {",
        "            TransactionRowView(transaction: t,",
        "                               viewModel: viewModel)",
        "        }",
        "    }",
        "    .onDelete(perform: deleteTransactions)",
        "}",
        "",
        "func deleteTransactions(at offsets: IndexSet) {",
        "    offsets.forEach { idx in",
        "        viewModel.deleteTransaction(",
        "            filteredTransactions[idx])",
        "    }",
        "}",
        "",
        "// Filter — segmented picker at top",
        "@State private var filterType = \"All\"",
        "",
        "var filteredTransactions: [Transaction] {",
        "    switch filterType {",
        "    case \"Income\":  return transactions.filter { $0.type == \"income\" }",
        "    case \"Expense\": return transactions.filter { $0.type == \"expense\" }",
        "    default:        return Array(transactions)",
        "    }",
        "}",
        "",
        "Picker(\"Filter\", selection: $filterType) {",
        "    Text(\"ទាំងអស់\").tag(\"All\")",
        "    Text(\"ចំណូល\").tag(\"Income\")",
        "    Text(\"ចំណាយ\").tag(\"Expense\")",
        "}.pickerStyle(SegmentedPickerStyle())",
    ],.35,1.1,7.5,6.25,size=9)
    notes=[
        (GREEN, ".onDelete","swipe left → delete\n→ IndexSet → viewModel\n→ auto removes from list"),
        (BLUE,  "filteredTransactions","computed from @FetchRequest\n→ filter + All modes\n→ live with CoreData"),
        (ORANGE,"SegmentedPickerStyle","All / Income / Expense\n→ top of Finance tab\n→ familiar iOS pattern"),
        (PURPLE,"@State filterType","drives filteredTransactions\nchanges → view re-renders\n→ no extra fetch"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,8.0,cy,5.1,1.38,CARD); add_rect(s,8.0,cy,.08,1.38,c)
        add_text(s,t,8.22,cy+.1,4.6,.4,size=12,bold=True,color=c)
        add_text(s,d,8.22,cy+.55,4.6,.72,size=11,color=WHITE)
    snum(s,8)

# ── Slide 9 — Project Structure ───────────────────────────────────────────────
def s9(prs):
    s=add_slide(prs)
    hdr(s,"🏗️  Project Structure · Finance Module Files")
    tree=[
        "SmartFarmerAssistant/",
        "├── ViewModels/",
        "│   ├── FarmViewModel.swift         # CRUD + formatCurrency/Date",
        "│   └── FinanceCoordinator.swift    # navigation state",
        "└── Views/",
        "    ├── MainTabView.swift           # TabView + @StateObject owners",
        "    ├── FinanceTabView.swift        # summary cards + filter + list",
        "    ├── AddTransactionView.swift    # sheet: add new",
        "    ├── EditTransactionView.swift   # sheet: edit existing",
        "    ├── TransactionDetailView.swift # push: full detail",
        "    ├── TransactionRowView.swift    # list row",
        "    └── CalendarTabView.swift       # farm activities",
    ]
    add_code(s,tree,.35,1.1,12.6,3.0,size=11)
    screens=[
        (GREEN, "FinanceTabView","Summary cards + segmented filter\n+ filterable list + add button"),
        (BLUE,  "AddTransactionView","Sheet · amount, type,\ncategory, note form"),
        (ORANGE,"EditTransactionView","Sheet · same form\npre-filled with existing data"),
        (PURPLE,"TransactionDetailView","Push · full detail view\n+ Edit button → EditSheet"),
        (TEAL,  "TransactionRowView","List row · amount color\ntype badge + date"),
    ]
    add_text(s,"Screens at a Glance",.35,4.25,12.6,.4,size=13,bold=True,color=GREEN)
    for i,(c,t,d) in enumerate(screens):
        col=i%3; row=i//3; cx=.35+col*4.35; cy=4.75+row*1.3
        add_rect(s,cx,cy,4.1,1.15,CARD); add_rect(s,cx,cy,.08,1.15,c)
        add_text(s,t,cx+.22,cy+.08,3.7,.38,size=12,bold=True,color=c)
        add_text(s,d,cx+.22,cy+.52,3.7,.55,size=10,color=GREY)
    snum(s,9)

# ── Slide 10 — Data Flow ──────────────────────────────────────────────────────
def s10(prs):
    s=add_slide(prs)
    hdr(s,"🔄  Data Flow · Full Transaction Lifecycle")
    flow=[
        (GREEN, "User taps +","AddTransactionView opens as sheet"),
        (BLUE,  "Fill form","amount, type, category, note → tap រក្សាទុក"),
        (ORANGE,"viewModel.addTransaction()","CoreData creates + saves new Transaction"),
        (PURPLE,"@FetchRequest refreshes","FinanceTabView list + totals update instantly"),
        (TEAL,  "Tap row","NavigationLink → TransactionDetailView pushed"),
        (GREEN, "Tap Edit → save","viewModel.updateTransaction() → detail refreshes"),
    ]
    for i,(c,kh,en) in enumerate(flow):
        col=i%2; row=i//2; cx=.35+col*6.5; cy=1.1+row*2.0
        add_rect(s,cx,cy,6.1,1.75,CARD); add_rect(s,cx,cy,.08,1.75,c)
        add_rect(s,cx+.15,cy+.4,.55,.55,c)
        add_text(s,str(i+1),cx+.15,cy+.4,.55,.55,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        add_text(s,kh,cx+.88,cy+.2,5.0,.5,size=14,bold=True,color=WHITE)
        add_text(s,en,cx+.88,cy+.82,5.0,.6,size=11,color=GREY)
    snum(s,10)

# ── Slide 11 — Mini-Project ───────────────────────────────────────────────────
def s11(prs):
    s=add_slide(prs)
    hdr(s,"🏠  Mini-Project · Finance Module ពេញ")
    cols=[
        ("📋  Requirements",GREEN,[
            "Summary cards: income/expense/balance",
            "Add Transaction form (4 fields)",
            "Category Picker adapts to type",
            "Edit Transaction — pre-filled form",
            "Transaction Detail — full view",
            "Swipe to delete + segmented filter",
        ]),
        ("✅  Checklist",BLUE,[
            "FarmViewModel CRUD methods ✓",
            "addTransaction() works ✓",
            "updateTransaction() reflects ✓",
            "deleteTransaction() removes ✓",
            "Real-time totals update ✓",
            "Filter: All / Income / Expense ✓",
        ]),
        ("🎯  Grading",ORANGE,[
            "Add Transaction form: 20%",
            "Category system: 15%",
            "Edit + update: 20%",
            "Real-time totals: 20%",
            "Delete + filter: 15%",
            "Bonus: balance color alert",
        ]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*4.35
        add_rect(s,cx,1.1,4.1,5.8,CARD); add_rect(s,cx,1.1,4.1,.55,c)
        add_text(s,t,cx+.12,1.15,3.85,.45,size=13,bold=True,color=WHITE)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.15,cy,3.7,.72,DARK2)
            add_text(s,item,cx+.28,cy+.1,3.4,.52,size=11,color=WHITE)
    add_rect(s,.35,7.0,12.6,.42,GREEN)
    add_text(s,"💰  Week 04 ចប់! — Finance module ពេញ ready · Next: Calendar & Reminders →",
             .5,7.02,12.3,.38,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    snum(s,11)

# ── Slide 12 — Cheat Sheet ───────────────────────────────────────────────────
def s12(prs):
    s=add_slide(prs)
    hdr(s,"📌  សង្ខេប Week 04 · Cheat Sheet",accent=GREEN)
    cols=[
        ("🧠  ViewModel",GREEN,["addTransaction()","updateTransaction()","deleteTransaction()","formatCurrency()"]),
        ("➕  AddView",BLUE,["@State amountText","SegmentedPickerStyle","categories computed","pm.wrappedValue.dismiss()"]),
        ("📊  Totals",ORANGE,["@FetchRequest live","totalIncome .filter","totalExpense .filter","balance = in - exp"]),
        ("🗑️  Delete/Filter",PURPLE,[".onDelete(perform:)","IndexSet → delete","filterType @State","filteredTransactions"]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*3.25
        add_rect(s,cx,1.1,3.0,5.8,CARD); add_rect(s,cx,1.1,3.0,.55,c)
        add_text(s,t,cx+.1,1.15,2.8,.45,size=12,bold=True,color=WHITE)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.1,cy,2.8,.72,DARK2)
            add_text(s,item,cx+.18,cy+.1,2.6,.52,size=10,color=TEAL)
            s.shapes[-1].text_frame.paragraphs[0].runs[0].font.name="Courier New"
    add_rect(s,.35,7.0,12.6,.42,GREEN)
    add_text(s,"💰  Finance module = CRUD + live totals + filter · @FetchRequest does the heavy lifting!",
             .5,7.02,12.3,.38,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    snum(s,12)

def build():
    prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H)
    s1(prs);s2(prs);s3(prs)
    # slide 4
    s=add_slide(prs)
    hdr(s,"➕  AddTransactionView · Form បញ្ចូលប្រតិបត្តិការ")
    add_code(s,[
        "struct AddTransactionView: View {",
        "    @Environment(\\.presentationMode) var pm",
        "    @EnvironmentObject private var viewModel: FarmViewModel",
        "    @State private var amountText = \"\"",
        "    @State private var note      = \"\"",
        "    @State private var type      = \"expense\"",
        "    @State private var category  = \"Seeds\"",
        "",
        "    var categories: [String] {",
        "        type == \"expense\"",
        "            ? [\"Seeds\",\"Fertiliser\",\"Equipment\",\"Labour\",\"Other\"]",
        "            : [\"Crop Sales\",\"Livestock\",\"Subsidy\",\"Other\"]",
        "    }",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            Form {",
        '                Section(header: Text("ចំនួនទឹកប្រាក់")) {',
        "                    TextField(\"0.00\", text: $amountText)",
        "                        .keyboardType(.decimalPad)",
        "                }",
        '                Section(header: Text("ប្រភេទ")) {',
        "                    Picker(\"ប្រភេទ\", selection: $type) {",
        "                        Text(\"ចំណាយ\").tag(\"expense\")",
        "                        Text(\"ចំណូល\").tag(\"income\")",
        "                    }.pickerStyle(SegmentedPickerStyle())",
        "                }",
        '                Section(header: Text("ប្រភេទរង")) {',
        "                    Picker(\"ប្រភេទរង\", selection: $category) {",
        "                        ForEach(categories, id: \\.self) {",
        "                            Text($0).tag($0)",
        "                        }",
        "                    }",
        "                }",
        '                Section(header: Text("កំណត់ចំណាំ")) {',
        "                    TextField(\"optional\", text: $note)",
        "                }",
        "            }",
        '            .navigationTitle("ប្រតិបត្តិការថ្មី")',
        "        }",
        "    }",
        "}",
    ],.35,1.1,7.5,6.25,size=8.8)
    notes=[
        (GREEN, "Form sections","amount / type / category / note\n→ clean grouped UI"),
        (BLUE,  "Segmented Picker","expense vs income\n→ drives category list"),
        (ORANGE,".keyboardType","decimalPad for amount\n→ numeric-only keyboard"),
        (PURPLE,"categories computed","adapts when type changes\n→ always correct options"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,8.0,cy,5.1,1.38,CARD); add_rect(s,8.0,cy,.08,1.38,c)
        add_text(s,t,8.22,cy+.1,4.6,.4,size=12,bold=True,color=c)
        add_text(s,d,8.22,cy+.55,4.6,.72,size=11,color=WHITE)
    snum(s,4)
    # slide 5 clean
    s5_clean(prs)
    s6(prs);s7(prs);s8(prs);s9(prs);s10(prs);s11(prs);s12(prs)
    out=os.path.join(os.path.dirname(__file__),"Week04_Finance_Complete_KH.pptx")
    prs.save(out); print(f"✅  Saved → {out}")

if __name__=="__main__": build()
