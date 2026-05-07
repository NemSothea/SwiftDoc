#!/usr/bin/env python3
"""Generate Week 07 — Daily Journal slide deck (Khmer)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

PURPLE = RGBColor(0x8E,0x44,0xAD)
GREEN  = RGBColor(0x1B,0xB8,0x89)
BLUE   = RGBColor(0x28,0x7D,0xFA)
ORANGE = RGBColor(0xF3,0x96,0x20)
TEAL   = RGBColor(0x00,0xC9,0xC8)
DARK   = RGBColor(0x1A,0x1A,0x2E)
DARK2  = RGBColor(0x16,0x21,0x3E)
CARD   = RGBColor(0x0F,0x2A,0x45)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GREY   = RGBColor(0xAA,0xAA,0xBB)
RED    = RGBColor(0xE5,0x47,0x47)
YELLOW = RGBColor(0xFF,0xD7,0x00)
W,H    = 13.33,7.5

def add_slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    for ph in s.placeholders: ph._element.getparent().remove(ph._element)
    f=s.background.fill; f.solid(); f.fore_color.rgb=DARK; return s

def add_text(s,text,x,y,w,h,size=16,bold=False,italic=False,
             color=WHITE,align=PP_ALIGN.LEFT,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; return tb

def add_rect(s,x,y,w,h,fc,lc=None):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if lc: sh.line.color.rgb=lc
    else: sh.line.fill.background(); return sh

def add_bullets(s,items,x,y,w,h,size=14,color=WHITE,bc=PURPLE):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r1=p.add_run(); r1.text="• "; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=bc
        r2=p.add_run(); r2.text=item; r2.font.size=Pt(size); r2.font.color.rgb=color

def add_code(s,lines,x,y,w,h,size=10):
    add_rect(s,x,y,w,h,CARD)
    tb=s.shapes.add_textbox(Inches(x+.18),Inches(y+.15),Inches(w-.36),Inches(h-.3))
    tf=tb.text_frame; tf.word_wrap=False
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.color.rgb=TEAL; r.font.name="Courier New"

def hdr(s,title,accent=PURPLE):
    add_rect(s,0,0,W,.9,DARK2); add_rect(s,0,0,.06,.9,accent)
    add_text(s,title,.25,.1,12.5,.7,size=22,bold=True,color=accent)

def snum(s,n,total=12):
    add_text(s,f"{n}/{total}",12.3,7.1,.9,.3,size=10,color=GREY,align=PP_ALIGN.RIGHT)

def s1(prs):
    s=add_slide(prs)
    add_rect(s,0,0,.12,H,PURPLE)
    add_rect(s,.4,.4,2.2,.42,PURPLE)
    add_text(s,"Week 07 · SmartFarmer",.4,.4,2.2,.42,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(s,"កំណត់ហេតុប្រចាំថ្ងៃ",.4,1.0,12.0,1.0,size=44,bold=True,color=PURPLE)
    add_text(s,"Daily Journal · Project 4",.4,1.95,12.0,.7,size=28,bold=True,color=WHITE)
    add_text(s,"iOS 13+  ·  CoreData  ·  Photos  ·  Weather  ·  UIKit Bridge",.4,2.85,12.0,.45,size=14,italic=True,color=GREY)
    add_rect(s,.4,3.45,10.5,.04,PURPLE)
    icons=[("📖","Timeline"),("🌤️","Weather"),("📷","Photos"),("🔍","Search"),("📱","UIKit Bridge"),("🌾","Khmer UI")]
    for i,(e,l) in enumerate(icons):
        cx=.4+i*2.1; add_rect(s,cx,3.65,1.9,1.0,CARD)
        add_text(s,e,cx,3.7,1.9,.45,size=20,align=PP_ALIGN.CENTER)
        add_text(s,l,cx,4.18,1.9,.4,size=10,color=GREY,align=PP_ALIGN.CENTER)
    add_text(s,"SmartFarmer Assistant · ភ្នំពេញ 2026",.4,6.9,12.0,.4,size=11,italic=True,color=GREY)
    snum(s,1)

def s2(prs):
    s=add_slide(prs)
    hdr(s,"📋  មាតិកា Week 07 · Agenda")
    topics=[
        (PURPLE,"01","JournalEntry Entity","7 CoreData attributes + photoData: Data?"),
        (GREEN, "02","Timeline UI","reverse chronological list + date headers"),
        (BLUE,  "03","Weather System","enum WeatherType + SF Symbols"),
        (ORANGE,"04","Photo Picker","UIImagePickerController UIKit bridge"),
        (TEAL,  "05","Search & Filter","text + weather filter combined"),
        (PURPLE,"06","Mini-Project","Journal ពេញ + photos + weather + search"),
    ]
    for i,(c,n,kh,en) in enumerate(topics):
        col=i%2; row=i//2; cx=.35+col*6.5; cy=1.1+row*2.0
        add_rect(s,cx,cy,6.1,1.75,CARD); add_rect(s,cx,cy,.08,1.75,c)
        add_rect(s,cx+.18,cy+.38,.6,.6,c)
        add_text(s,n,cx+.18,cy+.38,.6,.6,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        add_text(s,kh,cx+.95,cy+.25,4.9,.5,size=15,bold=True,color=WHITE)
        add_text(s,en,cx+.95,cy+.85,4.9,.45,size=11,color=GREY)
    snum(s,2)

def s3(prs):
    s=add_slide(prs)
    hdr(s,"🗂️  JournalEntry Entity · CoreData Model")
    attrs=[
        ("id","UUID","Identifiable conformance"),
        ("date","Date?","កាលបរិច្ឆេទ — sort descending"),
        ("content","String?","ខ្លឹមសារ — unlimited text"),
        ("weather","String?","rawValue: sunny/rainy/cloudy/windy"),
        ("photoData","Data?","Binary Data — allows external storage"),
        ("location","String?","ទីកន្លែង optional"),
        ("mood","String?","optional — ស្ថានភាពអារម្មណ៍"),
    ]
    for i,(n,t,d) in enumerate(attrs):
        cy=1.1+i*.75; bg=CARD if i%2==0 else DARK2
        add_rect(s,.35,cy,12.6,.65,bg); add_rect(s,.35,cy,.06,.65,PURPLE)
        add_code(s,[n],.55,cy+.07,3.0,.5,size=12)
        add_code(s,[t],3.75,cy+.07,2.2,.5,size=12)
        add_text(s,d,6.1,cy+.12,6.7,.42,size=11,color=GREY)
    add_rect(s,.35,7.1,12.6,.32,DARK2)
    add_text(s,"⚠️  photoData: Binary Data → check 'Allows External Storage' in Xcode (large images stored outside SQLite)",
             .5,7.12,12.2,.28,size=10,color=GREY)
    snum(s,3)

def s4(prs):
    s=add_slide(prs)
    hdr(s,"📖  Timeline UI · Reverse Chronological List")
    add_bullets(s,[
        "Sort @FetchRequest by date ascending: false → newest first",
        "Section headers: group by Calendar.current.startOfDay(for: entry.date)",
        "JournalRowView: date headline + weather icon + content preview (2 lines)",
        "NavigationLink → JournalDetailView — full content + photo",
        "Photo: if let data = entry.photoData { Image(uiImage: UIImage(data:)) }",
        "Empty state: if entries.isEmpty { Text(\"មិនទាន់មានកំណត់ហេតុ\") }",
    ],.35,1.1,6.3,4.0,size=13,bc=PURPLE)
    add_code(s,[
        "struct JournalTabView: View {",
        "    @FetchRequest(",
        "        entity: JournalEntry.entity(),",
        "        sortDescriptors: [NSSortDescriptor(",
        "            keyPath: \\JournalEntry.date,",
        "            ascending: false)]",
        "    ) var entries: FetchedResults<JournalEntry>",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            List {",
        "                ForEach(entries, id: \\.self) { entry in",
        "                    NavigationLink(",
        "                        destination: JournalDetailView(entry: entry))",
        "                    {",
        "                        JournalRowView(entry: entry)",
        "                    }",
        "                }",
        "                .onDelete(perform: deleteEntries)",
        "            }",
        "            .navigationTitle(\"កំណត់ហេតុកសិកម្ម\")",
        "        }",
        "    }",
        "}",
    ],.35,5.3,6.3,2.05,size=9.5)
    add_text(s,"JournalRowView Preview",7.0,1.05,6.0,.4,size=12,bold=True,color=PURPLE)
    add_code(s,[
        "struct JournalRowView: View {",
        "    let entry: JournalEntry",
        "    var body: some View {",
        "        VStack(alignment: .leading, spacing: 6) {",
        "            HStack {",
        "                if let date = entry.date {",
        "                    Text(date, style: .date)",
        "                        .font(.headline)",
        "                }",
        "                Spacer()",
        "                // Weather icon",
        "                Text(weatherEmoji(entry.weather ?? \"\"))",
        "            }",
        "            if let content = entry.content {",
        "                Text(content)",
        "                    .font(.subheadline)",
        "                    .foregroundColor(.gray)",
        "                    .lineLimit(2)",
        "            }",
        "            if entry.photoData != nil {",
        "                Label(\"📷 មានរូបភាព\", systemImage: \"photo\")",
        "                    .font(.caption)",
        "                    .foregroundColor(.blue)",
        "            }",
        "        }",
        "    }",
        "}",
    ],7.0,1.55,6.0,5.8,size=9.5)
    snum(s,4)

def s5(prs):
    s=add_slide(prs)
    hdr(s,"🌤️  Weather System · Enum + SF Symbols")
    add_code(s,[
        "enum WeatherType: String, CaseIterable {",
        '    case sunny  = "sunny"',
        '    case rainy  = "rainy"',
        '    case cloudy = "cloudy"',
        '    case windy  = "windy"',
        "",
        "    var displayName: String {",
        "        switch self {",
        '        case .sunny:  return "☀️ ថ្ងៃរះ"',
        '        case .rainy:  return "🌧️ ភ្លៀង"',
        '        case .cloudy: return "⛅ ពពក"',
        '        case .windy:  return "💨 ខ្យល់"',
        "        }",
        "    }",
        "",
        "    var sfSymbol: String {",
        "        switch self {",
        '        case .sunny:  return "sun.max.fill"',
        '        case .rainy:  return "cloud.rain.fill"',
        '        case .cloudy: return "cloud.fill"',
        '        case .windy:  return "wind"',
        "        }",
        "    }",
        "}",
        "",
        "// In AddJournalEntryView:",
        "Picker(\"អាកាសធាតុ\", selection: $weather) {",
        "    ForEach(WeatherType.allCases, id: \\.rawValue) {",
        "        Label($0.displayName,",
        "              systemImage: $0.sfSymbol).tag($0.rawValue)",
        "    }",
        "}",
    ],.35,1.1,6.5,6.25,size=9.5)
    notes=[
        (PURPLE,"rawValue in CoreData","Store \"sunny\" string — not enum\n→ CoreData = String attribute\n→ easier migration later"),
        (ORANGE,"CaseIterable","ForEach(WeatherType.allCases)\n→ Picker gets all options\nautomatically"),
        (GREEN, "SF Symbols","sun.max.fill, cloud.rain.fill\n→ consistent Apple icons\n→ always in iOS 13+"),
        (BLUE,  "displayName","Khmer + emoji label\n→ farmer-friendly UI\n→ keep rawValue English"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,7.05,cy,6.0,1.38,CARD); add_rect(s,7.05,cy,.08,1.38,c)
        add_text(s,t,7.28,cy+.1,5.5,.4,size=12,bold=True,color=c)
        add_text(s,d,7.28,cy+.55,5.5,.72,size=11,color=WHITE)
    snum(s,5)

def s6(prs):
    s=add_slide(prs)
    hdr(s,"📷  Photo Picker · UIKit Bridge (iOS 13+)")
    add_code(s,[
        "// UIViewControllerRepresentable bridge",
        "struct ImagePickerView: UIViewControllerRepresentable {",
        "    @Binding var selectedImage: UIImage?",
        "    @Environment(\\.presentationMode) var presentMode",
        "",
        "    func makeUIViewController(context: Context)",
        "            -> UIImagePickerController {",
        "        let picker = UIImagePickerController()",
        "        picker.delegate = context.coordinator",
        "        picker.sourceType = .photoLibrary",
        "        return picker",
        "    }",
        "",
        "    func updateUIViewController(_ uiVC: UIImagePickerController,",
        "                               context: Context) {}",
        "",
        "    func makeCoordinator() -> Coordinator {",
        "        Coordinator(self)",
        "    }",
        "",
        "    class Coordinator: NSObject,",
        "        UIImagePickerControllerDelegate,",
        "        UINavigationControllerDelegate {",
        "        let parent: ImagePickerView",
        "        init(_ p: ImagePickerView) { parent = p }",
        "",
        "        func imagePickerController(_ picker:",
        "            UIImagePickerController,",
        "            didFinishPickingMediaWithInfo info:",
        "            [UIImagePickerController.InfoKey: Any]) {",
        "            if let img = info[.originalImage] as? UIImage {",
        "                parent.selectedImage = img",
        "            }",
        "            parent.presentMode.wrappedValue.dismiss()",
        "        }",
        "    }",
        "}",
    ],.35,1.1,7.2,6.25,size=9.5)
    notes=[
        (PURPLE,"3 Required Methods","makeUIViewController\nupdateUIViewController\nmakeCoordinator"),
        (ORANGE,"Coordinator pattern","NSObject subclass\nimplements UIKit delegates\nbridges back to SwiftUI"),
        (GREEN, "Store as Data","selectedImage → jpegData(0.8)\n→ entry.photoData: Data?\n→ CoreData stores"),
        (BLUE,  "iOS 13+ safe","UIImagePickerController\n= iOS 2+\n❌ PHPickerViewController (iOS 14+)"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*1.55
        add_rect(s,7.55,cy,5.5,1.38,CARD); add_rect(s,7.55,cy,.08,1.38,c)
        add_text(s,t,7.78,cy+.1,5.0,.4,size=12,bold=True,color=c)
        add_text(s,d,7.78,cy+.55,5.0,.72,size=11,color=WHITE)
    snum(s,6)

def s7(prs):
    s=add_slide(prs)
    hdr(s,"🔍  Search & Filter · Text + Weather Combined")
    add_code(s,[
        "struct JournalTabView: View {",
        "    @State private var searchText   = \"\"",
        "    @State private var filterWeather = \"\"  // \"\" = all",
        "",
        "    var displayedEntries: [JournalEntry] {",
        "        entries.filter { entry in",
        "            let matchText = searchText.isEmpty ||",
        "                (entry.content ?? \"\")",
        "                    .localizedCaseInsensitiveContains(searchText)",
        "            let matchWeather = filterWeather.isEmpty ||",
        "                entry.weather == filterWeather",
        "            return matchText && matchWeather",
        "        }",
        "    }",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            VStack {",
        "                // Search TextField",
        "                TextField(\"ស្វែងរក...\", text: $searchText)",
        "                    .padding(8)",
        "                    .background(Color(.systemGray6))",
        "                    .cornerRadius(10)",
        "                    .padding(.horizontal)",
        "",
        "                // Weather filter",
        "                ScrollView(.horizontal) {",
        "                    HStack {",
        "                        WeatherFilterChip(label: \"ទាំងអស់\",",
        "                            isSelected: filterWeather.isEmpty,",
        "                            action: { filterWeather = \"\" })",
        "                        ForEach(WeatherType.allCases,id: \\.rawValue){w in",
        "                            WeatherFilterChip(label: w.displayName,",
        "                                isSelected: filterWeather==w.rawValue,",
        "                                action: { filterWeather=w.rawValue })",
        "                        }",
        "                    }",
        "                }",
        "            }",
        "        }",
        "    }",
        "}",
    ],.35,1.1,7.5,6.25,size=9)
    notes=[
        (PURPLE,"Combined filter","searchText && filterWeather\n→ both conditions must match"),
        (GREEN, "Weather chips","ScrollView(.horizontal)\n→ filter buttons row"),
        (ORANGE,"localizedCaseInsensitiveContains","ស្វែងរក Khmer + English\ncase-insensitive"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*2.1
        add_rect(s,8.0,cy,5.1,1.9,CARD); add_rect(s,8.0,cy,.08,1.9,c)
        add_text(s,t,8.22,cy+.12,4.6,.45,size=13,bold=True,color=c)
        add_text(s,d,8.22,cy+.65,4.6,.9,size=12,color=WHITE)
    snum(s,7)

def s8(prs):
    s=add_slide(prs)
    hdr(s,"⚠️  Common Mistakes · ❌ vs ✅")
    rows=[
        ("UIKit bridge","UIViewControllerRepresentable missing","Implement all 3: make/update/makeCoordinator"),
        ("Photo storage","UIImage stored directly","Convert: jpegData(0.8) → Data → CoreData"),
        ("PHPickerVC","PHPickerViewController (iOS 14+)","UIImagePickerController (iOS 2+)"),
        ("External storage","photoData without external storage","Check 'Allows External Storage' in model"),
        ("Weather filter","SwiftUI Picker vs custom chips","ScrollView HStack chips (better UX)"),
        ("Search bar",".searchable(text:) (iOS 15+)","Custom TextField in VStack (iOS 13+)"),
    ]
    hdrs=["Pattern","❌  Wrong","✅  Correct (iOS 13+)"]
    widths=[2.8,4.5,5.1]; starts=[.35,3.3,7.95]
    for c,(h,cw,cs) in enumerate(zip(hdrs,widths,starts)):
        add_rect(s,cs,1.1,cw,.55,PURPLE if c==2 else RED if c==1 else DARK2)
        add_text(s,h,cs+.1,1.15,cw-.15,.45,size=12,bold=True,color=WHITE)
    for r,(p,w,g) in enumerate(rows):
        cy=1.75+r*.67; bg=CARD if r%2==0 else DARK2
        for c,(cell,cw,cs) in enumerate(zip([p,w,g],widths,starts)):
            add_rect(s,cs,cy,cw,.62,bg)
            col=GREY if c==0 else (RED if c==1 else GREEN)
            add_text(s,cell,cs+.1,cy+.1,cw-.15,.42,size=10,color=col)
    snum(s,8)

def s9(prs):
    s=add_slide(prs)
    hdr(s,"🔄  AddJournalEntryView · Complete Form")
    add_code(s,[
        "struct AddJournalEntryView: View {",
        "    @Environment(\\.presentationMode) var pm",
        "    @Environment(\\.managedObjectContext) private var ctx",
        "    @State private var content  = \"\"",
        "    @State private var weather  = WeatherType.sunny.rawValue",
        "    @State private var location = \"\"",
        "    @State private var date     = Date()",
        "    @State private var selectedImage: UIImage?",
        "    @State private var showPicker = false",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            Form {",
        "                Section(header: Text(\"កាលបរិច្ឆេទ\")) {",
        "                    DatePicker(\"\", selection: $date,",
        "                        displayedComponents: .date).labelsHidden()",
        "                }",
        "                Section(header: Text(\"អាកាសធាតុ\")) {",
        "                    Picker(\"Weather\", selection: $weather) {",
        "                        ForEach(WeatherType.allCases, id: \\.rawValue) {",
        "                            Text($0.displayName).tag($0.rawValue)",
        "                        }",
        "                    }",
        "                }",
        "                Section(header: Text(\"ទីតាំង\")) {",
        "                    TextField(\"ឧ. វាលស្រែ ១\", text: $location)",
        "                }",
        "                Section(header: Text(\"មាតិកា\")) {",
        "                    TextEditor(text: $content)",
        "                        .frame(minHeight: 120)",
        "                }",
        "                Section(header: Text(\"រូបភាព\")) {",
        "                    Button(\"ជ្រើសរូបភាព\") { showPicker = true }",
        "                    if let img = selectedImage {",
        "                        Image(uiImage: img).resizable()",
        "                            .scaledToFit().frame(maxHeight: 200)",
        "                    }",
        "                }",
        "            }",
        "            .sheet(isPresented: $showPicker) { ImagePickerView(...) }",
        "        }",
        "    }",
        "}",
    ],.35,1.1,7.5,6.25,size=9)
    notes=[
        (PURPLE,"TextEditor","Multi-line text input\niOS 13+ ✅\n(UITextView bridge)"),
        (GREEN, "Form sections","Group: date, weather,\nlocation, content, photo\n→ clean UX"),
        (ORANGE,"Sheet for picker","@State showPicker = false\n.sheet(isPresented:)\n→ open ImagePickerView"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.1+i*2.1
        add_rect(s,8.0,cy,5.1,1.9,CARD); add_rect(s,8.0,cy,.08,1.9,c)
        add_text(s,t,8.22,cy+.12,4.6,.45,size=13,bold=True,color=c)
        add_text(s,d,8.22,cy+.65,4.6,.9,size=12,color=WHITE)
    snum(s,9)

def s10(prs):
    s=add_slide(prs)
    hdr(s,"🔄  Data Flow · JournalEntry Full Lifecycle")
    flow=[
        (PURPLE,"User taps + button","AddJournalEntryView opens as sheet"),
        (GREEN, "Fill form","date, weather, location, content, photo"),
        (BLUE,  "Photo picked","UIImagePickerController → UIImage → jpegData → Data"),
        (ORANGE,"Tap Save","JournalEntry(context:) → set all properties → saveContext()"),
        (TEAL,  "@FetchRequest refreshes","JournalTabView list updates automatically"),
        (PURPLE,"Tap row","NavigationLink → JournalDetailView — full content + photo"),
    ]
    for i,(c,kh,en) in enumerate(flow):
        col=i%2; row=i//2; cx=.35+col*6.5; cy=1.1+row*2.0
        add_rect(s,cx,cy,6.1,1.75,CARD); add_rect(s,cx,cy,.08,1.75,c)
        add_rect(s,cx+.15,cy+.4,.55,.55,c)
        add_text(s,str(i+1),cx+.15,cy+.4,.55,.55,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        add_text(s,kh,cx+.88,cy+.2,5.0,.5,size=14,bold=True,color=WHITE)
        add_text(s,en,cx+.88,cy+.82,5.0,.6,size=11,color=GREY)
    snum(s,10)

def s11(prs):
    s=add_slide(prs)
    hdr(s,"🏠  Mini-Project · Daily Journal ពេញ")
    cols=[
        ("📋  Requirements",PURPLE,[
            "Timeline list — reverse chronological",
            "Add journal entry form",
            "Weather picker (4 options + emoji)",
            "Photo from library — stored as Data",
            "Search by content text",
            "Filter by weather type",
        ]),
        ("✅  Checklist",GREEN,[
            "@FetchRequest sorted by date ✓",
            "WeatherType enum + displayName ✓",
            "ImagePickerView UIKit bridge ✓",
            "jpegData → photoData: Data ✓",
            "Search + weather filter ✓",
            "JournalDetailView shows photo ✓",
        ]),
        ("🎯  Grading",BLUE,[
            "Timeline UI: 20%",
            "CoreData CRUD: 20%",
            "Weather system: 15%",
            "Photo picker UIKit bridge: 25%",
            "Search + filter: 20%",
            "Bonus: mood tracking",
        ]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*4.35
        add_rect(s,cx,1.1,4.1,5.8,CARD); add_rect(s,cx,1.1,4.1,.55,c)
        add_text(s,t,cx+.12,1.15,3.85,.45,size=13,bold=True,color=DARK if c!=PURPLE else WHITE)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.15,cy,3.7,.72,DARK2)
            add_text(s,item,cx+.28,cy+.1,3.4,.52,size=11,color=WHITE)
    add_rect(s,.35,7.0,12.6,.42,PURPLE)
    add_text(s,"📖  Week 07 ចប់! — Journal with photos + weather ready · Next: Dashboard →",
             .5,7.02,12.3,.38,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    snum(s,11)

def s12(prs):
    s=add_slide(prs)
    hdr(s,"📌  សង្ខេប Week 07 · Cheat Sheet",accent=PURPLE)
    cols=[
        ("📖  Journal",PURPLE,["JournalEntry entity","date/content/weather/photo","reverse sort @FetchRequest","JournalDetailView"]),
        ("🌤️  Weather",GREEN,["WeatherType enum","rawValue: String","displayName: KH emoji","CaseIterable Picker"]),
        ("📷  Photos",BLUE,["UIViewControllerRepresentable","makeUIViewController","makeCoordinator","jpegData → Data"]),
        ("🔍  Search",ORANGE,["@State searchText","@State filterWeather","localizedCaseInsensitiveContains","ScrollView chips"]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*3.25
        add_rect(s,cx,1.1,3.0,5.8,CARD); add_rect(s,cx,1.1,3.0,.55,c)
        add_text(s,t,cx+.1,1.15,2.8,.45,size=12,bold=True,color=DARK if c!=PURPLE else WHITE)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.1,cy,2.8,.72,DARK2)
            add_text(s,item,cx+.18,cy+.1,2.6,.52,size=10,color=TEAL)
            s.shapes[-1].text_frame.paragraphs[0].runs[0].font.name="Courier New"
    add_rect(s,.35,7.0,12.6,.42,PURPLE)
    add_text(s,"📖  Journal = ការ record ជីវិត farmer! · Photos + weather = memory richer than text alone.",
             .5,7.02,12.3,.38,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    snum(s,12)

def build():
    prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H)
    s1(prs);s2(prs);s3(prs);s4(prs);s5(prs);s6(prs)
    s7(prs);s8(prs);s9(prs);s10(prs);s11(prs);s12(prs)
    out=os.path.join(os.path.dirname(__file__),"Week07_Daily_Journal_KH.pptx")
    prs.save(out); print(f"✅  Saved → {out}")

if __name__=="__main__": build()
