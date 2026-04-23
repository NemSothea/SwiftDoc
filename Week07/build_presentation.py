"""Generate Keynote/PowerPoint-ready deck for Week 7 (Daily Journal).

Produces: Week07_Daily_Journal.pptx

Design system (adapted from Week 6; module accent swapped from Leaf Green
to Journal Purple):
- Palette:
    Primary Blue    #0A84FF   (iOS system blue — nav / accents)
    Journal Purple  #A052D4   (module theme — replaces GREEN from Week 6)
    Green (ship)    #34C759   (kept for success chips on the summary slide)
    Alert Orange    #FF9500   (warnings, photo-picker flow)
    Background      #F2F2F7
    Surface         #FFFFFF
    Text Primary    #1C1C1E
    Text Secondary  #8E8E93
- Fonts: Title = SF Pro Display / Helvetica Neue Bold; Body = SF Pro Text
- Layout: 2-column (concept left, UI / code right)
- Icons: SF Symbols (book, sun.max.fill, cloud.rain.fill, cloud.fill, wind,
    photo.on.rectangle, camera, magnifyingglass)
- Animations: Fade (build in), Slide-from-right for code, Zoom for the
    photo-picker flow diagram
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# --- Design tokens -----------------------------------------------------------

PRIMARY = RGBColor(0x0A, 0x84, 0xFF)   # iOS blue
PURPLE  = RGBColor(0xA0, 0x52, 0xD4)   # journal purple (module theme)
GREEN   = RGBColor(0x34, 0xC7, 0x59)   # success / ship
ORANGE  = RGBColor(0xFF, 0x95, 0x00)   # alerts / photo flow
BG      = RGBColor(0xF2, 0xF2, 0xF7)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT    = RGBColor(0x1C, 0x1C, 0x1E)
MUTED   = RGBColor(0x8E, 0x8E, 0x93)
DIVIDER = RGBColor(0xD1, 0xD1, 0xD6)
CODE_BG = RGBColor(0x1C, 0x1C, 0x1E)
CODE_FG = RGBColor(0xF2, 0xF2, 0xF7)
TIP_BG  = RGBColor(0xFF, 0xF4, 0xE0)

# Weather tints (for slide 7)
WEATHER_YELLOW = RGBColor(0xFF, 0xCC, 0x00)
WEATHER_BLUE   = RGBColor(0x0A, 0x84, 0xFF)
WEATHER_GRAY   = RGBColor(0x8E, 0x8E, 0x93)
WEATHER_TEAL   = RGBColor(0x30, 0xB0, 0xC7)

TITLE_FONT = "SF Pro Display"
BODY_FONT  = "SF Pro Text"
MONO_FONT  = "Menlo"

# 16:9 deck
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


# --- Helpers -----------------------------------------------------------------

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=18, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(40000)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT,
                font=BODY_FONT, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(80000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(40000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_card(slide, x, y, w, h, *, fill=SURFACE, border=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = border
        card.line.width = Pt(0.75)
    card.shadow.inherit = False
    return card


def slide_header(slide, eyebrow, title, *, accent=PRIMARY):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.55), Inches(0.55),
                                 Inches(0.08), Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    add_text(slide, Inches(0.75), Inches(0.45), Inches(12), Inches(0.35),
             eyebrow, font=BODY_FONT, size=12, bold=True, color=accent)
    add_text(slide, Inches(0.75), Inches(0.75), Inches(12), Inches(0.75),
             title, font=TITLE_FONT, size=30, bold=True, color=TEXT)


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_code_block(slide, x, y, w, h, code_lines, *, size=11):
    add_card(slide, x, y, w, h, fill=CODE_BG)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(140000)
    tf.margin_right = Emu(80000)
    tf.margin_top = Emu(100000)
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


def add_chip(slide, x, y, text, *, fill=PRIMARY, text_color=SURFACE, w=1.7, h=0.32):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    tf = shp.text_frame
    tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
    tf.margin_top = Emu(20000);  tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = BODY_FONT; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.55), Inches(7.1), Inches(6), Inches(0.3),
             "Week 7 · Daily Journal · SmartFarmerAssistant",
             size=10, color=MUTED)
    add_text(slide, Inches(10.5), Inches(7.1), Inches(2.3), Inches(0.3),
             f"{idx} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def set_transition(slide, kind="fade"):
    xml_map = {
        "fade": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        "push": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        "zoom": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
    }
    node = etree.fromstring(xml_map.get(kind, xml_map["fade"]))
    slide.element.append(node)


TOTAL = 12


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = add_slide()

band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(4.5))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = PURPLE

for (cx, cy, r, color) in [
    (11.0, 0.6, 1.6, PRIMARY),
    (12.3, 3.4, 0.9, ORANGE),
    (1.1, 3.2, 0.55, SURFACE),
]:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(cx), Inches(cy),
                              Inches(r), Inches(r))
    circ.line.fill.background()
    circ.fill.solid(); circ.fill.fore_color.rgb = color

add_chip(s, 0.6, 0.6, "iOS · SwiftUI · Week 7",
         fill=SURFACE, text_color=PURPLE, w=2.4, h=0.38)

add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(1.5),
         "Daily Journal",
         font=TITLE_FONT, size=52, bold=True, color=SURFACE)
add_text(s, Inches(0.6), Inches(2.25), Inches(12), Inches(1.0),
         "Your Digital Notebook",
         font=TITLE_FONT, size=44, bold=True, color=SURFACE)
add_text(s, Inches(0.6), Inches(3.35), Inches(12), Inches(0.6),
         "SwiftUI  ·  Core Data  ·  Photos  ·  Weather",
         font=BODY_FONT, size=22, color=SURFACE)

add_card(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.6))
add_text(s, Inches(0.95), Inches(5.15), Inches(6), Inches(0.45),
         "AUDIENCE", size=11, bold=True, color=MUTED)
add_text(s, Inches(0.95), Inches(5.45), Inches(6), Inches(0.6),
         "Beginner iOS Developers", font=TITLE_FONT, size=20, bold=True)
add_text(s, Inches(0.95), Inches(5.95), Inches(7), Inches(0.5),
         "JournalEntry · Timeline list · Weather picker · UIKit bridge · Custom search",
         size=13, color=MUTED)

add_text(s, Inches(7.8), Inches(5.15), Inches(4.8), Inches(0.45),
         "PROJECT 4 — SMART FARMER ASSISTANT",
         size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
add_text(s, Inches(7.8), Inches(5.45), Inches(4.8), Inches(0.6),
         "Journal Module", font=TITLE_FONT, size=20, bold=True,
         align=PP_ALIGN.RIGHT)
add_text(s, Inches(7.8), Inches(5.95), Inches(4.8), Inches(0.5),
         "iOS 13+ — ObservableObject · NavigationView · UIImagePickerController",
         size=13, color=MUTED, align=PP_ALIGN.RIGHT)

add_footer(s, 1, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Welcome to Week 7. This week we ship the Daily Journal — an Apple Notes-style "
    "digital notebook that lives inside the Smart Farmer Assistant. By the end, "
    "students will have a JournalEntry Core Data entity, a reverse-chronological "
    "timeline, a weather-picker UI, a UIKit-bridged photo picker, and a custom "
    "search + filter bar. Everything works offline — no weather API, no cloud sync. "
    "UI suggestion: open with the finished Journal tab on screen — show the day-grouped "
    "timeline, tap a row to see the weather badge and photo gallery — so learners see "
    "the target before we write code."
))


# ============================================================================
# SLIDE 2 — Agenda / Learning Objectives
# ============================================================================
s = add_slide()
slide_header(s, "AGENDA", "What you'll build this week", accent=PURPLE)

add_text(s, Inches(0.75), Inches(1.9), Inches(5.8), Inches(0.4),
         "LEARNING OBJECTIVES", size=12, bold=True, color=PURPLE)
add_bullets(s, Inches(0.75), Inches(2.25), Inches(6.1), Inches(4.5), [
    "Model a JournalEntry entity in Core Data (date, content, weather, photos)",
    "Build a reverse-chronological timeline UI grouped by day",
    "Add a 4-option weather picker with SF Symbols",
    "Bridge UIImagePickerController via UIViewControllerRepresentable",
    "Combine a custom SearchBar with a weather filter chip row (no .searchable)",
], size=17)

lessons = [
    ("7.1", "JournalEntry Core Data Entity"),
    ("7.2", "Timeline list (newest first)"),
    ("7.3", "Weather selection UI"),
    ("7.4", "Photo picker — UIKit bridge"),
    ("7.5", "Custom search + weather filter"),
]
for i, (num, title) in enumerate(lessons):
    top = Inches(1.9 + i * 0.95)
    add_card(s, Inches(7.1), top, Inches(5.6), Inches(0.8))
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(7.25), top + Inches(0.15),
                              Inches(0.5), Inches(0.5))
    chip.line.fill.background()
    chip.fill.solid(); chip.fill.fore_color.rgb = PURPLE
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = TITLE_FONT; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = SURFACE
    add_text(s, Inches(7.95), top + Inches(0.2), Inches(4.5), Inches(0.5),
             title, font=TITLE_FONT, size=16, bold=True)

add_footer(s, 2, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Walk the five lessons at a high level. The payoff is Lesson 7.4 (the UIKit "
    "bridge) — it's a pattern students will reuse every time SwiftUI can't "
    "do something natively on iOS 13. Stress that we pick iOS 13+ APIs all the way: "
    "no .searchable, no PhotosPicker, no NavigationStack. "
    "UI suggestion: objectives on the left, numbered purple lesson chips on the "
    "right — build in the chips one-by-one with a fade."
))


# ============================================================================
# SLIDE 3 — Architecture
# ============================================================================
s = add_slide()
slide_header(s, "ARCHITECTURE", "MVVM — clean SwiftUI layers", accent=PRIMARY)

layers = [
    ("View",       "JournalTabView · JournalRowView · JournalDetailView · Add/Edit sheets · PhotoPicker · WeatherPickerView", PRIMARY),
    ("ViewModel",  "JournalViewModel — @Published searchText + weatherFilter + filter(_:)",  PURPLE),
    ("Service",    "JournalPhotoStore — UIImage ↔ JPEG Data helpers",                         ORANGE),
    ("Model",      "JournalEntry (NSManagedObject) + Weather (enum)",                         RGBColor(0x55, 0xB5, 0x5C)),
    ("Resource",   "(empty — entries are user-authored, no seed JSON)",                        MUTED),
]
y = 1.9
for role, desc, color in layers:
    add_card(s, Inches(0.75), Inches(y), Inches(11.8), Inches(0.85))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.75), Inches(y),
                             Inches(0.18), Inches(0.85))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    add_text(s, Inches(1.1), Inches(y + 0.12), Inches(2.2), Inches(0.4),
             role, font=TITLE_FONT, size=16, bold=True, color=color)
    add_text(s, Inches(3.4), Inches(y + 0.2), Inches(9.0), Inches(0.5),
             desc, size=13, color=TEXT)
    y += 1.00

add_footer(s, 3, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "This is the same layering students saw in Finance, CalendarReminders, and "
    "PestDisease. The new box this week is 'Service' — JournalPhotoStore — because "
    "UIImage and Data are not the same thing, and we want that translation isolated "
    "from views. Notice the Resource layer is intentionally empty: journal entries "
    "come from the user, not from a bundled file. "
    "UI suggestion: five stacked cards with a colored accent bar on the left. "
    "Animation: fade in top-to-bottom, 120ms delay between cards."
))


# ============================================================================
# SLIDE 4 — Folder Structure
# ============================================================================
s = add_slide()
slide_header(s, "CODE LAYOUT", "Journal module — folder by folder", accent=PRIMARY)

add_bullets(s, Inches(0.75), Inches(1.9), Inches(5.5), Inches(4.5), [
    "Mirrors Finance/, CalendarReminders/, PestDisease/ — same shape across modules",
    "Models/ holds the NSManagedObject class + the Weather enum",
    "Services/ isolates the UIImage ↔ Data conversion",
    "ViewModels/ owns search + weather filter + day grouping",
    "SearchBar and ExpandableSection are reused from PestDisease — no duplicates",
], size=16)

tree = [
    "SmartFarmerAssistantFinish/",
    "└── Journal/",
    "    ├── Models/",
    "    │   ├── JournalEntry+CoreDataClass.swift",
    "    │   ├── JournalEntry+CoreDataProperties.swift",
    "    │   └── Weather.swift",
    "    ├── ViewModels/",
    "    │   └── JournalViewModel.swift",
    "    ├── Services/",
    "    │   └── JournalPhotoStore.swift",
    "    ├── Views/",
    "    │   ├── JournalTabView.swift",
    "    │   ├── JournalRowView.swift",
    "    │   ├── JournalDetailView.swift",
    "    │   ├── AddJournalEntryView.swift",
    "    │   ├── EditJournalEntryView.swift",
    "    │   ├── WeatherPickerView.swift",
    "    │   ├── PhotoPicker.swift",
    "    │   └── PhotoGalleryView.swift",
    "    └── Resources/",
    "        └── (empty)",
]
add_code_block(s, Inches(6.55), Inches(1.85), Inches(6.2), Inches(5.2),
               tree, size=11)

add_footer(s, 4, TOTAL)
set_transition(s, "push")
speaker_notes(s, (
    "Walk the folder tree top to bottom. Emphasise that every week we open a new "
    "folder alongside Finance/, CalendarReminders/, and PestDisease/ — not a new "
    "project. Shared Core Data model, shared CoreDataManager, shared MainTabView. "
    "Point out that the Resources folder is empty — this is the first module "
    "where we don't ship seed data. "
    "UI suggestion: bullets on the left describing WHY each folder exists, the "
    "actual tree rendered as a dark code block on the right. "
    "Animation: slide-from-right for the code block."
))


# ============================================================================
# SLIDE 5 — Core Data JournalEntry Entity
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 7.1", "Core Data — JournalEntry Entity", accent=PRIMARY)

add_text(s, Inches(0.75), Inches(1.85), Inches(6.0), Inches(0.4),
         "ATTRIBUTES", size=12, bold=True, color=PRIMARY)

attrs = [
    ("id",       "UUID",          "Stable identifier"),
    ("date",     "Date",          "Timeline sort key + day grouping"),
    ("title",    "String?",       "Optional — first line of content otherwise"),
    ("content",  "String",        "Body text (edited in TextEditor)"),
    ("weather",  "String",        "Raw value of the Weather enum"),
    ("location", "String?",       "Optional bonus field"),
    ("photos",   "Transformable", "[Data] — JPEGs via NSSecureUnarchiveFromData"),
]
row_y = 2.25
for name, typ, desc in attrs:
    add_card(s, Inches(0.75), Inches(row_y), Inches(6.0), Inches(0.55))
    add_text(s, Inches(0.95), Inches(row_y + 0.07), Inches(1.7), Inches(0.4),
             name, font=MONO_FONT, size=12, bold=True, color=PRIMARY)
    add_text(s, Inches(2.65), Inches(row_y + 0.07), Inches(1.3), Inches(0.4),
             typ, font=MONO_FONT, size=11, color=MUTED)
    add_text(s, Inches(4.0), Inches(row_y + 0.07), Inches(2.65), Inches(0.4),
             desc, size=10, color=TEXT)
    row_y += 0.62

code = [
    "// Journal/Models/JournalEntry+CoreDataProperties.swift",
    "extension JournalEntry {",
    "  @NSManaged public var id: UUID?",
    "  @NSManaged public var date: Date?",
    "  @NSManaged public var title: String?",
    "  @NSManaged public var content: String?",
    "  @NSManaged public var weather: String?",
    "  @NSManaged public var location: String?",
    "  @NSManaged public var photos: NSObject?",
    "}",
    "",
    "extension JournalEntry: Identifiable {}",
]
add_code_block(s, Inches(7.25), Inches(2.1), Inches(5.45), Inches(3.5), code, size=12)

add_card(s, Inches(7.25), Inches(5.75), Inches(5.45), Inches(1.3), fill=TIP_BG)
add_text(s, Inches(7.5), Inches(5.85), Inches(5.0), Inches(0.4),
         "⚠️  Codegen = Manual/None  ·  Transformer = NSSecureUnarchiveFromData",
         font=TITLE_FONT, size=12, bold=True, color=ORANGE)
add_text(s, Inches(7.5), Inches(6.2), Inches(5.0), Inches(0.85),
         "Set both in the Data Model Inspector — otherwise Xcode "
         "auto-generates the class (duplicate-symbol error) or the "
         "photos attribute fails to decode at fetch time.",
         size=11, color=TEXT)

add_footer(s, 5, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Walk each attribute. The four core ones (date, content, weather, photos) are "
    "required by the spec; title is optional and falls back to the first line of "
    "content (same as Apple Notes); location is a bonus field we kept from an "
    "earlier draft. Show the Data Model Inspector live: point to Codegen → "
    "Manual/None, and to the Transformer → NSSecureUnarchiveFromData on the "
    "photos attribute. These two settings are the #1 cause of fetch-time crashes "
    "for learners. "
    "UI suggestion: attribute cards left, code right, orange warning card "
    "bottom-right with both configuration gotchas. "
    "Animation: slide-from-left for cards, slide-from-right for the code block."
))


# ============================================================================
# SLIDE 6 — Timeline UI
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 7.2", "Timeline UI — newest first, grouped by day", accent=PURPLE)

add_bullets(s, Inches(0.75), Inches(1.9), Inches(5.5), Inches(4.8), [
    "@FetchRequest(sort: date, ascending: false)",
    "Dictionary(grouping:by:) keyed by Calendar.startOfDay",
    "Each day → its own List Section with a formatted header",
    "Row: time · title · snippet · weather icon · photo count",
    "PlainListStyle for Apple Notes feel (no inset gaps)",
    ".onDelete deletes from viewContext; fetch re-runs automatically",
], size=15)

# iPhone-style mock
mock_x = 7.3; mock_y = 1.85
add_card(s, Inches(mock_x), Inches(mock_y), Inches(5.4), Inches(5.0))

add_text(s, Inches(mock_x + 0.3), Inches(mock_y + 0.25), Inches(5.0), Inches(0.35),
         "TUESDAY, APRIL 23", font=TITLE_FONT, size=12, bold=True, color=MUTED)

rows1 = [
    ("08:14", "Morning fog over the rice", WEATHER_YELLOW, 2),
    ("17:02", "Harvest prep — met the team",  WEATHER_GRAY,  0),
]
ry = mock_y + 0.7
for time, title, color, photos in rows1:
    add_text(s, Inches(mock_x + 0.3), Inches(ry), Inches(0.6), Inches(0.35),
             time, font=MONO_FONT, size=11, color=MUTED)
    add_text(s, Inches(mock_x + 0.95), Inches(ry), Inches(3.6), Inches(0.35),
             title, font=TITLE_FONT, size=13, bold=True, color=TEXT)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(mock_x + 4.65), Inches(ry + 0.05),
                             Inches(0.3), Inches(0.3))
    dot.line.fill.background()
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    if photos > 0:
        add_text(s, Inches(mock_x + 5.0), Inches(ry), Inches(0.4), Inches(0.35),
                 f"📷 {photos}", size=10, color=MUTED)
    ry += 0.55

add_text(s, Inches(mock_x + 0.3), Inches(ry + 0.2), Inches(5.0), Inches(0.35),
         "MONDAY, APRIL 22", font=TITLE_FONT, size=12, bold=True, color=MUTED)
ry += 0.65
rows2 = [
    ("09:30", "Heavy rain — stayed indoors",  WEATHER_BLUE, 1),
    ("20:11", "Market prices update",          WEATHER_TEAL, 0),
]
for time, title, color, photos in rows2:
    add_text(s, Inches(mock_x + 0.3), Inches(ry), Inches(0.6), Inches(0.35),
             time, font=MONO_FONT, size=11, color=MUTED)
    add_text(s, Inches(mock_x + 0.95), Inches(ry), Inches(3.6), Inches(0.35),
             title, font=TITLE_FONT, size=13, bold=True, color=TEXT)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(mock_x + 4.65), Inches(ry + 0.05),
                             Inches(0.3), Inches(0.3))
    dot.line.fill.background()
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    if photos > 0:
        add_text(s, Inches(mock_x + 5.0), Inches(ry), Inches(0.4), Inches(0.35),
                 f"📷 {photos}", size=10, color=MUTED)
    ry += 0.55

add_footer(s, 6, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "The #1 learner mistake is ascending: true — it puts the oldest entry on top. "
    "A journal must read newest-first. Walk the grouping step: Dictionary(grouping:) "
    "keyed by startOfDay buckets every entry into a date. Then sort the buckets "
    "descending. The view just renders List { ForEach(groups) { Section { ForEach(entries) } } }. "
    "UI suggestion: bullets left, iPhone-shape mock right with two day sections "
    "and four rows; color dots represent the weather tint. "
    "Animation: fade in bullets; slide-from-top for each section."
))


# ============================================================================
# SLIDE 7 — Weather System
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 7.3", "Weather system — enum + SF Symbol + tint", accent=PURPLE)

weathers = [
    ("sun.max.fill",    "Sunny",  WEATHER_YELLOW),
    ("cloud.rain.fill", "Rainy",  WEATHER_BLUE),
    ("cloud.fill",      "Cloudy", WEATHER_GRAY),
    ("wind",            "Windy",  WEATHER_TEAL),
]
chip_y = 1.95
chip_w = 1.4
gap = 0.15
start_x = 0.75
for i, (sym, label, color) in enumerate(weathers):
    x = start_x + i * (chip_w + gap)
    add_card(s, Inches(x), Inches(chip_y), Inches(chip_w), Inches(1.3),
             fill=SURFACE, border=color)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(x + chip_w / 2 - 0.3), Inches(chip_y + 0.2),
                              Inches(0.6), Inches(0.6))
    circ.line.fill.background()
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    add_text(s, Inches(x), Inches(chip_y + 0.85), Inches(chip_w), Inches(0.3),
             label, font=TITLE_FONT, size=12, bold=True, color=TEXT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(x), Inches(chip_y + 1.05), Inches(chip_w), Inches(0.25),
             sym, font=MONO_FONT, size=9, color=MUTED, align=PP_ALIGN.CENTER)

code = [
    "enum Weather: String, CaseIterable, Identifiable {",
    "    case sunny, rainy, cloudy, windy",
    "    var id: String { rawValue }",
    "    var symbolName: String {",
    "        switch self {",
    "        case .sunny:  return \"sun.max.fill\"",
    "        case .rainy:  return \"cloud.rain.fill\"",
    "        case .cloudy: return \"cloud.fill\"",
    "        case .windy:  return \"wind\"",
    "        }",
    "    }",
    "    var tint: Color { /* yellow / blue / gray / teal */ }",
    "}",
]
add_code_block(s, Inches(0.75), Inches(3.6), Inches(7.5), Inches(3.3), code, size=12)

add_card(s, Inches(8.5), Inches(3.6), Inches(4.2), Inches(3.3), fill=TIP_BG)
add_text(s, Inches(8.7), Inches(3.7), Inches(4.0), Inches(0.4),
         "WHY rawValue IN CORE DATA", size=12, bold=True, color=ORANGE)
add_bullets(s, Inches(8.7), Inches(4.1), Inches(4.0), Inches(2.7), [
    "@NSManaged can't hold a Swift enum",
    "String round-trips cleanly through SQLite",
    "entry.weatherTag is the one place the string → enum mapping happens",
    "New case? One file to edit — no literals scattered across the codebase",
], size=11, color=TEXT)

add_footer(s, 7, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Three takeaways. First, the enum owns both the SF Symbol and the tint — "
    "so the picker and the timeline read from the same source of truth. Second, "
    "Core Data stores the rawValue String because @NSManaged can't persist a "
    "Swift enum. Third, the adapter — entry.weatherTag — is the only place the "
    "conversion happens, so you never hand-roll the literal 'sunny' anywhere "
    "else. Adding a fifth case (e.g. .stormy) is literally one line in the enum "
    "plus one line in each switch. "
    "UI suggestion: four colored weather chips across the top, enum code block "
    "left, amber why-card right. "
    "Animation: pop each weather chip in sequence; then fade in the code and "
    "the why-card together."
))


# ============================================================================
# SLIDE 8 — Photo Picker Flow
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 7.4", "Photo picker — UIKit bridge flow", accent=ORANGE)

stages = [
    ("Add photo button",       "SwiftUI .sheet(isPresented:)", PURPLE),
    ("PhotoPicker (sheet)",    "UIViewControllerRepresentable", ORANGE),
    ("UIImagePickerController", "makeUIViewController(context:)", ORANGE),
    ("didFinishPicking",       "Coordinator delegate method",    PRIMARY),
    ("UIImage → JPEG Data",    "jpegData(compressionQuality: 0.8)", PRIMARY),
    ("entry.photos append",    "[Data] as NSArray + save()",      GREEN),
]
count = len(stages)
gap_x = 0.20
box_w = (13.333 - 1.5 - (count - 1) * gap_x) / count
box_h = 2.1
y = 2.9
for i, (title, desc, color) in enumerate(stages):
    x = 0.75 + i * (box_w + gap_x)
    add_card(s, Inches(x), Inches(y), Inches(box_w), Inches(box_h),
             fill=SURFACE, border=color)
    add_text(s, Inches(x), Inches(y + 0.25), Inches(box_w), Inches(0.6),
             title, font=TITLE_FONT, size=13, bold=True, color=color,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(x + 0.08), Inches(y + 0.9), Inches(box_w - 0.16), Inches(1.1),
             desc, size=10, color=TEXT, align=PP_ALIGN.CENTER)

    if i < count - 1:
        arr_x = x + box_w + 0.01
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(arr_x), Inches(y + box_h / 2 - 0.12),
                                 Inches(gap_x - 0.02), Inches(0.24))
        arr.line.fill.background()
        arr.fill.solid(); arr.fill.fore_color.rgb = MUTED

add_text(s, Inches(0.75), Inches(5.6), Inches(11.8), Inches(0.5),
         "Same bridge pattern works for any UIKit controller — camera, document picker, share sheet.",
         size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

add_card(s, Inches(0.75), Inches(2.0), Inches(11.8), Inches(0.7), fill=TIP_BG)
add_text(s, Inches(0.95), Inches(2.1), Inches(11.4), Inches(0.5),
         "PhotosPicker is iOS 16+ — on iOS 13 we wrap UIImagePickerController. "
         "Three methods: makeUIViewController, updateUIViewController, makeCoordinator.",
         size=12, bold=True, color=ORANGE)

add_footer(s, 8, TOTAL)
set_transition(s, "zoom")
speaker_notes(s, (
    "This is the key lesson of Week 7 — the UIKit → SwiftUI bridge. Walk the six "
    "stages left to right. Emphasise two details: (1) the Coordinator class exists "
    "because SwiftUI structs are value types and can't conform to @objc delegate "
    "protocols; (2) we JPEG-encode at quality 0.8 before storing — a single camera "
    "photo at raw UIImage bytes is ~24MB; at JPEG 0.8 it's ~3MB, visually identical. "
    "Same bridge pattern reappears for the camera, the document picker, and share sheets. "
    "UI suggestion: 6 boxes left-to-right with arrows between them; amber why-card "
    "on top with the three lifecycle methods. "
    "Animation: zoom + fade for the flow; slide-from-top for each stage."
))


# ============================================================================
# SLIDE 9 — Search & Filter
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 7.5", "Search + weather filter — one ViewModel", accent=PRIMARY)

add_bullets(s, Inches(0.75), Inches(1.9), Inches(5.5), Inches(4.8), [
    "Reuse the Week 6 SearchBar (custom TextField, iOS 13+)",
    "Horizontal chip row: 'All' + 4 weather chips under the search bar",
    "Tap the same chip twice → toggles back to 'All' (nil)",
    "Single filter(_:) on the VM — composes text + weather",
    "localizedCaseInsensitiveContains works in Khmer and Latin scripts",
    "Swift-side filter on fetched results — instant, no SQLite per keystroke",
], size=15)

mock_x = 7.3; mock_y = 1.85
add_card(s, Inches(mock_x), Inches(mock_y), Inches(5.4), Inches(5.0))

pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(mock_x + 0.25), Inches(mock_y + 0.3),
                          Inches(4.9), Inches(0.45))
pill.adjustments[0] = 0.5
pill.line.fill.background()
pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0xE5, 0xE5, 0xEA)
add_text(s, Inches(mock_x + 0.4), Inches(mock_y + 0.37), Inches(4.7), Inches(0.35),
         "🔍  rice", size=12, color=TEXT)

chip_defs = [("All", MUTED), ("Sunny", WEATHER_YELLOW),
             ("Rainy", WEATHER_BLUE), ("Cloudy", WEATHER_GRAY), ("Windy", WEATHER_TEAL)]
cx = mock_x + 0.25
for label, col in chip_defs:
    is_selected = (label == "Sunny")
    w = 0.95
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(cx), Inches(mock_y + 0.95),
                              Inches(w), Inches(0.34))
    chip.adjustments[0] = 0.5
    chip.line.fill.background()
    chip.fill.solid()
    chip.fill.fore_color.rgb = col if is_selected else RGBColor(0xE5, 0xE5, 0xEA)
    tf = chip.text_frame
    tf.margin_left = Emu(30000); tf.margin_right = Emu(30000)
    tf.margin_top = Emu(15000); tf.margin_bottom = Emu(15000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = BODY_FONT; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = SURFACE if is_selected else TEXT
    cx += w + 0.05

rows = [
    ("Morning fog over the rice",    WEATHER_YELLOW),
    ("Rice harvest — day 2",         WEATHER_YELLOW),
    ("Rice prices spiked at market", WEATHER_YELLOW),
]
ry = mock_y + 1.7
for title, color in rows:
    add_text(s, Inches(mock_x + 0.4), Inches(ry), Inches(4.0), Inches(0.35),
             title, font=TITLE_FONT, size=13, bold=True, color=TEXT)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(mock_x + 4.7), Inches(ry + 0.05),
                             Inches(0.3), Inches(0.3))
    dot.line.fill.background()
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    ry += 0.55

add_text(s, Inches(mock_x + 0.4), Inches(ry + 0.1), Inches(4.6), Inches(0.3),
         "3 results · query: 'rice' · weather: Sunny",
         size=10, color=MUTED)

add_footer(s, 9, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "The ViewModel holds TWO pieces of state — searchText and weatherFilter — and "
    "a single filter(_:) combines them. Keep the combination logic in one function "
    "so you don't iterate the entries twice. Why nil-biased filtering? Because 'All' "
    "isn't a weather — it's the absence of a filter. Using Weather? instead of a "
    "five-case enum keeps both the UI and the tests simpler. "
    "UI suggestion: bullets left, iPhone mock right showing the search pill, a "
    "horizontal chip row with 'Sunny' selected, and three filtered result rows. "
    "Animation: fade in bullets; pop the chip row; fade in rows top-to-bottom."
))


# ============================================================================
# SLIDE 10 — Data Flow Diagram
# ============================================================================
s = add_slide()
slide_header(s, "DATA FLOW", "From '+' tap to re-rendered timeline", accent=PRIMARY)

stages = [
    ("User taps +",            "Toolbar + button on JournalTabView",      PURPLE),
    ("AddJournalEntryView",    "Form — title, weather, content, photo",    PRIMARY),
    ("JournalEntry(context:)", "New NSManagedObject attached to context",  PRIMARY),
    ("context.save()",         "Commits to SQLite store",                  GREEN),
    ("@FetchRequest re-renders","Timeline auto-updates — no manual refresh", GREEN),
]
count = len(stages)
gap_x = 0.35
box_w = (13.333 - 1.5 - (count - 1) * gap_x) / count
box_h = 2.2
y = 2.8
for i, (title, desc, color) in enumerate(stages):
    x = 0.75 + i * (box_w + gap_x)
    add_card(s, Inches(x), Inches(y), Inches(box_w), Inches(box_h),
             fill=SURFACE, border=color)
    add_text(s, Inches(x), Inches(y + 0.25), Inches(box_w), Inches(0.5),
             title, font=TITLE_FONT, size=15, bold=True, color=color,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(x + 0.1), Inches(y + 0.85), Inches(box_w - 0.2), Inches(1.3),
             desc, size=12, color=TEXT, align=PP_ALIGN.CENTER)

    if i < count - 1:
        arr_x = x + box_w + 0.02
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(arr_x), Inches(y + box_h / 2 - 0.12),
                                 Inches(gap_x - 0.04), Inches(0.24))
        arr.line.fill.background()
        arr.fill.solid(); arr.fill.fore_color.rgb = MUTED

add_text(s, Inches(0.75), Inches(5.5), Inches(11.8), Inches(0.5),
         "One save() at step 4 → every view bound to @FetchRequest<JournalEntry> re-renders automatically.",
         size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 10, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Read left-to-right. The key insight is the jump from step 4 to step 5: a "
    "single context.save() triggers every @FetchRequest in the app. No manual "
    "refresh, no NotificationCenter post — Core Data does the work. This is the "
    "same pattern from Finance, CalendarReminders, and PestDisease, but the "
    "payload here is a full entry plus a photo blob. "
    "UI suggestion: five color-coded cards in a row with arrows between them; "
    "caption at the bottom summarises the auto-re-render behaviour. "
    "Animation: slide-from-left for each box with a 100ms delay."
))


# ============================================================================
# SLIDE 11 — Offline-First
# ============================================================================
s = add_slide()
slide_header(s, "CONCEPT", "Offline-first — reframed for a journal", accent=ORANGE)

add_text(s, Inches(0.75), Inches(1.85), Inches(5.8), Inches(0.4),
         "❌ NETWORK-FIRST", size=12, bold=True, color=RGBColor(0xD0, 0x30, 0x30))
add_text(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.4),
         "✅ OFFLINE-FIRST", size=12, bold=True, color=GREEN)

rows = [
    ("First launch",
     "Blocked until cloud sign-in completes",
     "Opens empty, ready for the first entry"),
    ("Weather",
     "Fetches live forecast from an API",
     "User picks — it's the felt weather"),
    ("Photos",
     "Upload to cloud, store URL",
     "JPEG Data lives inside the Core Data store"),
    ("Privacy",
     "Every entry leaks to a server",
     "Entries never touch the network"),
]
y = 2.3
for label, bad, good in rows:
    add_text(s, Inches(0.75), Inches(y), Inches(1.5), Inches(0.4),
             label, size=13, bold=True, color=TEXT)
    add_card(s, Inches(2.35), Inches(y - 0.05), Inches(4.5), Inches(0.9),
             fill=RGBColor(0xFF, 0xEC, 0xEC))
    add_text(s, Inches(2.55), Inches(y + 0.08), Inches(4.2), Inches(0.8),
             bad, size=12, color=TEXT)
    add_card(s, Inches(7.0), Inches(y - 0.05), Inches(5.6), Inches(0.9),
             fill=RGBColor(0xE8, 0xF7, 0xEB))
    add_text(s, Inches(7.2), Inches(y + 0.08), Inches(5.3), Inches(0.8),
             good, size=12, color=TEXT)
    y += 1.1

add_card(s, Inches(0.75), Inches(6.55), Inches(11.8), Inches(0.55),
         fill=TIP_BG)
add_text(s, Inches(0.95), Inches(6.6), Inches(11.4), Inches(0.45),
         "Rule of thumb: a journal's weather is the *felt* weather, not the meteorological truth — user-picked beats API-fetched.",
         size=13, bold=True, color=ORANGE)

add_footer(s, 11, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Don't read every row — pick the Weather row. It's the most counter-intuitive: "
    "students assume a Journal 'should' pull live weather from an API. Explain why "
    "that's wrong — the user is remembering what THEY experienced, not what a "
    "satellite recorded. Live-demo airplane mode: the journal keeps working with "
    "zero latency. "
    "UI suggestion: four rows of red-vs-green cards; callout footer carries the "
    "rule of thumb. "
    "Animation: fade in rows top-to-bottom; pop the rule-of-thumb card last."
))


# ============================================================================
# SLIDE 12 — Summary + Q&A
# ============================================================================
s = add_slide()

band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(3.2))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = PURPLE

add_chip(s, 0.6, 0.6, "WEEK 7 · WRAP",
         fill=SURFACE, text_color=PURPLE, w=1.8, h=0.38)

add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
         "You built a digital notebook.",
         font=TITLE_FONT, size=38, bold=True, color=SURFACE)
add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.8),
         "JournalEntry · Timeline list · Weather picker · UIKit bridge · Custom search",
         font=BODY_FONT, size=18, color=SURFACE)

add_card(s, Inches(0.6), Inches(3.5), Inches(6.0), Inches(3.3))
add_text(s, Inches(0.85), Inches(3.65), Inches(5.6), Inches(0.4),
         "WHAT YOU SHIPPED", size=12, bold=True, color=GREEN)
bullets = [
    "JournalEntry Core Data entity with 7 attributes",
    "JournalViewModel — search + weather filter + day grouping",
    "JournalTabView — timeline list grouped by day",
    "Weather enum + WeatherPickerView (4 options)",
    "PhotoPicker — UIViewControllerRepresentable bridge",
    "AddJournalEntryView & EditJournalEntryView sheets",
]
for i, b in enumerate(bullets):
    add_text(s, Inches(0.85), Inches(4.05 + i * 0.42), Inches(5.7), Inches(0.4),
             "•  " + b, size=13, color=TEXT)

add_card(s, Inches(6.75), Inches(3.5), Inches(6.0), Inches(3.3))
add_text(s, Inches(7.0), Inches(3.65), Inches(5.6), Inches(0.4),
         "QUESTIONS TO DRIVE DISCUSSION", size=12, bold=True, color=PRIMARY)
qs = [
    "Why store Weather as a String instead of the raw enum?",
    "What breaks if we forget Codegen = Manual/None?",
    "Why JPEG-encode photos before saving?",
    "Why filter in Swift instead of via NSPredicate?",
    "What single change would make this network-first?",
]
for i, q in enumerate(qs):
    add_text(s, Inches(7.0), Inches(4.05 + i * 0.42), Inches(5.7), Inches(0.4),
             "?  " + q, size=13, color=TEXT)

add_text(s, Inches(7.0), Inches(6.3), Inches(5.6), Inches(0.4),
         "NEXT WEEK →  Water Tracker (Week 8)",
         size=12, bold=True, color=PRIMARY)

add_footer(s, 12, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Recap what the class shipped. Use the questions on the right to drive the Q&A — "
    "don't lecture, invite answers. The single most important takeaway: the UIKit "
    "bridge pattern from Lesson 7.4 is the same shape you'll use for the camera, the "
    "document picker, and every share sheet you ever build — SwiftUI + UIKit is a "
    "dance students will do for years. Tease Week 8 (Water Tracker) briefly — it's "
    "a stats/charting module that reuses the same Core Data plumbing. "
    "UI suggestion: purple hero band up top, two cards below — 'What you shipped' "
    "and 'Questions'. "
    "Animation: fade in the cards after the hero band; reveal the questions one at a time."
))


# ============================================================================
# Save
# ============================================================================
out_path = "Week07_Daily_Journal.pptx"
prs.save(out_path)
print(f"Wrote {out_path} — {len(prs.slides)} slides")
