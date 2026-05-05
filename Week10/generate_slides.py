#!/usr/bin/env python3
"""Generate Week 10 Data Export & Reports slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ──────────────────────────────────────────────────────────
GREEN   = RGBColor(0x1B, 0xB8, 0x89)
BLUE    = RGBColor(0x28, 0x7D, 0xFA)
PURPLE  = RGBColor(0x8E, 0x44, 0xAD)
ORANGE  = RGBColor(0xF3, 0x96, 0x20)
GOLD    = RGBColor(0xFF, 0xD7, 0x00)
RED     = RGBColor(0xE5, 0x47, 0x47)
TEAL    = RGBColor(0x00, 0xC9, 0xC8)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
DARK2   = RGBColor(0x16, 0x21, 0x3E)
CARD    = RGBColor(0x0F, 0x2A, 0x45)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xAA, 0xAA, 0xBB)
PINK    = RGBColor(0xE9, 0x4F, 0x97)

W = 13.33
H = 7.5


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide  = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, items, x, y, w, h, size=15, color=WHITE,
                bullet_color=GREEN, indent=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        dot = p.add_run()
        dot.text = "  • " if indent else "• "
        dot.font.size = Pt(size)
        dot.font.color.rgb = bullet_color
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_code(slide, code_lines, x, y, w, h, size=11):
    add_rect(slide, x, y, w, h, CARD)
    tb = slide.shapes.add_textbox(
        Inches(x + 0.18), Inches(y + 0.15),
        Inches(w - 0.36), Inches(h - 0.3)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = TEAL
        run.font.name  = "Courier New"


def slide_number(slide, n, total=12):
    add_text(slide, f"{n} / {total}", W - 1.2, H - 0.42, 1.0, 0.35,
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
def slide_01(prs):
    s = add_slide(prs)

    add_rect(s, 0, 0, 0.12, H, ORANGE)

    add_rect(s, 1.0, 1.2, 2.2, 0.5, ORANGE)
    add_text(s, "WEEK 10", 1.0, 1.2, 2.2, 0.5,
             size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    add_text(s, "Data Export", 1.0, 1.9, 11.0, 1.0,
             size=46, bold=True, color=WHITE)
    add_text(s, "& Reports", 1.0, 2.75, 11.0, 1.0,
             size=46, bold=True, color=ORANGE)

    add_text(s, "SmartFarmerAssistant — Generating & Sharing Reports for Farmers",
             1.0, 3.8, 11.0, 0.5, size=16, italic=True, color=GREY)

    icons = [
        ("📊", "Bar Charts"),
        ("📄", "CSV Export"),
        ("🗒️", "PDFKit"),
        ("📤", "UIActivityVC"),
        ("💰", "P&L Reports"),
        ("🌾", "Farmer UX"),
    ]
    for i, (icon, label) in enumerate(icons):
        bx = 1.0 + i * 2.05
        add_rect(s, bx, 4.7, 1.85, 1.35, DARK2)
        add_text(s, icon, bx, 4.75, 1.85, 0.6,
                 size=22, align=PP_ALIGN.CENTER)
        add_text(s, label, bx, 5.35, 1.85, 0.45,
                 size=10, color=GREY, align=PP_ALIGN.CENTER)

    slide_number(s, 1)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
def slide_02(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "📋  Agenda — Week 10", 0.4, 0.12, 12.5, 0.7,
             size=26, bold=True, color=ORANGE)

    topics = [
        ("01", ORANGE, "Custom Bar Charts",
         "GeometryReader + Shape — iOS 13+, zero dependencies"),
        ("02", BLUE,   "CSV Export",
         "Build comma-separated reports — shareable anywhere"),
        ("03", GREEN,  "PDFKit Basics",
         "PDFDocument + UIGraphicsPDFRenderer — iOS 11+"),
        ("04", PURPLE, "PDF with Custom Drawing",
         "Render text, lines, and bar charts into a PDF page"),
        ("05", TEAL,   "UIActivityViewController",
         "Share sheets — AirDrop, Mail, Files, Messages"),
        ("06", GOLD,   "End-to-End Report Flow",
         "Generate → Preview → Share in one farmer-friendly tap"),
    ]

    for i, (num, color, title, sub) in enumerate(topics):
        row_y = 1.1 + i * 1.0
        add_rect(s, 0.4, row_y, 0.55, 0.62, color)
        add_text(s, num, 0.4, row_y, 0.55, 0.62,
                 size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.1, row_y, 5.5, 0.38,
                 size=16, bold=True, color=WHITE)
        add_text(s, sub,   1.1, row_y + 0.36, 5.5, 0.3,
                 size=12, color=GREY)

    slide_number(s, 2)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Custom Bar Chart with GeometryReader + Shape
# ════════════════════════════════════════════════════════════════════════════
def slide_03(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "📊  Custom Bar Chart — GeometryReader + Shape",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=ORANGE)

    add_text(s, "Why GeometryReader + Shape instead of Swift Charts?",
             0.4, 1.0, 12.5, 0.38, size=14, bold=True, color=WHITE)
    add_bullets(s, [
        "Swift Charts requires iOS 16+ — our target is iOS 13",
        "GeometryReader gives proportional sizing inside any layout",
        "Shape protocol lets us draw filled rectangles precisely",
        "Full control over colours, spacing, labels, and animations",
    ], 0.4, 1.45, 6.0, 1.85, size=13, bullet_color=ORANGE)

    add_code(s, [
        "struct BarShape: Shape {",
        "    var heightFraction: CGFloat   // 0.0 … 1.0",
        "",
        "    func path(in rect: CGRect) -> Path {",
        "        var p = Path()",
        "        let barH = rect.height * heightFraction",
        "        p.addRect(CGRect(",
        "            x: 0,",
        "            y: rect.height - barH,",
        "            width: rect.width,",
        "            height: barH",
        "        ))",
        "        return p",
        "    }",
        "}",
    ], 6.55, 0.95, 6.5, 4.1, size=11)

    add_text(s, "BarShape draws a bottom-anchored filled rectangle — the fraction scales the height:",
             0.4, 3.4, 6.0, 0.5, size=12, color=GREY, italic=True)

    # Mini visual bar chart demo
    bar_data = [0.4, 0.7, 0.55, 0.85, 0.3, 0.65]
    bar_labels = ["ម.ក", "ក.փ", "មីនា", "មេសា", "ឧសភា", "មិថុ"]
    bar_colors = [GREEN, GREEN, RED, GREEN, RED, GREEN]
    chart_x = 0.4
    chart_y = 3.95
    chart_w = 5.8
    chart_h = 2.8
    add_rect(s, chart_x, chart_y, chart_w, chart_h, CARD)
    bar_w = 0.58
    gap   = 0.18
    for i, (frac, label, color) in enumerate(zip(bar_data, bar_labels, bar_colors)):
        bx = chart_x + 0.2 + i * (bar_w + gap)
        max_bar_h = 1.9
        actual_h  = max_bar_h * frac
        by = chart_y + chart_h - 0.5 - actual_h
        add_rect(s, bx, by, bar_w, actual_h, color)
        add_text(s, label, bx, chart_y + chart_h - 0.45, bar_w, 0.35,
                 size=9, color=GREY, align=PP_ALIGN.CENTER)

    slide_number(s, 3)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BarChartView full implementation
# ════════════════════════════════════════════════════════════════════════════
def slide_04(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "📊  BarChartView — Full Implementation",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=ORANGE)

    add_code(s, [
        "struct BarChartView: View {",
        "    let entries: [MonthlyReport]   // [(month, income, expense)]",
        "",
        "    private var maxValue: Double {",
        "        entries.flatMap { [$0.income, $0.expense] }.max() ?? 1",
        "    }",
        "",
        "    var body: some View {",
        "        GeometryReader { geo in",
        "            let barW = (geo.size.width / CGFloat(entries.count)) * 0.35",
        "            let spacing = (geo.size.width / CGFloat(entries.count)) * 0.12",
        "            HStack(alignment: .bottom, spacing: 0) {",
        "                ForEach(entries) { entry in",
        "                    VStack(spacing: 2) {",
        "                        HStack(alignment: .bottom, spacing: spacing) {",
        "                            BarShape(heightFraction: entry.income / maxValue)",
        "                                .fill(Color.green)",
        "                                .frame(width: barW)",
        "                            BarShape(heightFraction: entry.expense / maxValue)",
        "                                .fill(Color.red.opacity(0.8))",
        "                                .frame(width: barW)",
        "                        }",
        "                        .frame(height: geo.size.height - 24)",
        "                        Text(entry.monthLabel)",
        "                            .font(.system(size: 9))",
        "                            .foregroundColor(.secondary)",
        "                    }",
        "                    .frame(maxWidth: .infinity)",
        "                }",
        "            }",
        "        }",
        "    }",
        "}",
    ], 0.3, 0.95, 7.6, 6.1, size=9.5)

    # Right — design notes
    add_text(s, "Design Decisions", 8.1, 1.0, 5.0, 0.35,
             size=15, bold=True, color=WHITE)

    notes = [
        (GREEN,   "maxValue computed",
                  "Bars scale relative to the largest value — chart always fills the available height"),
        (BLUE,    "GeometryReader",
                  "Width adapts to any screen size — no hardcoded pixel values"),
        (ORANGE,  "Paired bars",
                  "Income (green) + Expense (red) side by side for instant P&L comparison"),
        (TEAL,    "MonthLabel below",
                  "ខ្មែរ month abbreviations keep the UI compact on small screens"),
    ]

    for i, (color, title, desc) in enumerate(notes):
        ny = 1.45 + i * 1.4
        add_rect(s, 8.1, ny, 5.0, 1.28, CARD)
        add_rect(s, 8.1, ny, 0.08, 1.28, color)
        add_text(s, title, 8.28, ny + 0.1, 4.7, 0.38,
                 size=13, bold=True, color=color)
        add_text(s, desc, 8.28, ny + 0.5, 4.7, 0.68,
                 size=11, color=GREY)

    slide_number(s, 4)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CSV Export
# ════════════════════════════════════════════════════════════════════════════
def slide_05(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "📄  CSV Export — Shareable Spreadsheet Data",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=BLUE)

    add_code(s, [
        "// Reports/CSVExporter.swift",
        "struct CSVExporter {",
        "",
        "    static func makeTransactionCSV(",
        "        _ transactions: [Transaction]",
        "    ) -> String {",
        '        var lines = ["Date,Type,Category,Amount,Note"]',
        "        let df = DateFormatter()",
        '        df.dateStyle = .short',
        "        for t in transactions {",
        "            let date   = df.string(from: t.date ?? Date())",
        "            let type_  = t.isIncome ? \"Income\" : \"Expense\"",
        "            let cat    = t.category ?? \"\"",
        "            let amount = String(format: \"%.2f\", t.amount)",
        "            let note   = (t.note ?? \"\").replacingOccurrences(",
        "                             of: \",\", with: \";\")  // escape commas",
        '            lines.append("\\(date),\\(type_),\\(cat),\\(amount),\\(note)")',
        "        }",
        '        return lines.joined(separator: "\\n")',
        "    }",
        "",
        "    static func writeToTempFile(_ csv: String,",
        "                                named: String) -> URL? {",
        "        let url = FileManager.default",
        "            .temporaryDirectory",
        '            .appendingPathComponent(named + ".csv")',
        "        try? csv.write(to: url, atomically: true,",
        '                       encoding: .utf8)',
        "        return url",
        "    }",
        "}",
    ], 0.3, 0.95, 7.1, 6.1, size=9.5)

    # Right column
    add_text(s, "CSV Usage", 7.6, 1.0, 5.5, 0.35, size=15, bold=True, color=WHITE)
    add_code(s, [
        "let csv = CSVExporter",
        "    .makeTransactionCSV(transactions)",
        "guard let url = CSVExporter",
        '    .writeToTempFile(csv, named: "farm_report")',
        "else { return }",
        "// → pass url to UIActivityViewController",
    ], 7.6, 1.42, 5.5, 1.9, size=11)

    add_text(s, "Key Rules", 7.6, 3.45, 5.5, 0.35, size=15, bold=True, color=WHITE)
    add_bullets(s, [
        'Escape commas in text fields (",") → replace with ";"',
        "Use .temporaryDirectory — no permission needed",
        "atomically: true prevents partial writes on crash",
        "UTF-8 ensures Khmer text survives the round trip",
        "Open in Excel/Google Sheets on desktop immediately",
    ], 7.6, 3.88, 5.5, 2.7, size=12, bullet_color=BLUE)

    slide_number(s, 5)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PDFKit Basics
# ════════════════════════════════════════════════════════════════════════════
def slide_06(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🗒️  PDFKit — Building a PDF Report (iOS 11+)",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=GREEN)

    # Left: approach overview
    add_text(s, "Two ways to create PDFs on iOS:", 0.4, 1.0, 6.0, 0.38,
             size=14, bold=True, color=WHITE)

    approaches = [
        (ORANGE, "UIGraphicsPDFRenderer  ✅ Recommended",
         ["Draw into a CGContext — text, lines, images",
          "Works on iOS 10+ — no PDFKit import needed",
          "Best for custom-designed report pages"]),
        (BLUE,   "PDFDocument + PDFPage",
         ["Compose pages from UIView snapshots",
          "Better for existing SwiftUI views → PDF",
          "Requires PDFKit (iOS 11+) import"]),
    ]

    for i, (color, title, points) in enumerate(approaches):
        ay = 1.45 + i * 2.5
        add_rect(s, 0.35, ay, 6.0, 2.3, CARD)
        add_rect(s, 0.35, ay, 0.08, 2.3, color)
        add_text(s, title, 0.55, ay + 0.1, 5.7, 0.38,
                 size=13, bold=True, color=color)
        for j, pt in enumerate(points):
            add_text(s, "• " + pt, 0.55, ay + 0.55 + j * 0.5, 5.7, 0.44,
                     size=11.5, color=WHITE)

    add_code(s, [
        "// Reports/PDFGenerator.swift",
        "import PDFKit",
        "",
        "struct PDFGenerator {",
        "",
        "    static func makeReportPDF(",
        "        title: String,",
        "        rows: [(String, String)]",
        "    ) -> Data {",
        "        let pageRect = CGRect(x: 0, y: 0,",
        "                             width: 595, height: 842)  // A4",
        "        let renderer = UIGraphicsPDFRenderer(bounds: pageRect)",
        "",
        "        return renderer.pdfData { ctx in",
        "            ctx.beginPage()",
        "            drawTitle(title, in: pageRect)",
        "            drawTable(rows, in: pageRect)",
        "        }",
        "    }",
        "}",
    ], 6.65, 0.95, 6.4, 4.95, size=10)

    add_text(s, "A4 = 595 × 842 pt (72 dpi).  US Letter = 612 × 792 pt.",
             0.4, 6.4, 12.5, 0.38, size=11, italic=True, color=GREY)

    slide_number(s, 6)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PDF Custom Drawing Detail
# ════════════════════════════════════════════════════════════════════════════
def slide_07(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🖊️  Drawing Text & Tables into a PDF Page",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=TEAL)

    add_code(s, [
        "// Draw a title at the top of the page",
        "private static func drawTitle(_ title: String,",
        "                              in rect: CGRect) {",
        "    let attrs: [NSAttributedString.Key: Any] = [",
        "        .font: UIFont.boldSystemFont(ofSize: 22),",
        "        .foregroundColor: UIColor.black",
        "    ]",
        "    title.draw(at: CGPoint(x: 40, y: 40), withAttributes: attrs)",
        "}",
        "",
        "// Draw a two-column table",
        "private static func drawTable(",
        "    _ rows: [(String, String)], in rect: CGRect) {",
        "    let rowH: CGFloat = 24",
        "    var y: CGFloat = 90",
        "    let labelAttrs: [NSAttributedString.Key: Any] = [",
        "        .font: UIFont.systemFont(ofSize: 13),",
        "        .foregroundColor: UIColor.darkGray",
        "    ]",
        "    let valueAttrs: [NSAttributedString.Key: Any] = [",
        "        .font: UIFont.boldSystemFont(ofSize: 13),",
        "        .foregroundColor: UIColor.black",
        "    ]",
        "    for (label, value) in rows {",
        "        label.draw(at: CGPoint(x: 40, y: y),",
        "                   withAttributes: labelAttrs)",
        "        value.draw(at: CGPoint(x: 300, y: y),",
        "                   withAttributes: valueAttrs)",
        "        y += rowH",
        "    }",
        "}",
    ], 0.3, 0.95, 7.6, 6.1, size=9.5)

    add_text(s, "Output Preview", 8.1, 1.0, 5.0, 0.35,
             size=14, bold=True, color=WHITE)

    # Simulated A4 page preview
    page_x, page_y = 8.1, 1.42
    page_w, page_h = 5.0, 5.5
    add_rect(s, page_x, page_y, page_w, page_h, WHITE)
    add_text(s, "Farm Monthly Report", page_x + 0.25, page_y + 0.18,
             page_w - 0.4, 0.4, size=13, bold=True, color=DARK)
    add_rect(s, page_x + 0.15, page_y + 0.65, page_w - 0.3, 0.02,
             RGBColor(0xCC, 0xCC, 0xCC))

    table_rows = [
        ("ខែ:", "ឧសភា 2026"),
        ("ចំណូលសរុប:", "$1,850.00"),
        ("ចំណាយសរុប:", "$1,120.00"),
        ("ចំណេញ:", "$730.00"),
        ("ប្រតិបត្តិការ:", "24 entries"),
    ]
    for i, (label, value) in enumerate(table_rows):
        ry = page_y + 0.82 + i * 0.62
        bg = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE
        add_rect(s, page_x + 0.15, ry, page_w - 0.3, 0.55, bg)
        add_text(s, label, page_x + 0.25, ry + 0.1,
                 2.0, 0.35, size=10, color=DARK)
        vcolor = GREEN if "ចំណេញ" in label else DARK
        add_text(s, value, page_x + 2.35, ry + 0.1,
                 2.5, 0.35, size=10, bold=True, color=vcolor)

    slide_number(s, 7)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — UIActivityViewController
# ════════════════════════════════════════════════════════════════════════════
def slide_08(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "📤  UIActivityViewController — Sharing Reports",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=PURPLE)

    add_code(s, [
        "// Sharing a file URL (CSV or PDF)",
        "func shareFile(url: URL) {",
        "    let avc = UIActivityViewController(",
        "        activityItems: [url],",
        "        applicationActivities: nil",
        "    )",
        "    // iPad needs a popover source",
        "    if let popover = avc.popoverPresentationController {",
        "        popover.sourceView = UIApplication.shared",
        "            .windows.first?.rootViewController?.view",
        "        popover.sourceRect = CGRect(x: 0, y: 0,",
        "                                   width: 1, height: 1)",
        "    }",
        "    UIApplication.shared.windows.first?",
        "        .rootViewController?",
        "        .present(avc, animated: true)",
        "}",
        "",
        "// Share raw Data (PDF bytes) directly",
        "func sharePDFData(_ data: Data, filename: String) {",
        "    let url = FileManager.default",
        "        .temporaryDirectory",
        "        .appendingPathComponent(filename + \".pdf\")",
        "    try? data.write(to: url)",
        "    shareFile(url: url)",
        "}",
    ], 0.3, 0.95, 7.1, 6.1, size=10)

    add_text(s, "What Farmers See", 7.6, 1.0, 5.5, 0.35,
             size=15, bold=True, color=WHITE)
    add_bullets(s, [
        "Share sheet slides up automatically",
        "AirDrop → send to laptop instantly",
        "Mail → attach to email for agronomist",
        "Files app → save to iCloud / local",
        "WhatsApp / Telegram (if installed)",
        "Print → AirPrint-enabled printers",
    ], 7.6, 1.45, 5.5, 3.2, size=12.5, bullet_color=PURPLE)

    add_text(s, "SwiftUI Bridge (iOS 13+):", 7.6, 4.75, 5.5, 0.35,
             size=13, bold=True, color=WHITE)
    add_code(s, [
        "struct ShareSheet: UIViewControllerRepresentable {",
        "    let items: [Any]",
        "    func makeUIViewController(context: Context)",
        "        -> UIActivityViewController {",
        "        UIActivityViewController(",
        "            activityItems: items,",
        "            applicationActivities: nil)",
        "    }",
        "    func updateUIViewController(_ vc:",
        "        UIActivityViewController, context: Context) {}",
        "}",
    ], 7.55, 5.15, 5.55, 1.95, size=9.5)

    slide_number(s, 8)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ReportTabView  (end-to-end flow)
# ════════════════════════════════════════════════════════════════════════════
def slide_09(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🌾  ReportTabView — End-to-End Farmer Flow",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=GOLD)

    add_code(s, [
        "struct ReportTabView: View {",
        "    @FetchRequest(entity: Transaction.entity(),",
        "        sortDescriptors: [NSSortDescriptor(",
        "            keyPath: \\Transaction.date, ascending: false)])",
        "    private var transactions: FetchedResults<Transaction>",
        "",
        "    @StateObject private var vm = ReportViewModel()",
        "    @State private var showShare   = false",
        "    @State private var shareItems: [Any] = []",
        "    @State private var isExporting = false",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            ScrollView {",
        "                VStack(spacing: 20) {",
        "                    MonthPickerView(",
        "                        selectedMonth: $vm.selectedMonth,",
        "                        canGoNext: vm.canGoToNextMonth,",
        "                        onPrevious: vm.goToPreviousMonth,",
        "                        onNext: vm.goToNextMonth)",
        "                    FarmCard(title: \"ចំណូល / ចំណាយ ប្រចាំខែ\",",
        "                             icon: \"chart.bar.fill\",",
        "                             iconColor: .orange) {",
        "                        BarChartView(entries:",
        "                            Array(monthlyReports.suffix(6)))",
        "                            .frame(height: 200)",
        "                    }",
        "                    SummaryCardView(report: selectedReport)",
        "                    exportButtons   // PrimaryButton CSV + PDF",
        "                }",
        "            }",
        "            .navigationTitle(\"របាយការណ៍\")",
        "            .sheet(isPresented: $showShare) {",
        "                ShareSheet(items: shareItems) }",
        "            .loadingOverlay(isLoading: isExporting)",
        "        }",
        "    }",
        "}",
    ], 0.3, 0.95, 7.6, 6.1, size=9)

    # Right — flow diagram
    add_text(s, "Flow", 8.1, 1.0, 5.0, 0.35, size=15, bold=True, color=WHITE)
    steps = [
        (ORANGE, "1. Select Month",    "MonthPickerView ‹/› → vm.selectedMonth → SummaryCardView updates"),
        (BLUE,   "2. View Chart",      "BarChartView shows suffix(6) months; empty state if no transactions"),
        (GREEN,  "3. Tap CSV / PDF",   "PrimaryButton triggers exportCSV() / exportPDF() — background thread"),
        (PURPLE, "4. Build & Write",   "CSVExporter/PDFGenerator → PDFGenerator.writeTempFile → URL"),
        (TEAL,   "5. Share Sheet",     "ShareSheet(items: [url]) wraps UIActivityViewController"),
        (GOLD,   "6. Farmer Shares",   "AirDrop / WhatsApp / Mail / Files / Print"),
    ]
    for i, (color, step, desc) in enumerate(steps):
        sy = 1.45 + i * 0.98
        add_rect(s, 8.1, sy, 5.0, 0.88, CARD)
        add_rect(s, 8.1, sy, 0.08, 0.88, color)
        add_text(s, step, 8.28, sy + 0.06, 4.7, 0.32,
                 size=12, bold=True, color=color)
        add_text(s, desc, 8.28, sy + 0.42, 4.7, 0.38,
                 size=10, color=GREY)

    slide_number(s, 9)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ReportViewModel & MonthlyReport model
# ════════════════════════════════════════════════════════════════════════════
def slide_10(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🧠  ReportViewModel & MonthlyReport Model",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=TEAL)

    add_code(s, [
        "// Reports/Models/MonthlyReport.swift",
        "struct MonthlyReport: Identifiable {",
        "    let id = UUID()",
        "    let month: Date; let income: Double; let expense: Double",
        "    var profit: Double { income - expense }",
        "    var monthLabel: String { /* MMM */ ... }",
        "    var monthYearLabel: String { /* MMMM yyyy */ ... }",
        "}",
        "",
        "// Reports/ViewModels/ReportViewModel.swift",
        "class ReportViewModel: ObservableObject {",
        "    // Canonical month start — avoids boundary mismatches",
        "    @Published var selectedMonth: Date = {",
        "        Calendar.current",
        "            .dateInterval(of: .month, for: Date())!.start",
        "    }()",
        "",
        "    func buildMonthlyReports(_ txns: [Transaction]) -> [MonthlyReport] {",
        "        let cal = Calendar.current",
        "        let grouped = Dictionary(grouping: txns) {",
        "            cal.dateInterval(of: .month, for: $0.date ?? Date())!.start",
        "        }",
        "        return grouped.map { month, txns in",
        "            MonthlyReport(month: month,",
        "                income:  txns.filter { $0.isIncome }.map(\\.amount).reduce(0,+),",
        "                expense: txns.filter { $0.isExpense }.map(\\.amount).reduce(0,+))",
        "        }.sorted { $0.month < $1.month }",
        "    }",
        "",
        "    func selectedReport(from r: [MonthlyReport]) -> MonthlyReport {",
        "        r.first { Calendar.current.isDate($0.month,",
        "                    equalTo: selectedMonth, toGranularity: .month) }",
        "        ?? MonthlyReport(month: selectedMonth, income: 0, expense: 0)",
        "    }",
        "",
        "    func goToPreviousMonth() { selectedMonth =",
        "        Calendar.current.date(byAdding: .month, value: -1, to: selectedMonth)!",
        "    }",
        "    var canGoToNextMonth: Bool {",
        "        Calendar.current.date(byAdding: .month, value:1, to: selectedMonth)! <= Date()",
        "    }",
        "    // summaryRows uses Double.formattedCurrency from Formatters.swift",
        "}",
    ], 0.3, 0.95, 7.6, 6.1, size=9.5)

    add_text(s, "Key Design Choices", 8.1, 1.0, 5.0, 0.35,
             size=14, bold=True, color=WHITE)

    choices = [
        (GREEN,  "Month start not Date()",
                 "dateInterval(of:for:).start avoids off-by-one near midnight / month boundary"),
        (ORANGE, "$0.isExpense not !$0.isIncome",
                 "Transaction has both computed properties — explicit is clearer and safer"),
        (BLUE,   "selectedReport returns zero report",
                 "Chart & summary never crash on an empty month — $0.00 displays gracefully"),
        (TEAL,   "canGoToNextMonth guard",
                 "Prevents navigating into the future where no data will ever exist"),
    ]
    for i, (color, title, desc) in enumerate(choices):
        cy = 1.42 + i * 1.45
        add_rect(s, 8.1, cy, 5.0, 1.3, CARD)
        add_rect(s, 8.1, cy, 0.08, 1.3, color)
        add_text(s, title, 8.28, cy + 0.1, 4.7, 0.35,
                 size=13, bold=True, color=color)
        add_text(s, desc, 8.28, cy + 0.52, 4.7, 0.6,
                 size=11, color=GREY)

    slide_number(s, 10)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — iOS API Quick Reference
# ════════════════════════════════════════════════════════════════════════════
def slide_11(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "⚠️  iOS API Rules — Week 10 Quick Reference",
             0.4, 0.12, 12.5, 0.7, size=24, bold=True, color=GOLD)

    headers = ["Feature", "❌ Wrong (newer API)", "✅ Correct (iOS 13+)"]
    col_x   = [0.35, 4.3, 8.7]
    col_w   = [3.75, 4.2, 4.4]

    add_rect(s, 0.3, 1.0, 12.8, 0.42, DARK2)
    for j, h in enumerate(headers):
        add_text(s, h, col_x[j] + 0.1, 1.04, col_w[j] - 0.1, 0.35,
                 size=12, bold=True, color=TEAL)

    rows = [
        ("Bar Charts",
         "Swift Charts (iOS 16+)",
         "GeometryReader + BarShape: Shape (iOS 13+)"),
        ("Date grouping",
         "dateComponents(_:from:to:) careless",
         "Calendar.dateInterval(of: .month, for:)"),
        ("File sharing",
         "ShareLink (iOS 16+)",
         "UIActivityViewController wrapped in UIViewControllerRepresentable"),
        ("Photo in PDF",
         "UIGraphicsImageRenderer (fine, but check bounds)",
         "UIImage.draw(in:) inside UIGraphicsPDFRenderer context"),
        ("Navigation",
         "NavigationStack (iOS 16+)",
         "NavigationView (existing convention)"),
        ("PDF from SwiftUI view",
         "ImageRenderer (iOS 16+)",
         "UIGraphicsPDFRenderer + Core Text / NSString.draw(at:)"),
    ]

    for i, (feat, wrong, correct) in enumerate(rows):
        row_y = 1.46 + i * 0.98
        bg = CARD if i % 2 == 0 else DARK2
        add_rect(s, 0.3, row_y, 12.8, 0.92, bg)
        add_text(s, feat, col_x[0] + 0.1, row_y + 0.12,
                 col_w[0] - 0.1, 0.7, size=11, bold=True, color=WHITE)
        add_text(s, wrong, col_x[1] + 0.1, row_y + 0.12,
                 col_w[1] - 0.1, 0.7, size=10.5, color=RED)
        add_text(s, correct, col_x[2] + 0.1, row_y + 0.12,
                 col_w[2] - 0.1, 0.7, size=10.5, color=GREEN)

    slide_number(s, 11)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary & What's Next
# ════════════════════════════════════════════════════════════════════════════
def slide_12(prs):
    s = add_slide(prs)
    add_rect(s, 0, 0, W, 0.9, DARK2)
    add_text(s, "🎯  Summary & What's Next",
             0.4, 0.12, 12.5, 0.7, size=26, bold=True, color=GREEN)

    add_text(s, "Week 10 Deliverables", 0.4, 1.0, 6.2, 0.38,
             size=16, bold=True, color=WHITE)
    deliverables = [
        "Reports/Models/MonthlyReport.swift — data model",
        "Reports/ViewModels/ReportViewModel.swift — grouping + P&L",
        "Reports/Services/CSVExporter.swift — CSV generation + temp file",
        "Reports/Services/PDFGenerator.swift — UIGraphicsPDFRenderer report",
        "Reports/Views/BarChartView.swift — GeometryReader + BarShape",
        "Reports/Views/ReportTabView.swift — combined UI + share flow",
        "Shared/ShareSheet.swift — UIViewControllerRepresentable wrapper",
    ]
    add_bullets(s, deliverables, 0.4, 1.45, 6.3, 4.2,
                size=12, bullet_color=GREEN)

    add_text(s, "Coming Up — Week 11 / 12", 7.0, 1.0, 5.9, 0.38,
             size=16, bold=True, color=WHITE)
    next_items = [
        "CloudKit sync — share farm data across devices",
        "WidgetKit — farm summary on the Home Screen",
        "App Store submission checklist",
        "TestFlight internal distribution",
        "Performance profiling with Instruments",
    ]
    add_bullets(s, next_items, 7.0, 1.45, 5.9, 3.5,
                size=12.5, bullet_color=BLUE)

    add_rect(s, 0.3, 5.85, 12.75, 1.25, ORANGE)
    add_text(s, "💡  Key Takeaway",
             0.55, 5.9, 12.3, 0.38, size=14, bold=True, color=DARK)
    add_text(s,
             "A farmer who can export their own P&L report and share it in WhatsApp "
             "trusts the app. CSV + PDF + UIActivityViewController — three tools, "
             "one tap, maximum farmer value.",
             0.55, 6.28, 12.3, 0.65, size=12.5, color=DARK)

    slide_number(s, 12)
    return s


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)

    out = "Week10_Data_Export_Reports.pptx"
    prs.save(out)
    print(f"✅  Saved → {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
