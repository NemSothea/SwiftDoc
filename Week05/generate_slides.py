#!/usr/bin/env python3
"""Generate Week 05 — Calendar & Reminders slide deck (Khmer)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BLUE   = RGBColor(0x28, 0x7D, 0xFA)
GREEN  = RGBColor(0x1B, 0xB8, 0x89)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
ORANGE = RGBColor(0xF3, 0x96, 0x20)
TEAL   = RGBColor(0x00, 0xC9, 0xC8)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
DARK2  = RGBColor(0x16, 0x21, 0x3E)
CARD   = RGBColor(0x0F, 0x2A, 0x45)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xAA, 0xAA, 0xBB)
RED    = RGBColor(0xE5, 0x47, 0x47)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
W, H   = 13.33, 7.5

def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in s.placeholders: ph._element.getparent().remove(ph._element)
    f = s.background.fill; f.solid(); f.fore_color.rgb = DARK
    return s

def add_text(s, text, x, y, w, h, size=16, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def add_rect(s, x, y, w, h, fc, lc=None):
    sh = s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if lc: sh.line.color.rgb = lc
    else: sh.line.fill.background()
    return sh

def add_bullets(s, items, x, y, w, h, size=14, color=WHITE, bc=BLUE):
    tb = s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r1=p.add_run(); r1.text="• "; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=bc
        r2=p.add_run(); r2.text=item; r2.font.size=Pt(size); r2.font.color.rgb=color

def add_code(s, lines, x, y, w, h, size=10):
    add_rect(s, x, y, w, h, CARD)
    tb = s.shapes.add_textbox(Inches(x+.18),Inches(y+.15),Inches(w-.36),Inches(h-.3))
    tf = tb.text_frame; tf.word_wrap = False
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.color.rgb=TEAL; r.font.name="Courier New"

def hdr(s, title, accent=BLUE):
    add_rect(s,0,0,W,.9,DARK2); add_rect(s,0,0,.06,.9,accent)
    add_text(s,title,.25,.1,12.5,.7,size=22,bold=True,color=accent)

def snum(s, n, total=12):
    add_text(s,f"{n}/{total}",12.3,7.1,.9,.3,size=10,color=GREY,align=PP_ALIGN.RIGHT)

# ── Slide 1: Title ──────────────────────────────────────────────────────────
def s1(prs):
    s = add_slide(prs)
    add_rect(s,0,0,.12,H,BLUE)
    add_rect(s,.4,.4,2.2,.42,BLUE)
    add_text(s,"Week 05 · SmartFarmer",.4,.4,2.2,.42,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(s,"ប្រតិទិន & ការរំលឹក",.4,1.0,12.0,1.0,size=46,bold=True,color=BLUE)
    add_text(s,"Calendar & Reminders · Project 2",.4,1.95,12.0,.7,size=28,bold=True,color=WHITE)
    add_text(s,"iOS 13+  ·  CoreData  ·  UNUserNotificationCenter  ·  DatePicker",.4,2.85,12.0,.45,size=14,italic=True,color=GREY)
    add_rect(s,.4,3.45,10.5,.04,BLUE)
    icons=[("📅","DatePicker"),("🔔","Notifications"),("🗂️","FarmActivity"),("🔗","Deep-link"),("🌾","Khmer UI"),("⏰","Scheduling")]
    for i,(e,l) in enumerate(icons):
        cx=.4+i*2.1; add_rect(s,cx,3.65,1.9,1.0,CARD)
        add_text(s,e,cx,3.7,1.9,.45,size=20,align=PP_ALIGN.CENTER)
        add_text(s,l,cx,4.18,1.9,.4,size=10,color=GREY,align=PP_ALIGN.CENTER)
    add_text(s,"SmartFarmer Assistant · ភ្នំពេញ 2026",.4,6.9,12.0,.4,size=11,italic=True,color=GREY)
    snum(s,1)

# ── Slide 2: Agenda ─────────────────────────────────────────────────────────
def s2(prs):
    s = add_slide(prs)
    hdr(s,"📋  មាតិកា Week 05 · Agenda")
    topics=[
        (BLUE,  "01","Calendar UI","DatePicker + LazyVGrid + DayCellView"),
        (GREEN, "02","FarmActivity Entity","CoreData model + 6 attributes"),
        (PURPLE,"03","ការអនុញ្ញាត Notification","UNUserNotificationCenter.requestAuthorization"),
        (ORANGE,"04","ការកំណត់ Notification","UNTimeIntervalNotificationTrigger + UNCalendarTrigger"),
        (TEAL,  "05","Deep Linking","NotificationCenter observer → open activity"),
        (BLUE,  "06","Mini-Project","ប្រតិទិន + ការជូនដំណឹងពេញ"),
    ]
    for i,(c,n,kh,en) in enumerate(topics):
        col=i%2; row=i//2
        cx=.35+col*6.5; cy=1.1+row*2.0
        add_rect(s,cx,cy,6.1,1.75,CARD); add_rect(s,cx,cy,.08,1.75,c)
        add_rect(s,cx+.18,cy+.38,.6,.6,c)
        add_text(s,n,cx+.18,cy+.38,.6,.6,size=13,bold=True,color=DARK,align=PP_ALIGN.CENTER)
        add_text(s,kh,cx+.95,cy+.25,4.9,.5,size=15,bold=True,color=WHITE)
        add_text(s,en,cx+.95,cy+.85,4.9,.45,size=11,color=GREY)
    snum(s,2)

# ── Slide 3: Calendar UI ────────────────────────────────────────────────────
def s3(prs):
    s = add_slide(prs)
    hdr(s,"📅  Calendar UI · DatePicker + LazyVGrid")
    add_bullets(s,[
        "DatePicker → ជ្រើសសាររបស់ user (ខែ/ថ្ងៃ)",
        "LazyVGrid(columns: [GridItem(.flexible()) × 7]) → ក្រឡា ៧ ថ្ងៃ/ជួរ",
        "DayCellView → dot indicator: ● activity  ● today  ● selected",
        "selectedDate @State → drive @FetchRequest filter",
        "NavigationView + NavigationLink → detail screen (iOS 13+)",
        "Khmer weekday labels: អា ច អ ពុ ព្រ សុ ស",
    ],.35,1.1,6.3,5.6,bc=BLUE)
    # spacer line removed
    add_text(s,"ឧទាហរណ៍ · Calendar Layout",6.85,1.1,6.1,.4,size=12,bold=True,color=BLUE)
    add_code(s,[
        "struct CalendarTabView: View {",
        "    @State private var selectedDate = Date()",
        "    @FetchRequest(... predicate for date)",
        "    var activities: FetchedResults<FarmActivity>",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            VStack {",
        "                MonthHeaderView(date: $selectedDate)",
        "                WeekdayHeaderView()  // អា ច អ ពុ ព្រ សុ ស",
        "                LazyVGrid(columns: columns) {",
        "                    ForEach(daysInMonth, id: \\.self) { day in",
        "                        DayCellView(date: day,",
        "                            isSelected: day == selectedDate,",
        "                            hasActivity: activities",
        "                              .contains { sameDay($0.date, day) })",
        "                        .onTapGesture { selectedDate = day }",
        "                    }",
        "                }",
        "                ActivityListView(date: selectedDate)",
        "            }",
        "        }",
        "    }",
        "}",
    ],6.85,1.6,6.1,5.75,size=9.5)
    snum(s,3)

# ── Slide 4: FarmActivity Entity ────────────────────────────────────────────
def s4(prs):
    s = add_slide(prs)
    hdr(s,"🗂️  FarmActivity Entity · CoreData Model")
    attrs=[
        ("id","UUID","អត្តសញ្ញាណតែមួយ"),
        ("title","String?","ចំណងជើងសកម្មភាព"),
        ("activityType","String?","ប្រភេទ: ដាំ, ស្រោច, ច្រូត..."),
        ("date","Date?","កាលបរិច្ឆេទ + ម៉ោង"),
        ("notes","String?","កំណត់ចំណាំ optional"),
        ("isCompleted","Bool","ស្ថានភាព: ធ្វើរួច / មិនទាន់"),
        ("reminderEnabled","Bool","ត្រូវការ notification?"),
    ]
    for i,(name,typ,desc) in enumerate(attrs):
        cy=1.1+i*.85
        add_rect(s,.35,cy,12.6,.75,CARD if i%2==0 else DARK2)
        add_rect(s,.35,cy,.06,.75,BLUE)
        add_code(s,[name],  .55,cy+.1,3.2,.55,size=12)
        add_code(s,[typ],   3.9,cy+.1,2.4,.55,size=12)
        add_text(s,desc,   6.5,cy+.15,6.3,.45,size=12,color=GREY)
    add_rect(s,.35,7.1,12.6,.32,DARK2)
    add_text(s,"⚙️  Codegen → Manual/None   |   reminderEnabled → triggers UNUserNotificationCenter scheduling",
             .5,7.12,12.2,.28,size=10,color=GREY)
    snum(s,4)

# ── Slide 5: Notification Permission ────────────────────────────────────────
def s5(prs):
    s = add_slide(prs)
    hdr(s,"🔔  ការអនុញ្ញាត Notification · Permission Flow")
    steps=[
        (GREEN, "①","App ចាប់ផ្ដើម","AppDelegate / .task modifier → request permission"),
        (BLUE,  "②","requestAuthorization","UNUserNotificationCenter.current().requestAuthorization(options: [.alert,.badge,.sound])"),
        (PURPLE,"③","iOS Dialog","system shows permission dialog → user taps Allow/Don't Allow"),
        (ORANGE,"④","ផ្ទៀងផ្ទាត់","check granted Bool → save to UserDefaults"),
        (TEAL,  "⑤","Schedule","if granted → call scheduleNotification(for: activity)"),
    ]
    for i,(c,n,kh,en) in enumerate(steps):
        cy=1.1+i*1.22
        add_rect(s,.35,cy,12.6,1.08,CARD); add_rect(s,.35,cy,.08,1.08,c)
        add_rect(s,.55,cy+.25,.55,.55,c)
        add_text(s,n,.55,cy+.25,.55,.55,size=14,bold=True,color=DARK,align=PP_ALIGN.CENTER)
        add_text(s,kh,1.25,cy+.1,11.3,.42,size=14,bold=True,color=WHITE)
        add_text(s,en,1.25,cy+.56,11.3,.42,size=11,color=GREY)
    add_code(s,[
        "UNUserNotificationCenter.current().requestAuthorization(",
        "    options: [.alert, .badge, .sound]",
        ") { granted, error in",
        "    DispatchQueue.main.async {",
        "        self.notificationsEnabled = granted",
        "    }",
        "}",
    ],.35,7.15,12.6,.8,size=10) # below steps
    snum(s,5)

# ── Slide 6: Scheduling Notifications ───────────────────────────────────────
def s6(prs):
    s = add_slide(prs)
    hdr(s,"⏰  ការកំណត់ Notification · Scheduling")
    add_text(s,"scheduleNotification(for:) · Function",.35,1.05,7.5,.4,size=12,bold=True,color=BLUE)
    add_code(s,[
        "func scheduleNotification(for activity: FarmActivity) {",
        "    guard activity.reminderEnabled,",
        "          let date = activity.date,",
        "          let title = activity.title else { return }",
        "",
        "    // Content",
        "    let content = UNMutableNotificationContent()",
        '    content.title = "🌾 SmartFarmer"',
        "    content.body  = title",
        "    content.sound = .default",
        "",
        "    // Trigger — 1 day before",
        "    let triggerDate = Calendar.current.date(",
        "        byAdding: .day, value: -1, to: date)!",
        "    let components = Calendar.current",
        "        .dateComponents([.year,.month,.day,.hour,.minute],",
        "                        from: triggerDate)",
        "    let trigger = UNCalendarNotificationTrigger(",
        "        dateMatching: components, repeats: false)",
        "",
        "    // Request",
        "    let id = activity.id?.uuidString ?? UUID().uuidString",
        "    let request = UNNotificationRequest(",
        "        identifier: id, content: content, trigger: trigger)",
        "    UNUserNotificationCenter.current().add(request)",
        "}",
    ],.35,1.55,7.5,5.8,size=9.5)
    notes=[
        (BLUE,  "UNMutableNotificationContent","ដាក់ title, body, sound\nstrings — ខ្លីច្បាស់"),
        (GREEN, "UNCalendarNotificationTrigger","Fire ត្រឹមត្រូវតាម date\n(ត្រូវ dateComponents ពេញ)"),
        (PURPLE,"UNNotificationRequest","identifier = activity UUID\n→ អាច cancel បាន"),
        (ORANGE,"1 day before","ជូនដំណឹងមុន 1 ថ្ងៃ\nដើម្បីឲ្យ farmer រៀបចំ"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.55+i*1.45
        add_rect(s,8.05,cy,5.1,1.3,CARD); add_rect(s,8.05,cy,.08,1.3,c)
        add_text(s,t,8.28,cy+.1,4.6,.4,size=12,bold=True,color=c)
        add_text(s,d,8.28,cy+.55,4.6,.65,size=11,color=WHITE)
    snum(s,6)

# ── Slide 7: Deep Linking ───────────────────────────────────────────────────
def s7(prs):
    s = add_slide(prs)
    hdr(s,"🔗  Deep Linking · Tap Notification → Open Activity")
    add_bullets(s,[
        "UNUserNotificationCenterDelegate.userNotificationCenter(_:didReceive:) — called when user taps notification",
        "Extract activityID from notification.request.identifier",
        "Post NotificationCenter notification with the ID",
        "CalendarTabView listens via .onReceive(NotificationCenter.default.publisher(for:))",
        "Set selectedDate = activity.date → SelectedActivity = activity → sheet opens",
        "iOS 13+ pattern: NotificationCenter.default.post(name:object:) + .onReceive()",
    ],.35,1.1,12.6,4.2,size=13,bc=BLUE)
    add_code(s,[
        "// AppDelegate / SceneDelegate",
        "func userNotificationCenter(_ center: UNUserNotificationCenter,",
        "    didReceive response: UNNotificationResponse,",
        "    withCompletionHandler completionHandler: @escaping () -> Void) {",
        "    let id = response.notification.request.identifier",
        "    NotificationCenter.default.post(",
        '        name: .openActivity, object: id)',
        "    completionHandler()",
        "}",
        "",
        "// CalendarTabView",
        ".onReceive(NotificationCenter.default.publisher(",
        "    for: .openActivity)) { note in",
        "    if let id = note.object as? String {",
        "        openActivity(id: id)",
        "    }",
        "}",
    ],.35,5.5,12.6,1.85,size=10)
    snum(s,7)

# ── Slide 8: End-to-End Flow ─────────────────────────────────────────────────
def s8(prs):
    s = add_slide(prs)
    hdr(s,"🔄  End-to-End Flow · User → Notification → Deep-link")
    nodes=[
        (GREEN, "User saves\nActivity",      .35,1.1 ),
        (BLUE,  "CoreData\nstores",          3.0, 1.1 ),
        (PURPLE,"Schedule\nNotification",    5.65,1.1 ),
        (ORANGE,"iOS delivers\nat trigger",  8.3, 1.1 ),
        (TEAL,  "User taps\nnotification",   10.95,1.1),
    ]
    for c,lbl,cx,cy in nodes:
        add_rect(s,cx,cy,2.4,1.3,CARD); add_rect(s,cx,cy,2.4,.08,c)
        add_text(s,lbl,cx+.1,cy+.2,2.2,.9,size=11,color=WHITE,align=PP_ALIGN.CENTER)
        if cx<10.95:
            add_text(s,"→",cx+2.4,cy+.55,.5,.4,size=16,bold=True,color=GREY,align=PP_ALIGN.CENTER)
    add_rect(s,10.95,2.4,.12,1.5,TEAL)
    add_rect(s,.35,2.4,10.7,.08,DARK2)
    add_text(s,"↓ deep-link",10.5,2.4,1.8,.4,size=11,color=TEAL)
    # bottom explanation
    add_bullets(s,[
        "saveContext() → CoreData writes to SQLite → trigger schedule",
        "UNUserNotificationCenter fires at exact dateComponents → iOS delivers",
        "Tap → UNUserNotificationCenterDelegate → NotificationCenter.post",
        "CalendarTabView .onReceive → opens correct FarmActivity detail sheet",
    ],.35,3.15,12.6,3.0,size=13,bc=BLUE)
    snum(s,8)

# ── Slide 9: Code Walkthrough ────────────────────────────────────────────────
def s9(prs):
    s = add_slide(prs)
    hdr(s,"📖  Code Walkthrough · scheduleNotification")
    steps=[
        (GREEN, "1","guard reminderEnabled, let date","ចៀសវាង schedule ដែលគ្មានការរំលឹក"),
        (BLUE,  "2","UNMutableNotificationContent()","បង្កើត notification payload"),
        (PURPLE,"3","content.title / body / sound","ដាក់ข้อchunk ទូរស័ព្ទ"),
        (ORANGE,"4","byAdding: .day value: -1","1 day before = farmer has time"),
        (TEAL,  "5","dateComponents from triggerDate","Calendar precision — year/month/day/hour/minute"),
        (GREEN, "6","UNCalendarNotificationTrigger","Fire ត្រឹមមួយដង repeats: false"),
        (BLUE,  "7","UNNotificationRequest + .add()","register ជាមួយ iOS system"),
    ]
    for i,(c,n,code,kh) in enumerate(steps):
        col=i%2; row=i//2
        cx=.35+col*6.5; cy=1.1+row*1.6
        if i==6: cx=.35
        add_rect(s,cx,cy,6.1,1.45,CARD); add_rect(s,cx,cy,.08,1.45,c)
        add_rect(s,cx+.15,cy+.4,.5,.5,c)
        add_text(s,n,cx+.15,cy+.4,.5,.5,size=13,bold=True,color=DARK,align=PP_ALIGN.CENTER)
        add_code(s,[code],cx+.78,cy+.1,5.1,.55,size=11)
        add_text(s,kh,cx+.78,cy+.72,5.1,.55,size=11,color=GREY)
    snum(s,9)

# ── Slide 10: Common Mistakes ────────────────────────────────────────────────
def s10(prs):
    s = add_slide(prs)
    hdr(s,"⚠️  Common Mistakes · ❌ vs ✅")
    rows=[
        ("Duplicate container","NSPersistentContainer() in App","CoreDataManager.shared.context"),
        ("NavigationStack","NavigationStack {} (iOS 16+)","NavigationView {} (iOS 13+)"),
        ("Missing requestAuthorization","Schedule without permission","Always request first, check granted"),
        ("Wrong trigger type","UNTimeIntervalTrigger for exact date","UNCalendarTrigger + dateComponents"),
        ("Deep-link pattern","@Environment(\\.dismiss) (iOS 15+)","NotificationCenter + .onReceive"),
        ("Fetch without context",".environment missing on sheet","Always .environment(\\.managedObjectContext, ctx)"),
    ]
    hdrs=["Pattern","❌  Wrong","✅  Correct (iOS 13+)"]
    widths=[2.8,4.5,5.1]; starts=[.35,3.3,7.95]
    for c,(h,cw,cs) in enumerate(zip(hdrs,widths,starts)):
        add_rect(s,cs,1.1,cw,.55,BLUE if c==2 else RED if c==1 else DARK2)
        add_text(s,h,cs+.1,1.15,cw-.15,.45,size=12,bold=True,color=WHITE)
    for r,(p,w,g) in enumerate(rows):
        cy=1.75+r*.67; bg=CARD if r%2==0 else DARK2
        for c,(cell,cw,cs) in enumerate(zip([p,w,g],widths,starts)):
            add_rect(s,cs,cy,cw,.62,bg)
            col=GREY if c==0 else (RED if c==1 else GREEN)
            add_text(s,cell,cs+.1,cy+.1,cw-.15,.42,size=10,color=col)
    snum(s,10)

# ── Slide 11: Mini-Project Assignment ───────────────────────────────────────
def s11(prs):
    s = add_slide(prs)
    hdr(s,"🏠  Mini-Project · ប្រតិទិន + ការជូនដំណឹង")
    cols=[
        ("📋  Requirements",BLUE,[
            "Calendar grid shows month view",
            "Tap day → see activities for that day",
            "Add/Edit/Delete FarmActivity",
            "Toggle reminderEnabled per activity",
            "Notification fires 1 day before",
            "Tap notification → opens activity",
        ]),
        ("✅  Checklist",GREEN,[
            "LazyVGrid calendar renders",
            "DayCellView has activity dot",
            "Add activity sheet works",
            "@FetchRequest filters by date",
            "requestAuthorization called",
            "Deep-link opens correct activity",
        ]),
        ("🎯  Grading",PURPLE,[
            "Calendar UI: 25%",
            "CoreData CRUD: 25%",
            "Notification permission: 20%",
            "Scheduling: 20%",
            "Deep-linking: 10%",
            "Bonus: monthly view toggle",
        ]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*4.35
        add_rect(s,cx,1.1,4.1,5.8,CARD); add_rect(s,cx,1.1,4.1,.55,c)
        add_text(s,t,cx+.12,1.15,3.85,.45,size=13,bold=True,color=DARK)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.15,cy,3.7,.72,DARK2)
            add_text(s,item,cx+.28,cy+.1,3.4,.52,size=11,color=WHITE)
    add_rect(s,.35,7.0,12.6,.42,BLUE)
    add_text(s,"📅  Week 05 ចប់! — ប្រតិទិន + notification ready! · Next: Pest & Disease Guide →",
             .5,7.02,12.3,.38,size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    snum(s,11)

# ── Slide 12: Summary ───────────────────────────────────────────────────────
def s12(prs):
    s = add_slide(prs)
    hdr(s,"📌  សង្ខេប Week 05 · Cheat Sheet",accent=BLUE)
    cols=[
        ("⚙️  Setup",BLUE,["UNUserNotificationCenter","requestAuthorization(options:)","FarmActivity CoreData entity","LazyVGrid calendar grid"]),
        ("📅  Calendar",GREEN,["DatePicker selection","DayCellView + dot","@FetchRequest by date","NavigationLink to detail"]),
        ("🔔  Notifications",PURPLE,["UNMutableNotificationContent","UNCalendarNotificationTrigger","UNNotificationRequest + .add()","cancel: removePendingRequests"]),
        ("🔗  Deep-link",ORANGE,["UNUserNotificationCenterDelegate","NotificationCenter.post","CalendarTabView .onReceive","open activity sheet"]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*3.25
        add_rect(s,cx,1.1,3.0,5.8,CARD); add_rect(s,cx,1.1,3.0,.55,c)
        add_text(s,t,cx+.1,1.15,2.8,.45,size=12,bold=True,color=DARK)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.1,cy,2.8,.72,DARK2)
            add_text(s,item,cx+.18,cy+.1,2.6,.52,size=10,color=TEAL)
            s.shapes[-1].text_frame.paragraphs[0].runs[0].font.name="Courier New"
    add_rect(s,.35,7.0,12.6,.42,BLUE)
    add_text(s,"🔔  Notification = ការសន្យារបស់ app ទៅ farmer — Schedule it right!",
             .5,7.02,12.3,.38,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    snum(s,12)

def build():
    prs = Presentation()
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    s1(prs);s2(prs);s3(prs);s4(prs);s5(prs);s6(prs)
    s7(prs);s8(prs);s9(prs);s10(prs);s11(prs);s12(prs)
    out = os.path.join(os.path.dirname(__file__),"Week05_Calendar_Reminders_KH.pptx")
    prs.save(out); print(f"✅  Saved → {out}")

if __name__ == "__main__": build()
