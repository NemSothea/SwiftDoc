"""Generate Keynote/PowerPoint-ready deck for Week 5 (Calendar & Reminders).

Produces: Week05_Calendar_Reminders.pptx

Design system (also listed on the summary slide):
- Palette:
    Primary Blue   #0A84FF   (iOS system blue — titles, accents)
    Farm Green     #34C759   (has-activity dot, checkmarks, success)
    Alert Orange   #FF9500   (reminder bell, notifications)
    Background     #F2F2F7   (slide background)
    Surface        #FFFFFF   (content cards)
    Text Primary   #1C1C1E
    Text Secondary #8E8E93
- Fonts: Title = SF Pro Display / Helvetica Neue Bold; Body = SF Pro Text / Helvetica Neue
- Layout: mostly 2-column (concept left, UI / code right)
- Icons: SF Symbols (calendar, bell.fill, drop.fill, leaf.fill, checkmark.circle.fill, ...)
- Animations: Fade (build in), Slide-from-right for code, Zoom for diagram
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# --- Design tokens -----------------------------------------------------------

PRIMARY = RGBColor(0x0A, 0x84, 0xFF)   # iOS blue
GREEN   = RGBColor(0x34, 0xC7, 0x59)   # farm green
ORANGE  = RGBColor(0xFF, 0x95, 0x00)   # alert orange
BG      = RGBColor(0xF2, 0xF2, 0xF7)   # slide background
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)   # card background
TEXT    = RGBColor(0x1C, 0x1C, 0x1E)   # primary text
MUTED   = RGBColor(0x8E, 0x8E, 0x93)   # secondary text
DIVIDER = RGBColor(0xD1, 0xD1, 0xD6)
CODE_BG = RGBColor(0x1C, 0x1C, 0x1E)
CODE_FG = RGBColor(0xF2, 0xF2, 0xF7)

TITLE_FONT = "SF Pro Display"
BODY_FONT  = "SF Pro Text"
MONO_FONT  = "Menlo"

# 16:9 deck
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


# --- Helpers -----------------------------------------------------------------

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def set_text(tf, text, *, font=BODY_FONT, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=18, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(40000)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT,
                font=BODY_FONT, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(80000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(40000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_card(slide, x, y, w, h, *, fill=SURFACE, border=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = border
        card.line.width = Pt(0.75)
    card.shadow.inherit = False
    return card


def slide_header(slide, eyebrow, title, *, accent=PRIMARY):
    # left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.55), Inches(0.55),
                                 Inches(0.08), Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    add_text(slide, Inches(0.75), Inches(0.45), Inches(12), Inches(0.35),
             eyebrow, font=BODY_FONT, size=12, bold=True, color=accent)
    add_text(slide, Inches(0.75), Inches(0.75), Inches(12), Inches(0.75),
             title, font=TITLE_FONT, size=30, bold=True, color=TEXT)


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_code_block(slide, x, y, w, h, code_lines, *, size=11):
    add_card(slide, x, y, w, h, fill=CODE_BG)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(140000)
    tf.margin_right = Emu(80000)
    tf.margin_top = Emu(100000)
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


def add_chip(slide, x, y, text, *, fill=PRIMARY, text_color=SURFACE, w=1.7, h=0.32):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    tf = shp.text_frame
    tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
    tf.margin_top = Emu(20000);  tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = BODY_FONT; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.55), Inches(7.1), Inches(6), Inches(0.3),
             "Week 5 · Calendar & Reminders · SmartFarmerAssistant",
             size=10, color=MUTED)
    add_text(slide, Inches(10.5), Inches(7.1), Inches(2.3), Inches(0.3),
             f"{idx} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def set_transition(slide, kind="fade"):
    """Attach a slide transition. Recognised by both Keynote and PowerPoint."""
    xml_map = {
        "fade":  '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        "push":  '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        "zoom":  '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
    }
    node = etree.fromstring(xml_map.get(kind, xml_map["fade"]))
    slide.element.append(node)


TOTAL = 10

# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = add_slide()

# full-bleed gradient-ish header band (two stacked rectangles)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(4.5))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = PRIMARY

# Decorative circles
for (cx, cy, r, color, alpha) in [
    (11.0, 0.6, 1.6, GREEN, None),
    (12.3, 3.4, 0.9, ORANGE, None),
    (1.1, 3.2, 0.55, SURFACE, None),
]:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(cx), Inches(cy),
                              Inches(r), Inches(r))
    circ.line.fill.background()
    circ.fill.solid(); circ.fill.fore_color.rgb = color

# "iOS" eyebrow chip
add_chip(s, 0.6, 0.6, "iOS · SwiftUI · Week 5", fill=SURFACE, text_color=PRIMARY, w=2.4, h=0.38)

add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(1.5),
         "Building Calendar & Reminder",
         font=TITLE_FONT, size=52, bold=True, color=SURFACE)
add_text(s, Inches(0.6), Inches(2.25), Inches(12), Inches(1.0),
         "Feature in iOS",
         font=TITLE_FONT, size=52, bold=True, color=SURFACE)
add_text(s, Inches(0.6), Inches(3.35), Inches(12), Inches(0.6),
         "SwiftUI  ·  Core Data  ·  Local Notifications",
         font=BODY_FONT, size=22, color=SURFACE)

# Bottom info card
add_card(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.6))
add_text(s, Inches(0.95), Inches(5.15), Inches(6), Inches(0.45),
         "AUDIENCE", size=11, bold=True, color=MUTED)
add_text(s, Inches(0.95), Inches(5.45), Inches(6), Inches(0.6),
         "Beginner iOS Developers", font=TITLE_FONT, size=20, bold=True)
add_text(s, Inches(0.95), Inches(5.95), Inches(7), Inches(0.5),
         "Month grid · FarmActivity entity · UNUserNotificationCenter · Deep-linking",
         size=13, color=MUTED)

add_text(s, Inches(7.8), Inches(5.15), Inches(4.8), Inches(0.45),
         "PROJECT 2 — SMART FARMER ASSISTANT",
         size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
add_text(s, Inches(7.8), Inches(5.45), Inches(4.8), Inches(0.6),
         "CalendarReminders Module", font=TITLE_FONT, size=20, bold=True,
         align=PP_ALIGN.RIGHT)
add_text(s, Inches(7.8), Inches(5.95), Inches(4.8), Inches(0.5),
         "iOS 13+ compatible — ObservableObject · NavigationView",
         size=13, color=MUTED, align=PP_ALIGN.RIGHT)

add_footer(s, 1, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Welcome! Today we build the Calendar & Reminders tab of the Smart Farmer Assistant. "
    "By the end, students will have a working month grid, a FarmActivity Core Data model, "
    "local notifications that survive app relaunches, and deep-link navigation from a "
    "notification tap. Everything is pure SwiftUI — no AppDelegate. "
    "UI suggestion: open with the finished app on screen so learners see the target."
))

# ============================================================================
# SLIDE 2 — Agenda / Learning Objectives
# ============================================================================
s = add_slide()
slide_header(s, "AGENDA", "What you'll build this week")

# Left column — learning objectives
add_text(s, Inches(0.75), Inches(1.9), Inches(5.8), Inches(0.4),
         "LEARNING OBJECTIVES", size=12, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.75), Inches(2.25), Inches(6.1), Inches(4.5), [
    "Build a month calendar with DatePicker + custom LazyVGrid",
    "Model FarmActivity in Core Data (date, type, notes, reminder)",
    "Request notification permission the right way",
    "Schedule local reminders with UNUserNotificationCenter",
    "Deep-link from a notification tap back to the activity",
], size=17)

# Right column — 5 lesson cards
lessons = [
    ("5.1", "Calendar View — DatePicker & Grid"),
    ("5.2", "FarmActivity Core Data Model"),
    ("5.3", "Notification Permission"),
    ("5.4", "Scheduling Local Notifications"),
    ("5.5", "Deep-linking on Notification Tap"),
]
for i, (num, title) in enumerate(lessons):
    top = Inches(1.9 + i * 0.95)
    add_card(s, Inches(7.1), top, Inches(5.6), Inches(0.8))
    # number bubble
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(7.25), top + Inches(0.15),
                              Inches(0.5), Inches(0.5))
    chip.line.fill.background()
    chip.fill.solid(); chip.fill.fore_color.rgb = PRIMARY
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = TITLE_FONT; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = SURFACE
    add_text(s, Inches(7.95), top + Inches(0.2), Inches(4.5), Inches(0.5),
             title, font=TITLE_FONT, size=16, bold=True)

add_footer(s, 2, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Walk through the five lessons briefly — do NOT deep-dive yet. "
    "Set expectations: by lesson 5.5 everything will click together via NotificationCenter. "
    "UI suggestion: 2-column layout — objectives on the left, scrollable lesson chips on the right. "
    "Animation: fade in the chips one-by-one."
))

# ============================================================================
# SLIDE 3 — Calendar UI (DatePicker + Custom Grid)
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 5.1", "Calendar UI — DatePicker + Custom Grid")

# Left — concept bullets
add_bullets(s, Inches(0.75), Inches(1.9), Inches(6.0), Inches(4.8), [
    "Month header with ‹ › chevrons (previousMonth / nextMonth)",
    "LazyVGrid with 7 columns — Sun → Sat",
    "firstWeekdayOffset renders blank leading cells",
    "DayCellView shows selection, today, and a green 'has activity' dot",
    "selectedDate is @Binding — taps drive the filtered list",
    "displayedMonth is @State local to the grid — chevrons never move the user's pick",
], size=16)

# Right — calendar UI mock
mock_x, mock_y, mock_w, mock_h = Inches(7.5), Inches(1.9), Inches(5.2), Inches(4.9)
add_card(s, mock_x, mock_y, mock_w, mock_h)

# Month header
add_text(s, mock_x, mock_y + Inches(0.2), mock_w, Inches(0.4),
         "‹     មេសា  2026     ›",
         font=TITLE_FONT, size=18, bold=True, align=PP_ALIGN.CENTER)

# day-of-week row
days = ["អា","ច","អ","ពុ","ព្រ","សុ","ស"]
cell_w = 0.66
start_x = 7.65
for i, d in enumerate(days):
    add_text(s, Inches(start_x + i * cell_w), mock_y + Inches(0.8),
             Inches(cell_w), Inches(0.3), d,
             size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# day cells (5 rows) — offset 3 blanks so the 1st is Wed
for row in range(5):
    for col in range(7):
        idx = row * 7 + col - 2  # 3 leading blanks
        if idx < 1 or idx > 30:
            continue
        cx = Inches(start_x + col * cell_w + 0.08)
        cy = mock_y + Inches(1.15 + row * 0.62)
        is_selected = (idx == 15)
        is_today    = (idx == 22)
        has_act     = idx in (3, 7, 15, 22, 25)
        # circle
        cell = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.5), Inches(0.5))
        cell.line.fill.background()
        cell.fill.solid()
        if is_selected:
            cell.fill.fore_color.rgb = PRIMARY
        elif is_today:
            cell.fill.fore_color.rgb = RGBColor(0xD6, 0xE9, 0xFF)
        else:
            cell.fill.fore_color.rgb = SURFACE
        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(idx)
        r.font.name = BODY_FONT
        r.font.size = Pt(12)
        r.font.bold = is_today or is_selected
        r.font.color.rgb = SURFACE if is_selected else (PRIMARY if is_today else TEXT)
        # green dot under cells with activity
        if has_act:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx + Inches(0.2),
                                     cy + Inches(0.52),
                                     Inches(0.1), Inches(0.1))
            dot.line.fill.background()
            dot.fill.solid(); dot.fill.fore_color.rgb = GREEN

# Legend
add_text(s, mock_x + Inches(0.3), mock_y + Inches(4.35), Inches(4.8), Inches(0.3),
         "●  has activity    ●  today    ●  selected",
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 3, TOTAL)
set_transition(s, "push")
speaker_notes(s, (
    "Explain WHY we do not use the stock DatePicker(.graphical): we need a custom "
    "'has-activity' dot under each day, Khmer month labels, and control over blank "
    "leading cells. Call out the two pieces of state — selectedDate is a @Binding "
    "owned by CalendarTabView, displayedMonth is @State private to the grid. "
    "UI suggestion: side-by-side — bullet list on the left, live mock calendar on the right. "
    "Animation: build-in day cells with a quick fade-cascade (Keynote: 'Appear' by object, 0.05s delay)."
))

# ============================================================================
# SLIDE 4 — Core Data Model (FarmActivity)
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 5.2", "Core Data Model — FarmActivity")

# Left — attributes table header
add_text(s, Inches(0.75), Inches(1.85), Inches(6.0), Inches(0.4),
         "ATTRIBUTES", size=12, bold=True, color=PRIMARY)

attrs = [
    ("id",              "UUID",   "Unique ID (also notification identifier)"),
    ("title",           "String", "Shown in the activity list"),
    ("activityType",    "String", "ដាំដំណាំ · ស្រោចទឹក · បាញ់ថ្នាំ · ច្រូតកាត់"),
    ("date",            "Date",   "Day + time — drives the trigger"),
    ("notes",           "String", "Optional description"),
    ("isCompleted",     "Bool",   "Checkbox state"),
    ("reminderEnabled", "Bool",   "Whether a local notification is scheduled"),
]
row_y = 2.25
for name, typ, desc in attrs:
    add_card(s, Inches(0.75), Inches(row_y), Inches(6.0), Inches(0.55))
    add_text(s, Inches(0.95), Inches(row_y + 0.07), Inches(1.9), Inches(0.4),
             name, font=MONO_FONT, size=12, bold=True, color=PRIMARY)
    add_text(s, Inches(2.85), Inches(row_y + 0.07), Inches(0.9), Inches(0.4),
             typ, font=MONO_FONT, size=11, color=MUTED)
    add_text(s, Inches(3.8), Inches(row_y + 0.07), Inches(2.85), Inches(0.4),
             desc, size=10, color=TEXT)
    row_y += 0.62

# Right — code snippet
code = [
    "// FarmActivity+CoreDataProperties.swift",
    "extension FarmActivity {",
    "  @NSManaged public var id: UUID?",
    "  @NSManaged public var title: String?",
    "  @NSManaged public var activityType: String?",
    "  @NSManaged public var date: Date?",
    "  @NSManaged public var notes: String?",
    "  @NSManaged public var isCompleted: Bool",
    "  @NSManaged public var reminderEnabled: Bool",
    "}",
    "",
    "extension FarmActivity: Identifiable {}",
]
add_code_block(s, Inches(7.25), Inches(2.1), Inches(5.45), Inches(3.2), code, size=12)

# Right — tip card
add_card(s, Inches(7.25), Inches(5.45), Inches(5.45), Inches(1.3),
         fill=RGBColor(0xFF, 0xF4, 0xE0))
add_text(s, Inches(7.5), Inches(5.55), Inches(5.0), Inches(0.4),
         "⚠️  Codegen: Manual / None",
         font=TITLE_FONT, size=13, bold=True, color=ORANGE)
add_text(s, Inches(7.5), Inches(5.9), Inches(5.0), Inches(0.8),
         "Set this in the Data Model Inspector — otherwise Xcode auto-generates the class and you'll get duplicate-symbol errors.",
         size=11, color=TEXT)

add_footer(s, 4, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Walk through each attribute and why it exists. Emphasise that `id: UUID?` does "
    "double duty as the notification request identifier — this is how we later cancel "
    "or replace the pending notification. Show Xcode's Data Model Inspector and point "
    "to Codegen → Manual/None. "
    "UI suggestion: image-left (attribute cards) + code-right. "
    "Animation: slide-from-left for cards, slide-from-right for the code block."
))

# ============================================================================
# SLIDE 5 — Notification Permission Flow
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 5.3", "Notification Permission Flow")

# Left — flow steps
steps = [
    ("1", "App launches → NotificationManager.shared as @StateObject"),
    ("2", "init() sets self as UNUserNotificationCenter delegate"),
    ("3", "User toggles 'បើកការរំលឹក' on an activity"),
    ("4", "requestPermission(options: [.alert, .badge, .sound])"),
    ("5", "Allow → isAuthorized = true   ·   Deny → open Settings alert"),
]
for i, (n, t) in enumerate(steps):
    top = Inches(1.95 + i * 0.85)
    # number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(0.75), top + Inches(0.08),
                              Inches(0.55), Inches(0.55))
    circ.line.fill.background()
    circ.fill.solid(); circ.fill.fore_color.rgb = PRIMARY
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = n
    r.font.name = TITLE_FONT; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = SURFACE
    add_text(s, Inches(1.45), top + Inches(0.15), Inches(5.4), Inches(0.6),
             t, size=14)

# Right — iPhone-ish permission mock
mx, my, mw, mh = Inches(8.3), Inches(1.9), Inches(3.9), Inches(4.9)
add_card(s, mx, my, mw, mh, fill=RGBColor(0x1C, 0x1C, 0x1E))

# inner white sheet
add_card(s, mx + Inches(0.25), my + Inches(0.9),
         mw - Inches(0.5), Inches(3.2), fill=SURFACE)

# bell icon placeholder
bell = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          mx + Inches(1.55), my + Inches(1.15),
                          Inches(0.8), Inches(0.8))
bell.line.fill.background()
bell.fill.solid(); bell.fill.fore_color.rgb = ORANGE
tf = bell.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🔔"
r.font.name = BODY_FONT; r.font.size = Pt(28)

add_text(s, mx + Inches(0.3), my + Inches(2.05),
         mw - Inches(0.6), Inches(0.45),
         '"SmartFarmer" Would Like to',
         font=TITLE_FONT, size=14, bold=True, align=PP_ALIGN.CENTER)
add_text(s, mx + Inches(0.3), my + Inches(2.4),
         mw - Inches(0.6), Inches(0.45),
         "Send You Notifications",
         font=TITLE_FONT, size=14, bold=True, align=PP_ALIGN.CENTER)
add_text(s, mx + Inches(0.3), my + Inches(2.85),
         mw - Inches(0.6), Inches(0.8),
         "Reminders for your farm activities will appear on the lock screen.",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)

# two buttons
btn_don = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             mx + Inches(0.3), my + Inches(3.65),
                             Inches(1.45), Inches(0.38))
btn_don.adjustments[0] = 0.5
btn_don.line.color.rgb = DIVIDER
btn_don.fill.solid(); btn_don.fill.fore_color.rgb = SURFACE
tf = btn_don.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Don't Allow"
r.font.name = BODY_FONT; r.font.size = Pt(11); r.font.color.rgb = PRIMARY

btn_ok = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            mx + Inches(1.85), my + Inches(3.65),
                            Inches(1.45), Inches(0.38))
btn_ok.adjustments[0] = 0.5
btn_ok.line.fill.background()
btn_ok.fill.solid(); btn_ok.fill.fore_color.rgb = PRIMARY
tf = btn_ok.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Allow"
r.font.name = BODY_FONT; r.font.size = Pt(11); r.font.bold = True
r.font.color.rgb = SURFACE

add_text(s, mx + Inches(0.3), my + Inches(4.3),
         mw - Inches(0.6), Inches(0.4),
         "System-level dialog — shows once",
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 5, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "The permission prompt fires ONCE. After that, requestAuthorization returns the "
    "stored answer without UI. This is why we ask only when the student first toggles "
    "'បើកការរំលឹក' — not at app launch. Also call out the main-thread hop: the callback "
    "runs on a background queue and @Published updates must happen on the main thread. "
    "UI suggestion: image-right with an iPhone-style mock of the permission dialog. "
    "Animation: zoom-in on the dialog card after the 5 flow steps fade in."
))

# ============================================================================
# SLIDE 6 — Scheduling Local Notifications
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 5.4", "Scheduling Local Notifications")

# Left — trio of cards for content / trigger / request
pieces = [
    ("UNMutableNotificationContent", "What the user sees",
     "title · body · sound · userInfo[\"activityID\"]"),
    ("UNCalendarNotificationTrigger", "When it fires",
     "dateMatching: [.year, .month, .day, .hour, .minute]"),
    ("UNNotificationRequest",         "Content + trigger",
     "identifier = activity.id.uuidString  (same ID replaces previous)"),
]
for i, (cls, what, detail) in enumerate(pieces):
    top = Inches(1.9 + i * 1.35)
    add_card(s, Inches(0.75), top, Inches(6.1), Inches(1.2))
    add_text(s, Inches(0.95), top + Inches(0.1), Inches(5.8), Inches(0.4),
             cls, font=MONO_FONT, size=12, bold=True, color=PRIMARY)
    add_text(s, Inches(0.95), top + Inches(0.42), Inches(5.8), Inches(0.35),
             what, size=11, bold=True, color=TEXT)
    add_text(s, Inches(0.95), top + Inches(0.72), Inches(5.8), Inches(0.45),
             detail, font=MONO_FONT, size=10, color=MUTED)

# Right — code
code = [
    "func scheduleNotification(for activity: FarmActivity) {",
    "  guard let id    = activity.id,",
    "        let title = activity.title,",
    "        let date  = activity.date else { return }",
    "  guard date > Date() else { return }     // skip past",
    "",
    "  let content = UNMutableNotificationContent()",
    "  content.title = \"🌾 កម្មវិធីកសិកម្ម\"",
    "  content.body  = title",
    "  content.sound = .default",
    "  content.userInfo = [\"activityID\": id.uuidString]",
    "",
    "  let comps = Calendar.current.dateComponents(",
    "    [.year,.month,.day,.hour,.minute], from: date)",
    "  let trigger = UNCalendarNotificationTrigger(",
    "    dateMatching: comps, repeats: false)",
    "",
    "  let request = UNNotificationRequest(",
    "    identifier: id.uuidString,",
    "    content: content, trigger: trigger)",
    "  UNUserNotificationCenter.current().add(request)",
    "}",
]
add_code_block(s, Inches(7.2), Inches(1.9), Inches(5.5), Inches(4.9), code, size=10)

add_footer(s, 6, TOTAL)
set_transition(s, "push")
speaker_notes(s, (
    "Three objects working together: Content (what), Trigger (when), Request (combine + add). "
    "Stress the guard on `date > Date()` — iOS silently drops past-date requests. "
    "Explain why identifier = UUID string: re-saving an activity REPLACES the old request; "
    "deleting an activity can cancel its reminder precisely. "
    "UI suggestion: 2-column — concept cards left, annotated Swift code right. "
    "Animation: slide-from-right for the code block to imply 'code as the destination'."
))

# ============================================================================
# SLIDE 7 — Handling Notification Tap (Deep Linking)
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 5.5", "Notification Tap → Deep Link")

# Left — the event-bus pattern
add_bullets(s, Inches(0.75), Inches(1.9), Inches(6.1), Inches(4.7), [
    "NotificationManager implements UNUserNotificationCenterDelegate",
    "didReceive(response:) reads activityID from userInfo",
    "Posts .didTapActivityNotification via Foundation NotificationCenter",
    "MainTabView.onReceive → selectedTab = 1 (Calendar)",
    "CalendarTabView.onReceive → selectedDate = activity.date",
    "Decoupled: the service has no SwiftUI dependencies",
], size=15)

# Tip callout
add_card(s, Inches(0.75), Inches(5.4), Inches(6.1), Inches(1.4),
         fill=RGBColor(0xE6, 0xF4, 0xFF))
add_text(s, Inches(0.95), Inches(5.5), Inches(5.7), Inches(0.4),
         "💡  Why NotificationCenter (not direct references)?",
         font=TITLE_FONT, size=13, bold=True, color=PRIMARY)
add_text(s, Inches(0.95), Inches(5.85), Inches(5.7), Inches(0.85),
         "NotificationManager doesn't know about views. An event bus keeps the service testable and the views independent subscribers.",
         size=11, color=TEXT)

# Right — code snippet (delegate)
code = [
    "func userNotificationCenter(",
    "  _ center: UNUserNotificationCenter,",
    "  didReceive response: UNNotificationResponse,",
    "  withCompletionHandler done: @escaping () -> Void) {",
    "",
    "  let info = response.notification.request",
    "                  .content.userInfo",
    "  if let s  = info[\"activityID\"] as? String,",
    "     let id = UUID(uuidString: s) {",
    "    DispatchQueue.main.async {",
    "      NotificationCenter.default.post(",
    "        name: .didTapActivityNotification,",
    "        object: nil,",
    "        userInfo: [\"activityID\": id])",
    "    }",
    "  }",
    "  done()",
    "}",
]
add_code_block(s, Inches(7.2), Inches(1.9), Inches(5.5), Inches(4.2), code, size=11)

# subscribers mini-code
sub_code = [
    ".onReceive(NotificationCenter.default",
    "  .publisher(for: .didTapActivityNotification)) { note in",
    "    if let id = note.userInfo?[\"activityID\"] as? UUID,",
    "       let a = activities.first(where: { $0.id == id }),",
    "       let d = a.date { selectedDate = d }",
    "}",
]
add_code_block(s, Inches(7.2), Inches(6.2), Inches(5.5), Inches(0.6),
               sub_code, size=9)

add_footer(s, 7, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "The delegate is the 'publisher' side; the two views are 'subscribers'. Foundation's "
    "NotificationCenter is a decoupled event bus — great fit because the service "
    "doesn't import SwiftUI. Remind students: delegate is already set in init() from "
    "Lesson 5.3, so no AppDelegate or UIApplicationDelegateAdaptor is needed. "
    "UI suggestion: 2-column — flow + callout left, delegate code right, subscriber code at bottom. "
    "Animation: fade the bullets, then slide-from-right for the delegate code, then fade the subscriber snippet."
))

# ============================================================================
# SLIDE 8 — Diagram: Notification Flow (full-width)
# ============================================================================
s = add_slide()
slide_header(s, "DIAGRAM", "End-to-end Notification Flow")

# A horizontal pipeline with 6 nodes and arrows
nodes = [
    ("User saves",      "AddActivityView"),
    ("Core Data",       "FarmActivity.save()"),
    ("Schedule",        "UNUserNotificationCenter"),
    ("Deliver",         "iOS at trigger time"),
    ("Tap",             "didReceive(response:)"),
    ("Deep-link",       ".didTapActivityNotification"),
]

node_w = 1.85
node_h = 1.0
start_x = 0.6
gap     = 0.18
y       = 2.3

colors = [PRIMARY, PRIMARY, ORANGE, ORANGE, GREEN, GREEN]
for i, ((label, sub), color) in enumerate(zip(nodes, colors)):
    x = start_x + i * (node_w + gap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y),
                              Inches(node_w), Inches(node_h))
    card.adjustments[0] = 0.15
    card.line.fill.background()
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE
    # top color strip
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y),
                               Inches(node_w), Inches(0.22))
    strip.adjustments[0] = 0.4
    strip.line.fill.background()
    strip.fill.solid(); strip.fill.fore_color.rgb = color

    # step number
    num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(x + 0.1), Inches(y + 0.08),
                             Inches(0.3), Inches(0.3))
    num.line.fill.background()
    num.fill.solid(); num.fill.fore_color.rgb = SURFACE
    tf = num.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1)
    r.font.name = TITLE_FONT; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = color

    add_text(s, Inches(x), Inches(y + 0.35),
             Inches(node_w), Inches(0.35),
             label, font=TITLE_FONT, size=13, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(x), Inches(y + 0.65),
             Inches(node_w), Inches(0.35),
             sub, font=MONO_FONT, size=9, color=MUTED,
             align=PP_ALIGN.CENTER)

    # arrow to the next
    if i < len(nodes) - 1:
        ax = x + node_w
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(ax + 0.01),
                                   Inches(y + 0.4),
                                   Inches(gap - 0.02),
                                   Inches(0.2))
        arrow.line.fill.background()
        arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED

# Second row — "what is in the payload"
add_card(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(1.15),
         fill=RGBColor(0xFF, 0xF4, 0xE0))
add_text(s, Inches(0.85), Inches(4.1), Inches(11), Inches(0.4),
         "📦  Payload travelling through every step",
         font=TITLE_FONT, size=13, bold=True, color=ORANGE)
add_text(s, Inches(0.85), Inches(4.45), Inches(11), Inches(0.65),
         "userInfo: [\"activityID\": UUID.uuidString]  →  carried on UNNotificationContent, "
         "read in didReceive, re-posted via NotificationCenter so the Calendar tab can find the row.",
         font=MONO_FONT, size=11)

# Third row — two subscriber bubbles
add_card(s, Inches(0.6), Inches(5.35), Inches(5.9), Inches(1.35))
add_text(s, Inches(0.85), Inches(5.45), Inches(5.3), Inches(0.4),
         "Subscriber 1 · MainTabView",
         font=TITLE_FONT, size=13, bold=True, color=PRIMARY)
add_text(s, Inches(0.85), Inches(5.75), Inches(5.5), Inches(0.85),
         "selectedTab = 1  →  switches to Calendar tab",
         font=MONO_FONT, size=11, color=TEXT)

add_card(s, Inches(6.8), Inches(5.35), Inches(5.9), Inches(1.35))
add_text(s, Inches(7.05), Inches(5.45), Inches(5.3), Inches(0.4),
         "Subscriber 2 · CalendarTabView",
         font=TITLE_FONT, size=13, bold=True, color=GREEN)
add_text(s, Inches(7.05), Inches(5.75), Inches(5.5), Inches(0.85),
         "selectedDate = activity.date  →  filters list to that day",
         font=MONO_FONT, size=11, color=TEXT)

add_footer(s, 8, TOTAL)
set_transition(s, "zoom")
speaker_notes(s, (
    "This is the money slide. Trace a single activityID from save → schedule → deliver → tap → "
    "both subscribers. Ask the class: 'what if we forgot to put activityID in userInfo?' "
    "(Answer: deep-linking breaks — we can switch tabs but can't select the right date.) "
    "UI suggestion: full-width horizontal pipeline with a payload strip underneath. "
    "Animation: appear each node left-to-right with a 0.15s delay; arrows wipe-right after each node."
))

# ============================================================================
# SLIDE 9 — Example Code Explanation
# ============================================================================
s = add_slide()
slide_header(s, "CODE WALKTHROUGH", "scheduleNotification(for:) — line by line")

# Left — numbered code
lines = [
    ("1",  "guard let id = activity.id, …,",
           "Reject if any required field is nil — a malformed request would crash later."),
    ("2",  "guard date > Date() else { return }",
           "Skip past-date activities. iOS would drop them anyway; we fail loudly in debug."),
    ("3",  "content.userInfo = [\"activityID\": id.uuidString]",
           "Breadcrumb the activity UUID so the tap handler can deep-link back to the row."),
    ("4",  "dateComponents([.year,.month,.day,.hour,.minute], from: date)",
           "Fire EXACTLY at the merged date+reminderTime. We ignore seconds on purpose."),
    ("5",  "UNCalendarNotificationTrigger(dateMatching:, repeats: false)",
           "One-shot reminder. repeats:true would require partial components — different use case."),
    ("6",  "UNNotificationRequest(identifier: id.uuidString, …)",
           "Same identifier = iOS replaces the old pending request — free idempotency."),
    ("7",  "UNUserNotificationCenter.current().add(request)",
           "Hands the request to the system; survives app relaunch and reboot."),
]
top_y = 1.95
for (num, code, note) in lines:
    add_card(s, Inches(0.6), Inches(top_y), Inches(12.15), Inches(0.68))
    # number
    nc = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(0.75), Inches(top_y + 0.14),
                            Inches(0.4), Inches(0.4))
    nc.line.fill.background()
    nc.fill.solid(); nc.fill.fore_color.rgb = PRIMARY
    tf = nc.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = TITLE_FONT; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = SURFACE

    add_text(s, Inches(1.3), Inches(top_y + 0.08),
             Inches(6.6), Inches(0.35),
             code, font=MONO_FONT, size=11, bold=True, color=PRIMARY)
    add_text(s, Inches(1.3), Inches(top_y + 0.36),
             Inches(6.6), Inches(0.32),
             note, size=10, color=MUTED)

    # right-hand "why" highlight
    add_text(s, Inches(8.1), Inches(top_y + 0.22),
             Inches(4.5), Inches(0.35),
             "✓ Beginner takeaway",
             size=10, bold=True, color=GREEN)
    top_y += 0.72

add_footer(s, 9, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Go line-by-line. Pause on line 6 (identifier = UUID string) — this is the trick "
    "that makes update/delete free. Also pause on line 1 to reinforce: early-return "
    "is the SwiftUI way, don't throw from here. "
    "UI suggestion: numbered list of code+commentary cards, 7 rows. "
    "Animation: cascading fade-in (0.1s stagger), Keynote calls this 'Appear by Object'."
))

# ============================================================================
# SLIDE 10 — Design system / Keynote handoff
# ============================================================================
s = add_slide()
slide_header(s, "HANDOFF", "Design System · Keynote Setup")

# --- Row 1: color palette ---
palette = [
    ("Primary Blue",  PRIMARY, "#0A84FF"),
    ("Farm Green",    GREEN,   "#34C759"),
    ("Alert Orange",  ORANGE,  "#FF9500"),
    ("Background",    BG,      "#F2F2F7"),
    ("Surface",       SURFACE, "#FFFFFF"),
    ("Text",          TEXT,    "#1C1C1E"),
]
add_text(s, Inches(0.75), Inches(1.85), Inches(12), Inches(0.35),
         "COLOR PALETTE", size=11, bold=True, color=MUTED)
for i, (name, rgb, hex_) in enumerate(palette):
    x = 0.75 + i * 2.05
    sw = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(2.2),
                            Inches(1.85), Inches(1.0))
    sw.adjustments[0] = 0.12
    sw.line.color.rgb = DIVIDER
    sw.fill.solid(); sw.fill.fore_color.rgb = rgb
    # caption card
    add_card(s, Inches(x), Inches(3.25), Inches(1.85), Inches(0.6),
             border=DIVIDER)
    add_text(s, Inches(x + 0.1), Inches(3.28), Inches(1.7), Inches(0.3),
             name, size=10, bold=True)
    add_text(s, Inches(x + 0.1), Inches(3.5), Inches(1.7), Inches(0.3),
             hex_, font=MONO_FONT, size=9, color=MUTED)

# --- Row 2: three cards (Fonts / Icons / Animations) ---
card_y = 4.05
card_h = 2.5
card_w = 4.05
gap_x  = 0.18
col_x  = [0.75, 0.75 + card_w + gap_x, 0.75 + 2*(card_w + gap_x)]

# Fonts card
add_card(s, Inches(col_x[0]), Inches(card_y), Inches(card_w), Inches(card_h))
add_text(s, Inches(col_x[0] + 0.25), Inches(card_y + 0.15), Inches(card_w), Inches(0.4),
         "TYPOGRAPHY", size=11, bold=True, color=PRIMARY)
add_text(s, Inches(col_x[0] + 0.25), Inches(card_y + 0.5), Inches(card_w - 0.5), Inches(0.5),
         "Title · SF Pro Display Bold",
         font=TITLE_FONT, size=18, bold=True)
add_text(s, Inches(col_x[0] + 0.25), Inches(card_y + 0.9), Inches(card_w - 0.5), Inches(0.5),
         "Fallback: Helvetica Neue Bold",
         size=11, color=MUTED)
add_text(s, Inches(col_x[0] + 0.25), Inches(card_y + 1.35), Inches(card_w - 0.5), Inches(0.5),
         "Body · SF Pro Text Regular",
         font=BODY_FONT, size=15)
add_text(s, Inches(col_x[0] + 0.25), Inches(card_y + 1.7), Inches(card_w - 0.5), Inches(0.5),
         "Mono · Menlo (for code blocks)",
         font=MONO_FONT, size=12, color=MUTED)
add_text(s, Inches(col_x[0] + 0.25), Inches(card_y + 2.05), Inches(card_w - 0.5), Inches(0.4),
         "Sizes: 32 / 20 / 16 / 11",
         size=10, color=MUTED)

# Icons card (SF Symbols)
add_card(s, Inches(col_x[1]), Inches(card_y), Inches(card_w), Inches(card_h))
add_text(s, Inches(col_x[1] + 0.25), Inches(card_y + 0.15), Inches(card_w), Inches(0.4),
         "SF SYMBOLS", size=11, bold=True, color=PRIMARY)
sym_rows = [
    ("calendar",               "Tab bar · empty state"),
    ("bell.fill",              "Reminder-enabled row"),
    ("checkmark.circle.fill",  "Completed activity"),
    ("leaf.fill",              "Activity type: ដាំដំណាំ"),
    ("drop.fill",              "Activity type: ស្រោចទឹក"),
    ("calendar.badge.plus",    "Empty-state prompt"),
]
for i, (name, use) in enumerate(sym_rows):
    y = card_y + 0.55 + i * 0.31
    add_text(s, Inches(col_x[1] + 0.25), Inches(y),
             Inches(2.0), Inches(0.3),
             name, font=MONO_FONT, size=10, color=TEXT, bold=True)
    add_text(s, Inches(col_x[1] + 2.2), Inches(y),
             Inches(card_w - 2.0), Inches(0.3),
             use, size=10, color=MUTED)

# Animations / Layout card
add_card(s, Inches(col_x[2]), Inches(card_y), Inches(card_w), Inches(card_h))
add_text(s, Inches(col_x[2] + 0.25), Inches(card_y + 0.15), Inches(card_w), Inches(0.4),
         "LAYOUT & MOTION", size=11, bold=True, color=PRIMARY)
layout_lines = [
    ("Layout", "2-column (concept · UI mock or code)"),
    ("Title",  "Full-bleed image-left · oversized type"),
    ("Diagram","Full-width horizontal pipeline"),
    ("",       ""),
    ("Transition", "Fade (default) · Push-left for code"),
    ("Build-in",   "Appear by object · 0.1s cascade"),
    ("Diagram",    "Zoom for the flow slide"),
]
for i, (k, v) in enumerate(layout_lines):
    y = card_y + 0.5 + i * 0.27
    if not k and not v:
        continue
    add_text(s, Inches(col_x[2] + 0.25), Inches(y),
             Inches(1.2), Inches(0.3),
             k, size=10, bold=True, color=TEXT)
    add_text(s, Inches(col_x[2] + 1.45), Inches(y),
             Inches(card_w - 1.55), Inches(0.3),
             v, size=10, color=MUTED)

# Footer bar with Keynote + PPTX note
add_card(s, Inches(0.75), Inches(6.65), Inches(12), Inches(0.4),
         fill=RGBColor(0xE6, 0xF4, 0xFF))
add_text(s, Inches(0.95), Inches(6.68), Inches(11.6), Inches(0.35),
         "Keynote: double-click Week05_Calendar_Reminders.pptx — it imports with transitions preserved. "
         "Export: File ▸ Export To ▸ PowerPoint to round-trip.",
         size=10, color=PRIMARY, bold=True)

add_footer(s, 10, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "This is the handoff slide. Students open the .pptx in Keynote (double-click). "
    "Fonts degrade gracefully to Helvetica Neue if SF Pro isn't installed. "
    "Keynote picks up the fade and push transitions automatically. "
    "For SF Symbols on slides, drag from the macOS SF Symbols app directly into Keynote."
))

# ----------------------------------------------------------------------------
out = "/Users/sothea007/Documents/Research/SwiftDoc/Week05/Week05_Calendar_Reminders.pptx"
prs.save(out)
print(f"Saved {out}")
