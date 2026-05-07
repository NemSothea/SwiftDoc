#!/usr/bin/env python3
"""Generate Git Fundamentals slide deck (Windows & macOS, Khmer language)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour palette ──────────────────────────────────────────────────────────
GIT    = RGBColor(0xF0, 0x50, 0x32)   # git brand orange-red
BLUE   = RGBColor(0x28, 0x7D, 0xFA)   # secondary blue
GREEN  = RGBColor(0x1B, 0xB8, 0x89)   # success green
TEAL   = RGBColor(0x00, 0xC9, 0xC8)   # code text
DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # slide background
DARK2  = RGBColor(0x16, 0x21, 0x3E)   # header bar
CARD   = RGBColor(0x0F, 0x2A, 0x45)   # code/card background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xAA, 0xAA, 0xBB)
RED    = RGBColor(0xE5, 0x47, 0x47)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

W = 13.33
H = 7.5


# ── Helpers ─────────────────────────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide  = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, items, x, y, w, h,
                size=15, color=WHITE, bullet_color=GREEN):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet
        r1 = p.add_run()
        r1.text = "• "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        # text
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_code(slide, code_lines, x, y, w, h, size=11):
    add_rect(slide, x, y, w, h, CARD)
    tb = slide.shapes.add_textbox(
        Inches(x + 0.18), Inches(y + 0.15),
        Inches(w - 0.36), Inches(h - 0.3))
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = TEAL
        run.font.name = "Courier New"
    return tb


def slide_number(slide, n, total=12):
    add_text(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.3,
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


def header_bar(slide, title, accent=GIT):
    add_rect(slide, 0, 0, W, 0.9, DARK2)
    add_rect(slide, 0, 0, 0.06, 0.9, accent)
    add_text(slide, title, 0.25, 0.1, 12.5, 0.7,
             size=24, bold=True, color=accent)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = add_slide(prs)

    # left accent bar
    add_rect(slide, 0, 0, 0.12, H, GIT)

    # badge
    add_rect(slide, 0.4, 0.4, 1.8, 0.42, GIT)
    add_text(slide, "Git Basics", 0.4, 0.4, 1.8, 0.42,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # main title
    add_text(slide, "Git", 0.4, 1.0, 12.0, 1.1,
             size=60, bold=True, color=GIT, align=PP_ALIGN.LEFT)
    add_text(slide, "សម្រាប់អ្នកចាប់ផ្ដើម", 0.4, 1.9, 12.0, 0.9,
             size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # subtitle
    add_text(slide, "Windows & macOS  ·  Basic Commands  ·  Branching  ·  GitHub Remote",
             0.4, 2.9, 12.0, 0.5, size=16, italic=True, color=GREY)

    # divider
    add_rect(slide, 0.4, 3.5, 10.0, 0.04, GIT)

    # icon row
    icons = [
        ("📁", "Repository"), ("🌿", "Branch"), ("⬆️", "Push"),
        ("⬇️", "Pull"),       ("🔀", "Merge"),  ("🐙", "GitHub"),
    ]
    for i, (emoji, label) in enumerate(icons):
        cx = 0.4 + i * 2.1
        add_rect(slide, cx, 3.7, 1.9, 1.0, CARD)
        add_text(slide, emoji, cx, 3.75, 1.9, 0.45,
                 size=22, align=PP_ALIGN.CENTER)
        add_text(slide, label, cx, 4.22, 1.9, 0.4,
                 size=11, color=GREY, align=PP_ALIGN.CENTER)

    # bottom tag
    add_text(slide, "ភាសាខ្មែរ · Complete Beginner Guide", 0.4, 6.9, 12.0, 0.45,
             size=12, italic=True, color=GREY)

    slide_number(slide, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
def slide_agenda(prs):
    slide = add_slide(prs)
    header_bar(slide, "📋  មាតិកា · Agenda")

    topics = [
        (GIT,   "01", "Git គឺជាអ្វី?",          "What is Git & why use it"),
        (BLUE,  "02", "ដំឡើង Windows",            "Install via Git for Windows"),
        (GREEN, "03", "ដំឡើង macOS",              "Xcode CLI Tools & Homebrew"),
        (TEAL,  "04", "ពាក្យបញ្ជាមូលដ្ឋាន",       "init · clone · add · commit · log"),
        (GIT,   "05", "Branch & Merge",           "branching workflow diagram"),
        (BLUE,  "06", "GitHub & Remote",          "push · pull · remote · PR"),
    ]

    for i, (color, num, kh, en) in enumerate(topics):
        col = i % 2
        row = i // 2
        cx = 0.35 + col * 6.5
        cy = 1.05 + row * 2.0

        add_rect(slide, cx, cy, 6.1, 1.7, CARD)
        add_rect(slide, cx, cy, 0.08, 1.7, color)

        # number badge
        add_rect(slide, cx + 0.18, cy + 0.35, 0.6, 0.6, color)
        add_text(slide, num, cx + 0.18, cy + 0.35, 0.6, 0.6,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, kh, cx + 0.95, cy + 0.25, 4.9, 0.55,
                 size=17, bold=True, color=WHITE)
        add_text(slide, en, cx + 0.95, cy + 0.82, 4.9, 0.45,
                 size=12, color=GREY)

    slide_number(slide, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is Git?
# ════════════════════════════════════════════════════════════════════════════
def slide_what_is_git(prs):
    slide = add_slide(prs)
    header_bar(slide, "🤔  Git គឺជាអ្វី?  (What is Git?)")

    # left bullets
    bullets = [
        "Git ជា Version Control System (VCS)",
        "រក្សាទុក ប្រវត្តិ ការផ្លាស់ប្ដូរ កូដ ទាំងអស់",
        "អាចត្រឡប់ទៅ version ចាស់ ពេលណាក៏បាន",
        "ធ្វើការបាន offline — មិនចាំបាច់ internet",
        "Git ≠ GitHub  (Git = tool, GitHub = cloud)",
        "ប្រើដោយ developer ជាង 90% នៅទូទាំងពិភពលោក",
    ]
    add_bullets(slide, bullets, 0.35, 1.05, 6.2, 5.5,
                size=15, bullet_color=GIT)

    # right: concept cards
    concepts = [
        (GIT,   "📸  Snapshot",    "រក្សាស្ថានភាពកូដ\nគ្រប់ commit"),
        (BLUE,  "🌿  Branch",      "ធ្វើការ feature\nផ្សេងៗដោយឯករាជ្យ"),
        (GREEN, "🔀  Merge",       "បញ្ចូលការផ្លាស់ប្ដូរ\nពី branch ផ្សេង"),
    ]
    for i, (color, title, desc) in enumerate(concepts):
        cy = 1.05 + i * 2.1
        add_rect(slide, 6.85, cy, 6.1, 1.85, CARD)
        add_rect(slide, 6.85, cy, 0.08, 1.85, color)
        add_text(slide, title, 7.05, cy + 0.15, 5.7, 0.5,
                 size=15, bold=True, color=color)
        add_text(slide, desc, 7.05, cy + 0.65, 5.7, 0.9,
                 size=13, color=WHITE)

    slide_number(slide, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Install on Windows
# ════════════════════════════════════════════════════════════════════════════
def slide_install_windows(prs):
    slide = add_slide(prs)
    header_bar(slide, "🪟  ដំឡើង Git នៅ Windows", accent=BLUE)

    steps = [
        ("01", "ចូល website",      "git-scm.com/download/win"),
        ("02", "ទាញ Installer",    "ចុច Download for Windows (.exe)"),
        ("03", "Run & Install",    "ចុច Next ទាំងអស់ → ចុច Install"),
        ("04", "ផ្ទៀងផ្ទាត់",       "បើក Git Bash ហើយវាយ: git --version"),
    ]

    for i, (num, kh, detail) in enumerate(steps):
        cy = 1.05 + i * 1.55
        add_rect(slide, 0.35, cy, 6.2, 1.35, CARD)
        add_rect(slide, 0.35, cy, 0.08, 1.35, BLUE)

        add_rect(slide, 0.55, cy + 0.38, 0.55, 0.55, BLUE)
        add_text(slide, num, 0.55, cy + 0.38, 0.55, 0.55,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, kh, 1.25, cy + 0.12, 5.1, 0.45,
                 size=15, bold=True, color=WHITE)
        add_text(slide, detail, 1.25, cy + 0.6, 5.1, 0.5,
                 size=12, color=GREY)

    # right: terminal mock
    add_text(slide, "Git Bash · Terminal Output", 6.85, 1.05, 6.1, 0.4,
             size=13, bold=True, color=BLUE)
    add_code(slide, [
        "Microsoft Windows [Version 11.0]",
        "",
        "C:\\Users\\you> git --version",
        "git version 2.44.0.windows.1",
        "",
        "C:\\Users\\you> git --help",
        "usage: git [--version] [--help]",
        "       git <command> [<args>]",
        "",
        "# Git Bash ត្រូវបានដំឡើងដោយ",
        "# ស្វ័យប្រវត្តិ ✅",
    ], 6.85, 1.55, 6.1, 3.8, size=11)

    # tip chip
    add_rect(slide, 6.85, 5.55, 6.1, 0.65, GREEN)
    add_text(slide, "✅  Git Bash ត្រូវបានដំឡើងដោយស្វ័យប្រវត្តិ — ប្រើជំនួស CMD!",
             7.0, 5.6, 5.8, 0.55, size=12, bold=True, color=DARK)

    slide_number(slide, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Install on macOS
# ════════════════════════════════════════════════════════════════════════════
def slide_install_macos(prs):
    slide = add_slide(prs)
    header_bar(slide, "🍎  ដំឡើង Git នៅ macOS", accent=GREEN)

    # Method 1 card
    add_rect(slide, 0.35, 1.05, 6.1, 4.2, CARD)
    add_rect(slide, 0.35, 1.05, 0.08, 4.2, BLUE)
    add_text(slide, "🔧  Method 1 · Xcode CLI Tools", 0.6, 1.15, 5.6, 0.5,
             size=15, bold=True, color=BLUE)
    add_text(slide, "(ណែនាំ — ងាយ និងរហ័ស)", 0.6, 1.65, 5.6, 0.35,
             size=12, italic=True, color=GREY)
    add_code(slide, [
        "# វាយក្នុង Terminal",
        "xcode-select --install",
        "",
        "# macOS នឹងបង្ហាញ dialog",
        "# ចុច Install → Agree",
        "",
        "# ចំណាយពេល 2-5 នាទី",
    ], 0.55, 2.1, 5.7, 2.7, size=11)

    # Method 2 card
    add_rect(slide, 6.85, 1.05, 6.1, 4.2, CARD)
    add_rect(slide, 6.85, 1.05, 0.08, 4.2, GREEN)
    add_text(slide, "🍺  Method 2 · Homebrew", 7.1, 1.15, 5.6, 0.5,
             size=15, bold=True, color=GREEN)
    add_text(slide, "(សម្រាប់ developer ដែលប្រើ Homebrew)", 7.1, 1.65, 5.6, 0.35,
             size=12, italic=True, color=GREY)
    add_code(slide, [
        "# ដំឡើង Homebrew មុន (ប្រសិនបើមិនទាន់មាន)",
        '/bin/bash -c "$(curl -fsSL',
        "  https://brew.sh/install.sh)",
        "",
        "# ដំឡើង Git",
        "brew install git",
        "",
        "# Update Git នៅថ្ងៃក្រោយ",
        "brew upgrade git",
    ], 7.05, 2.1, 5.7, 2.7, size=11)

    # verify row
    add_rect(slide, 0.35, 5.45, 12.6, 1.0, CARD)
    add_rect(slide, 0.35, 5.45, 0.08, 1.0, GIT)
    add_text(slide, "✅  ផ្ទៀងផ្ទាត់ (ទាំង Method 1 & 2):", 0.6, 5.5, 5.0, 0.4,
             size=13, bold=True, color=GIT)
    add_code(slide, ["$ git --version     →     git version 2.44.0"],
             5.8, 5.52, 7.0, 0.5, size=12)

    slide_number(slide, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — First-Time Setup
# ════════════════════════════════════════════════════════════════════════════
def slide_first_setup(prs):
    slide = add_slide(prs)
    header_bar(slide, "⚙️  កំណត់ Git លើកដំបូង · First-Time Setup", accent=GIT)

    add_text(slide, "ត្រូវធ្វើម្ដង — ទាំង Windows & macOS", 0.35, 1.0, 10.0, 0.4,
             size=13, italic=True, color=GREY)

    # large code block
    add_code(slide, [
        "# ── ដាក់ឈ្មោះ និង Email (ត្រូវដូចនឹង GitHub) ──────────────────",
        'git config --global user.name  "Your Name"',
        'git config --global user.email "you@example.com"',
        "",
        "# ── កំណត់ default branch ជា main ──────────────────────────────",
        "git config --global init.defaultBranch main",
        "",
        "# ── ផ្ទៀងផ្ទាត់ ──────────────────────────────────────────────",
        "git config --list",
        "",
        "# លទ្ធផល:",
        "# user.name=Your Name",
        "# user.email=you@example.com",
        "# init.defaultbranch=main",
    ], 0.35, 1.5, 7.8, 4.8, size=11)

    # explanation cards
    notes = [
        (GIT,   "user.name",   "ឈ្មោះនឹងបង្ហាញ\nក្នុង commit ទាំងអស់"),
        (BLUE,  "user.email",  "ត្រូវដូចនឹង\nGitHub account email"),
        (GREEN, "--global",    "អនុវត្តចំពោះ\nproject ទាំងអស់"),
    ]
    for i, (color, title, desc) in enumerate(notes):
        cy = 1.5 + i * 1.65
        add_rect(slide, 8.35, cy, 4.6, 1.45, CARD)
        add_rect(slide, 8.35, cy, 0.08, 1.45, color)
        add_text(slide, title, 8.58, cy + 0.1, 4.1, 0.45,
                 size=14, bold=True, color=color)
        add_text(slide, desc, 8.58, cy + 0.58, 4.1, 0.75,
                 size=12, color=WHITE)

    slide_number(slide, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Core Concepts (3-Area Diagram)
# ════════════════════════════════════════════════════════════════════════════
def slide_core_concepts(prs):
    slide = add_slide(prs)
    header_bar(slide, "🗂️  តំបន់ 3 ក្នុង Git · The Three Areas")

    # diagram boxes
    areas = [
        (GREY,  "Working Directory",  "ឯកសារដែល\nអ្នកកំពុងធ្វើការ",  "git add →"),
        (BLUE,  "Staging Area",       "ឯកសារ\nត្រៀមប្រគល់",         "git commit →"),
        (GIT,   "Repository (.git)",  "ប្រវត្តិ commit\nទាំងអស់",     ""),
    ]

    for i, (color, title, desc, arrow) in enumerate(areas):
        cx = 0.35 + i * 4.3
        add_rect(slide, cx, 1.2, 3.9, 3.2, CARD)
        add_rect(slide, cx, 1.2, 0.08, 3.2, color)
        add_text(slide, title, cx + 0.2, 1.35, 3.5, 0.55,
                 size=16, bold=True, color=color)
        add_text(slide, desc, cx + 0.2, 2.0, 3.5, 1.0,
                 size=14, color=WHITE)
        # arrow label
        if arrow:
            add_text(slide, arrow, cx + 3.6, 2.55, 1.5, 0.4,
                     size=13, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)

    # command arrows
    add_rect(slide, 4.15, 2.65, 0.2, 0.08, YELLOW)
    add_rect(slide, 8.45, 2.65, 0.2, 0.08, YELLOW)

    # workflow bullets
    add_bullets(slide, [
        "ធ្វើការកែប្រែ ក្នុង Working Directory (edit files)",
        "ជ្រើសរើសឯកសារ ដាក់ Staging Area: git add filename",
        "រក្សា snapshot ចូល Repository: git commit -m 'message'",
        "ប្រវត្តិ commit ទាំងអស់ ត្រូវបានរក្សាទុកជាអចិន្ត្រៃយ៍",
    ], 0.35, 4.65, 12.6, 2.5, size=14, bullet_color=GIT)

    slide_number(slide, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Basic Commands Part 1
# ════════════════════════════════════════════════════════════════════════════
def slide_basic_cmd1(prs):
    slide = add_slide(prs)
    header_bar(slide, "⌨️  ពាក្យបញ្ជាមូលដ្ឋាន · ផ្នែក ១  (Basic Commands)")

    cmds = [
        (GIT,   "git init",         "បង្កើត Git repository ថ្មី\nក្នុង folder បច្ចុប្បន្ន"),
        (BLUE,  "git clone <url>",  "ចម្លង repository\nពី GitHub មកលើ computer"),
        (GREEN, "git status",       "មើលស្ថានភាព — ឯកសារណា\nដែលបានផ្លាស់ប្ដូរ"),
        (TEAL,  "git add <file>",   "បន្ថែមឯកសារ ទៅ Staging Area\ngit add . = ទាំងអស់"),
    ]

    for i, (color, cmd, desc) in enumerate(cmds):
        cy = 1.05 + i * 1.6
        add_rect(slide, 0.35, cy, 6.1, 1.4, CARD)
        add_rect(slide, 0.35, cy, 0.08, 1.4, color)
        add_code(slide, [cmd], 0.55, cy + 0.12, 5.7, 0.55, size=13)
        add_text(slide, desc, 0.58, cy + 0.72, 5.6, 0.6,
                 size=12, color=GREY)

    # right code block
    add_text(slide, "ឧទាហរណ៍ · Example Workflow", 6.85, 1.05, 6.1, 0.4,
             size=13, bold=True, color=GIT)
    add_code(slide, [
        "# 1. បង្កើត project ថ្មី",
        "mkdir my-project",
        "cd my-project",
        "git init",
        "",
        "# 2. បន្ថែម file",
        "touch README.md",
        "git status",
        "# On branch main",
        "# Untracked files: README.md",
        "",
        "# 3. Add to staging",
        "git add README.md",
        "git status",
        "# Changes to be committed:",
        "#   new file: README.md",
    ], 6.85, 1.55, 6.1, 5.3, size=10)

    slide_number(slide, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Basic Commands Part 2
# ════════════════════════════════════════════════════════════════════════════
def slide_basic_cmd2(prs):
    slide = add_slide(prs)
    header_bar(slide, "⌨️  ពាក្យបញ្ជាមូលដ្ឋាន · ផ្នែក ២  (Basic Commands)")

    cmds = [
        (GIT,   'git commit -m "msg"', "រក្សា snapshot ជាមួយ\nmessage ពណ៌នា"),
        (BLUE,  "git log",             "មើលប្រវត្តិ commit\nទាំងអស់"),
        (GREEN, "git diff",            "មើលអ្វីដែល\nផ្លាស់ប្ដូរ (មុន add)"),
        (TEAL,  "git restore <file>",  "ត្រឡប់ file ទៅ\nស្ថានភាពចុងក្រោយ"),
    ]

    for i, (color, cmd, desc) in enumerate(cmds):
        cy = 1.05 + i * 1.6
        add_rect(slide, 0.35, cy, 6.1, 1.4, CARD)
        add_rect(slide, 0.35, cy, 0.08, 1.4, color)
        add_code(slide, [cmd], 0.55, cy + 0.12, 5.7, 0.55, size=13)
        add_text(slide, desc, 0.58, cy + 0.72, 5.6, 0.6,
                 size=12, color=GREY)

    # right code block
    add_text(slide, "git log --oneline · Output Example", 6.85, 1.05, 6.1, 0.4,
             size=13, bold=True, color=BLUE)
    add_code(slide, [
        "$ git log --oneline",
        "a3f9c12 fix: correct calculation bug",
        "8b2e001 feat: add login screen",
        "4d7fa3e chore: initial project setup",
        "",
        "$ git log --oneline --graph",
        "* a3f9c12 fix: correct bug",
        "* 8b2e001 feat: login screen",
        "* 4d7fa3e initial setup",
        "",
        "# commit message ល្អ:",
        "#  feat: add new feature",
        "#  fix: bug description",
        "#  docs: update README",
        "#  chore: update deps",
    ], 6.85, 1.55, 6.1, 5.3, size=10)

    slide_number(slide, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Branching & Merging
# ════════════════════════════════════════════════════════════════════════════
def slide_branching(prs):
    slide = add_slide(prs)
    header_bar(slide, "🌿  Branch & Merge · ការបំបែក និងបញ្ចូល")

    # branch diagram (visual)
    # main line
    add_rect(slide, 0.35, 2.5, 12.5, 0.06, GREEN)
    # feature branch line
    add_rect(slide, 3.5, 2.0, 0.06, 1.2, GIT)   # fork down
    add_rect(slide, 3.5, 2.0, 4.5, 0.06, GIT)    # feature line
    add_rect(slide, 7.94, 2.0, 0.06, 0.56, GIT)  # merge back

    # commit dots on main
    for cx in [1.0, 2.5, 8.5, 10.5, 12.2]:
        add_rect(slide, cx - 0.12, 2.38, 0.3, 0.3, GREEN)

    # commit dots on feature
    for cx in [4.5, 6.0, 7.5]:
        add_rect(slide, cx - 0.12, 1.88, 0.3, 0.3, GIT)

    add_text(slide, "main",          12.4, 2.3,  1.0, 0.35, size=12, bold=True, color=GREEN)
    add_text(slide, "feature/login", 7.7,  1.6,  2.5, 0.3,  size=11, color=GIT)
    add_text(slide, "① fork",        3.2,  1.5,  1.2, 0.3,  size=10, color=GREY)
    add_text(slide, "④ merge",        7.7,  2.65, 1.3, 0.3,  size=10, color=GREY)

    # code block
    add_text(slide, "Commands", 0.35, 3.0, 5.0, 0.4, size=13, bold=True, color=GIT)
    add_code(slide, [
        "# មើល branch ទាំងអស់",
        "git branch",
        "",
        "# បង្កើត branch ថ្មី & ប្ដូរទៅ",
        "git switch -c feature/login",
        "",
        "# ត្រឡប់ main branch",
        "git switch main",
        "",
        "# Merge feature branch",
        "git merge feature/login",
        "",
        "# លុប branch (បន្ទាប់ merge)",
        "git branch -d feature/login",
    ], 0.35, 3.45, 6.1, 4.0, size=11)

    # tips
    tips = [
        (GIT,   "Branch",  "branch = ច្រកការងារ\nឯករាជ្យ — main មិនប៉ះពាល់"),
        (BLUE,  "Switch",  "git switch (ថ្មី)\nឬ git checkout (ចាស់)"),
        (GREEN, "Merge",   "បញ្ចូល feature ត្រឡប់\nក្នុង main branch"),
    ]
    for i, (color, title, desc) in enumerate(tips):
        cy = 3.45 + i * 1.35
        add_rect(slide, 6.65, cy, 6.3, 1.2, CARD)
        add_rect(slide, 6.65, cy, 0.08, 1.2, color)
        add_text(slide, title, 6.9, cy + 0.1, 5.8, 0.4, size=14, bold=True, color=color)
        add_text(slide, desc, 6.9, cy + 0.55, 5.8, 0.55, size=12, color=WHITE)

    slide_number(slide, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — GitHub & Remote
# ════════════════════════════════════════════════════════════════════════════
def slide_github(prs):
    slide = add_slide(prs)
    header_bar(slide, "🐙  GitHub & Remote · ភ្ជាប់ Cloud")

    # step flow (4 cards)
    steps = [
        (GIT,   "① បង្កើត Repo",     "ចូល github.com → New → ដាក់ឈ្មោះ → Create"),
        (BLUE,  "② Remote Add",      "git remote add origin <url>"),
        (GREEN, "③ Push",            "git push -u origin main"),
        (TEAL,  "④ Pull",            "git pull origin main"),
    ]

    for i, (color, title, detail) in enumerate(steps):
        cx = 0.35 + i * 3.25
        add_rect(slide, cx, 1.05, 3.0, 1.3, CARD)
        add_rect(slide, cx, 1.05, 0.08, 1.3, color)
        add_text(slide, title, cx + 0.2, 1.12, 2.6, 0.45, size=13, bold=True, color=color)
        add_text(slide, detail, cx + 0.2, 1.6, 2.6, 0.65, size=10, color=GREY)

        if i < 3:
            add_text(slide, "→", cx + 3.0, 1.55, 0.35, 0.4,
                     size=18, bold=True, color=GREY, align=PP_ALIGN.CENTER)

    # code block
    add_text(slide, "ឧទាហរណ៍ · Full Workflow", 0.35, 2.5, 7.0, 0.4,
             size=13, bold=True, color=GIT)
    add_code(slide, [
        "# ភ្ជាប់ local repo ជាមួយ GitHub",
        "git remote add origin https://github.com/user/repo.git",
        "",
        "# Push ដំបូង (ចំពោះ main branch)",
        "git push -u origin main",
        "",
        "# Push ការផ្លាស់ប្ដូរ លើកបន្ទាប់",
        "git push",
        "",
        "# ទាញ ការផ្លាស់ប្ដូរ ពី GitHub",
        "git pull origin main",
        "",
        "# មើល remote ទាំងអស់",
        "git remote -v",
    ], 0.35, 3.0, 7.8, 4.4, size=11)

    # .gitignore tip
    add_text(slide, "📝  .gitignore · ឯកសារដែលមិន Track", 8.35, 2.5, 4.6, 0.4,
             size=13, bold=True, color=YELLOW)
    add_code(slide, [
        "# Xcode .gitignore",
        "*.xcuserstate",
        "xcuserdata/",
        "DerivedData/",
        ".DS_Store",
        "*.moved-aside",
        "Pods/",
        "",
        "# ចំណាំ: .gitignore ត្រូវ",
        "# commit ជាមួយ project",
    ], 8.35, 3.0, 4.6, 3.6, size=10)

    # pull vs fetch note
    add_rect(slide, 8.35, 6.7, 4.6, 0.65, CARD)
    add_rect(slide, 8.35, 6.7, 0.08, 0.65, BLUE)
    add_text(slide, "💡  git fetch = ទាញ ប៉ុន្ដែ មិន merge  |  git pull = fetch + merge",
             8.58, 6.75, 4.2, 0.55, size=10, color=WHITE)

    slide_number(slide, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary & Cheat Sheet
# ════════════════════════════════════════════════════════════════════════════
def slide_summary(prs):
    slide = add_slide(prs)
    header_bar(slide, "📌  សង្ខេប · Git Cheat Sheet", accent=GIT)

    columns = [
        ("⚙️  Setup", GIT, [
            "git init",
            "git clone <url>",
            "git config --global",
            "git remote add origin",
        ]),
        ("📅  Daily", BLUE, [
            "git status",
            "git add <file>",
            "git commit -m 'msg'",
            "git log --oneline",
            "git diff",
            "git restore <file>",
        ]),
        ("🤝  Team", GREEN, [
            "git branch",
            "git switch -c <name>",
            "git merge <branch>",
            "git push",
            "git pull",
            "git fetch",
        ]),
    ]

    for i, (title, color, cmds) in enumerate(columns):
        cx = 0.35 + i * 4.35
        add_rect(slide, cx, 1.05, 4.1, 5.8, CARD)
        add_rect(slide, cx, 1.05, 0.08, 5.8, color)

        # column header
        add_rect(slide, cx, 1.05, 4.1, 0.6, color)
        add_text(slide, title, cx + 0.15, 1.1, 3.8, 0.5,
                 size=15, bold=True, color=DARK)

        for j, cmd in enumerate(cmds):
            cy_cmd = 1.75 + j * 0.82
            add_rect(slide, cx + 0.15, cy_cmd, 3.7, 0.65, DARK2)
            add_text(slide, cmd, cx + 0.3, cy_cmd + 0.07, 3.4, 0.5,
                     size=12, color=TEAL)
            run_tb = slide.shapes[-1]
            run_tb.text_frame.paragraphs[0].runs[0].font.name = "Courier New"

    # bottom banner
    add_rect(slide, 0.35, 7.0, 12.6, 0.42, GIT)
    add_text(slide, "💪  Git រៀនដោយការអនុវត្ត! — Practice every day & it becomes second nature.",
             0.5, 7.03, 12.3, 0.38,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number(slide, 12)


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    slide_title(prs)
    slide_agenda(prs)
    slide_what_is_git(prs)
    slide_install_windows(prs)
    slide_install_macos(prs)
    slide_first_setup(prs)
    slide_core_concepts(prs)
    slide_basic_cmd1(prs)
    slide_basic_cmd2(prs)
    slide_branching(prs)
    slide_github(prs)
    slide_summary(prs)

    out = os.path.join(os.path.dirname(__file__), "Git_Fundamentals.pptx")
    prs.save(out)
    print(f"✅  Saved → {out}")


if __name__ == "__main__":
    build()
