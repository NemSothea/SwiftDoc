#!/usr/bin/env python3
"""Generate Week 06 — Pest & Disease Guide slide deck (Khmer)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

GREEN  = RGBColor(0x1B,0xB8,0x89)
BLUE   = RGBColor(0x28,0x7D,0xFA)
PURPLE = RGBColor(0x8E,0x44,0xAD)
ORANGE = RGBColor(0xF3,0x96,0x20)
TEAL   = RGBColor(0x00,0xC9,0xC8)
DARK   = RGBColor(0x1A,0x1A,0x2E)
DARK2  = RGBColor(0x16,0x21,0x3E)
CARD   = RGBColor(0x0F,0x2A,0x45)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GREY   = RGBColor(0xAA,0xAA,0xBB)
RED    = RGBColor(0xE5,0x47,0x47)
YELLOW = RGBColor(0xFF,0xD7,0x00)
W,H    = 13.33,7.5

def add_slide(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    for ph in s.placeholders: ph._element.getparent().remove(ph._element)
    f=s.background.fill; f.solid(); f.fore_color.rgb=DARK
    return s

def add_text(s,text,x,y,w,h,size=16,bold=False,italic=False,
             color=WHITE,align=PP_ALIGN.LEFT,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color
    return tb

def add_rect(s,x,y,w,h,fc,lc=None):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fc
    if lc: sh.line.color.rgb=lc
    else: sh.line.fill.background()
    return sh

def add_bullets(s,items,x,y,w,h,size=14,color=WHITE,bc=GREEN):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r1=p.add_run(); r1.text="• "; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=bc
        r2=p.add_run(); r2.text=item; r2.font.size=Pt(size); r2.font.color.rgb=color

def add_code(s,lines,x,y,w,h,size=10):
    add_rect(s,x,y,w,h,CARD)
    tb=s.shapes.add_textbox(Inches(x+.18),Inches(y+.15),Inches(w-.36),Inches(h-.3))
    tf=tb.text_frame; tf.word_wrap=False
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.color.rgb=TEAL; r.font.name="Courier New"

def hdr(s,title,accent=GREEN):
    add_rect(s,0,0,W,.9,DARK2); add_rect(s,0,0,.06,.9,accent)
    add_text(s,title,.25,.1,12.5,.7,size=22,bold=True,color=accent)

def snum(s,n,total=12):
    add_text(s,f"{n}/{total}",12.3,7.1,.9,.3,size=10,color=GREY,align=PP_ALIGN.RIGHT)

# ── Slide 1: Title ──────────────────────────────────────────────────────────
def s1(prs):
    s=add_slide(prs)
    add_rect(s,0,0,.12,H,GREEN)
    add_rect(s,.4,.4,2.2,.42,GREEN)
    add_text(s,"Week 06 · SmartFarmer",.4,.4,2.2,.42,size=12,bold=True,color=DARK,align=PP_ALIGN.CENTER)
    add_text(s,"មគ្គុទ្ទេសសត្វល្អិត",.4,1.0,12.0,1.0,size=44,bold=True,color=GREEN)
    add_text(s,"Pest & Disease Guide · Project 3",.4,1.95,12.0,.7,size=28,bold=True,color=WHITE)
    add_text(s,"iOS 13+  ·  CoreData  ·  JSON Preload  ·  Custom Search  ·  Offline",.4,2.85,12.0,.45,size=14,italic=True,color=GREY)
    add_rect(s,.4,3.45,10.5,.04,GREEN)
    icons=[("🐛","Pest Entity"),("🔍","Custom Search"),("📂","JSON Preload"),("📱","Expandable"),("🔌","Offline"),("🌿","Khmer UI")]
    for i,(e,l) in enumerate(icons):
        cx=.4+i*2.1; add_rect(s,cx,3.65,1.9,1.0,CARD)
        add_text(s,e,cx,3.7,1.9,.45,size=20,align=PP_ALIGN.CENTER)
        add_text(s,l,cx,4.18,1.9,.4,size=10,color=GREY,align=PP_ALIGN.CENTER)
    add_text(s,"SmartFarmer Assistant · ភ្នំពេញ 2026",.4,6.9,12.0,.4,size=11,italic=True,color=GREY)
    snum(s,1)

# ── Slide 2: Agenda ─────────────────────────────────────────────────────────
def s2(prs):
    s=add_slide(prs)
    hdr(s,"📋  មាតិកា Week 06 · Agenda")
    topics=[
        (GREEN, "01","Architecture & Folder","View → ViewModel → Service → Model → Resource"),
        (BLUE,  "02","Pest CoreData Entity","7 attributes + JSON preload flow"),
        (PURPLE,"03","Custom Search Bar","iOS 13+ TextField (no .searchable needed)"),
        (ORANGE,"04","Expandable Sections","@State toggles + detail view"),
        (TEAL,  "05","Offline-First Design","JSON → CoreData → always works"),
        (GREEN, "06","Mini-Project","Pest Guide ពេញ + search + categories"),
    ]
    for i,(c,n,kh,en) in enumerate(topics):
        col=i%2; row=i//2; cx=.35+col*6.5; cy=1.1+row*2.0
        add_rect(s,cx,cy,6.1,1.75,CARD); add_rect(s,cx,cy,.08,1.75,c)
        add_rect(s,cx+.18,cy+.38,.6,.6,c)
        add_text(s,n,cx+.18,cy+.38,.6,.6,size=13,bold=True,color=DARK,align=PP_ALIGN.CENTER)
        add_text(s,kh,cx+.95,cy+.25,4.9,.5,size=15,bold=True,color=WHITE)
        add_text(s,en,cx+.95,cy+.85,4.9,.45,size=11,color=GREY)
    snum(s,2)

# ── Slide 3: Architecture ───────────────────────────────────────────────────
def s3(prs):
    s=add_slide(prs)
    hdr(s,"🏗️  Architecture · 5-Layer MVVM")
    layers=[
        (BLUE,  "View",      "PestGuideTabView\nPestRowView, PestDetailView\nCustomSearchBar"),
        (GREEN, "ViewModel", "PestViewModel\nObservableObject\n@Published searchText"),
        (PURPLE,"Service",   "PestDataLoader\nJSON → CoreData seeding\n(first launch only)"),
        (ORANGE,"Model",     "Pest NSManagedObject\n7 @NSManaged properties\nPestType enum"),
        (TEAL,  "Resource",  "pests.json bundled\nAssets.xcassets\n.xcdatamodeld"),
    ]
    for i,(c,t,d) in enumerate(layers):
        cx=.35+i*2.58
        add_rect(s,cx,1.1,2.38,4.5,CARD); add_rect(s,cx,1.1,.08,4.5,c)
        add_text(s,t,cx+.18,1.2,2.0,.5,size=14,bold=True,color=c)
        add_text(s,d,cx+.18,1.78,2.0,2.6,size=11,color=WHITE)
        if i<4:
            add_text(s,"↓",cx+2.2,2.9,.4,.4,size=18,bold=True,color=GREY,align=PP_ALIGN.CENTER)
    add_bullets(s,[
        "View រៀបចំ UI — ViewModel គ្រប់គ្រង state — Service seed data — Model ជា CoreData entity",
        "PestDataLoader ត្រូវ run តែ ONCE (first launch) → check UserDefaults flag",
        "JSON bundle → ចម្លង → CoreData → app ធ្វើការ offline ជានិច្ច",
    ],.35,5.85,12.6,1.5,size=12,bc=GREEN)
    snum(s,3)

# ── Slide 4: Pest Entity ─────────────────────────────────────────────────────
def s4(prs):
    s=add_slide(prs)
    hdr(s,"🐛  Pest CoreData Entity · 8 Attributes")
    attrs=[
        ("id","UUID","Identifiable conformance"),
        ("name","String?","ឈ្មោះ: Rice Blast, Brown Planthopper..."),
        ("pestType","String?","ប្រភេទ: សត្វល្អិត / ផ្សិត / បាក់តេរី"),
        ("symptoms","String?","រោគសញ្ញា — ពណ៌នាលម្អិត"),
        ("treatment","String?","វិធីព្យាបាល — ជាប់ពាក្យ farmer"),
        ("prevention","String?","វិធីការពារ — long-term"),
        ("imageName","String?","Asset catalog name (optional)"),
        ("isFavorite","Bool","Bookmark — heart icon"),
    ]
    for i,(n,t,d) in enumerate(attrs):
        cy=1.1+i*.75; bg=CARD if i%2==0 else DARK2
        add_rect(s,.35,cy,12.6,.65,bg); add_rect(s,.35,cy,.06,.65,GREEN)
        add_code(s,[n],.55,cy+.07,3.0,.5,size=12)
        add_code(s,[t],3.75,cy+.07,2.2,.5,size=12)
        add_text(s,d,6.1,cy+.12,6.7,.42,size=11,color=GREY)
    add_rect(s,.35,7.1,12.6,.32,DARK2)
    add_text(s,"⚙️  Codegen → Manual/None  |  Create NSManagedObject Subclass  |  pestType ← filter ក្នុង NSPredicate",
             .5,7.12,12.2,.28,size=10,color=GREY)
    snum(s,4)

# ── Slide 5: JSON Preload Flow ───────────────────────────────────────────────
def s5(prs):
    s=add_slide(prs)
    hdr(s,"📂  JSON Preload Flow · pests.json → CoreData")
    flow=[
        (GREEN, "①","App launches","check UserDefaults: pestDataSeeded"),
        (BLUE,  "②","Not seeded?","load pests.json from Bundle.main"),
        (PURPLE,"③","JSONDecoder","decode [PestData] array"),
        (ORANGE,"④","Create entities","Pest(context:) loop → set all properties"),
        (TEAL,  "⑤","saveContext()","persist to SQLite"),
        (GREEN, "⑥","Set flag","UserDefaults.set(true, key: pestDataSeeded)"),
        (BLUE,  "⑦","@FetchRequest","UI auto-refreshes — shows all pests"),
    ]
    for i,(c,n,kh,en) in enumerate(flow):
        col=i%2; row=i//2; cx=.35+col*6.5; cy=1.1+row*1.55
        if i==6: cx=.35
        add_rect(s,cx,cy,6.1,1.35,CARD); add_rect(s,cx,cy,.08,1.35,c)
        add_rect(s,cx+.15,cy+.37,.52,.52,c)
        add_text(s,n,cx+.15,cy+.37,.52,.52,size=13,bold=True,color=DARK,align=PP_ALIGN.CENTER)
        add_text(s,kh,cx+.82,cy+.1,5.1,.45,size=13,bold=True,color=WHITE)
        add_text(s,en,cx+.82,cy+.62,5.1,.45,size=11,color=GREY)
    snum(s,5)

# ── Slide 6: Custom Search Bar ──────────────────────────────────────────────
def s6(prs):
    s=add_slide(prs)
    hdr(s,"🔍  Custom Search Bar · iOS 13+ (no .searchable)")
    add_text(s,"Search Implementation",.35,1.05,6.5,.4,size=12,bold=True,color=GREEN)
    add_code(s,[
        "// ✅ iOS 13+ custom search (no .searchable)",
        "struct PestGuideTabView: View {",
        "    @State private var searchText = \"\"",
        "",
        "    var displayedPests: [Pest] {",
        "        if searchText.isEmpty { return Array(pests) }",
        "        return pests.filter {",
        "            ($0.name ?? \"\")   ",
        "                .localizedCaseInsensitiveContains(searchText) ||",
        "            ($0.pestType ?? \"\")",
        "                .localizedCaseInsensitiveContains(searchText)",
        "        }",
        "    }",
        "",
        "    var body: some View {",
        "        NavigationView {",
        "            VStack {",
        "                // Custom search bar",
        "                HStack {",
        "                    Image(systemName: \"magnifyingglass\")",
        "                        .foregroundColor(.gray)",
        "                    TextField(\"ស្វែងរក...\", text: $searchText)",
        "                    if !searchText.isEmpty {",
        "                        Button(\"បោះបង់\") { searchText = \"\" }",
        "                    }",
        "                }",
        "                .padding(8)",
        "                .background(Color(.systemGray6))",
        "                .cornerRadius(10)",
        "                .padding(.horizontal)",
        "                List { ... }",
        "            }",
        "        }",
        "    }",
        "}",
    ],.35,1.55,6.5,5.8,size=9.5)
    notes=[
        (GREEN, "localizedCaseInsensitiveContains","ស្វែងរក case-insensitive\nខ្មែរ + Latin scripts"),
        (BLUE,  "searchText isEmpty check","return Array(pests) ភ្លាម\n→ fast path (no filter)"),
        (PURPLE,"Category filter","ប្រើ Picker → filter by\npestType (Insects/Fungal)"),
        (ORANGE,"iOS 13+ safe","❌ .searchable(text:) requires\niOS 15+ — avoid!"),
    ]
    for i,(c,t,d) in enumerate(notes):
        cy=1.55+i*1.45
        add_rect(s,7.05,cy,6.0,1.3,CARD); add_rect(s,7.05,cy,.08,1.3,c)
        add_text(s,t,7.28,cy+.1,5.5,.4,size=12,bold=True,color=c)
        add_text(s,d,7.28,cy+.55,5.5,.65,size=11,color=WHITE)
    snum(s,6)

# ── Slide 7: Expandable Sections ────────────────────────────────────────────
def s7(prs):
    s=add_slide(prs)
    hdr(s,"📱  Expandable Sections · Detail View Pattern")
    add_bullets(s,[
        "@State private var showSymptoms = false — ១ state per section",
        "Button → toggle state → chevron rotates .rotationEffect()",
        "if showSymptoms { Text(pest.symptoms) } — conditional content",
        "DisclosureGroup (iOS 14+) = simpler API — ប៉ុន្ដែ ❌ iOS 13",
        "Manual toggle pattern works iOS 13+ — ✅ ប្រើក្នុង course",
        "NavigationLink → PestDetailView(pest: pest) from list row",
    ],.35,1.1,6.3,4.0,size=13,bc=GREEN)
    add_code(s,[
        "struct PestDetailView: View {",
        "    let pest: Pest",
        "    @State private var showSymptoms  = false",
        "    @State private var showTreatment = false",
        "    @State private var showPrevention = false",
        "",
        "    var body: some View {",
        "        List {",
        "            // Expandable section",
        "            Button(action: { showSymptoms.toggle() }) {",
        "                HStack {",
        '                    Text("រោគសញ្ញា")',
        "                        .font(.headline)",
        "                    Spacer()",
        "                    Image(systemName:",
        "                        showSymptoms",
        '                            ? "chevron.up"',
        '                            : "chevron.down")',
        "                }",
        "            }",
        "            if showSymptoms {",
        "                Text(pest.symptoms ?? \"\")",
        "                    .padding(.leading)",
        "            }",
        "        }",
        "        .navigationTitle(pest.name ?? \"\")",
        "    }",
        "}",
    ],.35,5.3,6.3,2.05,size=9.5)
    add_text(s,"PestDetailView Code",7.0,1.05,6.0,.4,size=12,bold=True,color=PURPLE)
    add_code(s,[
        "// PestRowView — List row",
        "struct PestRowView: View {",
        "    let pest: Pest",
        "    var body: some View {",
        "        HStack(spacing: 12) {",
        "            Image(systemName: \"ant.fill\")",
        "                .foregroundColor(.red)",
        "                .frame(width:40, height:40)",
        "                .background(Color.red.opacity(.1))",
        "                .cornerRadius(8)",
        "            VStack(alignment: .leading) {",
        "                Text(pest.name ?? \"\")",
        "                    .font(.headline)",
        "                Text(pest.pestType ?? \"\")",
        "                    .font(.caption)",
        "                    .foregroundColor(.gray)",
        "            }",
        "            Spacer()",
        "            if pest.isFavorite {",
        "                Image(systemName: \"star.fill\")",
        "                    .foregroundColor(.yellow)",
        "            }",
        "        }",
        "    }",
        "}",
    ],7.0,1.55,6.0,5.8,size=9.5)
    snum(s,7)

# ── Slide 8: Offline-First ──────────────────────────────────────────────────
def s8(prs):
    s=add_slide(prs)
    hdr(s,"🔌  Offline-First Design · ធ្វើការដោយគ្មាន Internet")
    rows=[
        ("First launch","Download from API ❌","Load pests.json from Bundle ✅"),
        ("Day-to-day use","Network request each time ❌","@FetchRequest from CoreData ✅"),
        ("Airplane mode","App breaks ❌","App works perfectly ✅"),
        ("Privacy","User data on server ❌","All data stays on device ✅"),
    ]
    hdrs=["Scenario","❌  Network-First","✅  Offline-First"]
    widths=[2.8,4.5,4.9]; starts=[.35,3.3,7.95]
    for c,(h,cw,cs) in enumerate(zip(hdrs,widths,starts)):
        add_rect(s,cs,1.1,cw,.55,GREEN if c==2 else RED if c==1 else DARK2)
        add_text(s,h,cs+.1,1.15,cw-.15,.45,size=13,bold=True,color=WHITE)
    for r,(sc,w,g) in enumerate(rows):
        cy=1.75+r*.8; bg=CARD if r%2==0 else DARK2
        for c,(cell,cw,cs) in enumerate(zip([sc,w,g],widths,starts)):
            add_rect(s,cs,cy,cw,.7,bg)
            col=GREY if c==0 else (RED if c==1 else GREEN)
            add_text(s,cell,cs+.1,cy+.12,cw-.15,.46,size=11,color=col)
    add_rect(s,.35,5.08,12.6,.04,DARK2)
    add_bullets(s,[
        "pests.json ← bundle ទៀងទាត់ — ចម្លង data ទៅ CoreData លើកដំបូងប៉ុណ្ណោះ",
        "UserDefaults flag \"pestDataSeeded\" → ចៀសវាង duplicate seeding",
        "CoreData @FetchRequest → always reads from local SQLite (fast + offline)",
        "Update data: release new app version with updated pests.json",
    ],.35,5.3,12.6,2.0,size=12,bc=GREEN)
    snum(s,8)

# ── Slide 9: Common Mistakes ────────────────────────────────────────────────
def s9(prs):
    s=add_slide(prs)
    hdr(s,"⚠️  Common Mistakes · ❌ vs ✅")
    rows=[
        ("Seed every launch","Run PestDataLoader every app start","Check UserDefaults flag first"),
        ("iOS 16+ search",".searchable(text:) modifier","Custom HStack TextField (iOS 13+)"),
        ("DisclosureGroup","DisclosureGroup {} (iOS 14+)","Manual @State toggle + if/else"),
        ("Filter in memory","Swift .filter{} on FetchedResults","NSPredicate on @FetchRequest"),
        ("Missing context",".environment missing on sheet","Always pass managedObjectContext"),
        ("Duplicate records","No seeded check","UserDefaults pestDataSeeded flag"),
    ]
    hdrs=["Pattern","❌  Wrong","✅  Correct (iOS 13+)"]
    widths=[2.8,4.5,5.1]; starts=[.35,3.3,7.95]
    for c,(h,cw,cs) in enumerate(zip(hdrs,widths,starts)):
        add_rect(s,cs,1.1,cw,.55,GREEN if c==2 else RED if c==1 else DARK2)
        add_text(s,h,cs+.1,1.15,cw-.15,.45,size=12,bold=True,color=WHITE)
    for r,(p,w,g) in enumerate(rows):
        cy=1.75+r*.67; bg=CARD if r%2==0 else DARK2
        for c,(cell,cw,cs) in enumerate(zip([p,w,g],widths,starts)):
            add_rect(s,cs,cy,cw,.62,bg)
            col=GREY if c==0 else (RED if c==1 else GREEN)
            add_text(s,cell,cs+.1,cy+.1,cw-.15,.42,size=10,color=col)
    snum(s,9)

# ── Slide 10: Design System ──────────────────────────────────────────────────
def s10(prs):
    s=add_slide(prs)
    hdr(s,"🎨  Design System · Color + SF Symbols + Layout")
    palette=[(GREEN,"GREEN","#1BB889","Category: Insect"),
             (BLUE, "BLUE", "#287DFA","Category: Fungal"),
             (PURPLE,"PURPLE","#8E44AD","Category: Bacterial"),
             (ORANGE,"ORANGE","#F39620","Warning tips"),
             (RED,  "RED",  "#E54747","Danger / severe")]
    for i,(c,n,h,u) in enumerate(palette):
        cx=.35+i*2.55
        add_rect(s,cx,1.1,2.35,.7,c)
        add_text(s,n,cx+.08,1.15,2.2,.35,size=11,bold=True,color=DARK)
        add_text(s,h,cx+.08,1.5,2.2,.25,size=9,color=DARK)
    add_text(s,"SF Symbols Used",0.35,2.1,5.5,.4,size=12,bold=True,color=GREEN)
    symbols=[("ant.fill","Pest list row"),("star.fill","isFavorite"),
             ("magnifyingglass","Search bar"),("chevron.down","Expandable"),
             ("slider.horizontal.3","Filter picker"),("plus","Add new pest")]
    for i,(sym,desc) in enumerate(symbols):
        col=i%2; row=i//2; cx=.35+col*3.0; cy=2.6+row*.75
        add_rect(s,cx,cy,2.8,.65,CARD)
        add_code(s,[sym],cx+.1,cy+.07,1.6,.5,size=10)
        add_text(s,desc,cx+1.85,cy+.12,.85,.4,size=10,color=GREY)
    add_text(s,"Layout Rules",6.55,2.1,6.5,.4,size=12,bold=True,color=BLUE)
    add_bullets(s,[
        "NavigationView → NavigationLink to PestDetailView",
        "List → ForEach(displayedPests, id: \\.self)",
        "isFavorite toggle → pest.isFavorite.toggle() + save",
        "HStack search bar → top of VStack, padding(.horizontal)",
        "Category Picker (Segmented) → Insects | Fungal | Bacterial",
    ],6.55,2.6,6.5,3.5,size=11,bc=GREEN)
    snum(s,10)

# ── Slide 11: Mini-Project ───────────────────────────────────────────────────
def s11(prs):
    s=add_slide(prs)
    hdr(s,"🏠  Mini-Project · Pest & Disease Guide ពេញ")
    cols=[
        ("📋  Requirements",GREEN,[
            "JSON load 10+ pests on first launch",
            "List with custom search bar",
            "Category filter: Insects/Fungal/Bacterial",
            "Detail view with expandable sections",
            "isFavorite toggle + star icon",
            "AddPestView form",
        ]),
        ("✅  Checklist",BLUE,[
            "pests.json bundled + decoded ✓",
            "UserDefaults seed flag ✓",
            "@FetchRequest in PestGuideTabView ✓",
            "Search filters name + pestType ✓",
            "Symptoms/Treatment/Prevention expand ✓",
            "star.fill shows isFavorite ✓",
        ]),
        ("🎯  Grading",PURPLE,[
            "JSON seeding: 25%",
            "Search bar: 20%",
            "Category filter: 20%",
            "Expandable detail: 20%",
            "CRUD (add/delete): 15%",
            "Bonus: isFavorite filter",
        ]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*4.35
        add_rect(s,cx,1.1,4.1,5.8,CARD); add_rect(s,cx,1.1,4.1,.55,c)
        add_text(s,t,cx+.12,1.15,3.85,.45,size=13,bold=True,color=DARK)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.15,cy,3.7,.72,DARK2)
            add_text(s,item,cx+.28,cy+.1,3.4,.52,size=11,color=WHITE)
    add_rect(s,.35,7.0,12.6,.42,GREEN)
    add_text(s,"🐛  Week 06 ចប់! — Pest Guide offline + search ready · Next: Daily Journal →",
             .5,7.02,12.3,.38,size=12,bold=True,color=DARK,align=PP_ALIGN.CENTER)
    snum(s,11)

# ── Slide 12: Summary ───────────────────────────────────────────════════════
def s12(prs):
    s=add_slide(prs)
    hdr(s,"📌  សង្ខេប Week 06 · Cheat Sheet",accent=GREEN)
    cols=[
        ("⚙️  Setup",GREEN,["Pest.xcdatamodeld entity","pests.json in Bundle","PestDataLoader service","UserDefaults seed flag"]),
        ("🔍  Search",BLUE,["@State searchText","localizedCaseInsensitiveContains","Picker category filter","displayedPests computed var"]),
        ("📱  Detail",PURPLE,["NavigationLink destination","@State showSymptoms etc","if showX { Text(...) }","chevron icon rotation"]),
        ("🔌  Offline",ORANGE,["Bundle.main.url(forResource:)","JSONDecoder().decode","CoreData @FetchRequest","No network needed"]),
    ]
    for i,(t,c,items) in enumerate(cols):
        cx=.35+i*3.25
        add_rect(s,cx,1.1,3.0,5.8,CARD); add_rect(s,cx,1.1,3.0,.55,c)
        add_text(s,t,cx+.1,1.15,2.8,.45,size=12,bold=True,color=DARK)
        for j,item in enumerate(items):
            cy=1.75+j*.9; add_rect(s,cx+.1,cy,2.8,.72,DARK2)
            add_text(s,item,cx+.18,cy+.1,2.6,.52,size=10,color=TEAL)
            s.shapes[-1].text_frame.paragraphs[0].runs[0].font.name="Courier New"
    add_rect(s,.35,7.0,12.6,.42,GREEN)
    add_text(s,"🌿  Offline = farmer-first! · Data on device = always available even in the rice field.",
             .5,7.02,12.3,.38,size=13,bold=True,color=DARK,align=PP_ALIGN.CENTER)
    snum(s,12)

def build():
    prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H)
    s1(prs);s2(prs);s3(prs);s4(prs);s5(prs);s6(prs)
    s7(prs);s8(prs);s9(prs);s10(prs);s11(prs);s12(prs)
    out=os.path.join(os.path.dirname(__file__),"Week06_Pest_Disease_KH.pptx")
    prs.save(out); print(f"✅  Saved → {out}")

if __name__=="__main__": build()
