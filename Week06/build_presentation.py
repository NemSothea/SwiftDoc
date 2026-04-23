"""Generate Keynote/PowerPoint-ready deck for Week 6 (Pest & Disease Guide).

Produces: Week06_Pest_Disease_Guide.pptx

Design system (also listed on the summary slide):
- Palette (matches the course design tokens established in Week 5):
    Primary Blue   #0A84FF   (iOS system blue — nav / accents)
    Leaf Green     #34C759   (pest/disease theme, expandable section icons)
    Alert Orange   #FF9500   (warnings, favorites star)
    Background     #F2F2F7
    Surface        #FFFFFF
    Text Primary   #1C1C1E
    Text Secondary #8E8E93
- Fonts: Title = SF Pro Display / Helvetica Neue Bold; Body = SF Pro Text
- Layout: 2-column (concept left, UI / code right)
- Icons: SF Symbols (ladybug, magnifyingglass, stethoscope, cross.case, shield.lefthalf.filled)
- Animations: Fade (build in), Slide-from-right for code, Zoom for the JSON preload diagram
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# --- Design tokens -----------------------------------------------------------

PRIMARY = RGBColor(0x0A, 0x84, 0xFF)   # iOS blue
GREEN   = RGBColor(0x34, 0xC7, 0x59)   # leaf green
ORANGE  = RGBColor(0xFF, 0x95, 0x00)   # alert orange
BG      = RGBColor(0xF2, 0xF2, 0xF7)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT    = RGBColor(0x1C, 0x1C, 0x1E)
MUTED   = RGBColor(0x8E, 0x8E, 0x93)
DIVIDER = RGBColor(0xD1, 0xD1, 0xD6)
CODE_BG = RGBColor(0x1C, 0x1C, 0x1E)
CODE_FG = RGBColor(0xF2, 0xF2, 0xF7)
TIP_BG  = RGBColor(0xFF, 0xF4, 0xE0)

TITLE_FONT = "SF Pro Display"
BODY_FONT  = "SF Pro Text"
MONO_FONT  = "Menlo"

# 16:9 deck
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


# --- Helpers -----------------------------------------------------------------

def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=18, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(40000)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT,
                font=BODY_FONT, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(80000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(40000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_card(slide, x, y, w, h, *, fill=SURFACE, border=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = border
        card.line.width = Pt(0.75)
    card.shadow.inherit = False
    return card


def slide_header(slide, eyebrow, title, *, accent=PRIMARY):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.55), Inches(0.55),
                                 Inches(0.08), Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    add_text(slide, Inches(0.75), Inches(0.45), Inches(12), Inches(0.35),
             eyebrow, font=BODY_FONT, size=12, bold=True, color=accent)
    add_text(slide, Inches(0.75), Inches(0.75), Inches(12), Inches(0.75),
             title, font=TITLE_FONT, size=30, bold=True, color=TEXT)


def speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_code_block(slide, x, y, w, h, code_lines, *, size=11):
    add_card(slide, x, y, w, h, fill=CODE_BG)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(140000)
    tf.margin_right = Emu(80000)
    tf.margin_top = Emu(100000)
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


def add_chip(slide, x, y, text, *, fill=PRIMARY, text_color=SURFACE, w=1.7, h=0.32):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    tf = shp.text_frame
    tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
    tf.margin_top = Emu(20000);  tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = BODY_FONT; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = text_color


def add_footer(slide, idx, total):
    add_text(slide, Inches(0.55), Inches(7.1), Inches(6), Inches(0.3),
             "Week 6 · Pest & Disease Guide · SmartFarmerAssistant",
             size=10, color=MUTED)
    add_text(slide, Inches(10.5), Inches(7.1), Inches(2.3), Inches(0.3),
             f"{idx} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def set_transition(slide, kind="fade"):
    xml_map = {
        "fade": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        "push": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        "zoom": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
    }
    node = etree.fromstring(xml_map.get(kind, xml_map["fade"]))
    slide.element.append(node)


TOTAL = 12


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = add_slide()

band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(4.5))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = GREEN

for (cx, cy, r, color) in [
    (11.0, 0.6, 1.6, PRIMARY),
    (12.3, 3.4, 0.9, ORANGE),
    (1.1, 3.2, 0.55, SURFACE),
]:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(cx), Inches(cy),
                              Inches(r), Inches(r))
    circ.line.fill.background()
    circ.fill.solid(); circ.fill.fore_color.rgb = color

add_chip(s, 0.6, 0.6, "iOS · SwiftUI · Week 6",
         fill=SURFACE, text_color=GREEN, w=2.4, h=0.38)

add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(1.5),
         "Pest & Disease Guide",
         font=TITLE_FONT, size=52, bold=True, color=SURFACE)
add_text(s, Inches(0.6), Inches(2.25), Inches(12), Inches(1.0),
         "An Offline Reference Library",
         font=TITLE_FONT, size=44, bold=True, color=SURFACE)
add_text(s, Inches(0.6), Inches(3.35), Inches(12), Inches(0.6),
         "SwiftUI  ·  Core Data  ·  JSON Preload  ·  Custom Search",
         font=BODY_FONT, size=22, color=SURFACE)

add_card(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.6))
add_text(s, Inches(0.95), Inches(5.15), Inches(6), Inches(0.45),
         "AUDIENCE", size=11, bold=True, color=MUTED)
add_text(s, Inches(0.95), Inches(5.45), Inches(6), Inches(0.6),
         "Beginner iOS Developers", font=TITLE_FONT, size=20, bold=True)
add_text(s, Inches(0.95), Inches(5.95), Inches(7), Inches(0.5),
         "Pest entity · JSON preload · Custom search bar · Expandable sections",
         size=13, color=MUTED)

add_text(s, Inches(7.8), Inches(5.15), Inches(4.8), Inches(0.45),
         "PROJECT 2 — SMART FARMER ASSISTANT",
         size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
add_text(s, Inches(7.8), Inches(5.45), Inches(4.8), Inches(0.6),
         "PestDisease Module", font=TITLE_FONT, size=20, bold=True,
         align=PP_ALIGN.RIGHT)
add_text(s, Inches(7.8), Inches(5.95), Inches(4.8), Inches(0.5),
         "iOS 13+ — ObservableObject · NavigationView · custom SearchBar",
         size=13, color=MUTED, align=PP_ALIGN.RIGHT)

add_footer(s, 1, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Welcome to Week 6. This week we build an offline reference library — the Pest & "
    "Disease Guide tab of the Smart Farmer Assistant. By the end, students will have a "
    "Pest Core Data entity, a JSON file that seeds the database on first launch, a "
    "custom iOS-13 compatible search bar, and a detail screen with expandable sections. "
    "Everything works in airplane mode from launch 1. "
    "UI suggestion: open with the finished app on screen — show the list, a quick "
    "search, and a detail expand — so learners see the target before we write code."
))


# ============================================================================
# SLIDE 2 — Agenda / Learning Objectives
# ============================================================================
s = add_slide()
slide_header(s, "AGENDA", "What you'll build this week", accent=GREEN)

add_text(s, Inches(0.75), Inches(1.9), Inches(5.8), Inches(0.4),
         "LEARNING OBJECTIVES", size=12, bold=True, color=GREEN)
add_bullets(s, Inches(0.75), Inches(2.25), Inches(6.1), Inches(4.5), [
    "Model a Pest entity in Core Data (name, symptoms, treatment, imageName)",
    "Ship a JSON file in the app bundle and preload it on first launch",
    "Build a custom iOS 13+ search bar (no .searchable)",
    "Expand/collapse sections with local @State",
    "Design an offline-first feature — works with zero network",
], size=17)

lessons = [
    ("6.1", "Pest Core Data Entity"),
    ("6.2", "JSON Preload on First Launch"),
    ("6.3", "Custom iOS 13+ Search Bar"),
    ("6.4", "Expandable Sections with @State"),
    ("6.5", "Putting it together — PestGuideTabView"),
]
for i, (num, title) in enumerate(lessons):
    top = Inches(1.9 + i * 0.95)
    add_card(s, Inches(7.1), top, Inches(5.6), Inches(0.8))
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(7.25), top + Inches(0.15),
                              Inches(0.5), Inches(0.5))
    chip.line.fill.background()
    chip.fill.solid(); chip.fill.fore_color.rgb = GREEN
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = TITLE_FONT; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = SURFACE
    add_text(s, Inches(7.95), top + Inches(0.2), Inches(4.5), Inches(0.5),
             title, font=TITLE_FONT, size=16, bold=True)

add_footer(s, 2, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Walk the five lessons at a high level; the payoff is lesson 6.2 (JSON preload) — "
    "that single mechanism is what makes the app 'offline first'. Stress that we pick "
    "iOS 13+ APIs all the way — no .searchable, no DisclosureGroup-required, no "
    "AsyncImage. "
    "UI suggestion: objectives on the left, numbered lesson chips on the right — "
    "build in the chips one-by-one with a fade."
))


# ============================================================================
# SLIDE 3 — Architecture
# ============================================================================
s = add_slide()
slide_header(s, "ARCHITECTURE", "MVVM — clean SwiftUI layers", accent=PRIMARY)

layers = [
    ("View",       "PestGuideTabView · PestRowView · PestDetailView · SearchBar · ExpandableSection", PRIMARY),
    ("ViewModel",  "PestGuideViewModel — @Published searchText + filter(_:)",                         GREEN),
    ("Service",    "PestDataLoader — Bundle → JSONDecoder → Core Data (once)",                         ORANGE),
    ("Model",      "Pest (NSManagedObject) + PestDTO (Codable)",                                       RGBColor(0xA0, 0x52, 0xD4)),
    ("Resource",   "pests.json — committed to the repo, bundled with the app",                         MUTED),
]
y = 1.9
for role, desc, color in layers:
    add_card(s, Inches(0.75), Inches(y), Inches(11.8), Inches(0.85))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.75), Inches(y),
                             Inches(0.18), Inches(0.85))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    add_text(s, Inches(1.1), Inches(y + 0.12), Inches(2.2), Inches(0.4),
             role, font=TITLE_FONT, size=16, bold=True, color=color)
    add_text(s, Inches(3.4), Inches(y + 0.2), Inches(9.0), Inches(0.5),
             desc, size=13, color=TEXT)
    y += 1.00

add_footer(s, 3, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "This is the same layering students saw in Finance (Week 4) and CalendarReminders "
    "(Week 5). The new box here is 'Service' — PestDataLoader — because the data "
    "originates outside the runtime (JSON file shipped with the app) and has to be "
    "promoted into Core Data exactly once. "
    "UI suggestion: five stacked cards, colored accent bar on the left of each. "
    "Animation: fade in top-to-bottom, 120ms delay between cards."
))


# ============================================================================
# SLIDE 4 — Folder Structure
# ============================================================================
s = add_slide()
slide_header(s, "CODE LAYOUT", "PestDisease module — folder by folder", accent=PRIMARY)

add_bullets(s, Inches(0.75), Inches(1.9), Inches(5.5), Inches(4.5), [
    "Mirrors Finance/ and CalendarReminders/ — same shape across modules",
    "Models/ holds both the NSManagedObject class and the Codable DTO",
    "Services/ isolates the JSON-read + Core Data write",
    "ViewModels/ owns search text and filter logic",
    "Resources/ ships the seed JSON inside the app bundle",
], size=16)

tree = [
    "SmartFarmerAssistantFinish/",
    "└── PestDisease/",
    "    ├── Models/",
    "    │   ├── Pest+CoreDataClass.swift",
    "    │   ├── Pest+CoreDataProperties.swift",
    "    │   └── PestDTO.swift",
    "    ├── ViewModels/",
    "    │   └── PestGuideViewModel.swift",
    "    ├── Services/",
    "    │   └── PestDataLoader.swift",
    "    ├── Views/",
    "    │   ├── PestGuideTabView.swift",
    "    │   ├── PestRowView.swift",
    "    │   ├── PestDetailView.swift",
    "    │   ├── SearchBar.swift",
    "    │   └── ExpandableSection.swift",
    "    └── Resources/",
    "        └── pests.json",
]
add_code_block(s, Inches(6.55), Inches(1.85), Inches(6.2), Inches(5.0),
               tree, size=11)

add_footer(s, 4, TOTAL)
set_transition(s, "push")
speaker_notes(s, (
    "Walk the folder tree top to bottom. Emphasise that every week we've opened a new "
    "folder alongside Finance/ and CalendarReminders/ — not a new project. Shared "
    "Core Data model, shared CoreDataManager. "
    "UI suggestion: bullets on the left describing WHY each folder exists, the actual "
    "tree rendered as a dark code block on the right. "
    "Animation: slide-from-right for the code block."
))


# ============================================================================
# SLIDE 5 — Data Flow
# ============================================================================
s = add_slide()
slide_header(s, "DATA FLOW", "From JSON file to SwiftUI row", accent=PRIMARY)

# Stage boxes — 5 stages left-to-right
stages = [
    ("pests.json",         "in the app\nbundle",            MUTED,   "doc.text"),
    ("PestDataLoader",     "decode JSON\n→ Pest rows",      ORANGE,  "arrow.down.doc"),
    ("Core Data",          "local\npersistent\nstore",      PRIMARY, "cylinder"),
    ("@FetchRequest",      "auto-observes\nthe Pest table", GREEN,   "bolt.horizontal"),
    ("PestGuideTabView",   "List + Search\n+ Detail",       GREEN,   "list.bullet"),
]
count = len(stages)
gap_x = 0.35
box_w = (13.333 - 1.5 - (count - 1) * gap_x) / count
box_h = 2.2
y = 2.8
for i, (title, desc, color, _icon) in enumerate(stages):
    x = 0.75 + i * (box_w + gap_x)
    add_card(s, Inches(x), Inches(y), Inches(box_w), Inches(box_h),
             fill=SURFACE, border=color)
    add_text(s, Inches(x), Inches(y + 0.25), Inches(box_w), Inches(0.5),
             title, font=TITLE_FONT, size=15, bold=True, color=color,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(x + 0.1), Inches(y + 0.85), Inches(box_w - 0.2), Inches(1.3),
             desc, size=12, color=TEXT, align=PP_ALIGN.CENTER)

    # arrow between boxes
    if i < count - 1:
        arr_x = x + box_w + 0.02
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(arr_x), Inches(y + box_h / 2 - 0.12),
                                 Inches(gap_x - 0.04), Inches(0.24))
        arr.line.fill.background()
        arr.fill.solid(); arr.fill.fore_color.rgb = MUTED

# bottom caption
add_text(s, Inches(0.75), Inches(5.5), Inches(11.8), Inches(0.5),
         "ONE-TIME on first launch (steps 1→3) · EVERY session after (steps 3→5)",
         size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 5, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Read left-to-right. The key insight: steps 1→3 happen ONCE (on first launch). "
    "After that, the app only ever reads 3→5. This is what makes the experience "
    "feel native — no spinner, no 'loading', and no network dependency even on the "
    "very first tap. "
    "UI suggestion: five color-coded cards in a row with arrows between them; the "
    "caption under them clarifies the temporal split. "
    "Animation: slide-from-left for each box with a 100ms delay."
))


# ============================================================================
# SLIDE 6 — Core Data Pest Entity
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 6.1", "Core Data — Pest Entity", accent=PRIMARY)

add_text(s, Inches(0.75), Inches(1.85), Inches(6.0), Inches(0.4),
         "ATTRIBUTES", size=12, bold=True, color=PRIMARY)

attrs = [
    ("id",          "UUID",   "Stable identifier — also the row key"),
    ("name",        "String", "Shown in the list + detail title"),
    ("symptoms",    "String", "What the farmer sees on the plant"),
    ("treatment",   "String", "What to do once confirmed"),
    ("imageName",   "String", "Asset-catalog name (optional, SF Symbol fallback)"),
    ("pestType",    "String", "e.g. 'Disease (fungus)', 'Pest (insect)'"),
    ("prevention",  "String", "Optional — collapsible section"),
]
row_y = 2.25
for name, typ, desc in attrs:
    add_card(s, Inches(0.75), Inches(row_y), Inches(6.0), Inches(0.55))
    add_text(s, Inches(0.95), Inches(row_y + 0.07), Inches(1.9), Inches(0.4),
             name, font=MONO_FONT, size=12, bold=True, color=PRIMARY)
    add_text(s, Inches(2.85), Inches(row_y + 0.07), Inches(0.9), Inches(0.4),
             typ, font=MONO_FONT, size=11, color=MUTED)
    add_text(s, Inches(3.8), Inches(row_y + 0.07), Inches(2.85), Inches(0.4),
             desc, size=10, color=TEXT)
    row_y += 0.62

code = [
    "// PestDisease/Models/Pest+CoreDataProperties.swift",
    "extension Pest {",
    "  @NSManaged public var id: UUID?",
    "  @NSManaged public var name: String?",
    "  @NSManaged public var symptoms: String?",
    "  @NSManaged public var treatment: String?",
    "  @NSManaged public var imageName: String?",
    "  @NSManaged public var pestType: String?",
    "  @NSManaged public var prevention: String?",
    "  @NSManaged public var isFavorite: Bool",
    "}",
    "",
    "extension Pest: Identifiable {}",
]
add_code_block(s, Inches(7.25), Inches(2.1), Inches(5.45), Inches(3.5), code, size=12)

add_card(s, Inches(7.25), Inches(5.75), Inches(5.45), Inches(1.2), fill=TIP_BG)
add_text(s, Inches(7.5), Inches(5.85), Inches(5.0), Inches(0.4),
         "⚠️  Codegen: Manual / None",
         font=TITLE_FONT, size=13, bold=True, color=ORANGE)
add_text(s, Inches(7.5), Inches(6.2), Inches(5.0), Inches(0.7),
         "Set this in the Data Model Inspector — otherwise Xcode auto-generates the "
         "class and you'll see duplicate-symbol errors.",
         size=11, color=TEXT)

add_footer(s, 6, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Walk each attribute. The four core ones (name, symptoms, treatment, imageName) "
    "are required by the spec; the three extras (pestType, prevention, isFavorite) "
    "are optional enrichments that demonstrate how to extend a model without breaking "
    "previous features. Show the Data Model Inspector live and point to Codegen → "
    "Manual/None — this is the #1 cause of duplicate-class errors for learners. "
    "UI suggestion: attribute cards left, code right, orange warning card bottom-right. "
    "Animation: slide-from-left for cards, slide-from-right for the code block."
))


# ============================================================================
# SLIDE 7 — JSON Preload Flow (diagram)
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 6.2", "JSON preload flow — first launch only", accent=ORANGE)

# Vertical flowchart down the center
nodes = [
    ("App launches",                               PRIMARY),
    ("preloadIfNeeded(context:)",                  ORANGE),
    ("kPestsPreloaded == true?  →  RETURN",         MUTED),
    ("Bundle.main.url(\"pests.json\")",             PRIMARY),
    ("JSONDecoder → [PestDTO]",                    PRIMARY),
    ("for each DTO → Pest(context:)",              PRIMARY),
    ("context.save()",                              GREEN),
    ("UserDefaults.set(true, forKey: kPestsPreloaded)", GREEN),
]
box_w = 5.6
box_h = 0.55
gap  = 0.16
start_y = 1.85
x = (13.333 - box_w) / 2
for i, (label, color) in enumerate(nodes):
    y = start_y + i * (box_h + gap)
    add_card(s, Inches(x), Inches(y), Inches(box_w), Inches(box_h),
             fill=SURFACE, border=color)
    add_text(s, Inches(x), Inches(y + 0.1), Inches(box_w), Inches(0.4),
             label, font=MONO_FONT, size=13, bold=True, color=color,
             align=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        ay = y + box_h
        arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                 Inches(x + box_w / 2 - 0.1),
                                 Inches(ay), Inches(0.2), Inches(gap))
        arr.line.fill.background()
        arr.fill.solid(); arr.fill.fore_color.rgb = MUTED

# right-side notes
add_text(s, Inches(9.1), Inches(1.9), Inches(3.8), Inches(0.4),
         "WHY THIS SHAPE", size=12, bold=True, color=ORANGE)
add_bullets(s, Inches(9.1), Inches(2.3), Inches(3.9), Inches(4.0), [
    "Called once from App.init()",
    "UserDefaults flag is cheap + synchronous",
    "Flag NOT set on failure — retry next launch",
    "Static enum — no accidental double-load",
    "DTO decoupled from Core Data schema",
], size=13)

add_footer(s, 7, TOTAL)
set_transition(s, "zoom")
speaker_notes(s, (
    "This is the heart of Week 6. Walk every step of the flow and link it back to the "
    "PestDataLoader code. Emphasise two details: (1) the flag is NOT flipped when the "
    "load fails — that means a broken ship of pests.json can recover on the next launch, "
    "instead of silently leaving the user with an empty store forever; (2) it's called "
    "from App.init(), not from onAppear — onAppear fires every time the view appears, "
    "which would re-trigger work on every tab switch. "
    "UI suggestion: 8 stacked cards down the centre connected by down-arrows; "
    "reasoning notes on the right. "
    "Animation: zoom + fade for the flow, fade in the bullets last."
))


# ============================================================================
# SLIDE 8 — Custom Search Bar
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 6.3", "Custom search bar — iOS 13+", accent=PRIMARY)

add_bullets(s, Inches(0.75), Inches(1.9), Inches(5.5), Inches(4.8), [
    ".searchable is iOS 15+ — we build our own",
    "TextField in a rounded systemGray6 container",
    "Magnifying-glass leading icon + xmark.circle.fill clear button",
    "isEditing (@State Bool) driven by onEditingChanged",
    "Cancel button appears only while field has focus",
    "Filtering lives in the ViewModel — searchable fields: name, symptoms, treatment, type",
], size=15)

# iPhone-style mock
mock_x = 7.3; mock_y = 1.85
add_card(s, Inches(mock_x), Inches(mock_y), Inches(5.4), Inches(5.0))

# search bar pill
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(mock_x + 0.25), Inches(mock_y + 0.3),
                          Inches(4.0), Inches(0.45))
pill.adjustments[0] = 0.5
pill.line.fill.background()
pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0xE5, 0xE5, 0xEA)
add_text(s, Inches(mock_x + 0.4), Inches(mock_y + 0.37), Inches(3.8), Inches(0.35),
         "🔍  rice blast",
         size=12, color=TEXT)
# Cancel button
add_text(s, Inches(mock_x + 4.3), Inches(mock_y + 0.37), Inches(1.0), Inches(0.35),
         "បោះបង់", size=12, bold=True, color=PRIMARY)

# result rows
rows = [
    ("Rice blast",          "Disease (fungus)"),
    ("Brown planthopper",   "Pest (insect)"),
    ("Cassava mosaic",      "Disease (virus)"),
    ("Fall armyworm",       "Pest (caterpillar)"),
]
ry = mock_y + 1.0
for title, sub in rows:
    add_text(s, Inches(mock_x + 0.4), Inches(ry), Inches(4.6), Inches(0.35),
             title, font=TITLE_FONT, size=14, bold=True, color=TEXT)
    add_text(s, Inches(mock_x + 0.4), Inches(ry + 0.32), Inches(4.6), Inches(0.3),
             sub, size=11, color=MUTED)
    ry += 0.85

add_footer(s, 8, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Walk through each bullet. Emphasise WHY we don't use .searchable: the course "
    "minimum deployment target is iOS 13. Show that the 'Cancel' button animates in "
    "only while editing — a small UX detail that makes the bar feel like Apple's. "
    "Then show the ViewModel's filter(_:) — it's just a guard + .filter — and link "
    "back to the bullet that says 'filtering lives in the ViewModel'. "
    "UI suggestion: bullets left, iPhone-shape mock right with the live search pill + "
    "four result rows. "
    "Animation: fade in bullets; slide-from-top for the search pill; fade in rows."
))


# ============================================================================
# SLIDE 9 — Expandable Sections
# ============================================================================
s = add_slide()
slide_header(s, "LESSON 6.4", "Expandable sections with @State", accent=GREEN)

add_bullets(s, Inches(0.75), Inches(1.9), Inches(5.5), Inches(4.8), [
    "One generic struct — ExpandableSection<Content: View>",
    "Each section owns its own @State isExpanded",
    "Opening Treatment does NOT close Symptoms",
    "Chevron rotates 90° with .rotationEffect — no extra icons",
    "Content is a @ViewBuilder closure — pass anything",
    "iOS 13+ compatible — no DisclosureGroup required",
], size=15)

# Right — mock detail screen
mx = 7.3; my = 1.85
add_card(s, Inches(mx), Inches(my), Inches(5.4), Inches(5.0))

# title block
add_text(s, Inches(mx + 0.3), Inches(my + 0.2), Inches(5.0), Inches(0.4),
         "Rice blast", font=TITLE_FONT, size=18, bold=True, color=TEXT)
add_text(s, Inches(mx + 0.3), Inches(my + 0.55), Inches(5.0), Inches(0.3),
         "Disease (fungus)", size=12, color=MUTED)

sections = [
    ("Symptoms",   "stethoscope",         True,  "Diamond-shaped brown spots with gray centers …"),
    ("Treatment",  "cross.case",          True,  "Apply fungicides labeled for rice blast …"),
    ("Prevention", "shield.lefthalf",     False, "Use resistant varieties when available …"),
]
sy = my + 1.0
for title, _icon, expanded, body in sections:
    # header bar
    header_w = 5.0
    arrow = "▾" if expanded else "▸"
    add_text(s, Inches(mx + 0.3), Inches(sy), Inches(header_w), Inches(0.4),
             f"{arrow}  {title}",
             font=TITLE_FONT, size=15, bold=True, color=GREEN)
    sy += 0.42
    if expanded:
        add_text(s, Inches(mx + 0.55), Inches(sy), Inches(header_w - 0.3), Inches(0.7),
                 body, size=11, color=TEXT)
        sy += 0.75
    # divider
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(mx + 0.3), Inches(sy),
                              Inches(header_w), Inches(0.012))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = DIVIDER
    sy += 0.18

add_footer(s, 9, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Emphasise the 'each section owns its own @State' rule — the biggest mistake "
    "learners make here is to hoist one Bool for all sections into the detail view, "
    "which forces them open/close together. Show the ExpandableSection struct live "
    "and point to the private @State and the generic Content parameter. "
    "UI suggestion: bullets left, mock detail screen right showing two open sections "
    "('Symptoms' + 'Treatment') and one collapsed ('Prevention'). "
    "Animation: slide-up reveal for the body text when a section opens."
))


# ============================================================================
# SLIDE 10 — Offline-First Concept
# ============================================================================
s = add_slide()
slide_header(s, "CONCEPT", "Offline-first — what it really means", accent=ORANGE)

# comparison table (2 columns x 4 rows)
add_text(s, Inches(0.75), Inches(1.85), Inches(5.8), Inches(0.4),
         "❌ NETWORK-FIRST", size=12, bold=True, color=RGBColor(0xD0, 0x30, 0x30))
add_text(s, Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.4),
         "✅ OFFLINE-FIRST", size=12, bold=True, color=GREEN)

rows = [
    ("First launch",
     "Spinner while pulling JSON from the server",
     "Read bundled pests.json — zero latency"),
    ("Updates",
     "Poll an API on every open",
     "Ship a new pests.json with the next app release"),
    ("Airplane mode",
     "Blocked screen — 'no connection'",
     "Full app works exactly the same"),
    ("Privacy",
     "Every read leaks to a server",
     "No network traffic at all"),
]
y = 2.3
for label, bad, good in rows:
    add_text(s, Inches(0.75), Inches(y), Inches(1.5), Inches(0.4),
             label, size=13, bold=True, color=TEXT)
    add_card(s, Inches(2.35), Inches(y - 0.05), Inches(4.5), Inches(0.9),
             fill=RGBColor(0xFF, 0xEC, 0xEC))
    add_text(s, Inches(2.55), Inches(y + 0.08), Inches(4.2), Inches(0.8),
             bad, size=12, color=TEXT)
    add_card(s, Inches(7.0), Inches(y - 0.05), Inches(5.6), Inches(0.9),
             fill=RGBColor(0xE8, 0xF7, 0xEB))
    add_text(s, Inches(7.2), Inches(y + 0.08), Inches(5.3), Inches(0.8),
             good, size=12, color=TEXT)
    y += 1.1

# rule-of-thumb callout
add_card(s, Inches(0.75), Inches(6.55), Inches(11.8), Inches(0.45),
         fill=TIP_BG)
add_text(s, Inches(0.95), Inches(6.6), Inches(11.4), Inches(0.35),
         "Rule of thumb: if the user has to wait for the network on first tap, it isn't offline-first.",
         size=13, bold=True, color=ORANGE)

add_footer(s, 10, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "This is a concept slide — don't read every row. Pick one (airplane mode is the "
    "most visceral) and live-demo by toggling airplane mode on the device, then "
    "showing the app still works. "
    "UI suggestion: four rows of red-vs-green cards; callout footer carries the "
    "one-line rule of thumb. "
    "Animation: fade in rows top-to-bottom; pop the rule-of-thumb card last."
))


# ============================================================================
# SLIDE 11 — UI / UX Suggestions
# ============================================================================
s = add_slide()
slide_header(s, "DESIGN SYSTEM", "UI / UX — Apple-style recipe", accent=PRIMARY)

# Palette row
add_text(s, Inches(0.75), Inches(1.85), Inches(5.0), Inches(0.4),
         "PALETTE", size=12, bold=True, color=PRIMARY)
palette = [
    ("#0A84FF", "iOS Blue",    PRIMARY),
    ("#34C759", "Leaf Green",  GREEN),
    ("#FF9500", "Alert Orange", ORANGE),
    ("#F2F2F7", "Background",  BG),
    ("#1C1C1E", "Text",        TEXT),
]
px = 0.75
for hex_label, name, color in palette:
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(px), Inches(2.25),
                              Inches(2.2), Inches(0.65))
    chip.adjustments[0] = 0.12
    chip.line.fill.background()
    chip.fill.solid(); chip.fill.fore_color.rgb = color
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(80000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = hex_label
    r.font.name = MONO_FONT; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = SURFACE if color != BG else TEXT
    add_text(s, Inches(px), Inches(2.95), Inches(2.2), Inches(0.3),
             name, size=10, color=MUTED)
    px += 2.4

# SF Symbols card
add_card(s, Inches(0.75), Inches(3.85), Inches(6.0), Inches(3.15))
add_text(s, Inches(0.95), Inches(3.95), Inches(5.6), Inches(0.4),
         "SF SYMBOLS", size=12, bold=True, color=PRIMARY)
symbols = [
    ("ladybug.fill",        "Tab icon"),
    ("magnifyingglass",     "Search"),
    ("xmark.circle.fill",   "Clear search"),
    ("stethoscope",         "Symptoms"),
    ("cross.case",          "Treatment"),
    ("shield.lefthalf.filled", "Prevention"),
    ("star.fill",           "Favorite"),
]
for i, (sym, purpose) in enumerate(symbols):
    row_y = 4.35 + i * 0.36
    add_text(s, Inches(0.95), Inches(row_y), Inches(3.2), Inches(0.3),
             sym, font=MONO_FONT, size=11, bold=True, color=GREEN)
    add_text(s, Inches(4.15), Inches(row_y), Inches(2.6), Inches(0.3),
             purpose, size=11, color=TEXT)

# Layout recipe card
add_card(s, Inches(7.0), Inches(3.85), Inches(5.7), Inches(3.15))
add_text(s, Inches(7.2), Inches(3.95), Inches(5.3), Inches(0.4),
         "LAYOUT RECIPE", size=12, bold=True, color=PRIMARY)
recipe = [
    "Rounded systemGray6 search bar pinned to top",
    "PlainListStyle — no inset gaps between rows",
    "44pt thumbnail — asset if available, SF Symbol fallback",
    ".title2.bold() name · .subheadline muted type",
    "Detail header: full-width image OR pastel-green placeholder",
    "Sections: icon + title row, chevron rotates 90° on open",
    "Yellow star in the detail nav bar — favorites",
]
for i, line in enumerate(recipe):
    add_text(s, Inches(7.2), Inches(4.35 + i * 0.36), Inches(5.3), Inches(0.3),
             "•  " + line, size=11, color=TEXT)

add_footer(s, 11, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Design-system slide. Don't linger — point to the palette, read two or three SF "
    "Symbols, and mention one layout rule (usually 'pinned search bar'). Link it "
    "back to the course's ongoing design tokens from Weeks 4 and 5. "
    "UI suggestion: palette row top, two cards below (SF Symbols inventory + layout "
    "recipe bullets). "
    "Animation: fade in the palette chips sequentially, then pop both cards."
))


# ============================================================================
# SLIDE 12 — Summary + Q&A
# ============================================================================
s = add_slide()

# Hero band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(3.2))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = GREEN

add_chip(s, 0.6, 0.6, "WEEK 6 · WRAP",
         fill=SURFACE, text_color=GREEN, w=1.8, h=0.38)

add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
         "You built an offline reference library.",
         font=TITLE_FONT, size=38, bold=True, color=SURFACE)
add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.8),
         "Pest entity · JSON preload · Custom search · Expandable sections",
         font=BODY_FONT, size=18, color=SURFACE)

# What you shipped
add_card(s, Inches(0.6), Inches(3.5), Inches(6.0), Inches(3.3))
add_text(s, Inches(0.85), Inches(3.65), Inches(5.6), Inches(0.4),
         "WHAT YOU SHIPPED", size=12, bold=True, color=GREEN)
bullets = [
    "Pest Core Data entity with 8 attributes",
    "PestDataLoader — JSON → Core Data on first launch",
    "pests.json bundled in Resources/ (6 seeded entries)",
    "SearchBar (iOS 13+) + PestGuideViewModel filter",
    "ExpandableSection — reusable disclosure with @State",
    "PestGuideTabView — List + Search + Detail screen",
]
for i, b in enumerate(bullets):
    add_text(s, Inches(0.85), Inches(4.05 + i * 0.42), Inches(5.7), Inches(0.4),
             "•  " + b, size=13, color=TEXT)

# Q&A + next week
add_card(s, Inches(6.75), Inches(3.5), Inches(6.0), Inches(3.3))
add_text(s, Inches(7.0), Inches(3.65), Inches(5.6), Inches(0.4),
         "QUESTIONS TO DRIVE DISCUSSION", size=12, bold=True, color=PRIMARY)
qs = [
    "Why preload from App.init() instead of onAppear?",
    "What happens if pests.json has a typo on first launch?",
    "Why do we use a DTO instead of decoding straight into Pest?",
    "Why give each ExpandableSection its own @State?",
    "What single change would make this network-first?",
]
for i, q in enumerate(qs):
    add_text(s, Inches(7.0), Inches(4.05 + i * 0.42), Inches(5.7), Inches(0.4),
             "?  " + q, size=13, color=TEXT)

add_text(s, Inches(7.0), Inches(6.3), Inches(5.6), Inches(0.4),
         "NEXT WEEK →  Farm Journal with photos & weather (Week 7)",
         size=12, bold=True, color=PRIMARY)

add_footer(s, 12, TOTAL)
set_transition(s, "fade")
speaker_notes(s, (
    "Recap what the class shipped. Use the questions on the right to drive the Q&A — "
    "don't lecture, invite answers. The single most important takeaway: preload from "
    "App.init() so the work happens ONCE, and the UserDefaults flag is the gate. "
    "Tease Week 7 (Journal) only briefly — it reuses the same MVVM + Core Data pattern "
    "but adds photo capture. "
    "UI suggestion: green hero band up top, two cards below — 'What you shipped' and "
    "'Questions'. "
    "Animation: fade in the cards after the hero band; reveal the questions one at a time."
))

# ============================================================================
# Save
# ============================================================================
out_path = "Week06_Pest_Disease_Guide.pptx"
prs.save(out_path)
print(f"Wrote {out_path} — {len(prs.slides)} slides")
